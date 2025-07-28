# Standard library imports
import datetime
import logging
import os
import tempfile

import random
import re
import string
import subprocess
import sys
from collections import Counter
from decimal import Decimal
from io import BytesIO
from PyQt6.QtWidgets import QTabWidget
# Third-party imports
import matplotlib.dates as mdates
import matplotlib.font_manager as fm
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import matplotlib.ticker as ticker
import numpy as np
import openpyxl
import pandas as pd
import pymysql
import xlsxwriter
from cryptography.fernet import Fernet
from matplotlib.dates import DateFormatter, YearLocator
from matplotlib.ticker import AutoMinorLocator
from openpyxl import load_workbook
from openpyxl.chart import LineChart, Reference
from openpyxl.chart.label import DataLabelList
from openpyxl.styles import Font
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from PyQt6.QtCore import QThread, QTimer, Qt, pyqtSignal
from PyQt6.QtGui import QColor, QFont, QFontDatabase, QIcon, QIntValidator, QPalette, QPixmap
from PyQt6.QtWidgets import (
    QAbstractItemView, QApplication, QCheckBox, QComboBox, QDialog,
    QDialogButtonBox, QFileDialog, QGroupBox, QHBoxLayout, QHeaderView,
    QInputDialog, QLabel, QListWidget, QLineEdit, QMainWindow, QMessageBox,
    QProgressDialog, QPushButton, QScrollArea, QTableWidget, QTableWidgetItem,
    QVBoxLayout, QWidget
)
from scipy import stats
from sklearn.linear_model import LinearRegression
from tqdm import tqdm
from xlsxwriter.utility import xl_rowcol_to_cell

# 设置中文字体
plt.rcParams['font.sans-serif'] = ['SimHei']  # 使用黑体
plt.rcParams['axes.unicode_minus'] = False  # 正确显示负号

logging.basicConfig(level=logging.CRITICAL)
# 定义日志文件的路径
# 方案1：将日志保存到用户的Documents目录
log_file = os.path.join(tempfile.gettempdir(), "mating_app.log")
print("Current working directory:", os.getcwd())

# 配置日志记录
logging.basicConfig(
    filename='mating_app.log',
    level=logging.DEBUG,
    format='%(asctime)s - %(levelname)s - %(pathname)s:%(lineno)d - %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)

# 创建一个新的日志处理器,用于输出到控制台
console_handler = logging.StreamHandler()
console_handler.setLevel(logging.INFO)
formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
console_handler.setFormatter(formatter)

# 将控制台处理器添加到根日志器
logging.getLogger('').addHandler(console_handler)

# 配置异常日志记录
def exception_logger(exc_type, exc_value, exc_traceback):
    logging.exception("Uncaught exception", exc_info=(exc_type, exc_value, exc_traceback))

sys.excepthook = exception_logger

import time
from cryptography.fernet import Fernet

# 加密数据库连接信息
key = b'XGf7ZRXtj53qNCm9Ziuey22yXXEkzSq9FBTWZpfJiow='  # 实际密钥
cipher_suite = Fernet(key)

encoded_host = cipher_suite.encrypt(b'defectgene-new.mysql.polardb.rds.aliyuncs.com')
encoded_port = cipher_suite.encrypt(b'3306')
encoded_user = cipher_suite.encrypt(b'defect_genetic_checking')
encoded_password = cipher_suite.encrypt(b'Jaybz@890411')
encoded_db = cipher_suite.encrypt(b'bull_library')


class LoginDialog(QDialog):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("登录")
        self.setFixedSize(300, 150)
        self.setWindowFlags(self.windowFlags() | Qt.WindowType.WindowStaysOnTopHint)  # 设置窗口置顶
        self.setWindowIcon(QIcon("002.ico"))  # 设置应用程序图标

        self.layout = QVBoxLayout()

        self.username_label = QLabel("账号:")
        self.username_input = QLineEdit()
        self.layout.addWidget(self.username_label)
        self.layout.addWidget(self.username_input)

        self.password_label = QLabel("密码:")
        self.password_input = QLineEdit()
        self.password_input.setEchoMode(QLineEdit.EchoMode.Password)
        self.layout.addWidget(self.password_label)
        self.layout.addWidget(self.password_input)

        self.button_layout = QHBoxLayout()
        self.login_button = QPushButton("登录")
        self.login_button.clicked.connect(self.login)

        self.cancel_button = QPushButton("取消")
        self.cancel_button.clicked.connect(self.reject)
        self.button_layout.addWidget(self.login_button)
        self.button_layout.addWidget(self.cancel_button)
        self.layout.addLayout(self.button_layout)

        # 设置等待标签的字体和样式
        self.waiting_label_text = QLabel("  正在连接数据库......", self)
        self.waiting_label_hint = QLabel("（请等待2-5秒）", self)
        
        # 为“正在登陆中...”部分设置字体
        waiting_font = QFont("微软雅黑", 15, QFont.Weight.Medium)
        self.waiting_label_text.setFont(waiting_font)
        self.waiting_label_text.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.waiting_label_hint.setAlignment(Qt.AlignmentFlag.AlignCenter)
        
        self.waiting_layout = QVBoxLayout()
        self.waiting_layout.addWidget(self.waiting_label_text)
        self.waiting_layout.addWidget(self.waiting_label_hint)
        
        self.waiting_widget = QWidget()
        self.waiting_widget.setLayout(self.waiting_layout)
        self.waiting_widget.hide()
        
        self.layout.addWidget(self.waiting_widget)

        self.setLayout(self.layout)

    def show_waiting(self):
        for widget in [self.username_label, self.username_input, self.password_label, self.password_input, self.login_button, self.cancel_button]:
            widget.hide()
        self.waiting_widget.show()

    def login(self):
        username = self.username_input.text()
        password = self.password_input.text()

        self.show_waiting()

        QTimer.singleShot(100, lambda: self.process_login(username, password))

    def process_login(self, username, password):
        try:
            connection = pymysql.connect(
                host=cipher_suite.decrypt(encoded_host).decode(),
                port=int(cipher_suite.decrypt(encoded_port).decode()),
                user=cipher_suite.decrypt(encoded_user).decode(),
                password=cipher_suite.decrypt(encoded_password).decode(),
                database=cipher_suite.decrypt(encoded_db).decode(),
                charset='utf8mb4',
                cursorclass=pymysql.cursors.DictCursor
            )

            with connection.cursor() as cursor:
                sql = "SELECT * FROM `id-pw` WHERE ID=%s AND PW=%s"
                cursor.execute(sql, (username, password))
                result = cursor.fetchone()

            connection.close()

            if result:
                self.username = username  # 将用户名存储为实例属性
                self.accept()  # 直接接受对话框
            else:
                self.waiting_widget.hide()
                for widget in [self.username_label, self.username_input, self.password_label, self.password_input, self.login_button, self.cancel_button]:
                    widget.show()
                QMessageBox.warning(self, "登录失败", "账号或密码错误,请重试。")
                self.username_input.clear()
                self.password_input.clear()
                self.username_input.setFocus()

        except pymysql.Error as e:
            logging.error(f"数据库连接错误: {str(e)}")
            self.waiting_widget.hide()
            for widget in [self.username_label, self.username_input, self.password_label, self.password_input, self.login_button, self.cancel_button]:
                widget.show()
            QMessageBox.critical(self, "数据库连接错误", f"连接数据库时发生错误,请联系管理员。")
            self.username_input.clear()
            self.password_input.clear()
            self.username_input.setFocus()

    # 衔接登录窗口打开主窗口并延迟关闭登录对话框
    '''def open_main_window(self):
        self.main_window = MainWindow(self.username)
        self.main_window.show()
        QTimer.singleShot(500, self.accept)  # 延迟关闭登录对话框'''

class DataLoadingThread(QThread):
    finished = pyqtSignal()
    progress = pyqtSignal(int)

    def __init__(self, main_window):
        super().__init__()
        self.main_window = main_window

    def run(self):
        try:
            # 1. 连接数据库
            self.progress.emit(10)
            self.main_window.connect_to_database()
            
            # 2. 加载公牛库数据
            self.progress.emit(30)
            self.load_bull_library_data()
            
            # 3. 加载公牛号转换规则
            self.progress.emit(60)
            self.load_bull_rules()
            
            # 4. 初始化其他数据结构（如果有的话）
            self.progress.emit(80)
            self.initialize_other_data()

            self.progress.emit(100)
            self.finished.emit()
        except Exception as e:
            logging.error(f"数据加载过程中发生错误: {str(e)}")
            # 这里可以发射一个错误信号，让主窗口知道加载失败
            self.finished.emit()  # 即使出错也要发射完成信号，以便关闭加载窗口

    def load_bull_library_data(self):
        connection = pymysql.connect(
            host=cipher_suite.decrypt(encoded_host).decode(),
            port=int(cipher_suite.decrypt(encoded_port).decode()),
            user=cipher_suite.decrypt(encoded_user).decode(),
            password=cipher_suite.decrypt(encoded_password).decode(),
            database=cipher_suite.decrypt(encoded_db).decode(),
            charset='utf8mb4',
            cursorclass=pymysql.cursors.DictCursor
        )
        with connection.cursor() as cursor:
            sql = "SELECT `BULL NAAB`, `BULL REG` FROM bull_library"
            cursor.execute(sql)
            result = cursor.fetchall()
            self.main_window.bull_library_df = pd.DataFrame(result)
        connection.close()
        self.main_window.bull_library_reg_dict = {str(row['BULL REG']): str(row['BULL NAAB']) for _, row in
                                                  self.main_window.bull_library_df.iterrows() if pd.notna(row['BULL REG'])}
        self.main_window.bull_library_naab_dict = {str(row['BULL NAAB']): str(row['BULL NAAB']) for _, row in
                                                   self.main_window.bull_library_df.iterrows() if pd.notna(row['BULL NAAB'])}

    def load_bull_rules(self):
        connection = pymysql.connect(
            host=cipher_suite.decrypt(encoded_host).decode(),
            port=int(cipher_suite.decrypt(encoded_port).decode()),
            user=cipher_suite.decrypt(encoded_user).decode(),
            password=cipher_suite.decrypt(encoded_password).decode(),
            database=cipher_suite.decrypt(encoded_db).decode(),
            charset='utf8mb4',
            cursorclass=pymysql.cursors.DictCursor
        )
        with connection.cursor() as cursor:
            sql = "SELECT `NAAB-conventional`, `NAAB-sexed` FROM rules"
            cursor.execute(sql)
            result = cursor.fetchall()
            self.main_window.bull_rules = {row['NAAB-conventional']: row['NAAB-sexed'] for row in result}
        connection.close()

    def initialize_other_data(self):
        # 初始化其他可能的大型数据结构
        self.main_window.cow_df = pd.DataFrame()
        self.main_window.bull_df = pd.DataFrame()
        self.main_window.needed_bulls_naab = set()
        self.main_window.needed_bulls_reg = set()

class LoadingWidget(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowFlags(Qt.WindowType.Tool | Qt.WindowType.FramelessWindowHint)
        self.setAttribute(Qt.WidgetAttribute.WA_TranslucentBackground)
        self.setFixedSize(200, 50)
        # 添加一个尺寸一致的浅灰色背景
        self.background = QLabel(self)
        self.background.setStyleSheet("background-color: white; border-radius: 10px;")
        self.background.setGeometry(0, 0, 200, 50)

        # 保证背景在所有控件的底层
        self.background.lower()


        
        
        layout = QVBoxLayout()
        self.progress_bar = QProgressBar()
        self.label = QLabel("正在加载数据...")
        layout.addWidget(self.label)
        layout.addWidget(self.progress_bar)
        self.setLayout(layout)

    def set_progress(self, value):
        self.progress_bar.setValue(value)


from PyQt6.QtWidgets import QDialog, QVBoxLayout, QLabel, QProgressBar, QPushButton, QHBoxLayout
from PyQt6.QtCore import Qt, QTimer
import time



class TqdmProgressDialog(QDialog):
    def __init__(self, total, description="Processing", parent=None):
        super().__init__(parent)
        self.setWindowTitle("进度")
        self.setFixedSize(500, 200)
        
        # 使窗口总是置顶
        self.setWindowFlags(self.windowFlags() | Qt.WindowType.WindowStaysOnTopHint)
        
        layout = QVBoxLayout()
        
        self.description_label = QLabel(description)
        self.description_label.setStyleSheet("font-size: 16px; font-weight: bold;")
        layout.addWidget(self.description_label)
        
        self.progress_bar = QProgressBar()
        self.progress_bar.setRange(0, total)
        self.progress_bar.setStyleSheet("""
            QProgressBar {
                border: none;
                border-radius: 5px;
                text-align: right;
                background-color: #E0E0E0;
                height: 15px;
            }
            QProgressBar::chunk {
                background-color: #05B8CC;
                width: 20px;
            }
        """)
        layout.addWidget(self.progress_bar)
        
        self.info_label = QLabel()
        self.info_label.setStyleSheet("font-size: 14px;")
        layout.addWidget(self.info_label)
        
        self.cancel_button = QPushButton("取消")
        self.cancel_button.setStyleSheet("font-size: 14px; padding: 5px;")
        self.cancel_button.clicked.connect(self.cancel)
        layout.addWidget(self.cancel_button)
        
        self.setLayout(layout)
        
        self.total = total
        self.current = 0
        self.start_time = time.time()
        self.canceled = False
        
        self.timer = QTimer(self)
        self.timer.timeout.connect(self.update_info)
        self.timer.start(100)  # 每100毫秒更新一次
    
    def update_progress(self, n=1):
        self.current += n
        self.progress_bar.setValue(self.current)
        
    def update_info(self):
        if self.current == 0:
            return
        
        elapsed_time = time.time() - self.start_time
        rate = self.current / elapsed_time
        remaining = (self.total - self.current) / rate if rate > 0 else 0
        
        percentage = (self.current / self.total) * 100
        info = (f"{self.current}/{self.total} "
                f"[{percentage:.1f}%, "
                f"{elapsed_time:.1f}秒 已用, "
                f"{remaining:.1f}秒 剩余, "
                f"{rate:.2f}个/秒]")
        
        self.info_label.setText(info)
    
    def cancel(self):
        self.canceled = True
    
    def wasCanceled(self):
        return self.canceled
    
    def close(self):
        self.timer.stop()
        super().close()

def preprocess_cow_data(cow_df):
    # 替换表头中的中文列名为英文列名
    column_mapping = {
        "耳号": "cow_id",
        "品种": "breed",
        "性别": "sex",
        "父亲号": "sire",
        "外祖父": "mgs",
        "母亲号": "dam",
        "外曾外祖父": "mmgs",
        "胎次": "lac",
        "最近产犊日期": "calving_date",
        "牛只出生日期": "birth_date",
        "月龄": "age",

        "本胎次配次": "services_time",
        "本胎次奶厅高峰产量":"peak_milk",
        "305奶量":"milk_305",
        "泌乳天数":"DIM",
        "繁育状态": "repro_status",

    }
    cow_df.rename(columns=column_mapping, inplace=True)

    # 定义需要保留的列
    columns_to_keep = ["cow_id", "breed", "sex", "sire", "mgs", "dam", "mmgs", "lac", "calving_date", "birth_date", "age","services_time","DIM","peak_milk","milk_305","repro_status", "group", "是否在场"]

    # 添加不存在的列并填充为空值
    missing_columns = [col for col in columns_to_keep if col not in cow_df.columns]
    for col in missing_columns:
        cow_df[col] = ""

    # 调整列的顺序
    cow_df = cow_df[columns_to_keep]

    # 删除性别为“公”的记录
    cow_df = cow_df[cow_df['sex'] != '公']

    # 处理数值列，将无效值转换为空值
    numeric_columns = ['lac', 'age', 'services_time', 'DIM', 'peak_milk', 'milk_305']
    for column in numeric_columns:
        if column in cow_df.columns:
            # 转换为数值类型，无效值转为NaN
            cow_df[column] = pd.to_numeric(cow_df[column], errors='coerce')
            # 替换无限值为NaN
            cow_df[column] = cow_df[column].replace([np.inf, -np.inf], np.nan)
            # 将NaN转换为空字符串，保持与原代码一致
            cow_df[column] = cow_df[column].replace(np.nan, "")



    # 对cow_id, sire, mgs, dam, mmgs列进行空格清除
    columns_to_clean = ['cow_id', 'sire', 'mgs', 'dam', 'mmgs']
    for column in columns_to_clean:
        if column in cow_df.columns:
            cow_df.loc[:, column] = cow_df[column].astype(str).apply(lambda x: x.replace(' ', '').strip())

    # 对calving_date和birth_date列转换为日期格式
    date_columns = ['calving_date', 'birth_date']
    for column in date_columns:
        if column in cow_df.columns:
            cow_df[column] = pd.to_datetime(cow_df[column], errors='coerce')

    # 检查重复的cow_id
    duplicate_cows = cow_df[cow_df.duplicated(subset=['cow_id'], keep=False)]
    if not duplicate_cows.empty:
        duplicate_count = len(duplicate_cows['cow_id'].unique())
        msg = f"发现{duplicate_count}个重复的母牛号。将按以下优先级保留记录：\n1. 性别为母牛\n2. 在群状态\n3. 出生日期最近\n4. 胎次最小\n5. 随机选择"
        QMessageBox.warning(None, "警告", msg)
        
        # 定义一个函数来选择要保留的记录
        def select_record(group):
            # 1. 优先选择性别为母牛的记录
            females = group[group['sex'] == '母']
            if not females.empty:
                group = females
            
            # 2. 在性别相同的情况下，优先选择在群的记录
            in_herd = group[group['是否在场'] == '是']
            if not in_herd.empty:
                group = in_herd
            
            # 3. 如果还有多条记录，选择出生日期最近的
            if group['birth_date'].notna().any():
                return group.loc[group['birth_date'].idxmax()]
            
            # 4. 如果出生日期都缺失，选择胎次最小的
            elif group['lac'].notna().any():
                return group.loc[group['lac'].idxmin()]
            
            # 5. 如果以上条件都不满足，随机选择一条记录
            else:
                return group.sample(n=1).iloc[0]

        # 按cow_id分组，应用选择函数
        cow_df = cow_df.groupby('cow_id').apply(select_record).reset_index(drop=True)


    invalid_naab_numbers = set()

    # 对sire, mgs, mmgs列进行NAAB编号格式化
    naab_columns = ['sire', 'mgs', 'mmgs']
    for column in naab_columns:
        if column in cow_df.columns:
            cow_df[column] = cow_df[column].apply(lambda x: format_naab_number(x, invalid_naab_numbers))

    if invalid_naab_numbers:
        invalid_numbers_str = '\n'.join(invalid_naab_numbers)
        msg_box = CustomMessageBox(f"NAAB公牛号的牛号部分有误,\n以下牛号HO之后的数字超过5位:\n\n{invalid_numbers_str}\n\n（正确案例:123HO12345）\n\n点击确定或关闭弹窗,运算不会中断,将继续;\n\n如有更新,请重新上传")
        msg_box.exec()

    return cow_df

def preprocess_bull_data(bull_df):
    # 对bull_id列进行空格清除和修剪
    bull_df['bull_id'] = bull_df['bull_id'].astype(str).apply(lambda x: x.replace(' ', '').strip())

    # 对bull_id列进行NAAB编号格式化
    bull_df['bull_id'] = bull_df['bull_id'].apply(format_naab_number)

    return bull_df

def format_naab_number(naab_number, invalid_naab_numbers=None):
    if invalid_naab_numbers is None:
        invalid_naab_numbers = set()
    naab_number = str(naab_number).strip()

    if len(naab_number) < 15:
        # 识别"H",如果"H"后没有"O"则补1个"O"
        naab_number = re.sub(r'H(?!O)', 'HO', naab_number)

        # 删除从左起的以"0"开始的连续的"0"
        naab_number = naab_number.lstrip('0')

        # 对"HO"之前的数的位数进行判断和补零
        match = re.match(r'(\d+)HO', naab_number)
        if match:
            prefix = match.group(1)
            if len(prefix) < 3:
                naab_number = prefix.zfill(3) + naab_number[len(prefix):]
            elif len(prefix) > 3:
                raise ValueError(f"NAAB公牛号的站号有误,超过3位: {naab_number}")

        # 删除"HO"之后的以"0"开始的连续的"0"
        naab_number = re.sub(r'HO0+', 'HO', naab_number)

        # 对 "HO" 之后的数的位数进行判断
        match = re.search(r'HO(\d+)', naab_number)
        if match:
            suffix = match.group(1)
            if len(suffix) > 5:
                invalid_naab_numbers.add(naab_number)
            elif len(suffix) < 5:
                suffix = suffix.zfill(5)
                naab_number = re.sub(r'HO(\d+)', f'HO{suffix}', naab_number)

    return naab_number

class CustomMessageBox(QDialog):
    def __init__(self, message, parent=None):
        super().__init__(parent)
        self.setWindowTitle("警告")
        self.setMinimumSize(400, 300)  # 设置最小尺寸
        
        layout = QVBoxLayout(self)
        
        # 创建滚动区域
        scroll_area = QScrollArea(self)
        scroll_area.setWidgetResizable(True)
        scroll_area.setHorizontalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAlwaysOff)
        
        # 创建内容窗口部件
        content_widget = QWidget()
        content_layout = QVBoxLayout(content_widget)
        
        self.label = QLabel(message)
        self.label.setWordWrap(True)
        self.label.setTextInteractionFlags(
            Qt.TextInteractionFlag.TextSelectableByMouse | 
            Qt.TextInteractionFlag.TextSelectableByKeyboard
        )
        content_layout.addWidget(self.label)
        
        scroll_area.setWidget(content_widget)
        layout.addWidget(scroll_area)
        
        button = QPushButton("确定", self)
        button.clicked.connect(self.accept)
        layout.addWidget(button)

        self.setLayout(layout)


class MainWindow(QMainWindow):
    def __init__(self, username):
        super().__init__()
        self.username = username
        self.setWindowTitle("综合选配系统")
        self.setGeometry(100, 100, 800, 600)
        self.setWindowIcon(QIcon("002.ico"))

        # 初始化UI，但暂时禁用所有功能
        self.init_ui()
        self.disable_all_functions()

        # 创建加载窗口
        self.loading_widget = LoadingWidget(self)  # 将 self 作为父窗口传入
        self.update_loading_widget_position()

        # 开始加载数据
        self.loading_thread = DataLoadingThread(self)
        self.loading_thread.finished.connect(self.on_data_loaded)
        self.loading_thread.progress.connect(self.loading_widget.set_progress)
        self.loading_thread.start()

        # 显示加载窗口
        self.loading_widget.show()

    def init_ui(self):
        # 设置窗口背景色
        palette = QPalette()
        palette.setColor(QPalette.ColorRole.Window, QColor(240, 240, 240))
        self.setPalette(palette)

        # 创建主布局
        main_layout = QVBoxLayout()
        main_widget = QWidget()
        main_widget.setLayout(main_layout)
        self.setCentralWidget(main_widget)

        # 添加软件标题图片
        image_label = self.create_title_image()
        main_layout.addWidget(image_label)

        # 创建选项卡容器
        tab_widget = QTabWidget()
        tab_widget.setMinimumHeight(200)  # 设置最小高度以确保按钮显示完整
        
        # 创建并添加各个选项卡
        tab_widget.addTab(self.create_data_upload_tab(), "基础数据上传")
        tab_widget.addTab(self.create_ranking_tab(), "牛群排名")
        tab_widget.addTab(self.create_breeding_analysis_tab(), "配种记录分析")
        tab_widget.addTab(self.create_individual_matching_tab(), "个体选配")
        tab_widget.addTab(self.create_conformation_tab(), "体型外貌评定")
        tab_widget.addTab(self.create_automation_tab(), "自动化生成")
        
        main_layout.addWidget(tab_widget)

        # 创建状态标签和滚动区域
        scroll_area = self.create_status_scroll_area()
        main_layout.addWidget(scroll_area)

        # 添加管理公牛号转换规则按钮
        manage_rules_button = QPushButton("管理公牛号转换规则")
        manage_rules_button.clicked.connect(self.manage_bull_rules)
        main_layout.addWidget(manage_rules_button)

        # 根据用户名控制管理公牛号转换规则按钮的可见性
        if self.username != "jaybzbo":
            manage_rules_button.setVisible(False)

    def create_title_image(self):
        image_label = QLabel()
        if getattr(sys, 'frozen', False):
            base_path = sys._MEIPASS
        else:
            base_path = os.path.abspath(".")

        image_path = os.path.join(base_path, "软件名称.jpg")
        pixmap = QPixmap(image_path)
        if pixmap.isNull():
            """print("Failed to load image.")"""
            image_label.setText("图片加载失败")
        else:
            scaled_pixmap = pixmap.scaled(780, pixmap.height(), aspectRatioMode=Qt.AspectRatioMode.KeepAspectRatio, transformMode=Qt.TransformationMode.SmoothTransformation)
            image_label.setPixmap(scaled_pixmap)
            image_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        return image_label

    def create_data_upload_tab(self):
        tab = QWidget()
        layout = QVBoxLayout()

        # 创建按钮组并设置样式
        self.select_output_button = QPushButton("1. 选择输出文件夹")
        self.select_output_button.setMinimumHeight(50)
        self.select_output_button.clicked.connect(self.select_output_folder)
        layout.addWidget(self.select_output_button)

        self.upload_cow_button = QPushButton("2. 上传母牛信息\n可上传伊起牛-牛群信息\n（最好包括所有离群牛）")
        self.upload_cow_button.setMinimumHeight(50)
        self.upload_cow_button.clicked.connect(self.upload_cow_file)
        layout.addWidget(self.upload_cow_button)

        self.upload_bull_button = QPushButton("3. 上传选配公牛信息")
        self.upload_bull_button.setMinimumHeight(50)
        self.upload_bull_button.clicked.connect(self.upload_bull_file)
        layout.addWidget(self.upload_bull_button)

        self.upload_breeding_record_button = QPushButton("4. 上传配种记录")
        self.upload_breeding_record_button.setMinimumHeight(50)
        self.upload_breeding_record_button.clicked.connect(self.upload_breeding_record)
        layout.addWidget(self.upload_breeding_record_button)

        layout.addStretch()
        tab.setLayout(layout)
        return tab

    def create_ranking_tab(self):
        tab = QWidget()
        layout = QVBoxLayout()

        self.generate_ranking_button = QPushButton("在群母牛指数计算及排名")
        self.generate_ranking_button.setMinimumHeight(50)
        self.generate_ranking_button.clicked.connect(self.run_index_ranking)
        layout.addWidget(self.generate_ranking_button)

        self.calculate_key_traits_button = QPushButton("在群母牛关键性状计算")
        self.calculate_key_traits_button.setMinimumHeight(50)
        self.calculate_key_traits_button.clicked.connect(self.calculate_key_traits_index)
        layout.addWidget(self.calculate_key_traits_button)

        layout.addStretch()
        tab.setLayout(layout)
        return tab

    def create_breeding_analysis_tab(self):
        tab = QWidget()
        layout = QVBoxLayout()

        self.genetic_defect_screening_button = QPushButton("隐性基因筛查\n（将打开另一应用程序需要重新登陆, 并上传-processed_breeding_data.xlsx)")
        self.genetic_defect_screening_button.setMinimumHeight(50)
        self.genetic_defect_screening_button.clicked.connect(self.genetic_defect_screening)
        layout.addWidget(self.genetic_defect_screening_button)

        self.semen_usage_analysis_button = QPushButton("使用冻精分析")
        self.semen_usage_analysis_button.setMinimumHeight(50)
        self.semen_usage_analysis_button.clicked.connect(self.semen_usage_analysis)
        layout.addWidget(self.semen_usage_analysis_button)

        layout.addStretch()
        tab.setLayout(layout)
        return tab

    def create_individual_matching_tab(self):
        tab = QWidget()
        layout = QVBoxLayout()

        self.generate_1st_report_button = QPushButton("第1步（隐性基因、近交计算）")
        self.generate_1st_report_button.setMinimumHeight(50)
        self.generate_1st_report_button.clicked.connect(self.start_selection_process)
        layout.addWidget(self.generate_1st_report_button)

        self.generate_2nd_report_button = QPushButton("第2步（冻精分配系统）")
        self.generate_2nd_report_button.setMinimumHeight(50)
        self.generate_2nd_report_button.clicked.connect(self.select_group_for_breeding)
        layout.addWidget(self.generate_2nd_report_button)

        layout.addStretch()
        tab.setLayout(layout)
        return tab

    def create_conformation_tab(self):
        tab = QWidget()
        layout = QVBoxLayout()

        self.upload_conformation_button = QPushButton("体型外貌数据上传")
        self.upload_conformation_button.setMinimumHeight(50)
        self.upload_conformation_button.clicked.connect(self.upload_conformation_data)
        layout.addWidget(self.upload_conformation_button)

        self.process_conformation_button = QPushButton("体型外貌数据处理")
        self.process_conformation_button.setMinimumHeight(50)
        self.process_conformation_button.clicked.connect(self.process_conformation_data)
        layout.addWidget(self.process_conformation_button)

        layout.addStretch()
        tab.setLayout(layout)
        return tab


    def create_automation_tab(self):
        tab = QWidget()
        layout = QVBoxLayout()

        self.generate_ppt_button = QPushButton("生成PPT报告")
        self.generate_ppt_button.setMinimumHeight(50)
        self.generate_ppt_button.clicked.connect(self.generate_ppt_report)
        layout.addWidget(self.generate_ppt_button)

        layout.addStretch()
        tab.setLayout(layout)
        return tab




    def create_status_scroll_area(self):
        scroll_area = QScrollArea()
        scroll_area.setWidgetResizable(True)

        scroll_content = QWidget()
        scroll_layout = QVBoxLayout(scroll_content)

        self.status_label = QLabel("请上传在群母牛信息和备选公牛信息,然后选择输出文件夹。")
        self.status_label.setStyleSheet("background-color: white; color: black;")
        self.status_label.setFixedSize(780, 200)  # 设置固定大小
        self.status_label.setWordWrap(True)  # 自动换行
        scroll_layout.addWidget(self.status_label)

        scroll_area.setWidget(scroll_content)
        return scroll_area

    def disable_all_functions(self):
        # 禁用所有按钮
        for child in self.findChildren(QPushButton):
            child.setEnabled(False)
        
        # 启用指定的按钮
        self.select_output_button.setEnabled(True)
        self.upload_cow_button.setEnabled(True)
        self.upload_breeding_record_button.setEnabled(True)
        self.upload_bull_button.setEnabled(True)

    def enable_all_functions(self):
        # 启用所有按钮和功能
        for child in self.findChildren(QPushButton):
            child.setEnabled(True)

    def on_data_loaded(self):
        # 数据加载完成后的操作
        self.loading_widget.hide()
        self.enable_all_functions()  # 启用所有功能
        QMessageBox.information(self, "加载完成", "基础数据加载完成，您现在可以进行所有操作了。")

    '''def position_loading_widget(self):
        # 将加载窗口放置在主窗口的右上角
        geo = self.geometry()
        self.loading_widget.move(geo.right() - self.loading_widget.width() - 10, geo.top() + 10)'''

    def update_loading_widget_position(self):
        geo = self.geometry()
        self.loading_widget.move(geo.right() - self.loading_widget.width() - 10, geo.top() + 10)

    def moveEvent(self, event):
        super().moveEvent(event)
        self.update_loading_widget_position()

    def standardize_cow_id(self, cow_id):
        cow_id = str(cow_id)
        cow_id = cow_id.strip()
        cow_id = cow_id.split('.')[0]
        if cow_id.isdigit():
            cow_id = cow_id.zfill(10)
        return cow_id

    def generate_needed_bulls(self, progress_dialog):
        self.needed_bulls_naab = set()
        self.needed_bulls_reg = set()

        total_rows = len(self.bull_df) + len(self.cow_df) * 3  # 总行数
        processed_rows = 0  # 已处理的行数

        for bull_id in self.bull_df['bull_id']:
            if len(bull_id) > 11:
                self.needed_bulls_reg.add(bull_id)
            else:
                self.needed_bulls_naab.add(bull_id)
            processed_rows += 1
            progress = int(processed_rows / total_rows * 100)
            progress_dialog.setValue(progress)
            if progress_dialog.wasCanceled():
                break

        for _, row in self.cow_df.iterrows():
            for ancestor_type in ['sire', 'mgs', 'mmgs']:
                if ancestor_type in row:  # 检查列是否存在
                    ancestor_id = str(row[ancestor_type])
                    if len(ancestor_id) > 11:
                        self.needed_bulls_reg.add(ancestor_id)
                    else:
                        self.needed_bulls_naab.add(ancestor_id)
                processed_rows += 1
                progress = int(processed_rows / total_rows * 100)
                progress_dialog.setValue(progress)
                if progress_dialog.wasCanceled():
                    break

        self.needed_bulls_naab.discard('nan')  # 去除空值
        self.needed_bulls_reg.discard('nan')  # 去除空值
        self.needed_bulls_naab.add('999HO99999')  # 添加默认公牛ID

        progress_dialog.setValue(100)
        progress_dialog.close()

    def manage_bull_rules(self):
        dialog = RuleManagerDialog(self.username)
        dialog.exec()

    def initialize_login(self):
        login_dialog = LoginDialog(parent=self)
        if login_dialog.exec() == QDialog.DialogCode.Accepted:
            self.username = login_dialog.username
        else:
            sys.exit()

    def connect_to_database(self):
        try:
            connection = pymysql.connect(
                host=cipher_suite.decrypt(encoded_host).decode(),
                port=int(cipher_suite.decrypt(encoded_port).decode()),
                user=cipher_suite.decrypt(encoded_user).decode(),
                password=cipher_suite.decrypt(encoded_password).decode(),
                database=cipher_suite.decrypt(encoded_db).decode(),
                charset='utf8mb4',
                cursorclass=pymysql.cursors.DictCursor
            )
            with connection.cursor() as cursor:
                sql = "SELECT `BULL NAAB`, `BULL REG` FROM bull_library"
                cursor.execute(sql)
                result = cursor.fetchall()
                bull_library_df = pd.DataFrame(result)
            connection.close()
            self.bull_library_reg_dict = {str(row['BULL REG']): str(row['BULL NAAB']) for _, row in
                                          bull_library_df.iterrows() if pd.notna(row['BULL REG'])}
            self.bull_library_naab_dict = {str(row['BULL NAAB']): str(row['BULL NAAB']) for _, row in
                                           bull_library_df.iterrows() if pd.notna(row['BULL NAAB'])}

            # 读取公牛号转换规则
            connection = pymysql.connect(
                host=cipher_suite.decrypt(encoded_host).decode(),
                port=int(cipher_suite.decrypt(encoded_port).decode()),
                user=cipher_suite.decrypt(encoded_user).decode(),
                password=cipher_suite.decrypt(encoded_password).decode(),
                database=cipher_suite.decrypt(encoded_db).decode(),
                charset='utf8mb4',
                cursorclass=pymysql.cursors.DictCursor
            )
            with connection.cursor() as cursor:
                sql = "SELECT `NAAB-conventional`, `NAAB-sexed` FROM rules"
                cursor.execute(sql)
                result = cursor.fetchall()
                self.bull_rules = {row['NAAB-conventional']: row['NAAB-sexed'] for row in result}
            connection.close()

        except Exception as e:
            logging.error(f"连接数据库时发生错误: {e}")
            QMessageBox.critical(self, "数据库连接错误", f"连接数据库时发生错误：{e},请联系管理员。")
            sys.exit()

    def on_database_result_ready(self, bull_library_df):
        self.bull_library_reg_dict = {str(row['BULL REG']): str(row['BULL NAAB']) for _, row in
                                      bull_library_df.iterrows() if pd.notna(row['BULL REG'])}
        self.bull_library_naab_dict = {str(row['BULL NAAB']): str(row['BULL NAAB']) for _, row in
                                       bull_library_df.iterrows() if pd.notna(row['BULL NAAB'])}

        login_dialog = LoginDialog(parent=self)
        if login_dialog.exec() == QDialog.DialogCode.Accepted:
            self.username = login_dialog.username
        else:
            sys.exit()
        """print(bull_library_df.columns)"""      

    def login(self):
        login_dialog = LoginDialog(parent=self)
        if login_dialog.exec() == QDialog.DialogCode.Accepted:
            self.username = login_dialog.username
        else:
            sys.exit()

    def upload_cow_file(self):
        if not self.output_folder:
            QMessageBox.warning(self, "警告", "请先选择输出文件夹")
            return

        file_path, _ = QFileDialog.getOpenFileName(self, "上传在群母牛信息", "", "Excel Files (*.xlsx)")
        if file_path:
            try:
                original_cow_df = pd.read_excel(file_path)
                logging.info(f"原始母牛数据行数: {len(original_cow_df)}")
                
                self.cow_df = preprocess_cow_data(original_cow_df)  # 使用新的预处理函数
                logging.info(f"预处理后母牛数据行数: {len(self.cow_df)}")

                
                # 检查是否包含“是否在场”列以及是否有“是”的记录，不包括表头
                if '是否在场' not in self.cow_df.columns:
                    QMessageBox.warning(self, "警告", "上传牛只信息中缺少“是否在场”列，将无法生成“牛群关键性状年度进展图”，但不会影响牛只上传。")
                elif not any(self.cow_df['是否在场'] == '是'):  # 检查实际数据部分
                    QMessageBox.warning(self, "警告", "上传牛只信息中无“是否在场”信息，将无法生成“牛群关键性状年度进展图”，但不会影响牛只上传。")
                    

                # 继续执行后续逻辑，无论是否有“是否在场”列，数据都会上传
                self.cow_df.fillna('', inplace=True)  # 将缺失值填充为空字符串
                cow_output_path = os.path.join(self.output_folder, "processed_cow_data.xlsx")
                self.cow_df.to_excel(cow_output_path, index=False)  # 保存为Excel文件

                QMessageBox.information(self, "上传成功", f"在群母牛信息上传成功！\n原始数据行数: {len(original_cow_df)}\n处理后数据行数: {len(self.cow_df)}")
            except Exception as e:
                logging.error(f"读取在群母牛信息文件时发生错误: {e}")
                QMessageBox.critical(self, "错误", f"读取在群母牛信息文件时发生错误: {e},请检查文件格式。")

    def upload_bull_file(self):
        if not self.output_folder:

            QMessageBox.warning(self, "警告", "请先选择输出文件夹")
            return

        file_path, _ = QFileDialog.getOpenFileName(self, "上传备选公牛信息", "", "Excel Files (*.xlsx)")
        if file_path:
            try:
                self.bull_df = pd.read_excel(file_path)
                self.bull_df = preprocess_bull_data(self.bull_df)  # 对上传的公牛信息进行预处理
                bull_output_path = os.path.join(self.output_folder, 'processed_bull_data.xlsx')
                self.bull_df.to_excel(bull_output_path, index=False)  # 保存为Excel文件

                # 自动再上传一次生成的公牛信息文件
                self.bull_df = pd.read_excel(bull_output_path)
                
                
                QMessageBox.information(self, "上传成功", "备选公牛信息上传成功！")
                """print("Uploaded bull file:")
                print(self.bull_df.head())"""
            except Exception as e:
                logging.error(f"读取备选公牛信息文件时发生错误: {e}")
                QMessageBox.critical(self, "错误", f"读取备选公牛信息文件时发生错误,请检查文件格式。")

    def select_output_folder(self):
        folder_path = QFileDialog.getExistingDirectory(self, "选择输出文件夹")
        if folder_path:
            self.output_folder = folder_path
            QMessageBox.information(self, "选择成功", f"已选择输出文件夹: {self.output_folder}")

    def start_selection_process(self):
        if self.cow_df.empty or self.bull_df.empty or not self.output_folder:
            QMessageBox.warning(self, "错误", "请先上传在群母牛信息、备选公牛信息,并选择输出文件夹。")
            return

        # 创建进度条
        progress_dialog = QProgressDialog("正在准备个体选配数据，\n\n数据库连接中......", "取消", 0, 100, self)
        progress_dialog.setWindowModality(Qt.WindowModality.WindowModal)
        progress_dialog.setValue(0)
        progress_dialog.show()
        progress_dialog.raise_()
        QApplication.processEvents()
        
        # 检查是否有缺失公牛
        missing_bulls = self.check_missing_bulls(self.bull_df)
        if missing_bulls:
            self.update_missing_bulls(missing_bulls, "综合选配流程开始时发现缺失公牛")
        progress_dialog.setValue(22)
        QApplication.processEvents()
        progress_dialog.show()
        progress_dialog.raise_()
        QApplication.processEvents()
        

        # 指数排名
        self.run_index_ranking()
        progress_dialog.setValue(41)
        QApplication.processEvents()
        
           
        # 隐性基因筛查
        self.run_genetic_defect_screening()
        progress_dialog.setValue(56)
        QApplication.processEvents()  # 确保UI更新

        # 生成所需公牛列表
        self.generate_needed_bulls(progress_dialog)
        progress_dialog.setValue(78)
        

        """print("Generated needed bulls:")
        print("Needed NAAB bulls:", self.needed_bulls_naab)
        print("Needed REG bulls:", self.needed_bulls_reg)"""

        # 近交系数计算
        self.run_inbreeding_coefficient_calculation()
        progress_dialog.setValue(88)
        QApplication.processEvents()
        

        # 生成选配文件
        self.generate_selection_files()
        progress_dialog.setValue(100)
        

        QMessageBox.information(self, "完成", "综合选配流程已完成,请查看输出文件夹中的结果文件。")
        progress_dialog.close()

    def run_genetic_defect_screening(self):
        try:
            # 创建进度条
            progress_dialog = QProgressDialog("正在进行隐性基因筛查...", "取消", 0, 100, self)
            progress_dialog.setWindowModality(Qt.WindowModality.WindowModal)
            progress_dialog.setValue(0)

            # 检查是否有缺失公牛
            missing_bulls = self.check_missing_bulls(self.bull_df)
            if missing_bulls:
                self.update_missing_bulls(missing_bulls, "隐性基因筛查")
            progress_dialog.setValue(5)
            QApplication.processEvents()

            # 定义隐性基因列表
            defect_genes = ["MW", "HH1", "HH2", "HH3", "HH4", "HH5", "HH6", "BLAD", "Chondrodysplasia", "Citrullinemia",
                            "DUMPS", "Factor XI", "CVM", "Brachyspina", "Mulefoot", "Cholesterol deficiency"]

            self.status_label.setText(self.status_label.text() + "\n正在读取公牛信息...")
            # 读取公牛信息
            connection = pymysql.connect(
                host=cipher_suite.decrypt(encoded_host).decode(),
                port=int(cipher_suite.decrypt(encoded_port).decode()),
                user=cipher_suite.decrypt(encoded_user).decode(),
                password=cipher_suite.decrypt(encoded_password).decode(),
                database=cipher_suite.decrypt(encoded_db).decode(),
                charset='utf8mb4',
                cursorclass=pymysql.cursors.DictCursor
            )

            with connection.cursor() as cursor:
                sql = "SELECT `BULL NAAB`, `BULL REG`, `MW`, `HH1`, `HH2`, `HH3`, `HH4`, `HH5`, `HH6`, `BLAD`, `Chondrodysplasia`, `Citrullinemia`, `DUMPS`, `Factor XI`, `CVM`, `Brachyspina`, `Mulefoot`, `Cholesterol deficiency` FROM bull_library"
                cursor.execute(sql)
                result = cursor.fetchall()
                bull_library_df = pd.DataFrame(result)

            connection.close()
            progress_dialog.setValue(15)
            QApplication.processEvents()

            # 创建公牛信息字典
            bull_library_naab_dict = {}
            bull_library_reg_dict = {}
            for _, row in bull_library_df.iterrows():
                if pd.notna(row['BULL NAAB']):
                    bull_library_naab_dict[str(row['BULL NAAB'])] = {
                        'reg': str(row['BULL REG']) if pd.notna(row['BULL REG']) else None,
                        **{gene: str(row[gene]) for gene in defect_genes if pd.notna(row[gene])}
                    }
                if pd.notna(row['BULL REG']):
                    bull_library_reg_dict[str(row['BULL REG'])] = {
                        'naab': str(row['BULL NAAB']) if pd.notna(row['BULL NAAB']) else None,
                        **{gene: str(row[gene]) for gene in defect_genes if pd.notna(row[gene])}
                    }

            # 执行隐性基因筛查
            pair_info_list = []
            total_pairs = len(self.cow_df) * len(self.bull_df)
            processed_pairs = 0

            for _, cow_row in self.cow_df.iterrows():
                cow_id = str(cow_row['cow_id'])
                sire_id = str(cow_row.get('sire', ''))
                cow_genes = {gene: 'missing data' for gene in defect_genes}
                if len(sire_id) > 11:
                    sire_in_library = sire_id in bull_library_reg_dict
                else:
                    sire_in_library = sire_id in bull_library_naab_dict
                    if not sire_in_library:
                        for rule, replacement in self.bull_rules.items():
                            if sire_id.startswith(rule):
                                new_sire_id = replacement + sire_id[3:]
                                sire_in_library = new_sire_id in bull_library_naab_dict
                                if sire_in_library:
                                    sire_id = new_sire_id
                                    break
                            elif sire_id.startswith(replacement):
                                new_sire_id = rule + sire_id[3:]
                                sire_in_library = new_sire_id in bull_library_naab_dict
                                if sire_in_library:
                                    sire_id = new_sire_id
                                    break
                if sire_in_library:
                    if len(sire_id) > 11:
                        sire_data = bull_library_reg_dict[sire_id]
                    else:
                        sire_data = bull_library_naab_dict[sire_id]
                    for gene in defect_genes:
                        gene_data = sire_data.get(gene, 'missing data')
                        cow_genes[gene] = 'C' if gene_data == 'C' else 'Safe'
                for _, bull_row in self.bull_df.iterrows():
                    bull_id = str(bull_row['bull_id'])
                    if len(bull_id) > 11:
                        bull_in_library = bull_id in bull_library_reg_dict
                    else:
                        bull_in_library = bull_id in bull_library_naab_dict
                        if not bull_in_library:
                            for rule, replacement in self.bull_rules.items():
                                if bull_id.startswith(rule):
                                    new_bull_id = replacement + bull_id[3:]
                                    bull_in_library = new_bull_id in bull_library_naab_dict
                                    if bull_in_library:
                                        bull_id = new_bull_id
                                        break
                                elif bull_id.startswith(replacement):
                                    new_bull_id = rule + bull_id[3:]
                                    bull_in_library = new_bull_id in bull_library_naab_dict
                                    if bull_in_library:
                                        bull_id = new_bull_id
                                        break
                    bull_genes = {gene: 'missing data' for gene in defect_genes}
                    if bull_in_library:
                        if len(bull_id) > 11:
                            bull_data = bull_library_reg_dict[bull_id]
                        else:
                            bull_data = bull_library_naab_dict[bull_id]
                        for gene in defect_genes:
                            gene_data = bull_data.get(gene, 'missing data')
                            bull_genes[gene] = 'C' if gene_data == 'C' else 'Safe'
                    pair_info = {'Cow ID': cow_id, 'Bull ID': bull_id}
                    for gene in defect_genes:
                        cow_gene = cow_genes[gene]
                        bull_gene = bull_genes[gene]
                        if cow_gene == 'C' and bull_gene == 'C':
                            pair_info[gene] = 'NO safe'
                        elif cow_gene == 'missing data' and bull_gene == 'missing data':
                            pair_info[gene] = 'missing data'
                        elif cow_gene == 'missing data':
                            pair_info[gene] = 'missing cow data'
                        elif bull_gene == 'missing data':
                            pair_info[gene] = 'missing bull data'
                        else:
                            pair_info[gene] = 'safe'
                    pair_info_list.append(pair_info)

                    # 更新进度条
                    processed_pairs += 1
                    progress = int(processed_pairs / total_pairs * 100)
                    progress_dialog.setValue(progress_dialog.value() + int(70 * processed_pairs / total_pairs))
                    QApplication.processEvents()
                    if progress_dialog.wasCanceled():
                        return

            screening_result_df = pd.DataFrame(pair_info_list)
            progress_dialog.setValue(90)
            QApplication.processEvents()

            # 保存隐性基因筛查结果到文件
            screening_result_file = os.path.join(self.output_folder, "隐性基因筛查结果.xlsx").replace('\\', '/')
            screening_result_df.to_excel(screening_result_file, index=False)
            self.status_label.setText(self.status_label.text() + "\n隐性基因筛查完成。")
            progress_dialog.setValue(100)
            progress_dialog.close()

        except Exception as e:
            logging.error(f"隐性基因筛查时发生错误: {e}")
            QMessageBox.critical(self, "错误", f"隐性基因筛查时发生错误：{e},请按照模版上传母牛信息及备选公牛表。")

    def run_inbreeding_coefficient_calculation(self):
        try:
            # 创建进度条
            progress_dialog = QProgressDialog("正在进行近交系数计算...", "取消", 0, 100, self)
            progress_dialog.setWindowModality(Qt.WindowModality.WindowModal)
            progress_dialog.setValue(0)

            # 检查是否有缺失公牛
            missing_bulls = self.check_missing_bulls(self.bull_df)
            progress_dialog.setValue(13)
            QApplication.processEvents()
            if missing_bulls:
                self.update_missing_bulls(missing_bulls, "近交系数计算")
            progress_dialog.setValue(36)
            QApplication.processEvents()

            # 读取公牛信息
            connection = pymysql.connect(
                host=cipher_suite.decrypt(encoded_host).decode(),
                port=int(cipher_suite.decrypt(encoded_port).decode()),
                user=cipher_suite.decrypt(encoded_user).decode(),
                password=cipher_suite.decrypt(encoded_password).decode(),
                database=cipher_suite.decrypt(encoded_db).decode(),
                charset='utf8mb4',
                cursorclass=pymysql.cursors.DictCursor
            )

            with connection.cursor() as cursor:
                sql = "SELECT `BULL NAAB`, `BULL REG`, `SIRE REG`, `MGS REG` FROM bull_library"
                cursor.execute(sql)
                result = cursor.fetchall()
                bull_library_df = pd.DataFrame(result)

            connection.close()
            progress_dialog.setValue(53)
            QApplication.processEvents()

            # 创建公牛信息字典
            bull_dict = {}
            for _, row in bull_library_df.iterrows():
                bull_naab = str(row['BULL NAAB']) if pd.notna(row['BULL NAAB']) else None
                bull_reg = str(row['BULL REG'])
                bull_dict[bull_naab] = {
                    'reg': bull_reg,
                    'sire': str(row['SIRE REG']) if pd.notna(row['SIRE REG']) else None,
                    'mgs': str(row['MGS REG']) if pd.notna(row['MGS REG']) else None,
                }
                bull_dict[bull_reg] = bull_dict[bull_naab]

            # 计算近交系数
            results = []
            total_pairs = len(self.bull_df) * len(self.cow_df)
            processed_pairs = 0

            for _, bull_row in self.bull_df.iterrows():
                bull_id = bull_row['bull_id']
                bull_ancestors = self.get_ancestors(bull_id, bull_dict)
                bull_ancestors['bull_id'] = bull_id

                for _, cow_row in self.cow_df.iterrows():
                    cow_id = cow_row['cow_id']
                    cow_ancestors = self.get_ancestors(cow_row['sire'], bull_dict)
                    fx = self.calculate_inbreeding_coefficient(cow_ancestors, bull_ancestors)
                    fx_text = "{:.2%}".format(fx)
                    results.append({'Cow': cow_id, 'Bull': bull_id, 'Inbreeding Coefficient': fx_text})

                    # 更新进度条
                    processed_pairs += 1
                    progress = int(processed_pairs / total_pairs * 100)
                    progress_dialog.setValue(progress_dialog.value() + int(70 * processed_pairs / total_pairs))
                    QApplication.processEvents()
                    if progress_dialog.wasCanceled():
                        return

            inbreeding_result_df = pd.DataFrame(results)
            progress_dialog.setValue(90)
            QApplication.processEvents()

            # 保存近交系数计算结果到文件
            inbreeding_result_file = os.path.join(self.output_folder, "近交系数计算结果.xlsx").replace('\\', '/')
            inbreeding_result_df.to_excel(inbreeding_result_file, index=False)
            self.status_label.setText(self.status_label.text() + "\n近交系数计算完成。")
            progress_dialog.setValue(100)
            progress_dialog.close()

        except Exception as e:
            logging.error(f"近交系数计算时发生错误: {e}")
            QMessageBox.critical(self, "错误", f"近交系数计算时发生错误：{e},请按照模版上传母牛信息及备选公牛表。")

    def get_ancestors(self, bull_id, bull_dict):
        ancestors = {'sire': None, 'mgs': None, 'mmgs': None}
        bull_info = self.find_bull_in_library(bull_id, bull_dict)
        if bull_info is not None:
            ancestors['sire'] = bull_info['sire']
            ancestors['mgs'] = bull_info['mgs']
            if ancestors['sire'] is not None:
                sire_info = self.find_bull_in_library(ancestors['sire'], bull_dict)
                if sire_info is not None:
                    ancestors['mmgs'] = sire_info['mgs']
        return ancestors

    def find_bull_in_library(self, bull_id, bull_dict):
        if pd.isna(bull_id):
            return None
        bull_id = str(bull_id)
        if bull_id in bull_dict:
            return bull_dict[bull_id]
        return None

    def calculate_inbreeding_coefficient(self, cow_ancestors, bull_ancestors):
        if cow_ancestors['sire'] == bull_ancestors['bull_id']:
            return 0.25
        if cow_ancestors['sire'] == bull_ancestors['sire']:
            return 0.125
        if cow_ancestors['sire'] in [bull_ancestors['mgs'], bull_ancestors['mmgs']]:
            return 0.0625
        if cow_ancestors['mgs'] == bull_ancestors['bull_id'] or cow_ancestors['mmgs'] == bull_ancestors['bull_id']:
            return 0.125
        if cow_ancestors['mgs'] == bull_ancestors['sire'] or cow_ancestors['mmgs'] == bull_ancestors['sire']:
            return 0.0625
        return 0.0

    def check_weight_df_columns(self, weight_df):
        weight_df = weight_df[['traits', 'weight', 'SD']].dropna()
        weight_df['weight'] = weight_df['weight'].fillna(0)
        weight_df['SD'] = weight_df['SD'].fillna(1)
        return weight_df
    
    def check_weight_df_columns_correct(self,weight_df):
        required_columns = {'traits', 'weight', 'SD'}
        if not required_columns.issubset(weight_df.columns):
            missing_columns = required_columns - set(weight_df.columns)
            return False, missing_columns
        return True, None

    # 根据育种指数计算牛群的育种指数得分（并包括生成的xlsx的调用）
    def run_index_ranking(self):
        try:
            progress_dialog = QProgressDialog("正在进行指数排名，\n\n\n   数据库连接中......（1~3 mins）", "取消", 0, 100, self)
            progress_dialog.setWindowModality(Qt.WindowModality.WindowModal)
            progress_dialog.setValue(0)
            progress_dialog.show()
            progress_dialog.raise_()
            QApplication.processEvents()



            # 如果没有上传备选公牛信息,则生成一个默认的空表
            if self.bull_df.empty:
                self.bull_df = pd.DataFrame(columns=['bull_id', 'semen_type'])        
            if self.cow_df.empty:
                QMessageBox.warning(self, "警告", "请先上传在群母牛信息。")
                return
            progress_dialog.setValue(5)
            progress_dialog.show()
            progress_dialog.raise_()
            QApplication.processEvents()

            # 检查日期格式
            date_columns = [ 'birth_date']
            for column in date_columns:
                if column in self.cow_df.columns:
                    if not pd.api.types.is_datetime64_any_dtype(self.cow_df[column]):
                        QMessageBox.warning(self, "警告", f"{column}列的格式不是日期格式,请修改为日期格式后重新上传。")
                        return
            # 生成所需公牛列表
            self.generate_needed_bulls(progress_dialog)

            """print("Generated needed bulls:")
            print("Needed NAAB bulls:", self.needed_bulls_naab)
            print("Needed REG bulls:", self.needed_bulls_reg)"""
            progress_dialog.setValue(7)
            progress_dialog.show()
            progress_dialog.raise_()
            QApplication.processEvents()

            
            # 检查是否有缺失公牛
            missing_bulls = self.check_missing_bulls(self.bull_df)
            if missing_bulls:
                self.update_missing_bulls(missing_bulls, "指数排名")
            progress_dialog.setValue(10)
            progress_dialog.show()
            progress_dialog.raise_()
            QApplication.processEvents()
            while True:
                # 创建并设置消息框
                msg_box = QMessageBox(QMessageBox.Icon.Information, "提示", "请上传指数权重表。")
                msg_box.setWindowFlags(msg_box.windowFlags() | Qt.WindowType.WindowStaysOnTopHint)
                msg_box.exec()

                # 上传指数权重表
                weight_file_path, _ = QFileDialog.getOpenFileName(self, "上传指数权重表", "", "Excel Files (*.xlsx)")
                if not weight_file_path:
                    QMessageBox.warning(self, "错误", "未上传指数权重表,无法进行指数排名。")
                    progress_dialog.close()  # 关闭进度条
                    return
                weight_df = pd.read_excel(weight_file_path)
                valid, missing_columns = self.check_weight_df_columns_correct(weight_df)
                if valid:
                    break
                else:
                    QMessageBox.warning(self, "警告", f"上传的指数权重表缺少以下列: {', '.join(missing_columns)}")
            progress_dialog.setValue(15)
            QApplication.processEvents()

            """print("weight_df:")
            print(weight_df.head())"""
            progress_dialog.setValue(20)
            QApplication.processEvents()


            # 读取公牛信息
            connection = pymysql.connect(
                host=cipher_suite.decrypt(encoded_host).decode(),
                port=int(cipher_suite.decrypt(encoded_port).decode()),
                user=cipher_suite.decrypt(encoded_user).decode(),
                password=cipher_suite.decrypt(encoded_password).decode(),
                database=cipher_suite.decrypt(encoded_db).decode(),
                charset='utf8mb4',
                cursorclass=pymysql.cursors.DictCursor
            )
            with connection.cursor() as cursor:
                sql = "SELECT * FROM bull_library"
                cursor.execute(sql)
                result = cursor.fetchall()
                bull_library_df = pd.DataFrame(result)
            connection.close()

            """print("bull_library_df:")
            print(bull_library_df.head())"""
            progress_dialog.setValue(30)
            QApplication.processEvents()

            # 计算公牛指数
            bull_index_scores = self.calculate_bull_index_scores(bull_library_df, weight_df, progress_dialog)

            """print("bull_index_scores:")
            print(bull_index_scores.head())"""
            progress_dialog.setValue(40)
            QApplication.processEvents()

            # 保存公牛指数到文件
            bull_index_file = os.path.join(self.output_folder, "公牛指数.xlsx")
            bull_index_scores.to_excel(bull_index_file, index=False)
            self.status_label.setText(self.status_label.text() + "\n公牛指数计算完成。")
            progress_dialog.setValue(50)
            QApplication.processEvents()

            # 更新字典
            self.bull_scores_naab_dict = bull_index_scores.set_index('BULL NAAB')['Index Score'].to_dict()
            self.bull_scores_reg_dict = bull_index_scores.set_index('BULL REG')['Index Score'].to_dict()


            """print("self.bull_scores_naab_dict:")
            print(dict(list(self.bull_scores_naab_dict.items())[:5]))  # 打印前5个元素
            print("self.bull_scores_reg_dict:")
            print(dict(list(self.bull_scores_reg_dict.items())[:5]))  # 打印前5个元素"""
            progress_dialog.setValue(55)
            QApplication.processEvents()



            # 生成表a和表b
            table_a = self.generate_table_a(self.cow_df, bull_index_scores)
            """print("table_a:")
            print(table_a.head())"""


            table_b = self.generate_table_b(table_a)

            """print("table_b:")
            print(table_b)"""
            progress_dialog.setValue(60)
            QApplication.processEvents()


            """# 保存表a到文件
            table_a_file = os.path.join(self.output_folder, "表a.xlsx")
            table_a.to_excel(table_a_file, index=False)

            # 保存表b到文件
            table_b_file = os.path.join(self.output_folder, "表b.xlsx")
            table_b.to_excel(table_b_file, index=False)"""
            

            # 计算母牛指数
            default_bull_id = '999HO99999'
            default_bull_score = self.find_bull_index_score(default_bull_id, bull_index_scores, 0)
            sire_avg_scores = dict(zip(table_b['birth_year'], table_b['mean']))

            """print("default_bull_score:")
            print(default_bull_score)
            print("sire_avg_scores:")
            print(sire_avg_scores)"""
            progress_dialog.setValue(70)
            QApplication.processEvents()



            cow_index_scores = self.calculate_cow_index_scores(self.cow_df, bull_index_scores, sire_avg_scores,
                                                               default_bull_score)
            
            """print("cow_index_scores:")
            print(cow_index_scores.head())"""
            progress_dialog.setValue(80)
            QApplication.processEvents()
            

            # 保存母牛指数到文件
            cow_index_file = os.path.join(self.output_folder, "母牛指数.xlsx")
            cow_index_scores.to_excel(cow_index_file, index=False)
            self.status_label.setText(self.status_label.text() + "\n母牛指数计算完成。")

            # 计算母牛基因组指数
            geno_index_scores = pd.DataFrame()
            geno_library_file = os.path.join(self.output_folder, "geno_library.xlsx")
            if os.path.isfile(geno_library_file):
                geno_library_df = pd.read_excel(geno_library_file)
                geno_index_scores = self.calculate_geno_index_scores(geno_library_df, weight_df)
                # 保存基因组指数到文件
                geno_index_file = os.path.join(self.output_folder, "基因组指数.xlsx")
                geno_index_scores.to_excel(geno_index_file, index=False)
                self.status_label.setText(self.status_label.text() + "\n母牛基因组指数计算完成。")
            progress_dialog.setValue(95)
            QApplication.processEvents()


            # 计算综合指数排名
            ranking_result_df = self.calculate_combined_index_scores(self.cow_df, cow_index_scores)
            # 保存指数排名结果到文件
            ranking_result_file = os.path.join(self.output_folder, "结果-指数排名结果.xlsx")
            ranking_result_df.sort_values(by='Combine Index Score', ascending=False).to_excel(ranking_result_file,
                                                                                              index=False)
            self.status_label.setText(self.status_label.text() + "\n指数排名完成。")

            # 生成育种指数得分的年份正态分布图
            self.generate_normal_distribution_charts(ranking_result_df, 'Combine Index Score', self.output_folder, "育种指数得分")
            self.status_label.setText(self.status_label.text() + "\n育种指数得分的年份正态分布图完成。")

            # 生成育种指数得分的在群牛五等份分布表
            self.generate_quintile_distribution(ranking_result_df, 'Combine Index Score', self.output_folder, "育种指数得分")
            self.status_label.setText(self.status_label.text() + "\n育种指数得分的五等份分布表完成。")

            progress_dialog.setValue(100)
            QApplication.processEvents()


        except Exception as e:
            logging.error(f"指数排名时发生错误: {e}")
            QMessageBox.critical(self, "错误", f"指数排名时发生错误：{e},请按照模版上传母牛信息/权重分配表。")

    def calculate_bull_index_scores(self, bull_library_df, weight_df, progress_dialog):
        # 预处理 weight_df
        weight_df = weight_df[['traits', 'weight', 'SD']].dropna(subset=['weight'])
        weight_df['SD'] = weight_df['SD'].replace(0, 1).fillna(1)

        # 初始化结果列表
        bull_index_scores_data = []

        # 总的处理数
        total_bulls = len(self.needed_bulls_naab) + len(self.needed_bulls_reg)
        processed_bulls = 0

        # 处理 NAAB 公牛
        for bull in self.needed_bulls_naab:
            bull_data = bull_library_df[bull_library_df['BULL NAAB'] == bull]
            if bull_data.empty:
                for rule, replacement in self.bull_rules.items():
                    if bull.startswith(rule):
                        new_bull_naab = replacement + bull[3:]
                        bull_data = bull_library_df[bull_library_df['BULL NAAB'] == new_bull_naab]
                        if not bull_data.empty:
                            break
                    elif bull.startswith(replacement):
                        new_bull_naab = rule + bull[3:]
                        bull_data = bull_library_df[bull_library_df['BULL NAAB'] == new_bull_naab]
                        if not bull_data.empty:
                            break
            if not bull_data.empty:
                score = ((bull_data[weight_df['traits']].values.astype(float) * weight_df['weight'].values.astype(float)) / weight_df['SD'].values.astype(float)).sum()
                bull_index_scores_data.append(
                    [bull_data['BULL NAAB'].values[0], bull_data['BULL REG'].values[0], score])
            processed_bulls += 1
            progress_dialog.setValue(int((processed_bulls / total_bulls) * 100))

        # 处理 REG 公牛
        for bull in self.needed_bulls_reg:
            bull_data = bull_library_df[bull_library_df['BULL REG'] == bull]
            if not bull_data.empty:
                score = ((bull_data[weight_df['traits']].values.astype(float) * weight_df['weight'].values.astype(float)) / weight_df['SD'].values.astype(float)).sum()
                bull_index_scores_data.append(
                    [bull_data['BULL NAAB'].values[0], bull_data['BULL REG'].values[0], score])
            processed_bulls += 1
            progress_dialog.setValue(int((processed_bulls / total_bulls) * 100))

        # 将结果转换为 DataFrame
        bull_index_scores = pd.DataFrame(bull_index_scores_data, columns=['BULL NAAB', 'BULL REG', 'Index Score'])
        return bull_index_scores



    def find_bull_index_score(self, ancestor_id, bull_index_scores, *args):
        """
        在公牛指数得分中查找指定祖先的指数得分。

        :param ancestor_id: 要查找的祖先ID
        :param bull_index_scores: 包含公牛指数得分的字典或DataFrame
        :param args: 额外的参数（为了兼容性，但不使用）
        :return: 祖先的指数得分,如果找不到则返回np.nan
        """
        if pd.isna(ancestor_id) or str(ancestor_id).strip() == '':
            return np.nan

        ancestor_id = str(ancestor_id)

        try:
            if isinstance(bull_index_scores, dict):
                # 如果是字典，直接查找
                if len(ancestor_id) > 11:
                    return bull_index_scores.get(ancestor_id, np.nan)
                else:
                    score = bull_index_scores.get(ancestor_id, np.nan)
                    if np.isnan(score):
                        for rule, replacement in self.bull_rules.items():
                            if ancestor_id.startswith(rule):
                                new_ancestor_id = replacement + ancestor_id[3:]
                                score = bull_index_scores.get(new_ancestor_id, np.nan)
                                if not np.isnan(score):
                                    break
                            elif ancestor_id.startswith(replacement):
                                new_ancestor_id = rule + ancestor_id[3:]
                                score = bull_index_scores.get(new_ancestor_id, np.nan)
                                if not np.isnan(score):
                                    break
                    return score
            else:
                # 假设是DataFrame，使用原来的逻辑
                if len(ancestor_id) > 11:  # 如果祖先ID大于11位,直接查找
                    ancestor_score_row = bull_index_scores[bull_index_scores['BULL REG'] == ancestor_id]
                else:  # 如果祖先ID小于等于11位,先直接查找
                    ancestor_score_row = bull_index_scores[bull_index_scores['BULL NAAB'] == ancestor_id]

                    if ancestor_score_row.empty:  # 如果直接查找没有结果,根据规则替换前3位再查找
                        for rule, replacement in self.bull_rules.items():
                            if ancestor_id.startswith(rule):
                                new_ancestor_id = replacement + ancestor_id[3:]
                                ancestor_score_row = bull_index_scores[bull_index_scores['BULL NAAB'] == new_ancestor_id]
                                if not ancestor_score_row.empty:
                                    break
                            elif ancestor_id.startswith(replacement):
                                new_ancestor_id = rule + ancestor_id[3:]
                                ancestor_score_row = bull_index_scores[bull_index_scores['BULL NAAB'] == new_ancestor_id]
                                if not ancestor_score_row.empty:
                                    break

                if ancestor_score_row.empty:
                    logging.warning(f"No matching bull found for ID: {ancestor_id}")
                    return np.nan
                else:
                    return ancestor_score_row['Index Score'].values[0]

        except Exception as e:
            logging.error(f"Unexpected error in find_bull_index_score: {e}")
            return np.nan    

    def calculate_cow_index_scores(self, cow_df, bull_index_scores, sire_avg_scores, default_bull_score):
        cow_index_scores = pd.DataFrame(columns=['Cow ID', 'Index Score'])
        default_bull_id = '999HO99999'
        
        total_cows = len(cow_df)
        progress_dialog = TqdmProgressDialog(total_cows, "计算母牛指数得分", self)
        progress_dialog.show()

        try:
            for index, cow in cow_df.iterrows():
                score = 0
                ancestor_weights = {'sire': 0.5, 'mgs': 0.25, 'mmgs': 0.125}

                for ancestor_type, weight in ancestor_weights.items():
                    ancestor_id = cow[ancestor_type]

                    if ancestor_type == 'sire':
                        if pd.isna(ancestor_id) or (isinstance(ancestor_id, str) and ancestor_id.strip() == ''):
                            birth_year = cow['birth_date'].year
                            ancestor_score = 0.5 * sire_avg_scores.get(birth_year, default_bull_score) + 0.5 * self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                            # 因为算出的值比实际较高，预估是由于赋的值过高所致，因此使用999HO99999和不同年份均值（默认值）的均值进行尝试。
                            ancestor_score = ancestor_score if ancestor_score is not None else self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                        else:
                            ancestor_score = self.find_bull_index_score(ancestor_id, bull_index_scores, default_bull_score)
                            if pd.isna(ancestor_score):
                                birth_year = cow['birth_date'].year
                                ancestor_score = 0.5 * sire_avg_scores.get(birth_year, default_bull_score) + 0.5 * self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                                # 因为算出的值比实际较高，预估是由于赋的值过高所致，因此使用999HO99999和不同年份均值（默认值）的均值进行尝试。
                                ancestor_score = ancestor_score if ancestor_score is not None else self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)

                    elif ancestor_type == 'mgs':
                        if pd.isna(ancestor_id) or (isinstance(ancestor_id, str) and ancestor_id.strip() == ''):
                            dam_id = cow['dam']
                            if dam_id in cow_df['cow_id'].values:
                                dam_birth_year = cow_df.loc[cow_df['cow_id'] == dam_id, 'birth_date'].dt.year.iloc[0]
                                ancestor_score = 0.5 * sire_avg_scores.get(dam_birth_year, default_bull_score) + 0.5 * self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                                # 因为算出的值比实际较高，预估是由于赋的值过高所致，因此使用999HO99999和不同年份均值（默认值）的均值进行尝试。
                                ancestor_score = ancestor_score if ancestor_score is not None else self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                            else:
                                ancestor_score = self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                        else:
                            ancestor_score = self.find_bull_index_score(ancestor_id, bull_index_scores, default_bull_score)
                            if pd.isna(ancestor_score):
                                dam_id = cow['dam']
                                if dam_id in cow_df['cow_id'].values:
                                    dam_birth_year = cow_df.loc[cow_df['cow_id'] == dam_id, 'birth_date'].dt.year.iloc[0]
                                    ancestor_score =  0.5 * sire_avg_scores.get(dam_birth_year, default_bull_score) + 0.5 * self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                                # 因为算出的值比实际较高，预估是由于赋的值过高所致，因此使用999HO99999和不同年份均值（默认值）的均值进行尝试。
                                    ancestor_score = ancestor_score if ancestor_score is not None else self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                                else:
                                    ancestor_score = self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)

                    else:  # mmgs
                        if pd.isna(ancestor_id) or (isinstance(ancestor_id, str) and ancestor_id.strip() == ''):
                            dam_id = cow['dam']
                            if dam_id in cow_df['cow_id'].values:
                                dam_mother_id = cow_df.loc[cow_df['cow_id'] == dam_id, 'dam'].iloc[0]
                                if dam_mother_id in cow_df['cow_id'].values:
                                    dam_mother_birth_year = cow_df.loc[cow_df['cow_id'] == dam_mother_id, 'birth_date'].dt.year.iloc[0]
                                    ancestor_score =  0.5 * sire_avg_scores.get(dam_mother_birth_year, default_bull_score) + 0.5 * self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                                # 因为算出的值比实际较高，预估是由于赋的值过高所致，因此使用999HO99999和不同年份均值（默认值）的均值进行尝试。
                                    ancestor_score = ancestor_score if ancestor_score is not None else self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                                else:
                                    ancestor_score = self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                            else:
                                ancestor_score = self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                        else:
                            ancestor_score = self.find_bull_index_score(ancestor_id, bull_index_scores, default_bull_score)
                            if pd.isna(ancestor_score):
                                dam_id = cow['dam']
                                if dam_id in cow_df['cow_id'].values:
                                    dam_mother_id = cow_df.loc[cow_df['cow_id'] == dam_id, 'dam'].iloc[0]
                                    if dam_mother_id in cow_df['cow_id'].values:
                                        dam_mother_birth_year = cow_df.loc[cow_df['cow_id'] == dam_mother_id, 'birth_date'].dt.year.iloc[0]
                                        ancestor_score =  0.5 * sire_avg_scores.get(dam_mother_birth_year, default_bull_score) + 0.5 * self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                                        # 因为算出的值比实际较高，预估是由于赋的值过高所致，因此使用999HO99999和不同年份均值（默认值）的均值进行尝试。
                                        ancestor_score = ancestor_score if ancestor_score is not None else self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                                    else:
                                        ancestor_score = self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                                else:
                                    ancestor_score = self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)

                    score += float(ancestor_score) * weight  # 将 ancestor_score 转换为 float 类型

                # 添加默认得分的权重
                default_score = self.find_bull_index_score(default_bull_id, bull_index_scores, default_bull_score)
                score += float(default_score) * 0.125  # 将 default_score 转换为 float 类型

                new_row = pd.DataFrame({'Cow ID': [cow['cow_id']], 'Index Score': [score]})
                if not new_row.empty:
                    cow_index_scores = pd.concat([cow_index_scores, new_row], ignore_index=True)
                progress_dialog.update_progress()
                QApplication.processEvents()
                
                if progress_dialog.wasCanceled():
                    break

        finally:
            progress_dialog.close()

        return cow_index_scores

    def generate_table_a(self, cow_df, bull_index_scores):
        # 筛选出父号不为空的母牛
        table_a = cow_df[cow_df['sire'].notna()][['cow_id', 'birth_date', 'sire']]

        # 提取出生年份
        table_a['birth_year'] = table_a['birth_date'].dt.year

        # 根据父号在公牛指数表中查找对应的指数
        table_a['sire_index'] = table_a['sire'].apply(
            lambda x: self.find_bull_index_score(x, bull_index_scores, np.nan))

        # 重新排列列的顺序
        table_a = table_a[['cow_id', 'birth_year', 'sire', 'sire_index']]

        return table_a

    def generate_table_b(self, table_a):
        logging.info("开始生成表b")
        """print("Input table_a:")
        print(table_a.head())"""

        try:
            # 确保 birth_year 和 sire_index 列不包含 NaN 值
            table_a['birth_year'] = pd.to_numeric(table_a['birth_year'], errors='coerce')
            table_a['sire_index'] = pd.to_numeric(table_a['sire_index'], errors='coerce')
            table_a = table_a.dropna(subset=['birth_year', 'sire_index'])

            # 按出生年份分组,计算每年父号指数的平均值
            table_b = table_a.groupby('birth_year')['sire_index'].agg(['count', 'mean']).reset_index()

            """print("table_b after grouping and aggregating:")
            print(table_b.head())"""

            # 筛选出数据量大于等于10的年份
            table_b = table_b[table_b['count'] >= 10]

            """print("table_b after filtering:")
            print(table_b.head())"""

            # 获取母牛群的最小出生年份和表b的最大年份
            min_cow_birth_year = self.cow_df['birth_date'].dt.year.min()
            if pd.isna(min_cow_birth_year):
                min_cow_birth_year = 2010  # 设置一个默认值
            else:
                min_cow_birth_year = int(min_cow_birth_year)

            max_year = max(int(table_b['birth_year'].max() if not table_b.empty else min_cow_birth_year), datetime.datetime.now().year)
            """print("min_cow_birth_year:", min_cow_birth_year)
            print("max_year:", max_year)"""

            # 获取999HO99999的值作为默认值
            default_value = self.find_bull_index_score('999HO99999', self.bull_scores_naab_dict, self.bull_scores_reg_dict)

            # 如果只有一个数据点或没有数据点，使用特殊逻辑
            if len(table_b) <= 1:
                logging.warning("表b数据点不足，使用特殊逻辑。")
                min_cow_birth_year = 2010
                table_b = pd.DataFrame({'birth_year': [2010], 'count': [1], 'mean': [default_value]})

            # 创建完整的年份范围
            full_year_range = pd.DataFrame({'birth_year': range(min_cow_birth_year, max_year + 1)})
            table_b = pd.merge(full_year_range, table_b, on='birth_year', how='left')

            """print("table_b after merging with full year range:")
            print(table_b.head())"""

            # 添加一列来标记缺失值
            table_b['is_missing'] = table_b['mean'].isna()

            # 找出非空的年份和对应的值
            known_data = table_b[~table_b['mean'].isna()]
            known_years = known_data['birth_year'].values.reshape(-1, 1)
            known_values = known_data['mean'].values

            logging.info(f"表b - 已知年份: {known_years.flatten()}")
            logging.info(f"表b - 已知值: {known_values}")

            if len(known_years) >= 2:
                try:
                    # 创建并训练线性回归模型
                    model = LinearRegression()
                    model.fit(known_years, known_values)

                    # 预测缺失年份的值
                    missing_years = table_b.loc[table_b['mean'].isna(), 'birth_year'].values.reshape(-1, 1)
                    if len(missing_years) > 0:
                        predicted_values = model.predict(missing_years)
                        table_b.loc[table_b['mean'].isna(), 'mean'] = predicted_values
                    logging.info("表b - 线性回归完成")
                except Exception as e:
                    logging.error(f"表b在进行线性回归时发生错误: {str(e)}")
                    # 如果线性回归失败，使用已知值的平均值填充
                    mean_value = known_values.mean()
                    table_b.loc[table_b['mean'].isna(), 'mean'] = mean_value
                    logging.warning(f"表b - 使用平均值 {mean_value} 填充缺失值")
            else:
                # 如果没有足够的数据点进行回归，使用默认值填充
                table_b['mean'] = table_b['mean'].fillna(default_value)
                logging.warning(f"表b - 使用默认值 {default_value} 填充所有缺失值")

            # 添加一列来标记插值的结果
            table_b['interpolated'] = table_b['is_missing'].map({True: 'Yes', False: 'No'})

            """print("table_b after filling missing values:")
            print(table_b.head())"""

            # 重新排列列的顺序
            table_b = table_b[['birth_year', 'mean', 'interpolated']]

            """print("Final table_b:")
            print(table_b.head())"""

            logging.info("表b生成完成")
            return table_b
        except Exception as e:
            logging.error(f"生成表b时发生错误: {str(e)}", exc_info=True)
            # 如果发生错误，返回一个基本的DataFrame以避免程序崩溃
            default_min_year = 2010
            default_max_year = datetime.datetime.now().year
            return pd.DataFrame({
                'birth_year': range(default_min_year, default_max_year + 1),
                'mean': [default_value] * (default_max_year - default_min_year + 1),
                'interpolated': ['Yes'] * (default_max_year - default_min_year + 1)
            })

    def calculate_geno_index_scores(self, geno_library_df, weight_df):
        geno_index_scores = pd.DataFrame(columns=['Farm ID', 'Index Score'])

        weight_df = weight_df[['traits', 'weight', 'SD']].dropna(subset=['weight'])
        weight_df['SD'] = weight_df['SD'].replace(0, 1).fillna(1)

        geno_library_dict = geno_library_df.set_index('Farm ID').T.to_dict('dict')

        for farm_id, row in geno_library_dict.items():
            score = 0
            for _, weight_row in weight_df.iterrows():
                trait = weight_row['traits']
                if trait in row and pd.notnull(row[trait]):
                    trait_value = row[trait]
                    weight = weight_row['weight']
                    sd = weight_row['SD']
                    score += (trait_value * weight) / sd
            new_row = pd.DataFrame({
                'Farm ID': [farm_id],
                'Index Score': [score]
            })
            geno_index_scores = pd.concat([geno_index_scores, new_row], ignore_index=True)

        return geno_index_scores

    def calculate_combined_index_scores(self, cow_df, cow_index_scores):
        # 重命名 pedigree_index_scores_df 的列
        pedigree_index_scores_df = cow_index_scores.rename(columns={'Cow ID': 'cow_id', 'Index Score': 'Pedigree_Index Score'})
        
        # 检查是否存在基因组指数文件
        geno_index_file = os.path.join(self.output_folder, "基因组指数.xlsx")
        
        # 打印路径和文件是否存在的信息
        """print("geno_index_file path:", geno_index_file)
        print("os.path.exists(geno_index_file):", os.path.exists(geno_index_file))"""

        
        if os.path.exists(geno_index_file):
            # 读取基因组指数得分
            geno_index_scores_df = pd.read_excel(geno_index_file)
            geno_index_scores_df = geno_index_scores_df.rename(columns={'Farm ID': 'cow_id', 'Index Score': 'Geno_Index Score'})
        else:
            # 生成空的 geno_index_scores_df
            geno_index_scores_df = pd.DataFrame(columns=['cow_id', 'Geno_Index Score'], dtype='float64')

        # 打印调试信息
        """print("pedigree_index_scores_df:\n", pedigree_index_scores_df.head(10))
        print("geno_index_scores_df:\n", geno_index_scores_df.head(10))"""

        # 确保所有的 cow_id 都是字符串类型以避免合并问题
        cow_df['cow_id'] = cow_df['cow_id'].astype(str)
        pedigree_index_scores_df['cow_id'] = pedigree_index_scores_df['cow_id'].astype(str)
        geno_index_scores_df['cow_id'] = geno_index_scores_df['cow_id'].astype(str)

        # 合并 pedigree_index_scores_df 和 cow_df
        merged_cow_data = pd.merge(cow_df, pedigree_index_scores_df[['cow_id', 'Pedigree_Index Score']], on='cow_id', how='left')
        """print("merged_cow_data:\n", merged_cow_data.head(10))"""

        # 合并 geno_index_scores_df 和 merged_cow_data
        merged_cow_data = pd.merge(merged_cow_data, geno_index_scores_df[['cow_id', 'Geno_Index Score']], on='cow_id', how='left')
        """print("merged_cow_data_geno:\n", merged_cow_data.head(10))"""

        # 计算 Combine Index Score
        merged_cow_data['Combine Index Score'] = merged_cow_data.apply(
            lambda row: row['Geno_Index Score'] if pd.notna(row['Geno_Index Score']) else row['Pedigree_Index Score'], axis=1
        )

        # 重新排列列的顺序
        index_columns = ['Pedigree_Index Score', 'Geno_Index Score', 'Combine Index Score']
        columns = [col for col in merged_cow_data.columns if col not in index_columns] + index_columns
        merged_cow_data = merged_cow_data[columns]

        return merged_cow_data

    # 生成正态分布图的方法（关键性状及育种指数得分均调用此方法）
    def generate_normal_distribution_charts(self, df, value_column, output_folder, file_prefix):
        plt.rcParams['font.sans-serif'] = ['SimHei']
        plt.rcParams['axes.unicode_minus'] = False

        if '是否在场' not in df.columns or '是' not in df['是否在场'].values:
            QMessageBox.warning(self, "警告", "母牛信息中未标明在场牛只，将影响部分功能。")
            df_in_herd = df
        else:
            df_in_herd = df[df['是否在场'] == '是']

        if df_in_herd.empty:
            QMessageBox.warning(self, "警告", "没有在场的牛只数据，无法生成分布图。")
            return

        # 年份分布
        self.generate_year_distribution(df_in_herd, value_column, output_folder, file_prefix)
        
        # 胎次分布（只区分成母牛和后备牛）
        self.generate_parity_distribution(df_in_herd, value_column, output_folder, file_prefix)

    def generate_year_distribution(self, df, value_column, output_folder, file_prefix):
        max_birth_year = df['birth_date'].dt.year.max()
        df['group'] = pd.cut(
            df['birth_date'].dt.year,
            bins=[-float('inf')] + list(range(max_birth_year-4, max_birth_year+1)),
            labels=[f'{max_birth_year-4}年及以前'] + [f'{year}年' for year in range(max_birth_year-3, max_birth_year+1)]
        )
        
        self.plot_distribution(df, value_column, '年份', output_folder, file_prefix)

    def generate_parity_distribution(self, df, value_column, output_folder, file_prefix):
        df['group'] = df['lac'].apply(lambda x: '后备牛' if pd.isna(x) or x == 0 else '成母牛')
        
        self.plot_distribution(df, value_column, '胎次', output_folder, file_prefix)

    def plot_distribution(self, df, value_column, group_type, output_folder, file_prefix):
        plt.figure(figsize=(12, 8))
        
        cmap = plt.cm.get_cmap('Set1')
        unique_groups = df['group'].unique()
        colors = cmap(np.linspace(0, 1, len(unique_groups)))
        
        legend_text = []
        
        for i, (group, group_data) in enumerate(df.groupby('group', observed=True)):
            group_values = group_data[value_column].dropna()
            
            if len(group_values) > 0:
                mean = np.mean(group_values)
                std = np.std(group_values)
                
                x = np.linspace(mean - 3*std, mean + 3*std, 100)
                y = stats.norm.pdf(x, mean, std)
                
                plt.plot(x, y, color=colors[i], linestyle=['-', '--'][i % 2], label=group)
                
                legend_text.append(f"{group}\n均值: {mean:.1f}\n标准差: {std:.2f}\nN: {len(group_values)}")
        
        plt.title(f'{value_column}值{group_type}分布', fontsize=16)
        plt.xlabel(f'{value_column}值', fontsize=12)
        plt.ylabel('密度', fontsize=12)
        plt.grid(True, linestyle='--', alpha=0.7)
        
        plt.legend(legend_text, loc='upper right', bbox_to_anchor=(1.25, 1), fontsize=10)
        
        plt.tight_layout()
        
        plt.savefig(os.path.join(output_folder, f"{file_prefix}_{group_type}正态分布.png"), dpi=300, bbox_inches='tight')
        plt.close()
        
        excel_file = os.path.join(output_folder, f"{file_prefix}_{group_type}正态分布数据.xlsx")
        with pd.ExcelWriter(excel_file, engine='xlsxwriter') as writer:
            for group, group_data in df.groupby('group', observed=True):
                group_data.to_excel(writer, sheet_name=str(group)[:31], index=False)
        
        logging.info(f"{group_type}正态分布图和数据已保存: {excel_file}")

    # 五等份数据分析
    def generate_quintile_distribution(self, df, value_column, output_folder, file_prefix):
        excel_file = os.path.join(output_folder, f"{file_prefix}_5等份分布表.xlsx")
        
        workbook = xlsxwriter.Workbook(excel_file, {'nan_inf_to_errors': True})
        worksheet = workbook.add_worksheet("5等份分布")
        
        # 定义格式
        title_fmt = workbook.add_format({'bold': True, 'font_size': 14})
        header_fmt = workbook.add_format({'bold': True, 'border': 1, 'align': 'center'})
        border_fmt = workbook.add_format({'border': 1})
        number_fmt = workbook.add_format({'border': 1, 'num_format': '#,##0.00'})
        
        worksheet.write('A1', f"{value_column}5等份分布", title_fmt)
        
        current_row = 2
        
        for cattle_type in ['成母牛', '后备牛']:
            if cattle_type == '成母牛':
                data = df[(df['lac'] >= 1) & (df['是否在场'] != '否') & (df['sex'] != '公')][value_column]
            else:
                data = df[((df['lac'] == 0) | (df['lac'].isna())) & (df['是否在场'] != '否') & (df['sex'] != '公')][value_column]
            
           
            # 将数据转换为数值类型
            data = pd.to_numeric(data, errors='coerce')
            
            # 处理 NaN/INF 值
            data = data.replace([np.inf, -np.inf], np.nan)
            data = data.dropna()
            
            if data.empty:
                logging.warning(f"没有{cattle_type}的数据，跳过生成分布表。")
                continue

            quintiles = np.percentile(data, [20, 40, 60, 80])
            
            # 计算每个区间的头数和统计值
            intervals = [
                (float('-inf'), quintiles[0]),
                (quintiles[0], quintiles[1]),
                (quintiles[1], quintiles[2]),
                (quintiles[2], quintiles[3]),
                (quintiles[3], float('inf'))
            ]
            
            stats = []
            for start, end in intervals:
                if start == float('-inf'):
                    mask = (data <= end)
                elif end == float('inf'):
                    mask = (data > start)
                else:
                    mask = (data > start) & (data <= end)
                
                interval_data = data[mask]
                stats.append({
                    'count': len(interval_data),
                    'average': interval_data.mean() if not interval_data.empty else 0,
                    'min': interval_data.min() if not interval_data.empty else 0,
                    'max': interval_data.max() if not interval_data.empty else 0
                })
            
            distribution = pd.DataFrame({
                '指数分布': ['最差20%', '20%', '20%', '20%', '最优20%'],
                '头数': [s['count'] for s in stats],
                'AVERAGE': [s['average'] for s in stats],
                'MIN': [s['min'] for s in stats],
                'MAX': [s['max'] for s in stats]
            })
            
            # 写入小标题
            worksheet.write(current_row, 0, f"{cattle_type}分布", title_fmt)
            current_row += 1
            
            # 写入表头
            headers = ['指数分布', '头数', 'AVERAGE', 'MIN', 'MAX']
            for col, header in enumerate(headers):
                worksheet.write(current_row, col, header, header_fmt)
            current_row += 1
            
            # 写入数据
            for row in range(len(distribution)):
                # 写入指数分布
                worksheet.write(current_row + row, 0, distribution.iloc[row, 0], border_fmt)
                
                # 写入头数（使用整数格式）
                worksheet.write(current_row + row, 1, distribution.iloc[row, 1], border_fmt)
                
                # 写入其他统计值
                for col in range(2, 5):
                    value = distribution.iloc[row, col]
                    if pd.isna(value) or value in [np.inf, -np.inf]:
                        value = 0
                    worksheet.write(current_row + row, col, value, number_fmt)
            
            current_row += len(distribution) + 2  # 添加一个空行
        
        # 调整列宽
        worksheet.set_column('A:E', 15)  # 注意这里改为A:E，因为多了一列
        
        workbook.close()
        
        logging.info(f"{value_column}5等份分布表已保存: {excel_file}")


    # 计算牛群的关键性状值（并包括各种调用）
    def calculate_key_traits_index(self):
        try:
            progress_dialog = QProgressDialog("正在进行关键育种性状计算...\n\n数据库连接中......", "取消", 0, 100, self)
            progress_dialog.setWindowModality(Qt.WindowModality.WindowModal)
            progress_dialog.setValue(0)
            progress_dialog.show()
            progress_dialog.raise_()
            QApplication.processEvents()            
            # 如果没有上传备选公牛信息,则生成一个默认的空表
            if self.bull_df.empty:
                self.bull_df = pd.DataFrame(columns=['bull_id', 'semen_type'])

            if self.cow_df.empty:
                QMessageBox.warning(self, "警告", "请先上传在群母牛信息。")
                return
            progress_dialog.setValue(5)
            progress_dialog.show()
            progress_dialog.raise_()
            QApplication.processEvents()  
            # 检查日期格式
            date_columns = [ 'birth_date']
            for column in date_columns:
                if column in self.cow_df.columns:
                    if not pd.api.types.is_datetime64_any_dtype(self.cow_df[column]):
                        QMessageBox.warning(self, "警告", f"{column}列的格式不是日期格式,请修改为日期格式后重新上传。")
                        return
            progress_dialog.setValue(10)
            progress_dialog.show()
            progress_dialog.raise_()
            QApplication.processEvents()          
            # 生成所需公牛列表
            self.generate_needed_bulls(progress_dialog)
            """print("Generated needed bulls:")
            print("Needed NAAB bulls:", self.needed_bulls_naab)
            print("Needed REG bulls:", self.needed_bulls_reg)"""
            progress_dialog.setValue(15)
            progress_dialog.show()
            progress_dialog.raise_()
            QApplication.processEvents()
            # 检查是否有缺失公牛
            missing_bulls = self.check_missing_bulls(self.bull_df)
            if missing_bulls:
                self.update_missing_bulls(missing_bulls, "关键育种性状计算")
                """print("missing bull:",missing_bulls)"""
            


            default_traits = [ 'NM$','TPI', 'MILK', 'FAT', 'FAT %', 'PROT', 'PROT%', 'SCS', 'PL', 'DPR', 'PTAT', 'UDC',
                              'FLC', 'ST', 'TL']
            all_traits = ['TPI', 'NM$', 'CM$', 'FM$', 'GM$', 'MILK', 'FAT', 'PROT', 'FAT %', 'PROT%', 'SCS', 'DPR',
                          'HCR', 'CCR', 'PL', 'SCE', 'DCE', 'SSB', 'DSB', 'PTAT', 'UDC', 'FLC', 'BDC', 'ST', 'SG', 'BD',
                          'DF', 'RA', 'RW', 'LS', 'LR', 'FA', 'FLS', 'FU', 'UH', 'UW', 'UC', 'UD', 'FT', 'RT', 'TL',
                          'FE', 'FI', 'HI', 'LIV', 'GL', 'MAST', 'MET', 'RP', 'KET', 'DA', 'MFV', 'EFC', 'HLiv', 'FS',
                          'RFI', 'Milk Speed' ,'Eval Date']

            dialog = KeyTraitsDialog(default_traits, all_traits)
            progress_dialog.setValue(17)
            progress_dialog.show()
            progress_dialog.raise_()
            QApplication.processEvents()
            if dialog.exec() == QDialog.DialogCode.Accepted:
                selected_traits = dialog.get_selected_traits()
                if not selected_traits:
                    QMessageBox.warning(self, "警告", "请至少选择一个关键性状。")
                    return
                
                # 将selected_traits保存到selected_traits_key_traits.txt文件中
                selected_traits_file = os.path.join(self.output_folder, "selected_traits_key_traits.txt")
                with open(selected_traits_file, "w", encoding="utf-8") as f:
                    for trait in selected_traits:
                        f.write(trait + "\n")
                
                # 将selected_traits存储为实例属性
                self.selected_traits = selected_traits
                
                progress_dialog.setValue(20)
                # 获取所需的公牛列表
                needed_bulls_naab = set()
                needed_bulls_reg = set()
                for bull_id in self.bull_df['bull_id']:
                    if len(bull_id) > 11:
                        needed_bulls_reg.add(bull_id)
                    else:
                        needed_bulls_naab.add(bull_id)

                progress_dialog.setValue(25)
                for _, row in self.cow_df.iterrows():
                    for ancestor_type in ['sire', 'mgs', 'mmgs']:
                        if ancestor_type in row:  # 检查列是否存在
                            ancestor_id = str(row[ancestor_type])
                            if len(ancestor_id) > 11:
                                needed_bulls_reg.add(ancestor_id)
                            else:
                                needed_bulls_naab.add(ancestor_id)

                needed_bulls_naab.discard('nan')  # 去除空值
                needed_bulls_reg.discard('nan')  # 去除空值
                needed_bulls_naab.add('999HO99999')  # 添加默认公牛ID

                progress_dialog.setValue(28)
                # 生成所需公牛列表
                self.generate_needed_bulls(progress_dialog)

                """print("Generated needed bulls:")
                print("Needed NAAB bulls:", self.needed_bulls_naab)
                print("Needed REG bulls:", self.needed_bulls_reg)

                print("Needed NAAB bulls for key traits:")
                print(needed_bulls_naab)
                print("Needed REG bulls for key traits:")
                print(needed_bulls_reg)"""

                progress_dialog.setValue(32)
                # 读取公牛信息
                connection = pymysql.connect(
                    host=cipher_suite.decrypt(encoded_host).decode(),
                    port=int(cipher_suite.decrypt(encoded_port).decode()),
                    user=cipher_suite.decrypt(encoded_user).decode(),
                    password=cipher_suite.decrypt(encoded_password).decode(),
                    database=cipher_suite.decrypt(encoded_db).decode(),
                    charset='utf8mb4',
                    cursorclass=pymysql.cursors.DictCursor
                )
                
                progress_dialog.setValue(36)
                if needed_bulls_naab and needed_bulls_reg:
                    placeholders_naab = ','.join(['%s'] * len(needed_bulls_naab))
                    placeholders_reg = ','.join(['%s'] * len(needed_bulls_reg))
                    sql = f"SELECT * FROM bull_library WHERE `BULL NAAB` IN ({placeholders_naab}) OR `BULL REG` IN ({placeholders_reg})"
                    params = tuple(needed_bulls_naab) + tuple(needed_bulls_reg)
                elif needed_bulls_naab:
                    placeholders_naab = ','.join(['%s'] * len(needed_bulls_naab))
                    sql = f"SELECT * FROM bull_library WHERE `BULL NAAB` IN ({placeholders_naab})"
                    params = tuple(needed_bulls_naab)
                elif needed_bulls_reg:
                    placeholders_reg = ','.join(['%s'] * len(needed_bulls_reg))
                    sql = f"SELECT * FROM bull_library WHERE `BULL REG` IN ({placeholders_reg})"
                    params = tuple(needed_bulls_reg)
                else:
                    bull_library_df = pd.DataFrame()  # 返回一个空的 DataFrame
                    return
                progress_dialog.setValue(39)
                with connection.cursor() as cursor:
                    cursor.execute(sql, params)
                    result = cursor.fetchall()
                    bull_library_df = pd.DataFrame(result)
                progress_dialog.setValue(43)
                # 处理需要替换前3位的公牛
                for bull in needed_bulls_naab.copy():  # 创建一个副本,避免在迭代过程中修改原始集合
                    if bull not in bull_library_df['BULL NAAB'].values:
                        for rule, replacement in self.bull_rules.items():
                            if bull.startswith(rule):
                                new_bull = replacement + bull[3:]
                                needed_bulls_naab.add(new_bull)  # 将替换后的公牛ID添加到needed_bulls_naab
                                break
                            elif bull.startswith(replacement):
                                new_bull = rule + bull[3:]
                                needed_bulls_naab.add(new_bull)  # 将替换后的公牛ID添加到needed_bulls_naab
                                break
                progress_dialog.setValue(46)                    
                if needed_bulls_naab:
                    placeholders_naab = ','.join(['%s'] * len(needed_bulls_naab))
                    sql = f"SELECT * FROM bull_library WHERE `BULL NAAB` IN ({placeholders_naab})"
                    params = tuple(needed_bulls_naab)
                    with connection.cursor() as cursor:
                        cursor.execute(sql, params)
                        result = cursor.fetchall()
                        bull_library_df = pd.concat([bull_library_df, pd.DataFrame(result)], ignore_index=True)

                connection.close()

                progress_dialog.setValue(52)
                # 初始化 bull_key_traits_scores,包含所有需要的公牛
                bull_key_traits_scores = pd.DataFrame(columns=['BULL NAAB', 'BULL REG'] + selected_traits)

                for bull in needed_bulls_naab | needed_bulls_reg:
                    if len(bull) > 11:
                        bull_data = bull_library_df[bull_library_df['BULL REG'] == bull]
                        if bull_data.empty:
                            bull_data = pd.DataFrame({'BULL NAAB': [np.nan], 'BULL REG': [bull]})
                            for trait in selected_traits:
                                bull_data[trait] = np.nan
                    else:
                        bull_data = bull_library_df[bull_library_df['BULL NAAB'] == bull]
                        if bull_data.empty:
                            found = False
                            for rule, replacement in self.bull_rules.items():
                                if bull.startswith(rule):
                                    new_bull = replacement + bull[3:]
                                    bull_data = bull_library_df[bull_library_df['BULL NAAB'] == new_bull]
                                    if not bull_data.empty:
                                        found = True
                                        break
                                elif bull.startswith(replacement):
                                    new_bull = rule + bull[3:]
                                    bull_data = bull_library_df[bull_library_df['BULL NAAB'] == new_bull]
                                    if not bull_data.empty:
                                        found = True
                                        break
                            if not found:
                                # 如果替换前三位后仍未找到,则创建一个新的DataFrame行
                                bull_data = pd.DataFrame({'BULL NAAB': [bull], 'BULL REG': [np.nan]})
                                for trait in selected_traits:
                                    bull_data[trait] = np.nan
                            else:
                                # 如果替换前三位后找到了,则更新 BULL REG 和各性状指数
                                bull_data['BULL REG'] = \
                                bull_library_df[bull_library_df['BULL NAAB'] == new_bull]['BULL REG'].values[0]
                                for trait in selected_traits:
                                    bull_data[trait] = \
                                    bull_library_df[bull_library_df['BULL NAAB'] == new_bull][trait].values[0]
                    bull_key_traits_scores = pd.concat(
                        [bull_key_traits_scores, bull_data[['BULL NAAB', 'BULL REG'] + selected_traits]],
                        ignore_index=True)

                self.status_label.setText(self.status_label.text() + "\n公牛关键性状指数计算完成。")
                progress_dialog.setValue(64)
                # 保存公牛关键性状指数到文件
                bull_key_traits_file = os.path.join(self.output_folder, "结果-公牛关键性状指数.xlsx")
                bull_key_traits_scores.to_excel(bull_key_traits_file, index=False)

                self.status_label.setText(self.status_label.text() + "\n公牛关键性状指数计算完成。")
                logging.info("公牛关键性状指数计算完成。")

                self.status_label.setText(self.status_label.text() + "\n正在生成表c和表d...")

                progress_dialog.setValue(73)
                # 生成表c和表d
                table_c = self.generate_table_c(self.cow_df, bull_key_traits_scores, selected_traits)
                # 不需要将cow_df的内容合并到table_c中
                # table_c = self.cow_df.merge(table_c, on='cow_id', how='left')
                table_d = self.generate_table_d(table_c, selected_traits, bull_key_traits_scores)

                # 保存表c到文件
                table_c_file = os.path.join(self.output_folder, "表c.xlsx")
                table_c.to_excel(table_c_file, index=False)

                # 保存表d到文件
                table_d_file = os.path.join(self.output_folder, "表d.xlsx")
                with pd.ExcelWriter(table_d_file) as writer:
                    for trait, df in table_d.items():
                        df.to_excel(writer, sheet_name=trait, index=False)

                # 生成系谱识别情况分析表
                pedigree_analysis = self.generate_pedigree_identification_analysis(table_c, self.cow_df)

                # 保存系谱识别情况分析表到文件
                pedigree_analysis_file = os.path.join(self.output_folder, "结果-系谱识别情况分析.xlsx")
                pedigree_analysis.to_excel(pedigree_analysis_file, index=False)

                self.status_label.setText(self.status_label.text() + "\n系谱识别情况分析完成。")
               

                self.status_label.setText(self.status_label.text() + "\n正在计算母牛关键性状指数...")

                progress_dialog.setValue(83)
                # 计算母牛关键性状指数
                cow_key_traits_scores = self.calculate_cow_key_traits_scores(self.cow_df, bull_key_traits_scores, table_d, selected_traits)

                # 合并 cow_df 和计算出的关键性状指数
                merged_cow_key_traits_scores = pd.merge(self.cow_df, cow_key_traits_scores, left_on='cow_id', right_on='Cow ID', how='left')
                



                self.status_label.setText(self.status_label.text() + "\n母牛关键性状指数计算完成。")
                progress_dialog.setValue(93)
                
                
                # 保存母牛关键性状指数到文件
                cow_key_traits_file = os.path.join(self.output_folder, "结果-母牛关键性状指数.xlsx")
                merged_cow_key_traits_scores.to_excel(cow_key_traits_file, index=False)
                """print("merged结果-母牛关键性状指数:",merged_cow_key_traits_scores.head())
                print(list(selected_traits))"""

                self.status_label.setText(self.status_label.text() + "\n母牛关键性状指数计算完成。")
                logging.info("母牛关键性状指数计算完成。")

                self.status_label.setText(self.status_label.text() + "\n正在计算基因组关键性状指数...")

                # 计算基因组关键性状指数
                geno_library_file = os.path.join(self.output_folder, "geno_library.xlsx")

                if os.path.isfile(geno_library_file):
                    geno_library_df = pd.read_excel(geno_library_file)
                    geno_key_traits_scores = geno_library_df[['Farm ID'] + selected_traits].copy()
                    geno_key_traits_scores.fillna(0, inplace=True)

                    self.status_label.setText(self.status_label.text() + "\n基因组关键性状指数计算完成。")
                    # 保存基因组关键性状指数到文件
                    geno_key_traits_file = os.path.join(self.output_folder, "结果-基因组关键性状指数.xlsx")
                    geno_key_traits_scores.to_excel(geno_key_traits_file, index=False)
                    self.status_label.setText(self.status_label.text() + "\n基因组关键性状指数计算完成。")
                    logging.info("基因组关键性状指数计算完成。")
                else:
                    geno_key_traits_scores = pd.DataFrame()

                # 以下操作无论是否有基因组数据都应执行
                # 生成牛群关键性状年度变化表
                key_traits_annual_change = self.generate_key_traits_annual_change(merged_cow_key_traits_scores, selected_traits)

                # 牛群关键性状年度变化保存结果到Excel文件
                key_traits_annual_change_file = os.path.join(self.output_folder, "结果-牛群关键性状年度变化.xlsx")
                key_traits_annual_change.to_excel(key_traits_annual_change_file, index=False)

                self.status_label.setText(self.status_label.text() + "\n牛群关键性状年度变化表生成完成。")
                logging.info("牛群关键性状年度变化表生成完成。")    

                # 牛群关键性状年度变化生成Excel图表
                self.generate_excel_charts(key_traits_annual_change, selected_traits)

                self.status_label.setText(self.status_label.text() + "\n关键性状指数计算和牛群关键性状年度变化图生成完成。")

                # 生成NM$分布图表
                self.generate_nm_distribution_charts(merged_cow_key_traits_scores)
                self.status_label.setText(self.status_label.text() + "\nNM$正态分布图表生成完成。")


                # 生成NM$的年份分布图
                self.generate_normal_distribution_charts(merged_cow_key_traits_scores, 'NM$', self.output_folder, "NM$")
                self.status_label.setText(self.status_label.text() + "\nNM$的年份正态分布图生成完成。")


                # 生成NM$的五等份分布表
                self.generate_quintile_distribution(merged_cow_key_traits_scores, 'NM$', self.output_folder, "NM$")
                self.status_label.setText(self.status_label.text() + "\nNM$的五等份分布表生成完成。")


                progress_dialog.setValue(100)
                
        except Exception as e:
            logging.error(f"计算关键性状指数时发生错误: {e}", exc_info=True)
            QMessageBox.critical(self, "错误", f"计算关键性状指数时发生错误：{e},请按照模版上传母牛信息。")

    # 关键性状进展表
    def generate_key_traits_annual_change(self, merged_cow_key_traits_scores, selected_traits):
        # 确保birth_date是日期类型
        merged_cow_key_traits_scores['birth_date'] = pd.to_datetime(merged_cow_key_traits_scores['birth_date'])
        
        # 获取牛群中的最大出生年份
        max_birth_year = merged_cow_key_traits_scores['birth_date'].dt.year.max()

        # 创建出生年份分组
        merged_cow_key_traits_scores['birth_year_group'] = pd.cut(
            merged_cow_key_traits_scores['birth_date'].dt.year,
            bins=[-float('inf')] + list(range(max_birth_year-4, max_birth_year)) + [float('inf')],
            labels=[f'{max_birth_year-4}年及以前'] + [f'{year}年' for year in range(max_birth_year-3, max_birth_year+1)]
        )

        # 创建分析函数
        def analyze_group(group):
            total = len(group)
            result = {'头数': total}
            for trait in selected_traits:
                result[trait] = round(group[trait].mean(), 2)
            return pd.Series(result)

        # 总体分析
        results = []
        
        total_result = merged_cow_key_traits_scores.groupby('birth_year_group', observed=True).apply(analyze_group).reset_index()
        total_result['是否在场'] = '合计'
        results.append(total_result)
        
        # 添加总计的合计行
        total_summary = merged_cow_key_traits_scores.pipe(analyze_group)
        total_summary = pd.DataFrame(total_summary).T.reset_index(drop=True)
        total_summary['birth_year_group'] = '合计'
        total_summary['是否在场'] = '合计'
        results.append(total_summary)
        
        # 添加空行
        empty_row = pd.DataFrame({'是否在场': [''], 'birth_year_group': ['']})
        results.append(empty_row)

        # 按是否在场分组并创建结果表
        for in_herd in merged_cow_key_traits_scores['是否在场'].unique():
            df_in_herd = merged_cow_key_traits_scores[merged_cow_key_traits_scores['是否在场'] == in_herd]
            result = df_in_herd.groupby('birth_year_group').apply(analyze_group).reset_index()
            result['是否在场'] = in_herd
            results.append(result)
            
            # 添加小计行
            subtotal = df_in_herd.pipe(analyze_group)
            subtotal = pd.DataFrame(subtotal).T.reset_index(drop=True)
            subtotal['birth_year_group'] = '合计'
            subtotal['是否在场'] = in_herd
            results.append(subtotal)
            
            # 添加空行
            results.append(empty_row)

        # 合并所有结果
        final_result = pd.concat(results, ignore_index=True)

        # 重新排序列
        final_result = final_result[['是否在场', 'birth_year_group', '头数'] + selected_traits]

        return final_result

    # 关键性状进展表做成关键性状的进展图-处理年份结尾不同
    def process_year(self, year_str):
        import re
        # 去掉"年"和空格
        year_str = year_str.replace('年', '').replace(' ', '')
        # 如果包含"及以前"，只保留年份
        if '及以前' in year_str:
            match = re.search(r'\d{4}', year_str)
            if match:
                year_str = match.group()
        try:
            year = int(year_str)
        except ValueError:
            year = None  # 或者根据需要进行处理
        return year

    # 关键性状进展表做成关键性状的进展图
    def generate_excel_charts(self, key_traits_annual_change, selected_traits,):
        # 对照组数据
        comparison_groups = {
            '参测牛': pd.DataFrame({
                '出生日期': ['2020', '2021', '2022', '2023'],
                'TPI': [ 2086, 2229, 2306, 2488],
                'NM$': [209, 355, 423, 560],
                'MILK': [153, 356, 471, 680],
                'FAT %': [0.08, 0.09, 0.09, 0.1],
                'FAT': [21.5, 34.59, 42.06, 60.48],
                'PROT%': [0.04, 0.03, 0.04, 0.04],
                'PROT': [12.87, 20.82, 25.03, 37.87],
                'PL': [1.06, 1.8, 2.17, 3.2],
                'PTAT': [-0.65, -0.38, -0.21, 0.2],
                'UDC': [-0.44, -0.47, 0.01, 0.4],
                'FLC': [-0.54, -0.49, -0.4, -0.15],
                'ST': [-0.67, -0.14, -0.35, -0.15]
            }),
            '前20%参测牛': pd.DataFrame({
                '出生日期': ['2020年', '2021年', '2022年', '2023年'],
                'TPI': [2407, 2501, 2650, 2827],
                'NM$': [522, 623, 753, 899],
                'MILK': [1051, 1201, 1348, 1593],
                'FAT %': [0.29, 0.26, 0.21, 0.22],
                'FAT': [57.87, 68.13, 79.78, 95.97],
                'PROT%': [0.13, 0.1, 0.09, 0.09],
                'PROT': [35.53, 41.85, 48.71, 57.85],
                'PL': [3.31, 3.84, 4.44, 5.27],
                'PTAT': [0.4, 0.57, 0.83, 1.14],
                'UDC': [0.62, 0.8, 1.02, 1.32]
            }),
            '国内某牧场': pd.DataFrame({
                '出生日期': ['2020及以前年份', '2021', '2022', '2023', '2024'],
                'TPI': [1686, 1991, 2167, 2078, 2281],
                'NM$': [380, 455.19, 529.52, 583.28, 631.80],
                'MILK': [320, 499.45, 803.10, 891.25, 1004.94],
                'FAT %': [0.05, 0.07, 0.09, 0.11, 0.12],
                'FAT': [35.32, 43.22, 54.32, 62.83, 72.23],
                'PROT%': [0.03, 0.03, 0.04, 0.05, 0.057],
                'PROT': [28.22, 31.09, 42.34, 49.48, 50.12],
                'SCS': [2.78, 2.8, 2.77, 2.67, 2.76],
                'PL': [2.3, 2.9, 3.6, 4, 4.2],
                'DPR': [0, 0.1, 0.1, 0.2, 0.28],
                'FLC': [-0.3, -0.16, -0.2, 0, 0.12],
                'UDC': [0.12, 0.42, 0.30, 0.31, 0.36],
                'ST': [-0.35, -0.4, -0.4, -0.43, -0.42]
            })
        }

        # 处理年份标签
        key_traits_annual_change['birth_year_group'] = key_traits_annual_change['birth_year_group'].apply(self.process_year)
        key_traits_annual_change = key_traits_annual_change.dropna(subset=['birth_year_group'])
        key_traits_annual_change['birth_year_group'] = key_traits_annual_change['birth_year_group'].astype(int)

        for group in comparison_groups.values():
            group['出生日期'] = group['出生日期'].apply(self.process_year)
            group.dropna(subset=['出生日期'], inplace=True)
            group['出生日期'] = group['出生日期'].astype(int)

        # 只保留"是否在场"为"是"的数据
        key_traits_annual_change = key_traits_annual_change[key_traits_annual_change['是否在场'] == '是']

        # 移除"合计"行
        key_traits_annual_change = key_traits_annual_change[key_traits_annual_change['birth_year_group'] != '合计']

        # 获取本牧场数据的最小和最大年份
        min_year = key_traits_annual_change['birth_year_group'].min()
        max_year = key_traits_annual_change['birth_year_group'].max()
        years = list(range(min_year, max_year + 1))

        # 创建输出文件夹和Excel文件
        excel_file = os.path.join(self.output_folder, "结果-牛群关键性状进展表.xlsx")
        charts_folder = os.path.join(self.output_folder, "结果-牛群关键性状进展图")
        os.makedirs(charts_folder, exist_ok=True)

        # 设置matplotlib的中文字体
        plt.rcParams['font.sans-serif'] = ['SimHei']  # 使用黑体
        plt.rcParams['axes.unicode_minus'] = False  # 解决负号显示问题

        # 将数据写入Excel并生成图表
        with pd.ExcelWriter(excel_file) as writer:
            for trait in selected_traits:
                if trait not in key_traits_annual_change.columns:
                    print(f"Warning: Trait '{trait}' not found in data. Skipping.")
                    continue

                # 准备数据
                data = pd.DataFrame({'年份': years})
                data.set_index('年份', inplace=True)

                # 本牧场数据
                farm_data = key_traits_annual_change[['birth_year_group', trait]].set_index('birth_year_group')
                data = data.join(farm_data.rename(columns={trait: '本牧场'}), how='left')

                # 对照组数据
                for group_name, group_df in comparison_groups.items():
                    if trait in group_df.columns:
                        group_trait_data = group_df[['出生日期', trait]].set_index('出生日期')
                        data = data.join(group_trait_data.rename(columns={trait: group_name}), how='left')

                # 将索引重置为列
                data.reset_index(inplace=True)

                # 按年份排序
                data.sort_values(by='年份', inplace=True)

                # 写入Excel
                data.to_excel(writer, sheet_name=trait, index=False)

                # 计算Y轴范围
                all_values = data.drop(columns=['年份']).values.flatten()
                all_values = all_values[~np.isnan(all_values)]
                max_value = np.max(all_values)
                min_value = np.min(all_values)
                range_padding = 0.2 * (max_value - min_value) if max_value != min_value else 1
                y_axis_max = max_value + range_padding
                y_axis_min = min_value - range_padding if min_value < 0 else max(0, min_value - range_padding)

                # 绘制图表
                plt.figure(figsize=(12, 4.5))

                # 绘制本牧场数据
                plt.plot(data['年份'], data['本牧场'], label='本牧场', linewidth=2.25, marker='o')

                # 显示本牧场数据标签
                for x, y in zip(data['年份'], data['本牧场']):
                    if not np.isnan(y):
                        plt.text(x, y, f"{y:.2f}", fontsize=12, ha='center', va='bottom')

                # 绘制对照组数据
                for group_name in comparison_groups.keys():
                    if group_name in data.columns:
                        plt.plot(data['年份'], data[group_name], label=group_name, linestyle='--', marker='o')

                        # 显示对照组数据标签
                        for x, y in zip(data['年份'], data[group_name]):
                            if not np.isnan(y):
                                plt.text(x, y, f"{y:.2f}", fontsize=12, ha='center', va='bottom')

                translation_dict = {
                    'TPI': '育种综合指数', 'NM$': '净利润值', 'CM$': '奶酪净值', 'FM$': '液奶净值', 'GM$': '放牧净值',
                    'MILK': '产奶量', 'FAT': '乳脂量', 'PROT': '乳蛋白量', 'FAT %': '乳脂率', 'PROT%': '乳蛋白率',
                    'SCS': '体细胞指数', 'DPR': '女儿怀孕率', 'HCR': '青年牛受胎率', 'CCR': '成母牛受胎率', 'PL': '生产寿命',
                    'SCE': '配种产犊难易度', 'DCE': '女儿产犊难易度', 'SSB': '配种死胎率', 'DSB': '女儿死淘率',
                    'PTAT': '体型综合指数', 'UDC': '乳房综合指数', 'FLC': '肢蹄综合指数', 'BDC': '体重综合指数',
                    'ST': '体高', 'SG': '强壮度', 'BD': '体深', 'DF': '乳用特征', 'RA': '尻角度', 'RW': '尻宽',
                    'LS': '后肢侧视', 'LR': '后肢后视', 'FA': '蹄角度', 'FLS': '肢蹄评分', 'FU': '前方附着',
                    'UH': '后房高度', 'UW': '后方宽度', 'UC': '悬韧带', 'UD': '乳房深度', 'FT': '前乳头位置',
                    'RT': '后乳头位置', 'TL': '乳头长度', 'FE': '饲料效率指数', 'FI': '繁殖力指数', 'HI': '健康指数',
                    'LIV': '存活率', 'GL': '妊娠长度', 'MAST': '乳房炎抗病性', 'MET': '子宫炎抗病性', 'RP': '胎衣不下抗病性',
                    'KET': '酮病抗病性', 'DA': '变胃抗病性', 'MFV': '产后瘫抗病性', 'EFC': '首次产犊月龄',
                    'HLiv': '后备牛存活率', 'FS': '饲料节约指数', 'RFI': '剩余饲料采食量', 'Milk Speed': '排乳速度',
                    'Eval Date': '数据日期'
                }


                # 设置标题和标签
                chinese_trait = translation_dict.get(trait, trait)  # 如果没有翻译，就使用原来的trait
                plt.title(f"{chinese_trait} ({trait}) 进展情况", fontsize=18, fontweight='bold')
                """plt.title(f"{trait}进展情况", fontsize=18, fontweight='bold')"""
                plt.xlabel('出生年份', fontsize=14)
                plt.ylabel(trait, fontsize=14)
                plt.legend(fontsize=12)
                plt.ylim(y_axis_min, y_axis_max)
                plt.grid(True, linestyle='--', alpha=0.7)
                plt.tight_layout()

                # **设置横坐标刻度为整数年份**
                plt.xticks(years)  # 设置 x 轴刻度为年份列表
                plt.gca().xaxis.set_major_locator(ticker.FixedLocator(years))  # 确保刻度定位器使用整数年份

                plt.tight_layout()

                # 保存图表为 PNG 图片
                plt.savefig(os.path.join(charts_folder, f"{trait}进展情况.png"), dpi=300, bbox_inches='tight')
                plt.close()

        print("Excel charts and image files generated successfully.")

    # 生成表C sire\ mgs\ mmgs 各性状的值
    def generate_table_c(self, cow_df, bull_key_traits_scores, selected_traits):
        """print("Input data for generate_table_c:")
        print("cow_df:")
        print(cow_df.head())
        print("bull_key_traits_scores:")
        print(bull_key_traits_scores.head())
        print("selected_traits:")
        print(selected_traits)"""
    
        # 保留 cow_df 中的所有列
        table_c = cow_df.copy()

        # 提取出生年份
        table_c['birth_year'] = table_c['birth_date'].dt.year

        # 为每个祖先（sire, mgs, mmgs）创建性状列
        for ancestor in ['sire', 'mgs', 'mmgs']:
            for trait in selected_traits:
                table_c[f'{ancestor}_{trait}'] = table_c[ancestor].apply(
                    lambda x: self.find_bull_key_trait_score(x, bull_key_traits_scores, trait) if pd.notna(x) and str(x).strip() else np.nan
                )

        # 创建识别情况列，基于NM$的结果
        for ancestor in ['sire', 'mgs', 'mmgs']:
            table_c[f'{ancestor}_identified'] = table_c[f'{ancestor}_NM$'].apply(
                lambda x: '已识别' if pd.notna(x) else '未识别'
            )

        # 重新排列列的顺序
        original_columns = cow_df.columns.tolist()
        new_columns = ['birth_year', 'sire_identified', 'mgs_identified', 'mmgs_identified']
        for ancestor in ['sire', 'mgs', 'mmgs']:
            new_columns.extend([f'{ancestor}_{trait}' for trait in selected_traits])
        
        # 最终列顺序：原始列 + 新增列
        final_columns = original_columns + [col for col in new_columns if col not in original_columns]
        table_c = table_c[final_columns]

        return table_c

    # 生成"在群牛只净利润值分布.xlsx"-“merged_cow_key_traits_scores” -来自“def calculate_key_traits_index(self):”
    def generate_nm_distribution_charts(self, merged_cow_key_traits_scores):
        # 数据处理
        df = merged_cow_key_traits_scores
        
        # 添加“是否在场”限制条件，排除“是否在场”为“否”的记录
        df = df[df['是否在场'] != '否']


        if 'NM$' not in df.columns:
            raise ValueError("NM$ column not found in the dataframe")
        
        df['NM$'] = pd.to_numeric(df['NM$'], errors='coerce')
        
        # 动态确定区间
        min_value = df['NM$'].min()
        max_value = df['NM$'].max()
        
        # 根据数据范围确定步长
        range_value = max_value - min_value
        step = 300  # 默认步长
        
        if range_value > 3000:
            step = 500
        elif range_value <= 1500:
            step = 200
        
        # 创建区间
        lower_bound = np.floor(min_value / step) * step
        upper_bound = np.ceil(max_value / step) * step
        bins = np.arange(lower_bound, upper_bound + step, step)
        
        # 创建标签
        labels = [f'{bins[i]:.0f}-{bins[i+1]:.0f}' for i in range(len(bins)-1)]
        
        df['NM$_group'] = pd.cut(df['NM$'], bins=bins, labels=labels, include_lowest=True)
        nm_distribution = df['NM$_group'].value_counts().sort_index()
        
        nm_percentages = nm_distribution / len(df) * 100
        
        # 创建Excel文件
        excel_file = os.path.join(self.output_folder, "在群牛只净利润值分布.xlsx")
        workbook = xlsxwriter.Workbook(excel_file)
        worksheet = workbook.add_worksheet("NM$分布")
        
        # 写入数据
        data = [
            ['NM$区间', '头数', '占比'],
            *[[label, count, percent] for label, count, percent in 
            zip(nm_distribution.index, nm_distribution.values, nm_percentages)]
        ]
        
        for row, row_data in enumerate(data):
            for col, cell_data in enumerate(row_data):
                if col == 2 and row != 0:  # 对占比列进行特殊处理
                    worksheet.write(row, col, cell_data / 100, workbook.add_format({'num_format': '0.00%'}))
                else:
                    worksheet.write(row, col, cell_data)
        
        # 创建柱状图
        chart1 = workbook.add_chart({'type': 'column'})
        chart1.add_series({
            'name': '头数',
            'categories': ['NM$分布', 1, 0, len(data)-1, 0],
            'values': ['NM$分布', 1, 1, len(data)-1, 1],
            'data_labels': {'value': True},
        })
        chart1.set_title({'name': '在群牛只净利润值分布'})
        chart1.set_x_axis({'name': 'NM$ 区间'})
        chart1.set_y_axis({'name': '牛只数量'})
        worksheet.insert_chart('E1', chart1, {'x_scale': 1.5, 'y_scale': 1.5})
        
        # 创建饼图
        chart2 = workbook.add_chart({'type': 'pie'})
        chart2.add_series({
            'name': '占比',
            'categories': ['NM$分布', 1, 0, len(data)-1, 0],
            'values': ['NM$分布', 1, 2, len(data)-1, 2],
            'data_labels': {'percentage': True, 'leader_lines': True},
        })
        chart2.set_title({'name': '在群牛只不同净利润区间占比'})
        worksheet.insert_chart('E18', chart2, {'x_scale': 1.5, 'y_scale': 1.5})
        
        workbook.close()
        
        logging.info(f"NM$分布图表已保存至 {excel_file}")


    # 系谱识别情况表
    def generate_pedigree_identification_analysis(self, table_c, cow_df):
        merged_df = table_c

        # 添加性别限制，排除性别为“公”的记录
        merged_df = merged_df[merged_df['sex'] != '公']

        # 获取牛群中的最大出生年份
        max_birth_year = merged_df['birth_year'].max()

        # 创建出生年份分组
        merged_df['birth_year_group'] = pd.cut(
            merged_df['birth_year'],
            bins=[-float('inf')] + list(range(max_birth_year-4, max_birth_year)) + [float('inf')],
            labels=[f'{max_birth_year-4}年及以前'] + [str(year) for year in range(max_birth_year-3, max_birth_year+1)]
        )
        """print("max_birth_year:", max_birth_year)"""
        # 创建分析函数
        def analyze_group(group):
            total = len(group)
            sire_identified = (group['sire_identified'] == '已识别').sum()
            mgs_identified = (group['mgs_identified'] == '已识别').sum()
            mmgs_identified = (group['mmgs_identified'] == '已识别').sum()
            return pd.Series({
                '头数': total,
                '父号可识别头数': sire_identified,
                '父号识别率': sire_identified / total if total > 0 else 0,
                '外祖父可识别头数': mgs_identified,
                '外祖父识别率': mgs_identified / total if total > 0 else 0,
                '外曾外祖父可识别头数': mmgs_identified,
                '外曾外祖父识别率': mmgs_identified / total if total > 0 else 0,
            })

        # 总体分析
        results = []
        total_result = merged_df.groupby('birth_year_group').apply(analyze_group).reset_index()
        total_result['是否在场'] = '总计'
        results.append(total_result)
        
        # 添加总计的合计行
        total_summary = merged_df.pipe(analyze_group)
        total_summary = pd.DataFrame(total_summary).T.reset_index(drop=True)
        total_summary['birth_year_group'] = '合计'
        total_summary['是否在场'] = '总计'
        results.append(total_summary)
        
        # 添加空行
        empty_row = pd.DataFrame({'是否在场': [''], 'birth_year_group': ['']})
        results.append(empty_row)

        # 按是否在场分组并创建结果表
        for in_herd in merged_df['是否在场'].unique():
            df_in_herd = merged_df[merged_df['是否在场'] == in_herd]
            result = df_in_herd.groupby('birth_year_group').apply(analyze_group).reset_index()
            result['是否在场'] = in_herd
            results.append(result)
            
            # 添加小计行
            subtotal = df_in_herd.pipe(analyze_group)
            subtotal = pd.DataFrame(subtotal).T.reset_index(drop=True)
            subtotal['birth_year_group'] = '合计'
            subtotal['是否在场'] = in_herd
            results.append(subtotal)
            
            # 添加空行
            results.append(empty_row)

        # 合并所有结果
        final_result = pd.concat(results, ignore_index=True)

        # 重新排序列
        final_result = final_result[[
            '是否在场', 'birth_year_group', '头数', 
            '父号可识别头数', '父号识别率', 
            '外祖父可识别头数', '外祖父识别率', 
            '外曾外祖父可识别头数', '外曾外祖父识别率'
        ]]

        # 格式化百分比列
        percentage_columns = ['父号识别率', '外祖父识别率', '外曾外祖父识别率']
        for col in percentage_columns:
            final_result[col] = final_result[col].apply(lambda x: f"{x:.2%}" if pd.notnull(x) else '')

        return final_result
    
    def generate_table_d(self, table_c, selected_traits, bull_key_traits_scores):
        table_d = {}

        for trait in selected_traits:
            logging.info(f"开始处理特征: {trait}")
            
            # 将 'decimal.Decimal' 类型的数据转换为 'float' 类型，并确保 birth_year 为整数类型
            table_c[f'sire_{trait}'] = table_c[f'sire_{trait}'].astype(float)
            table_c['birth_year'] = table_c['birth_year'].astype(int)
            
            # 处理可能的空值
            table_c_clean = table_c.dropna(subset=['birth_year', f'sire_{trait}'])
            
            # 按出生年份分组,计算每年父号指数的平均值
            trait_data = table_c_clean.groupby('birth_year')[f'sire_{trait}'].agg(['count', 'mean']).reset_index()
            
            logging.info(f"特征 {trait} - 年份和数据量: \n{trait_data}")

            # 筛选出数据量大于等于10的年份
            trait_data = trait_data[trait_data['count'] >= 10]

            # 获取母牛群的最小出生年份和最大年份
            min_cow_birth_year = int(self.cow_df['birth_date'].dt.year.min())
            max_year = max(int(trait_data['birth_year'].max() if not trait_data.empty else min_cow_birth_year), datetime.datetime.now().year)

            logging.info(f"特征 {trait} - 最小出生年份: {min_cow_birth_year}, 最大年份: {max_year}")

            # 获取999HO99999的值作为默认值
            default_value = self.find_bull_key_trait_score('999HO99999', bull_key_traits_scores, trait)

            # 如果只有一个数据点或没有数据点，使用特殊逻辑
            if len(trait_data) <= 1:
                logging.warning(f"特征 {trait} 数据点不足，使用特殊逻辑。")
                min_cow_birth_year = 2010
                trait_data = pd.DataFrame({'birth_year': [2010], 'count': [1], 'mean': [default_value]})

            # 创建完整的年份范围
            full_year_range = pd.DataFrame({'birth_year': range(min_cow_birth_year, max_year + 1)})
            trait_data = pd.merge(full_year_range, trait_data, on='birth_year', how='left')

            # 添加一列来标记缺失值
            trait_data['is_missing'] = trait_data['mean'].isna()

            # 找出非空的年份和对应的值
            known_data = trait_data[~trait_data['mean'].isna()]
            known_years = known_data['birth_year'].values.reshape(-1, 1)
            known_values = known_data['mean'].values

            logging.info(f"特征 {trait} - 已知年份: {known_years.flatten()}")
            logging.info(f"特征 {trait} - 已知值: {known_values}")

            if len(known_years) >= 2:
                try:
                    # 创建并训练线性回归模型
                    model = LinearRegression()
                    model.fit(known_years, known_values)

                    # 预测缺失年份的值
                    missing_years = trait_data.loc[trait_data['mean'].isna(), 'birth_year'].values.reshape(-1, 1)
                    if len(missing_years) > 0:
                        predicted_values = model.predict(missing_years)
                        trait_data.loc[trait_data['mean'].isna(), 'mean'] = predicted_values
                    logging.info(f"特征 {trait} - 线性回归完成")
                except Exception as e:
                    logging.error(f"特征 {trait} 在进行线性回归时发生错误: {str(e)}")
                    # 如果线性回归失败，使用已知值的平均值填充
                    mean_value = known_values.mean()
                    trait_data.loc[trait_data['mean'].isna(), 'mean'] = mean_value
                    logging.warning(f"特征 {trait} - 使用平均值 {mean_value} 填充缺失值")
            else:
                # 如果没有足够的数据点进行回归，使用默认值填充
                trait_data['mean'] = trait_data['mean'].fillna(default_value)
                logging.warning(f"特征 {trait} - 使用默认值 {default_value} 填充所有缺失值")

            # 添加一列来标记插值的结果
            trait_data['interpolated'] = trait_data['is_missing'].map({True: 'Yes', False: 'No'})

            # 重新排列列的顺序
            trait_data = trait_data[['birth_year', 'mean', 'interpolated']]

            table_d[trait] = trait_data

            logging.info(f"特征 {trait} 处理完成")

        return table_d    
    
    def update_missing_bulls(self, missing_bulls, source):
        try:
            connection = pymysql.connect(
                host=cipher_suite.decrypt(encoded_host).decode(),
                port=int(cipher_suite.decrypt(encoded_port).decode()),
                user=cipher_suite.decrypt(encoded_user).decode(),
                password=cipher_suite.decrypt(encoded_password).decode(),
                database=cipher_suite.decrypt(encoded_db).decode(),
                charset='utf8mb4',
                cursorclass=pymysql.cursors.DictCursor
            )

            with connection.cursor() as cursor:
                for bull_id in missing_bulls:
                    sql = "INSERT INTO miss_bull (bull, source, time, user) VALUES (%s, %s, %s, %s)"
                    cursor.execute(sql, (bull_id, source, datetime.datetime.now(), self.username))

            connection.commit()
            connection.close()

        except Exception as e:
            logging.error(f"更新缺失公牛信息时发生错误: {e}")
            QMessageBox.critical(self, "错误", f"更新缺失公牛信息时发生错误：{e},请联系管理员。")

    def calculate_cow_key_traits_scores(self, cow_df, bull_key_traits_scores, table_d, selected_traits):
        cow_key_traits_scores = pd.DataFrame(columns=['Cow ID'] + selected_traits)
        cow_key_traits_scores_pedigree = pd.DataFrame(columns=['Cow ID'] + selected_traits)
        default_bull_id = '999HO99999'
        
        # 读取基因组数据
        geno_library_file = os.path.join(self.output_folder, "geno_library.xlsx")
        if os.path.isfile(geno_library_file):
            geno_library_df = pd.read_excel(geno_library_file)
            geno_library_df['Farm ID'] = geno_library_df['Farm ID'].apply(self.standardize_cow_id)
            geno_library_df.set_index('Farm ID', inplace=True)
        else:
            geno_library_df = None
        
        # 读取 "结果-公牛关键性状指数.xlsx" 文件
        default_bull_traits = pd.read_excel(os.path.join(self.output_folder, "结果-公牛关键性状指数.xlsx"))
        default_bull_traits = default_bull_traits[default_bull_traits['BULL NAAB'] == default_bull_id].iloc[0] if not default_bull_traits[default_bull_traits['BULL NAAB'] == default_bull_id].empty else None

        total_cows = len(cow_df)
        progress_dialog = TqdmProgressDialog(total_cows, "计算母牛关键性状指数", self)
        progress_dialog.show()

        try:
            for _, cow in cow_df.iterrows():
                score_dict = {'Cow ID': cow['cow_id']}
                score_dict_pedigree = {'Cow ID': cow['cow_id']}
                ancestor_weights = {'sire': 0.5, 'mgs': 0.25, 'mmgs': 0.125}

                for trait in selected_traits:
                    score = 0
                    for ancestor_type, weight in ancestor_weights.items():
                        ancestor_id = cow[ancestor_type]
                        ancestor_score = self.calculate_ancestor_score(cow, ancestor_id, ancestor_type, bull_key_traits_scores, table_d, trait, default_bull_id)
                        score += float(ancestor_score) * weight

                    # 添加默认得分的权重
                    default_score = self.find_bull_key_trait_score(default_bull_id, bull_key_traits_scores, trait)
                    score += float(default_score) * 0.125

                    # 如果计算结果仍为空值，使用替代方法
                    if pd.isna(score):
                        score = self.get_alternative_score(cow, trait, table_d, default_bull_traits)

                    # 保存系谱计算结果
                    score_dict_pedigree[trait] = score

                    # 如果有基因组数据，则使用基因组数据替换计算得到的指数
                    standardized_cow_id = self.standardize_cow_id(cow['cow_id'])
                    if geno_library_df is not None:
                        geno_index_lower = [idx.lower() for idx in geno_library_df.index]
                        if standardized_cow_id.lower() in geno_index_lower:
                            matched_index = geno_library_df.index[geno_index_lower.index(standardized_cow_id.lower())]
                            geno_score = geno_library_df.loc[matched_index, trait]
                            if pd.notna(geno_score):
                                score = geno_score
                                logging.info(f"Cow ID {standardized_cow_id} matched with genomic data for trait {trait}")
                        else:
                            logging.warning(f"No genomic data found for Cow ID {standardized_cow_id}")

                    score_dict[trait] = score

                cow_key_traits_scores = pd.concat([cow_key_traits_scores, pd.DataFrame([score_dict])], ignore_index=True)
                cow_key_traits_scores_pedigree = pd.concat([cow_key_traits_scores_pedigree, pd.DataFrame([score_dict_pedigree])], ignore_index=True)
                
                progress_dialog.update_progress()
                QApplication.processEvents()
                
                if progress_dialog.wasCanceled():
                    break

        finally:
            progress_dialog.close()

        # 保存系谱计算结果
        pedigree_result_file = os.path.join(self.output_folder, "结果-母牛关键性状指数（系谱计算结果）.xlsx")
        self.save_to_excel(cow_key_traits_scores_pedigree, pedigree_result_file)

        # 保存最终结果（包含基因组数据）
        # final_result_file = os.path.join(self.output_folder, "结果-母牛关键性状指数.xlsx")
        # self.save_to_excel(cow_key_traits_scores, final_result_file)

        return cow_key_traits_scores

    def save_to_excel(self, df, file_path):
        # 在写入Excel前处理Inf和NaN值
        df = df.replace([np.inf, -np.inf], np.nan)
        df = df.fillna("")  # 或者使用其他适当的值填充NaN
        
        with pd.ExcelWriter(file_path, engine='openpyxl') as writer:
            df.to_excel(writer, index=False)
            workbook = writer.book
            worksheet = writer.sheets['Sheet1']
            for idx, cell in enumerate(worksheet['A']):
                cell.number_format = '@'  # 设置为文本格式



    def get_alternative_score(self, cow, trait, table_d, default_bull_traits):
        # 首先尝试使用母牛出生年份在表d中查找对应性状值
        birth_year = cow['birth_date'].year
        year_score = table_d[trait].loc[table_d[trait]['birth_year'] == birth_year, 'mean'].values
        if len(year_score) > 0:
            return float(year_score[0])
        
        # 如果表d中没有，则使用默认公牛的性状值
        if default_bull_traits is not None and trait in default_bull_traits:
            return float(default_bull_traits[trait])
        
        # 如果还是没有，返回0或其他适当的默认值
        logging.warning(f"无法为母牛 {cow['cow_id']} 的 {trait} 特征找到有效的分数。使用默认值0。")
        return 0

    def calculate_ancestor_score(self, cow, ancestor_id, ancestor_type, bull_key_traits_scores, table_d, trait, default_bull_id):
        if pd.isna(ancestor_id) or (isinstance(ancestor_id, str) and ancestor_id.strip() == ''):
            if ancestor_type == 'sire':
                birth_year = cow['birth_date'].year
            elif ancestor_type == 'mgs':
                dam_id = cow['dam']
                birth_year = self.cow_df.loc[self.cow_df['cow_id'] == dam_id, 'birth_date'].dt.year.iloc[0] if dam_id in self.cow_df['cow_id'].values else None
            else:  # mmgs
                dam_id = cow['dam']
                if dam_id in self.cow_df['cow_id'].values:
                    dam_mother_id = self.cow_df.loc[self.cow_df['cow_id'] == dam_id, 'dam'].iloc[0]
                    birth_year = self.cow_df.loc[self.cow_df['cow_id'] == dam_mother_id, 'birth_date'].dt.year.iloc[0] if dam_mother_id in self.cow_df['cow_id'].values else None
                else:
                    birth_year = None

            if birth_year is not None:
                mean_score = table_d[trait].loc[table_d[trait]['birth_year'] == birth_year, 'mean'].values
                # 检查 mean_score 是否为空，并处理
                if len(mean_score) > 0:
                    mean_score = float(mean_score[0])
                else:
                    mean_score = float(self.find_bull_key_trait_score(default_bull_id, bull_key_traits_scores, trait))

                default_score = float(self.find_bull_key_trait_score(default_bull_id, bull_key_traits_scores, trait))
                # 计算加权得分
                ancestor_score = 0.5 * mean_score + 0.5 * default_score
                # 因为算出的值比实际较高，预估是由于赋的值过高所致，因此使用999HO99999和不同年份均值（默认值）的均值进行尝试。
                ancestor_score = Decimal(ancestor_score)  # 转换为 decimal.Decimal
            else:
                ancestor_score = Decimal(self.find_bull_key_trait_score(default_bull_id, bull_key_traits_scores, trait))
        else:
            ancestor_score = self.find_bull_key_trait_score(ancestor_id, bull_key_traits_scores, trait)
            if pd.isna(ancestor_score):
                if ancestor_type == 'sire':
                    birth_year = cow['birth_date'].year
                elif ancestor_type == 'mgs':
                    dam_id = cow['dam']
                    birth_year = self.cow_df.loc[self.cow_df['cow_id'] == dam_id, 'birth_date'].dt.year.iloc[0] if dam_id in self.cow_df['cow_id'].values else None
                else:  # mmgs
                    dam_id = cow['dam']
                    if dam_id in self.cow_df['cow_id'].values:
                        dam_mother_id = self.cow_df.loc[self.cow_df['cow_id'] == dam_id, 'dam'].iloc[0]
                        birth_year = self.cow_df.loc[self.cow_df['cow_id'] == dam_mother_id, 'birth_date'].dt.year.iloc[0] if dam_mother_id in self.cow_df['cow_id'].values else None
                    else:
                        birth_year = None

                if birth_year is not None:
                    mean_score = table_d[trait].loc[table_d[trait]['birth_year'] == birth_year, 'mean'].values
                    if len(mean_score) > 0:
                        mean_score = float(mean_score[0])
                    else:
                        mean_score = float(self.find_bull_key_trait_score(default_bull_id, bull_key_traits_scores, trait))

                    default_score = float(self.find_bull_key_trait_score(default_bull_id, bull_key_traits_scores, trait))
                    ancestor_score = 0.5 * mean_score + 0.5 * default_score
                    # 因为算出的值比实际较高，预估是由于赋的值过高所致，因此使用999HO99999和不同年份均值（默认值）的均值进行尝试。
                    ancestor_score = Decimal(ancestor_score)  # 转换为 decimal.Decimal
                else:
                    ancestor_score = Decimal(self.find_bull_key_trait_score(default_bull_id, bull_key_traits_scores, trait))

        return ancestor_score

    def find_bull_key_trait_score(self, bull_id, bull_key_traits_scores, trait):
        """
        在公牛关键性状指数DataFrame中查找指定公牛的关键性状指数。

        :param bull_id: 要查找的公牛ID
        :param bull_key_traits_scores: 包含公牛关键性状指数的DataFrame
        :param trait: 要查找的关键性状
        :return: 公牛的关键性状指数,如果找不到则返回np.nan
        """
        if pd.isna(bull_id) or (isinstance(bull_id, str) and bull_id.strip() == ''):
            return np.nan

        bull_id = str(bull_id)

        if len(bull_id) > 11:
            score = bull_key_traits_scores[bull_key_traits_scores['BULL REG'] == bull_id][trait].values
        else:
            score = bull_key_traits_scores[bull_key_traits_scores['BULL NAAB'] == bull_id][trait].values
            if len(score) == 0:
                for rule, replacement in self.bull_rules.items():
                    if bull_id.startswith(rule):
                        new_bull_id = replacement + bull_id[3:]
                        score = bull_key_traits_scores[bull_key_traits_scores['BULL NAAB'] == new_bull_id][trait].values
                        if len(score) > 0:
                            break
                    elif bull_id.startswith(replacement):
                        new_bull_id = rule + bull_id[3:]
                        score = bull_key_traits_scores[bull_key_traits_scores['BULL NAAB'] == new_bull_id][trait].values
                        if len(score) > 0:
                            break

        return score[0] if len(score) > 0 else np.nan

    def generate_selection_files(self):
        try:
            # 读取隐性基因筛查结果
            screening_result_file = os.path.join(self.output_folder, "隐性基因筛查结果.xlsx")
            screening_result_df = pd.read_excel(screening_result_file)

            # 读取近交系数计算结果
            inbreeding_result_file = os.path.join(self.output_folder, "近交系数计算结果.xlsx")
            inbreeding_result_df = pd.read_excel(inbreeding_result_file)

            # 读取指数排名结果
            ranking_result_file = os.path.join(self.output_folder, "结果-指数排名结果.xlsx")
            ranking_result_df = pd.read_excel(ranking_result_file)

            # 生成选配完整报告表
            complete_report_df = self.generate_complete_report(screening_result_df, inbreeding_result_df)

            if complete_report_df is None:
                return

            # 生成密钥并保存到文件（只需运行一次）
            key_file_path = os.path.join(self.output_folder, 'natiommont-oasthandler.key')
            if not os.path.exists(key_file_path):
                key = Fernet.generate_key()
                with open(key_file_path, 'wb') as key_file:
                    key_file.write(key)
            else:
                with open(key_file_path, 'rb') as key_file:
                    key = key_file.read()

            # 生成加密器
            cipher = Fernet(key)

            # 生成 Excel 文件字节流
            with BytesIO() as excel_buffer:
                complete_report_df.to_excel(excel_buffer, index=False, engine='openpyxl')
                excel_bytes = excel_buffer.getvalue()

            # 加密 Excel 内容
            encrypted_data = cipher.encrypt(excel_bytes)

            # 保存加密文件
            complete_report_file = os.path.join(self.output_folder, "sty_14.enc")
            with open(complete_report_file, 'wb') as file:
                file.write(encrypted_data)

            '''# 同时保存未加密的Excel文件,用于检查
            unencrypted_report_file = os.path.join(self.output_folder, "选配初始文件(未加密).xlsx")
            complete_report_df.to_excel(unencrypted_report_file, index=False)'''

            self.status_label.setText(self.status_label.text() + "\n选配初始文件生成完成。")


            # 生成额外的".key"文件和".enc"文件用于混淆视听
            def generate_random_filename(extension, base_name):
                random_part = ''.join(random.choices(string.ascii_letters + string.digits, k=10))
                return f'{base_name}_{random_part}{extension}'

            key_base_names = [
                'national-timezones',
                'shellcommon-broker',
                'switch-toasthandler',
                'system-userprofile',
                'taskflow-dataengine',
                'timezones.resources',
                'tional-chinese-core'
            ]

            enc_base_names = [
                'system_64', 'sys_64', 'data_file', 'dom_lac',
                'profile', 'backup_data', 'resource_pic',
                'config_data', 'temp_file', 'log_cec',
                'cache_data', 'misc_file'
            ]

            for base_name in key_base_names:
                fake_key = Fernet.generate_key()
                fake_key_path = os.path.join(self.output_folder, generate_random_filename('.key', base_name))
                with open(fake_key_path, 'wb') as fake_key_file:
                    fake_key_file.write(fake_key)
            
            for base_name in enc_base_names:
                fake_data = ''.join(random.choices(string.ascii_letters + string.digits, k=100)).encode()
                fake_encrypted_data = cipher.encrypt(fake_data)
                fake_enc_path = os.path.join(self.output_folder, generate_random_filename('.enc', base_name))
                with open(fake_enc_path, 'wb') as fake_enc_file:
                    fake_enc_file.write(fake_encrypted_data)

        except Exception as e:
            logging.error(f"生成选配文件时发生错误: {e}")
            QMessageBox.critical(self, "错误", f"生成选配文件时发生错误：{e},请按照模版上传母牛信息/备选公牛表/权重分配表。")

    def find_bull_score(self, bull_id):
        """
        在公牛指数得分DataFrame中查找指定公牛的指数得分。

        :param bull_id: 要查找的公牛ID
        :return: 公牛的指数得分,如果找不到则返回np.nan
        """
        if pd.isna(bull_id) or (isinstance(bull_id, str) and bull_id.strip() == ''):
            return np.nan

        bull_id = str(bull_id)
        if len(bull_id) > 11:
            return self.bull_scores_reg_dict.get(bull_id, np.nan)
        else:
            score = self.bull_scores_naab_dict.get(bull_id, np.nan)
            if np.isnan(score):
                for rule, replacement in self.bull_rules.items():
                    if bull_id.startswith(rule):
                        new_bull_id = replacement + bull_id[3:]
                        score = self.bull_scores_naab_dict.get(new_bull_id, np.nan)
                        if not np.isnan(score):
                            break
                    elif bull_id.startswith(replacement):
                        new_bull_id = rule + bull_id[3:]
                        score = self.bull_scores_naab_dict.get(new_bull_id, np.nan)
                        if not np.isnan(score):
                            break
            return score

    def generate_complete_report(self, screening_result_df, inbreeding_result_df):
        """print("screening_result_df 'Cow ID' 唯一性:", screening_result_df['Cow ID'].is_unique)
        print("inbreeding_result_df 'Cow' 唯一性:", inbreeding_result_df['Cow'].is_unique)"""
        # 根据公牛数据确定常规和性控冻精的种类数量
        regular_count, sex_controlled_count = self.get_semen_type_counts()

        # 创建报告表格的列名
        columns = ['在群母牛号']
        for i in range(1, regular_count + 1):
            columns += [f"推荐常规冻精{i}选", f"常规冻精{i}隐性基因情况", f"常规冻精{i}近交系数", f"常规冻精{i}得分"]
        for i in range(1, sex_controlled_count + 1):
            columns += [f"推荐性控冻精{i}选", f"性控冻精{i}隐性基因情况", f"性控冻精{i}近交系数", f"性控冻精{i}得分"]
        columns += ['breed', 'sire', 'mgs', 'dam', 'mmgs', 'lac', 'calving_date', 'birth_date', 'services_time',
                    'group', 'Combine Index Score']

        # 从 "指数排名结果.xlsx" 文件中读取数据
        ranking_result_file = os.path.join(self.output_folder, "结果-指数排名结果.xlsx")
        ranking_result_df = pd.read_excel(ranking_result_file)

        # 按'Combine Index Score'降序排列
        ranking_result_df.sort_values(by='Combine Index Score', ascending=False, inplace=True)

        # 创建一个包含所需列的DataFrame
        report_df = ranking_result_df[['cow_id', 'Combine Index Score']].rename(columns={'cow_id': '在群母牛号'})

        # 更新其他字段
        for column in ['breed', 'sire', 'mgs', 'dam', 'mmgs', 'lac', 'calving_date', 'birth_date', 'services_time',
                       'group']:
            report_df[column] = report_df['在群母牛号'].map(self.cow_df.set_index('cow_id')[column])

        # 加载公牛得分数据
        bull_scores_df = pd.read_excel(os.path.join(self.output_folder, "公牛指数.xlsx"))

        # 创建两个字典,分别存储 BULL NAAB 和 BULL REG 对应的得分
        bull_scores_naab_dict = bull_scores_df.set_index('BULL NAAB')['Index Score'].to_dict()
        bull_scores_reg_dict = bull_scores_df.set_index('BULL REG')['Index Score'].to_dict()

        # 加载近交系数数据
        inbreeding_dict = inbreeding_result_df.set_index(['Cow', 'Bull'])['Inbreeding Coefficient'].to_dict()

        # 加载隐性基因信息数据
        gene_info_dict = screening_result_df.set_index(['Cow ID', 'Bull ID'])

        missing_bulls = set()  # 用于记录未找到的公牛

        # 遍历报告DataFrame以填充得分
        for index, row in report_df.iterrows():
            cow_id = row['在群母牛号']
            cow_score = row['Combine Index Score']

            for semen_type in ['常规', '性控']:
                pairs_scores = []
                bulls = self.bull_df[self.bull_df['semen_type'] == semen_type]

                for bull_row in bulls.itertuples():
                    bull_id = getattr(bull_row, 'bull_id')
                    bull_score = self.find_bull_score(bull_id)

                    if np.isnan(bull_score):
                        missing_bulls.add(bull_id)
                    else:
                        offspring_score = 0.5 * (cow_score + bull_score)
                        inbreeding_coeff = inbreeding_dict.get((cow_id, bull_id), "N/A")
                        try:
                            gene_info = gene_info_dict.loc[(cow_id, bull_id)].values.tolist()[2:]

                        except KeyError:
                            gene_info = []

                        pairs_scores.append((bull_id, offspring_score, inbreeding_coeff, gene_info))

                # 按 offspring_score 降序排列
                pairs_scores.sort(key=lambda x: x[1], reverse=True)

                # 填充报告
                for i, (bull_id, score, inbreeding_coeff, gene_info) in enumerate(pairs_scores, start=1):
                    if semen_type == '常规' and i <= regular_count:
                        report_df.at[index, f"推荐常规冻精{i}选"] = bull_id
                        report_df.at[
                            index, f"常规冻精{i}隐性基因情况"] = "NO safe" if "NO safe" in gene_info else "Safe"
                        report_df.at[index, f"常规冻精{i}近交系数"] = inbreeding_coeff
                        report_df.at[index, f"常规冻精{i}得分"] = score
                    elif semen_type == '性控' and i <= sex_controlled_count:
                        report_df.at[index, f"推荐性控冻精{i}选"] = bull_id
                        report_df.at[
                            index, f"性控冻精{i}隐性基因情况"] = "NO safe" if "NO safe" in gene_info else "Safe"
                        report_df.at[index, f"性控冻精{i}近交系数"] = inbreeding_coeff
                        report_df.at[index, f"性控冻精{i}得分"] = score

        # 重新排列列的顺序
        new_columns = ['在群母牛号'] + [col for col in report_df.columns if '常规' in col] + [col for col in
                                                                                              report_df.columns if
                                                                                              '性控' in col] + [col for
                                                                                                                col in
                                                                                                                report_df.columns
                                                                                                                if
                                                                                                                col not in [
                                                                                                                    '在群母牛号',
                                                                                                                    'Combine Index Score'] and '常规' not in col and '性控' not in col] + [
                          'Combine Index Score']
        report_df = report_df[new_columns]

        # 如果有未找到的公牛,弹出提示框
        if missing_bulls:
            missing_bulls_str = ', '.join(missing_bulls)
            msg_box = CustomMessageBox(f"以下备选公牛在数据库中未找到:\n\n {missing_bulls_str}\n\n\n请仔细检查！\n\n如确认无误，可联系负责人更新至数据库.")
            msg_box.exec()


        return report_df

    def get_semen_type_counts(self):
        # 分别计算性控和常规冻精的数量
        semen_types = self.bull_df['semen_type'].value_counts()
        regular_count = semen_types.get('常规', 0)  # 如果没有常规冻精,则默认为0
        sex_controlled_count = semen_types.get('性控', 0)  # 如果没有性控冻精,则默认为0
        return regular_count, sex_controlled_count

    def generate_selection_result(self, complete_report_df, selected_groups, inbreeding_threshold,
                                  control_defect_genes, regular_percentages, sexed_percentages):
        # 根据选择的分组筛选母牛
        selected_cows = complete_report_df[complete_report_df['group'].isin(selected_groups)]

        # 根据近交系数标准和隐性基因标准筛选候选公牛
        candidate_bulls = {}
        for _, row in selected_cows.iterrows():
            cow_id = row['在群母牛号']
            candidate_bulls[cow_id] = []
            for semen_type in ['常规', '性控']:
                for i in range(1, 4):
                    bull_id = row[f"推荐{semen_type}冻精{i}选"]
                    inbreeding_coeff = float(row[f"{semen_type}冻精{i}近交系数"][:-1]) / 100
                    defect_gene = row[f"{semen_type}冻精{i}隐性基因情况"] != "Safe"
                    if inbreeding_coeff <= inbreeding_threshold and (not control_defect_genes or not defect_gene):
                        candidate_bulls[cow_id].append((bull_id, semen_type, i))

        # 根据冻精百分比分配候选公牛
        selection_result = []
        for cow_id, candidates in candidate_bulls.items():
            result = {'母牛号': cow_id}
            for semen_type, percentages in [('常规', regular_percentages), ('性控', sexed_percentages)]:
                bulls = [bull for bull, bull_semen_type, _ in candidates if bull_semen_type == semen_type]
                total_percentage = sum(percentages.values())
                if total_percentage > 0:
                    for i in range(1, 4):
                        bull_counts = {bull: int(percentages[bull] / total_percentage * len(bulls)) for bull in bulls}
                        bull_counts[bulls[-1]] += len(bulls) - sum(bull_counts.values())
                        selected_bull = max(bulls, key=lambda bull: (bull_counts[bull], complete_report_df.loc[
                            complete_report_df['在群母牛号'] == cow_id, f"{semen_type}冻精{i}得分"].values[0]))
                        result[f"{i}选{semen_type}"] = selected_bull
                        bulls.remove(selected_bull)
                else:
                    for i in range(1, 4):
                        result[f"{i}选{semen_type}"] = ""
            selection_result.append(result)

        selection_result_df = pd.DataFrame(selection_result)
        return selection_result_df

    def check_missing_bulls(self, bull_df):
        missing_bulls = set()

        for _, bull_row in bull_df.iterrows():
            bull_id = str(bull_row['bull_id'])
            if len(bull_id) > 10:  # 如果公牛ID大于10位,则与数据库中的BULL REG进行匹配
                if bull_id not in self.bull_library_reg_dict:
                    missing_bulls.add(bull_id)
            else:  # 如果公牛ID小于等于10位,则直接与数据库中的BULL NAAB进行匹配
                if bull_id not in self.bull_library_naab_dict:
                    for rule, replacement in self.bull_rules.items():
                        if bull_id.startswith(rule):
                            new_bull_id = replacement + bull_id[3:]
                            if new_bull_id not in self.bull_library_naab_dict:
                                missing_bulls.add(bull_id)
                                break
                        elif bull_id.startswith(replacement):
                            new_bull_id = rule + bull_id[3:]
                            if new_bull_id not in self.bull_library_naab_dict:
                                missing_bulls.add(bull_id)
                                break
                    else:
                        missing_bulls.add(bull_id)

        for bull_id in self.needed_bulls_naab:
            if bull_id not in self.bull_library_naab_dict:
                missing_bulls.add(bull_id)

        for bull_id in self.needed_bulls_reg:
            if bull_id not in self.bull_library_reg_dict:
                missing_bulls.add(bull_id)

        """print(list(missing_bulls))
        print("Needed NAAB bulls:", self.needed_bulls_naab)
        print("Needed REG bulls:", self.needed_bulls_reg)

        print("First 5 items in bull_library_naab_dict:")"""
        for i, (bull_id, bull_info) in enumerate(self.bull_library_naab_dict.items()):
            """print(f"{bull_id}: {bull_info}")"""
            if i >= 4:
                break

        """print("First 5 items in bull_library_reg_dict:")"""
        for i, (bull_id, bull_info) in enumerate(self.bull_library_reg_dict.items()):
            """print(f"{bull_id}: {bull_info}")"""
            if i >= 4:
                break

        if missing_bulls:
            self.update_missing_bulls(missing_bulls, "check_missing_bulls")

        return list(missing_bulls)

    def select_group_for_breeding(self):
        if not self.output_folder:
            QMessageBox.warning(self, "警告", "请先选择输出文件夹")
            return
        if self.cow_df.empty:
            QMessageBox.warning(self, "警告", "请先上传在群母牛信息")
            return
        if self.bull_df.empty:
            QMessageBox.warning(self, "警告", "请先上传备选公牛信息")
            return
        dialog = GroupSelectionDialog(self.cow_df, self.bull_df, self.output_folder, parent=self)
        dialog.exec()

    # 配种记录分析部分的代码：
    def upload_breeding_record(self):
        if not self.output_folder:
            QMessageBox.warning(self, "警告", "请先选择输出文件夹")
            return

        file_path, _ = QFileDialog.getOpenFileName(self, "上传配种记录", "", "Excel Files (*.xlsx)")
        if file_path:
            try:
                breeding_df = pd.read_excel(file_path)
                required_columns = ['耳号', '配种日期', '冻精编号', '冻精类型']
                if not all(col in breeding_df.columns for col in required_columns):
                    QMessageBox.warning(self, "格式错误", f"上传的配种记录表必须包含以下列: {', '.join(required_columns)}")
                    return

                breeding_df = breeding_df[required_columns]
                breeding_df['冻精编号'] = breeding_df['冻精编号'].apply(format_naab_number)
                breeding_df['配种日期'] = pd.to_datetime(breeding_df['配种日期'], errors='coerce')
                breeding_df.fillna('', inplace=True)

                # 从母牛数据中获取父号信息
                if not self.cow_df.empty:
                    # 确保cow_df中的cow_id列为字符串类型
                    self.cow_df['cow_id'] = self.cow_df['cow_id'].astype(str)
                    # 创建耳号到父号的映射
                    sire_dict = dict(zip(self.cow_df['cow_id'], self.cow_df['sire']))
                    # 将耳号转换为字符串类型
                    breeding_df['耳号'] = breeding_df['耳号'].astype(str)
                    # 添加父号列
                    breeding_df['父号'] = breeding_df['耳号'].map(sire_dict)
                else:
                    breeding_df['父号'] = ''  # 如果没有母牛数据，添加空的父号列

                # 重新排列列的顺序
                breeding_df = breeding_df[['耳号', '父号', '冻精编号', '配种日期', '冻精类型']]

                processed_breeding_file = os.path.join(self.output_folder, "processed_breeding_data.xlsx")
                breeding_df.to_excel(processed_breeding_file, index=False)

                self.breeding_df = breeding_df  # 保存处理后的 DataFrame
                QMessageBox.information(self, "上传成功", "配种记录已上传并处理成功！")
            except Exception as e:
                logging.error(f"读取配种记录文件时发生错误: {e}")
                QMessageBox.critical(self, "错误", f"读取配种记录文件时发生错误: {str(e)}")  

    # 隐性基因筛查-配种记录
    def genetic_defect_screening(self):
        if self.cow_df.empty or self.breeding_df is None or not self.output_folder:
            QMessageBox.warning(self, "错误", "请先上传在群母牛信息、配种记录,并选择输出文件夹。")
            return

        try:
            # 显示加载数据的进度对话框
            progress_dialog = QProgressDialog("正在加载公牛库数据...", "取消", 0, 100, self)
            progress_dialog.setWindowModality(Qt.WindowModality.WindowModal)
            progress_dialog.show()

            # 加载公牛库数据
            self.load_bull_library_data(progress_dialog)

            if not hasattr(self, 'bull_library_df') or self.bull_library_df.empty:
                QMessageBox.critical(self, "错误", "加载公牛库数据失败，无法进行隐性基因筛查。")
                return

            # 从cow_df中获取耳号和对应的父号
            sire_dict = dict(zip(self.cow_df['cow_id'], self.cow_df['sire']))
            
            # 将父号添加到breeding_df中
            self.breeding_df['父号'] = self.breeding_df['耳号'].map(sire_dict)
            
            # 重新排列列的顺序
            self.breeding_df = self.breeding_df[['耳号', '父号', '冻精编号', '配种日期', '冻精类型']]
            
            # 更新processed_breeding_data.xlsx文件
            processed_breeding_file = os.path.join(self.output_folder, "processed_breeding_data.xlsx")
            self.breeding_df.to_excel(processed_breeding_file, index=False)
            
            # 执行隐性基因筛查
            self.start_screening()

            QMessageBox.information(self, "筛查完成", "隐性基因筛查已完成,结果已保存。")

        except Exception as e:
            logging.error(f"执行隐性基因筛查时发生错误: {str(e)}")
            QMessageBox.critical(self, "错误", f"执行隐性基因筛查时发生错误,请联系管理员。")

    # 隐性基因筛查-配种记录（主方法）
    def start_screening(self):
        try:
            # 检查缺失公牛并更新到数据库
            missing_bulls_set = set()
            for _, mating_row in self.breeding_df.iterrows():
                sire_id = str(mating_row.get('父号', ''))
                if sire_id != 'nan':
                    if sire_id not in self.bull_library_naab_dict and sire_id not in self.bull_library_reg_dict:
                        missing_bulls_set.add(sire_id)
                ai_id = str(mating_row.get('冻精编号', ''))
                if ai_id != 'nan':
                    if ai_id not in self.bull_library_naab_dict and ai_id not in self.bull_library_reg_dict:
                        missing_bulls_set.add(ai_id)
            
            if missing_bulls_set:
                connection = pymysql.connect(
                    host=cipher_suite.decrypt(encoded_host).decode(),
                    port=int(cipher_suite.decrypt(encoded_port).decode()),
                    user=cipher_suite.decrypt(encoded_user).decode(),
                    password=cipher_suite.decrypt(encoded_password).decode(),
                    database=cipher_suite.decrypt(encoded_db).decode(),
                    charset='utf8mb4',
                    cursorclass=pymysql.cursors.DictCursor
                )
                with connection.cursor() as cursor:
                    for bull_id in missing_bulls_set:
                        sql = "INSERT INTO miss_bull (bull, source, time, user) VALUES (%s, %s, %s, %s)"
                        cursor.execute(sql, (bull_id, '隐性基因筛查-配种记录', datetime.datetime.now(), self.username))
                connection.commit()
                connection.close()
            
            self.screened_df = self.breeding_df.apply(self.screen_record, axis=1)
            self.screened_df = pd.concat([self.breeding_df, self.screened_df], axis=1)
            
            # 保存筛查结果
            screening_result_file = os.path.join(self.output_folder, "隐性基因筛查结果-配种记录.xlsx")
            self.screened_df.to_excel(screening_result_file, index=False)
            
        except Exception as e:
            logging.error(f"执行筛查时发生异常:{str(e)}")
            QMessageBox.critical(self, "错误", f"执行筛查时发生错误,请联系管理员。")

    # 隐性基因筛查-配种记录（标注方法）
    def screen_record(self, record):
        sire_id = str(record.get('父号', ''))
        ai_id = str(record.get('冻精编号', ''))
        sire_in_library = sire_id in self.bull_library_naab_dict or sire_id in self.bull_library_reg_dict
        ai_in_library = ai_id in self.bull_library_naab_dict or ai_id in self.bull_library_reg_dict
        
        result = {}
        unsafe_count = 0
        for gene in self.defect_genes:
            sire_gene = ''
            ai_gene = ''
            if sire_in_library:
                sire_data = self.bull_library_naab_dict.get(sire_id) or self.bull_library_reg_dict.get(sire_id, {})
                sire_gene = str(sire_data.get(gene, ''))
            if ai_in_library:
                ai_data = self.bull_library_naab_dict.get(ai_id) or self.bull_library_reg_dict.get(ai_id, {})
                ai_gene = str(ai_data.get(gene, ''))
            if sire_gene == 'C' and ai_gene == 'C':
                result[gene] = '纯合高风险'
                unsafe_count += 1
            else:
                result[gene] = ''
        
        if not sire_in_library and not ai_in_library:
            result['备注'] = '库中缺少: 母牛父号/与配公牛信息'
        elif not sire_in_library:
            result['备注'] = '库中缺少: 母牛父号信息'
        elif not ai_in_library:
            result['备注'] = '库中缺少: 与配公牛信息'
        else:
            result['备注'] = ''
        result['汇总'] = unsafe_count
        return pd.Series(result)

    # 加载隐性基因相关数据库
    def load_bull_library_data(self, progress_dialog):
        try:
            connection = pymysql.connect(
                host=cipher_suite.decrypt(encoded_host).decode(),
                port=int(cipher_suite.decrypt(encoded_port).decode()),
                user=cipher_suite.decrypt(encoded_user).decode(),
                password=cipher_suite.decrypt(encoded_password).decode(),
                database=cipher_suite.decrypt(encoded_db).decode(),
                charset='utf8mb4',
                cursorclass=pymysql.cursors.DictCursor
            )
            
            with connection.cursor() as cursor:
                sql = "SELECT `BULL NAAB`, `BULL REG`, `MW`, `HH1`, `HH2`, `HH3`, `HH4`, `HH5`, `HH6`, `BLAD`, `Chondrodysplasia`, `Citrullinemia`, `DUMPS`, `Factor XI`, `CVM`, `Brachyspina`, `Mulefoot`, `Cholesterol deficiency` FROM bull_library"
                cursor.execute(sql)
                result = cursor.fetchall()
                self.bull_library_df = pd.DataFrame(result)
            
            connection.close()
            
            progress_dialog.setValue(50)  # 更新进度
            
            self.defect_genes = ["MW", "HH1", "HH2", "HH3", "HH4", "HH5", "HH6", "BLAD", "Chondrodysplasia", "Citrullinemia",
                                "DUMPS", "Factor XI", "CVM", "Brachyspina", "Mulefoot", "Cholesterol deficiency"]
            
            self.bull_library_naab_dict = {}
            self.bull_library_reg_dict = {}
            total_rows = len(self.bull_library_df)
            for i, (_, row) in enumerate(self.bull_library_df.iterrows()):
                if pd.notna(row['BULL NAAB']):
                    self.bull_library_naab_dict[str(row['BULL NAAB'])] = {
                        'naab': str(row['BULL NAAB']),
                        **{gene: str(row[gene]) for gene in self.defect_genes if pd.notna(row[gene])}
                    }
                if pd.notna(row['BULL REG']):
                    self.bull_library_reg_dict[str(row['BULL REG'])] = {
                        'naab': str(row['BULL NAAB']) if pd.notna(row['BULL NAAB']) else None,
                        **{gene: str(row[gene]) for gene in self.defect_genes if pd.notna(row[gene])}
                    }
                progress_dialog.setValue(50 + int((i + 1) / total_rows * 50))  # 更新进度

            logging.info("公牛库数据已成功加载！")
        
        except Exception as e:
            logging.error(f"加载公牛库数据时发生错误: {e}")
            QMessageBox.critical(self, "错误", "加载公牛库数据时发生错误，请联系管理员。")
        finally:
            progress_dialog.close()

    # 冻精使用情况分析
    def semen_usage_analysis(self):
        try:
            progress_dialog = QProgressDialog("正在进行配种记录分析...\n\n数据库连接中......", "取消", 0, 100, self)
            progress_dialog.setWindowModality(Qt.WindowModality.WindowModal)
            progress_dialog.setValue(0)
            progress_dialog.show()
            progress_dialog.raise_()
            QApplication.processEvents()            

            if self.breeding_df is None or self.breeding_df.empty or not self.output_folder:
                QMessageBox.warning(self, "警告", "请先选择储存文件夹，并上传配种记录。")
                return
            progress_dialog.setValue(5)
            progress_dialog.show()
            progress_dialog.raise_()
            QApplication.processEvents()  

            # 检查点1：打印原始的breeding_df的行数和列名
            """print("检查点1 - 原始breeding_df:")
            print(f"行数: {len(self.breeding_df)}")
            print(f"列名: {self.breeding_df.columns.tolist()}")
            print(self.breeding_df.head())"""

            # 检查配种日期格式
            if '配种日期' in self.breeding_df.columns:
                if not pd.api.types.is_datetime64_any_dtype(self.breeding_df['配种日期']):
                    QMessageBox.warning(self, "警告", "配种日期列的格式不是日期格式,请修改为日期格式后重新上传。")
                    return
            progress_dialog.setValue(10)
            progress_dialog.show()
            progress_dialog.raise_()
            QApplication.processEvents()          



            # 生成所需公牛列表
            needed_bulls = set(self.breeding_df['冻精编号'].unique())
            needed_bulls.discard('nan')  # 去除空值
            needed_bulls = {format_naab_number(bull) for bull in needed_bulls}  # 格式化NAAB号

            # 检查点2：打印needed_bulls
            """print("检查点2 - needed_bulls:")
            print(needed_bulls)"""

            progress_dialog.setValue(15)
            progress_dialog.show()
            progress_dialog.raise_()
            QApplication.processEvents()

            # 检查是否有缺失公牛
            missing_bulls = self.check_missing_bulls(pd.DataFrame({'bull_id': list(needed_bulls)}))
            if missing_bulls:
                self.update_missing_bulls(missing_bulls, "关键育种性状计算-配种记录")
                """print("missing bull:", missing_bulls)"""

            default_traits = ['NM$', 'TPI', 'MILK', 'FAT', 'FAT %', 'PROT', 'PROT%', 'SCS', 'PL', 'DPR', 'PTAT', 'UDC',
                            'FLC', 'ST', 'TL']
            all_traits = ['TPI', 'NM$', 'CM$', 'FM$', 'GM$', 'MILK', 'FAT', 'PROT', 'FAT %', 'PROT%', 'SCS', 'DPR',
                        'HCR', 'CCR', 'PL', 'SCE', 'DCE', 'SSB', 'DSB', 'PTAT', 'UDC', 'FLC', 'BDC', 'ST', 'SG', 'BD',
                        'DF', 'RA', 'RW', 'LS', 'LR', 'FA', 'FLS', 'FU', 'UH', 'UW', 'UC', 'UD', 'FT', 'RT', 'TL',
                        'FE', 'FI', 'HI', 'LIV', 'GL', 'MAST', 'MET', 'RP', 'KET', 'DA', 'MFV', 'EFC', 'HLiv', 'FS',
                        'RFI', 'Milk Speed','Eval Date']

            dialog = KeyTraitsDialog(default_traits, all_traits)
            progress_dialog.setValue(17)
            progress_dialog.show()
            progress_dialog.raise_()
            QApplication.processEvents()
            if dialog.exec() == QDialog.DialogCode.Accepted:
                selected_traits = dialog.get_selected_traits()
                if not selected_traits:
                    QMessageBox.warning(self, "警告", "请至少选择一个关键性状。")
                    return
                progress_dialog.setValue(20)

                progress_dialog.setValue(32)
                # 读取公牛信息
                connection = pymysql.connect(
                    host=cipher_suite.decrypt(encoded_host).decode(),
                    port=int(cipher_suite.decrypt(encoded_port).decode()),
                    user=cipher_suite.decrypt(encoded_user).decode(),
                    password=cipher_suite.decrypt(encoded_password).decode(),
                    database=cipher_suite.decrypt(encoded_db).decode(),
                    charset='utf8mb4',
                    cursorclass=pymysql.cursors.DictCursor
                )
                
                progress_dialog.setValue(36)
                placeholders = ','.join(['%s'] * len(needed_bulls))
                sql = f"SELECT * FROM bull_library WHERE `BULL NAAB` IN ({placeholders}) OR `BULL REG` IN ({placeholders})"
                params = tuple(needed_bulls) * 2

                progress_dialog.setValue(39)
                with connection.cursor() as cursor:
                    cursor.execute(sql, params)
                    result = cursor.fetchall()
                    bull_library_df = pd.DataFrame(result)
                progress_dialog.setValue(43)

                # 处理需要替换前3位的公牛
                for bull in needed_bulls.copy():  # 创建一个副本,避免在迭代过程中修改原始集合
                    if bull not in bull_library_df['BULL NAAB'].values:
                        for rule, replacement in self.bull_rules.items():
                            if bull.startswith(rule):
                                new_bull = replacement + bull[3:]
                                needed_bulls.add(new_bull)  # 将替换后的公牛ID添加到needed_bulls
                                break
                            elif bull.startswith(replacement):
                                new_bull = rule + bull[3:]
                                needed_bulls.add(new_bull)  # 将替换后的公牛ID添加到needed_bulls
                                break
                progress_dialog.setValue(46)                    

                placeholders = ','.join(['%s'] * len(needed_bulls))
                sql = f"SELECT * FROM bull_library WHERE `BULL NAAB` IN ({placeholders})"
                params = tuple(needed_bulls)
                with connection.cursor() as cursor:
                    cursor.execute(sql, params)
                    result = cursor.fetchall()
                    bull_library_df = pd.concat([bull_library_df, pd.DataFrame(result)], ignore_index=True)

                connection.close()

                progress_dialog.setValue(52)
                # 初始化 bull_key_traits_scores,包含所有需要的公牛
                bull_key_traits_scores = pd.DataFrame(columns=['BULL NAAB', 'BULL REG'] + selected_traits)

                for bull in needed_bulls:
                    bull_data = bull_library_df[bull_library_df['BULL NAAB'] == bull]
                    if bull_data.empty:
                        found = False
                        for rule, replacement in self.bull_rules.items():
                            if bull.startswith(rule):
                                new_bull = replacement + bull[3:]
                                bull_data = bull_library_df[bull_library_df['BULL NAAB'] == new_bull]
                                if not bull_data.empty:
                                    found = True
                                    break
                            elif bull.startswith(replacement):
                                new_bull = rule + bull[3:]
                                bull_data = bull_library_df[bull_library_df['BULL NAAB'] == new_bull]
                                if not bull_data.empty:
                                    found = True
                                    break
                        if not found:
                            # 如果替换前三位后仍未找到,则创建一个新的DataFrame行
                            bull_data = pd.DataFrame({'BULL NAAB': [bull], 'BULL REG': [np.nan]})
                            for trait in selected_traits:
                                bull_data[trait] = np.nan
                        else:
                            # 如果替换前三位后找到了,则更新 BULL REG 和各性状指数
                            bull_data['BULL REG'] = bull_library_df[bull_library_df['BULL NAAB'] == new_bull]['BULL REG'].values[0]
                            for trait in selected_traits:
                                bull_data[trait] = bull_library_df[bull_library_df['BULL NAAB'] == new_bull][trait].values[0]
                    bull_key_traits_scores = pd.concat(
                        [bull_key_traits_scores, bull_data[['BULL NAAB', 'BULL REG'] + selected_traits]],
                        ignore_index=True)

                self.status_label.setText(self.status_label.text() + "\n配种记录中公牛关键性状指数计算完成。")
                logging.info("配种记录中公牛关键性状指数计算完成。")

                # 在这里添加生成冻精使用情况报告的代码
                self.generate_semen_usage_report(bull_key_traits_scores, selected_traits, progress_dialog)

                progress_dialog.setValue(100)
                QMessageBox.information(self, "完成", "配种记录中公牛关键性状指数计算和冻精使用情况报告生成完成。")

        except Exception as e:
            logging.error(f"计算配种记录中公牛关键性状指数时发生错误: {e}", exc_info=True)
            QMessageBox.critical(self, "错误", f"计算配种记录中公牛关键性状指数时发生错误：{e},请按照模版上传配种记录。")

    # 生成冻精使用不同年份各性状的数据
    def generate_semen_usage_report(self, bull_key_traits_scores, selected_traits, progress_dialog):
        try:
            # 处理breeding_df，添加年份列并计算使用次数
            self.breeding_df['年份'] = pd.to_datetime(self.breeding_df['配种日期']).dt.year
            
            # 根据'耳号'、'配种日期'、'冻精编号'和'冻精类型'去重
            unique_breeding_df = self.breeding_df.drop_duplicates(subset=['耳号', '配种日期', '冻精编号', '冻精类型'])

            usage_count = unique_breeding_df.groupby(['年份', '冻精编号', '冻精类型']).size().reset_index(name='使用次数')

            # 确保bull_key_traits_scores中每个BULL NAAB只有一条记录
            bull_key_traits_scores_unique = bull_key_traits_scores.drop_duplicates(subset='BULL NAAB')

            # 合并usage_count和bull_key_traits_scores_unique
            merged_df = pd.merge(usage_count, bull_key_traits_scores_unique, left_on='冻精编号', right_on='BULL NAAB', how='left')
            
            # 将所选特征转换为float类型
            for trait in selected_traits:
                merged_df[trait] = pd.to_numeric(merged_df[trait], errors='coerce')

            # 重新排列列的顺序，将 'BULL NAAB' 和 'BULL REG' 移到最后
            columns_order = ['年份', '冻精编号', '冻精类型', '使用次数'] + selected_traits + ['BULL NAAB', 'BULL REG']
            merged_df = merged_df[columns_order]

            # 创建两个输出文件名
            """output_file_original = os.path.join(self.output_folder, "冻精使用情况报告.xlsx")"""
            output_file_with = os.path.join(self.output_folder, "冻精使用情况报告.xlsx")
            charts_folder = os.path.join(self.output_folder, "结果-冻精关键性状进展图")
            os.makedirs(charts_folder, exist_ok=True)

            # 设置matplotlib的中文字体
            plt.rcParams['font.sans-serif'] = ['SimHei']
            plt.rcParams['axes.unicode_minus'] = False

            # 生成原始报告 - 包含不同年份的冻精使用情况
            with pd.ExcelWriter(output_file_with, engine='xlsxwriter') as writer:
                workbook = writer.book
                red_format = workbook.add_format({'font_color': 'red'})

                # 生成表1：年份汇总
                yearly_summary = merged_df.groupby('年份').agg({
                    '使用次数': 'sum',
                    **{trait: lambda x: self.weighted_average(x, merged_df.loc[x.index, '使用次数']) for trait in selected_traits}
                }).reset_index()
                yearly_summary.to_excel(writer, sheet_name='年份汇总', index=False)

                # 生成表2：按年份和冻精类型的详细信息
                for year in merged_df['年份'].unique():
                    year_data = merged_df[merged_df['年份'] == year]
                    
                    # 计算年度均值（加权平均）
                    year_avg = year_data.agg({
                        '使用次数': 'sum',
                        **{trait: lambda x: self.weighted_average(x, year_data.loc[x.index, '使用次数']) for trait in selected_traits}
                    }).to_frame().T
                    year_avg['年份'] = year
                    year_avg['冻精类型'] = '均值'
                    year_avg = year_avg.reindex(columns=year_data.columns)

                    # 对每种冻精类型计算加权平均
                    type_avg = year_data.groupby('冻精类型').agg({
                        '使用次数': 'sum',
                        **{trait: lambda x: self.weighted_average(x, year_data.loc[x.index, '使用次数']) for trait in selected_traits}
                    }).reset_index()
                    type_avg['年份'] = year
                    type_avg = type_avg.reindex(columns=year_data.columns)

                    # 合并年度均值、类型均值和详细数据
                    final_data = pd.concat([year_avg, type_avg, year_data])
                    final_data = final_data.sort_values(['冻精类型', '使用次数'], ascending=[True, False])

                    # 写入Excel
                    sheet_name = f'{year}年冻精使用情况'
                    final_data.to_excel(writer, sheet_name=sheet_name, index=False)

                    # 此处可以添加原始报告的其他表格或内容

            self.status_label.setText(self.status_label.text() + f"\n冻精使用情况报告已生成: {output_file_with}")
            logging.info(f"冻精使用情况报告已生成: {output_file_with}")

            # 生成带趋势图和散点图的报告
            with pd.ExcelWriter(output_file_with, engine='xlsxwriter') as writer:
                workbook = writer.book
                red_format = workbook.add_format({'font_color': 'red'})

                # 生成表1：年份汇总
                yearly_summary = merged_df.groupby('年份').agg({
                    '使用次数': 'sum',
                    **{trait: lambda x: self.weighted_average(x, merged_df.loc[x.index, '使用次数']) for trait in selected_traits}
                }).reset_index()
                yearly_summary.to_excel(writer, sheet_name='年份汇总', index=False)

                # 生成表2：按年份和冻精类型的详细信息，和趋势图
                for year in merged_df['年份'].unique():
                    year_data = merged_df[merged_df['年份'] == year]
                    
                    # 计算年度均值（加权平均）
                    year_avg = year_data.agg({
                        '使用次数': 'sum',
                        **{trait: lambda x: self.weighted_average(x, year_data.loc[x.index, '使用次数']) for trait in selected_traits}
                    }).to_frame().T
                    year_avg['年份'] = year
                    year_avg['冻精类型'] = '均值'
                    year_avg = year_avg.reindex(columns=year_data.columns)

                    # 对每种冻精类型计算加权平均
                    type_avg = year_data.groupby('冻精类型').agg({
                        '使用次数': 'sum',
                        **{trait: lambda x: self.weighted_average(x, year_data.loc[x.index, '使用次数']) for trait in selected_traits}
                    }).reset_index()
                    type_avg['年份'] = year
                    type_avg = type_avg.reindex(columns=year_data.columns)

                    # 合并年度均值、类型均值和详细数据
                    final_data = pd.concat([year_avg, type_avg, year_data])
                    final_data = final_data.sort_values(['冻精类型', '使用次数'], ascending=[True, False])

                    # 写入Excel
                    sheet_name = f'{year}年冻精使用情况'
                    final_data.to_excel(writer, sheet_name=sheet_name, index=False)

                # 生成趋势图
                for trait in selected_traits:
                    plt.figure(figsize=(14, 6))  # 增加图表大小以适应标签

                    # 准备数据
                    trait_data = merged_df.groupby(['年份', '冻精类型']).apply(
                        lambda x: pd.Series({
                            trait: self.weighted_average(x[trait], x['使用次数']),
                            '使用次数': x['使用次数'].sum()
                        })
                    ).reset_index()

                    # 绘制每种冻精类型的趋势线
                    for semen_type in trait_data['冻精类型'].unique():
                        type_data = trait_data[trait_data['冻精类型'] == semen_type]
                        line = plt.plot(type_data['年份'], type_data[trait], marker='o', label=semen_type)
                        
                        # 添加数据标签
                        for x, y in zip(type_data['年份'], type_data[trait]):
                            plt.annotate(f'{y:.2f}', (x, y), textcoords="offset points", 
                                        xytext=(0,10), ha='center', va='bottom',
                                        bbox=dict(boxstyle='round,pad=0.5', fc='yellow', alpha=0.5),
                                        fontsize=8)

                    # 绘制总体均值线
                    overall_avg = trait_data.groupby('年份').apply(
                        lambda x: self.weighted_average(x[trait], x['使用次数'])
                    ).reset_index()
                    plt.plot(overall_avg['年份'], overall_avg[0], label='总体均值', linewidth=2, color='red')
                    
                    # 为总体均值添加标签
                    for x, y in zip(overall_avg['年份'], overall_avg[0]):
                        plt.annotate(f'{y:.2f}', (x, y), textcoords="offset points", 
                                    xytext=(0,10), ha='center', va='bottom',
                                    bbox=dict(boxstyle='round,pad=0.5', fc='red', alpha=0.5),
                                    fontsize=8, color='white')

                    plt.title(f"{trait} 冻精使用趋势", fontsize=18, fontweight='bold')
                    plt.xlabel('年份', fontsize=14)
                    plt.ylabel(trait, fontsize=14)
                    plt.legend(fontsize=10, loc='center left', bbox_to_anchor=(1, 0.5))
                    plt.grid(True, linestyle='--', alpha=0.7)
                    plt.tight_layout()

                    # 保存图表
                    chart_path = os.path.join(charts_folder, f"{trait}冻精使用趋势.png")
                    plt.savefig(chart_path, dpi=300, bbox_inches='tight')
                    plt.close()

                    # 将数据写入Excel
                    trait_data.pivot(index='年份', columns='冻精类型', values=trait).reset_index().to_excel(writer, sheet_name=trait, index=False)

                self.status_label.setText(self.status_label.text() + f"\n带趋势图的冻精使用情况报告已生成: {output_file_with}")
                logging.info(f"带趋势图的冻精使用情况报告已生成: {output_file_with}")

            # 生成散点图（单值图）
            scatter_plot_files = self.generate_semen_usage_scatter_plots(self.breeding_df, self.output_folder)

            if scatter_plot_files:
                self.status_label.setText(self.status_label.text() + f"\n冻精使用散点图（单值图）已生成: {', '.join(scatter_plot_files)}")
                logging.info(f"冻精使用散点图（单值图）已生成: {', '.join(scatter_plot_files)}")
            else:
                self.status_label.setText(self.status_label.text() + "\n未能生成冻精使用散点图（单值图），可能是因为没有相关数据。")
                logging.warning("未能生成冻精使用散点图（单值图），可能是因为没有相关数据。")

        except Exception as e:
            logging.error(f"生成冻精使用情况报告时发生错误: {e}", exc_info=True)
            QMessageBox.critical(self, "错误", f"生成冻精使用情况报告时发生错误：{e}")


    # 加权平均方法
    def weighted_average(self, values, weights):
        """
        计算加权平均值，忽略NaN值，并保留两位小数
        """
        valid = ~np.isnan(values)
        return round(np.average(values[valid], weights=weights[valid]), 2) if np.any(valid) else np.nan

    # 冻精使用趋势图
    def generate_semen_usage_trend_charts(self, merged_df, selected_traits, yearly_summary):
        excel_file = os.path.join(self.output_folder, "结果-冻精使用趋势图.xlsx")
        workbook = xlsxwriter.Workbook(excel_file, {'nan_inf_to_errors': True})
        worksheet = workbook.add_worksheet()
        chart_row = 0
        chart_col = 0

        # 获取冻精类型
        semen_types = merged_df['冻精类型'].unique().tolist()
        """print("检查点1：冻精类型 - semen_types:", semen_types)"""

        # 为每个性状生成图表
        for trait in selected_traits:
            print(f"\n检查点2：开始处理性状 {trait}")
            
            # 计算每个冻精类型的加权平均值
            type_avg = merged_df.groupby(['年份', '冻精类型']).apply(
                lambda x: pd.Series({
                    trait: self.weighted_average(x[trait], x['使用次数']),
                    '使用次数': x['使用次数'].sum()
                })
            ).reset_index()
            
            # 数据透视
            pivot_data = type_avg.pivot(index='年份', columns='冻精类型', values=trait).reset_index()
            
            # 合并总体数据
            yearly_trait_data = yearly_summary[['年份', trait]].rename(columns={trait: f'{trait}_总体'})
            pivot_data = pd.merge(pivot_data, yearly_trait_data, on='年份', how='left')
            
            # 确保列的顺序
            columns_order = ['年份'] + semen_types + [f'{trait}_总体']
            pivot_data = pivot_data.reindex(columns=columns_order)
            
            print(f"检查点3：{trait} 的 pivot_data:")
            print(pivot_data)

            # 写入数据到 Excel 表中
            start_row = chart_row
            worksheet.write(start_row, chart_col, f"{trait}进展图数据")
            for i, col_name in enumerate(pivot_data.columns):
                worksheet.write(start_row + 1, chart_col + i, col_name)
            for i, row in enumerate(pivot_data.values):
                for j, value in enumerate(row):
                    worksheet.write(start_row + 2 + i, chart_col + j, value)

            # 计算 Y 轴范围
            all_values = pivot_data.iloc[:, 1:].values.flatten()
            all_values = all_values[~pd.isna(all_values)]  # 排除 NaN 值
            if len(all_values) == 0:
                print(f"{trait} 没有可用的数据")
                continue

            max_value = np.max(all_values)
            min_value = np.min(all_values)
            range_padding = 0.2 * (max_value - min_value)

            # 动态调整 Y 轴的最大和最小值，确保最小值不会总是 0
            y_axis_max = max_value + range_padding
            # 如果最小值和最大值都为负数，允许最小值为负
            y_axis_min = min_value - range_padding

            print(f"检查点：{trait} Y轴范围 - min: {y_axis_min}, max: {y_axis_max}")



            # 创建图表
            chart = workbook.add_chart({'type': 'line'})

            # 添加数据系列
            for i, semen_type in enumerate(semen_types + [f'{trait}_总体']):
                col = chr(ord('B') + i)
                series = {
                    'name': f'=Sheet1!${col}${start_row + 2}',
                    'categories': f'=Sheet1!$A${start_row + 3}:$A${start_row + 2 + len(pivot_data)}',
                    'values': f'=Sheet1!${col}${start_row + 3}:${col}${start_row + 2 + len(pivot_data)}',
                    'data_labels': {'value': True, 'font': {'size': 10}},
                    'marker': {'type': 'circle', 'size': 3},
                }
                
                # 为总体数据设置特殊样式
                if semen_type == f'{trait}_总体':
                    series['data_labels']['font']['size'] = 12
                    series['data_labels']['font']['bold'] = True
                    series['line'] = {'color': 'red', 'width': 2.25}
                    series['marker']['fill'] = {'color': 'red'}
                
                chart.add_series(series)

            # 设置图表标题和轴标签
            chart.set_title({'name': f'{trait}冻精使用趋势'})
            chart.set_x_axis({'name': '年份'})
            chart.set_y_axis({'name': trait, 'min': y_axis_min, 'max': y_axis_max})

            # 设置图表大小和位置
            chart.set_size({'width': 720, 'height': 360})
            worksheet.insert_chart(start_row + len(pivot_data) + 5, chart_col, chart)

            # 更新下一个图表的位置
            chart_row = start_row + len(pivot_data) + 35
            chart_col = 0

            print(f"检查点5：{trait} 图表已生成")

        workbook.close()
        """print(f"冻精使用趋势图已生成：{excel_file}")"""

    # 生成冻精使用的单值图
    def generate_semen_usage_scatter_plots(self, breeding_df, output_folder):
        if breeding_df.empty:
            logging.warning("配种记录为空，无法生成冻精使用单值图。")
            return None

        plt.rcParams['font.sans-serif'] = ['SimHei']
        plt.rcParams['axes.unicode_minus'] = False

        breeding_df['配种日期'] = pd.to_datetime(breeding_df['配种日期'])
        semen_types = breeding_df['冻精类型'].unique()

        if len(semen_types) == 0:
            logging.warning("配种记录中没有冻精类型信息，无法生成冻精使用单值图。")
            return None

        output_files = []

        for semen_type in semen_types:
            df_filtered = breeding_df[breeding_df['冻精类型'] == semen_type]
            
            if df_filtered.empty:
                logging.warning(f"没有{semen_type}冻精的数据，跳过生成图表。")
                continue

            fig, ax = plt.subplots(figsize=(24, 12))  # 增加图表大小

            unique_bulls = sorted(df_filtered['冻精编号'].unique())
            bull_indices = {bull: i for i, bull in enumerate(unique_bulls)}

            for bull in unique_bulls:
                bull_data = df_filtered[df_filtered['冻精编号'] == bull]
                x = []
                y = []
                for _, row in bull_data.iterrows():
                    x.append(bull_indices[bull] + random.uniform(-0.3, 0.3))
                    y.append(row['配种日期'])
                ax.scatter(x, y, s=40, alpha=0.6)  # 增加点的大小

            ax.set_xticks(range(len(unique_bulls)))
            ax.set_xticklabels(unique_bulls, rotation=90)
            
            # 增加刻度标签的字体大小
            ax.tick_params(axis='both', which='major', labelsize=14)
            
            ax.yaxis.set_major_locator(YearLocator())
            ax.yaxis.set_major_formatter(DateFormatter('%Y-%m-%d'))
            ax.yaxis.set_minor_locator(AutoMinorLocator())

            ax.set_title(f'{semen_type}冻精配种日期的单值图', fontsize=32)
            ax.set_xlabel('冻精编号', fontsize=24)
            ax.set_ylabel('配种日期', fontsize=24)

            # 增加坐标轴标签的字体大小
            ax.xaxis.label.set_size(32)
            ax.yaxis.label.set_size(32)

            ax.grid(True, linestyle='--', alpha=0.7)

            plt.tight_layout()

            output_file = os.path.join(output_folder, f'{semen_type}配种日期单值图.png')
            plt.savefig(output_file, dpi=300, bbox_inches='tight')
            plt.close()

            output_files.append(output_file)
            logging.info(f"{semen_type}配种日期单值图已保存: {output_file}")

        if not output_files:
            logging.warning("没有生成任何冻精使用单值图。")
            return None

        return output_files




    # 体型外貌鉴定-功能实现部分-上传数据
    def upload_conformation_data(self):
        if not hasattr(self, 'output_folder') or not self.output_folder:
            QMessageBox.warning(self, "警告", "请先选择输出文件夹。")
            return

        file_path, _ = QFileDialog.getOpenFileName(
            self, 
            "上传体型外貌数据", 
            "", 
            "Excel Files (*.xlsx *.xls);;CSV Files (*.csv)"
        )
        
        if not file_path:
            return

        try:
            # 根据文件扩展名读取数据
            if file_path.endswith('.csv'):
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)

            # 检查必需的列
            required_columns = [
                "牧场", "牛号", "体高", "胸宽", "体深", "腰强度", "尻角度", "尻宽", 
                "蹄角度", "蹄踵深度", "骨质地", "后肢侧视", "后肢后视", "乳房深度", 
                "中央悬韧带", "前乳房附着", "前乳头位置", "前乳头长度", "后乳房附着高度", 
                "后乳房附着宽度", "后乳头位置", "棱角性"
            ]
            
            missing_columns = [col for col in required_columns if col not in df.columns]
            if missing_columns:
                QMessageBox.warning(self, "数据格式错误", 
                                f"缺少以下列：\n{', '.join(missing_columns)}")
                return

            # 保存处理后的数据
            self.conformation_df = df
            
            # 保存到输出文件夹
            output_file = os.path.join(self.output_folder, "processed_conformation_data.xlsx")
            df.to_excel(output_file, index=False)

            QMessageBox.information(self, "上传成功", f"成功上传 {len(df)} 条体型外貌数据。")
            
        except Exception as e:
            QMessageBox.critical(self, "错误", f"数据上传失败：{str(e)}")

    def process_conformation_data(self):
        if not hasattr(self, 'conformation_df') or self.conformation_df.empty:
            QMessageBox.warning(self, "警告", "请先上传体型外貌数据。")
            return

        try:
            # 创建输出文件夹
            output_dir = os.path.join(self.output_folder, "体型外貌鉴定")
            os.makedirs(output_dir, exist_ok=True)
            
            # 生成线性评分图
            self.generate_linear_score_chart(output_dir)
            
            # 生成缺陷性状占比图
            self.generate_defect_ratio_chart(output_dir)
            
            # 生成各性状缺陷占比饼图
            self.generate_defect_pie_charts(output_dir)
            
            QMessageBox.information(self, "处理完成", "体型外貌数据处理完成。")
            
        except Exception as e:
            QMessageBox.critical(self, "错误", f"数据处理失败：{str(e)}")

    def generate_linear_score_chart(self, output_dir):
        # 设置中文字体
        plt.rcParams['font.sans-serif'] = ['SimHei']  
        plt.rcParams['axes.unicode_minus'] = False

        # 定义评分标准
        traits = {
            '体高': 8, '胸宽': 9, '体深': 7, '腰强度': 9, '尻角度': 5, '尻宽': 9,
            '蹄角度': 7, '蹄踵深度': 9, '骨质地': 9, '后肢侧视': 5, '后肢后视': 9,
            '乳房深度': 5, '中央悬韧带': 9, '前乳房附着': 9, '前乳头位置': 6,
            '前乳头长度': 5, '后乳房附着高度': 8, '后乳房附着宽度': 9, '后乳头位置': 6,
            '棱角性': 9
        }

        df = self.conformation_df
        # 确保数据为数值型并计算平均值
        for trait in traits.keys():
            df[trait] = pd.to_numeric(df[trait], errors='coerce')
        means = {trait: df[trait].mean() for trait in traits.keys()}
        
        # 创建图表
        plt.figure(figsize=(15, 10))
        y_pos = np.arange(len(traits))
        
        # 绘制均值条形图
        bars = plt.barh(y_pos, list(means.values()), color='royalblue', alpha=0.6)
        
        # 添加最佳得分点
        plt.plot(list(traits.values()), y_pos, 'ro', markersize=10, label='最佳得分')
        
        # 添加数值标签
        for i, bar in enumerate(bars):
            width = bar.get_width()
            plt.text(width, bar.get_y() + bar.get_height()/2,
                    f'{width:.1f}', ha='left', va='center')
        
        # 设置坐标轴
        plt.yticks(y_pos, list(traits.keys()), fontsize=12)
        plt.xlabel('评分', fontsize=12)
        plt.title('牛群体型外貌线性评分', fontsize=14, pad=20)
        
        # 添加网格和图例
        plt.grid(True, alpha=0.3)
        plt.legend(loc='lower right')
        
        # 调整布局和保存
        plt.tight_layout()
        plt.savefig(os.path.join(output_dir, '牛群体型外貌线性评分.png'), 
                    dpi=300, bbox_inches='tight')
        plt.close()


    def generate_defect_ratio_chart(self, output_dir):
        # 设置中文字体
        plt.rcParams['font.sans-serif'] = ['SimHei']  
        plt.rcParams['axes.unicode_minus'] = False

        df = self.conformation_df
        defect_ratios = {}

        # 确保数据为数值型  
        for col in df.columns:
            if col not in ['牧场', '牛号']:
                df[col] = pd.to_numeric(df[col], errors='coerce')
                
        # 定义完整的缺陷规则并直接计算缺陷比例
        defect_rules = {
            '体高': {'values': [1, 2, 3, 9]},
            '胸宽': {'values': [1, 2, 3]},
            '体深': {'values': [1, 2, 3]},
            '腰强度': {'values': [1, 2, 3]},
            '尻角度': {'values': [1, 2, 3, 7, 8, 9]},
            '尻宽': {'values': [1, 2, 3]},
            '蹄角度': {'values': [1, 2, 3]},
            '蹄踵深度': {'values': [1, 2, 3]},
            '骨质地': {'values': [1, 2, 3]},
            '后肢侧视': {'values': [1, 2, 3, 7, 8, 9]},
            '后肢后视': {'values': [1, 2, 3]},
            '乳房深度': {'values': [1, 2, 8, 9]},
            '中央悬韧带': {'values': [1, 2, 3]},
            '前乳房附着': {'values': [1, 2, 3]},
            '前乳头位置': {'values': [1, 2, 8, 9]},
            '前乳头长度': {'values': [1, 2, 3, 7, 8, 9]},
            '后乳房附着高度': {'values': [1, 2, 3]},
            '后乳房附着宽度': {'values': [1, 2, 3]},
            '后乳头位置': {'values': [1, 2, 8, 9]},
            '棱角性': {'values': [1, 2, 3]}
        }
        
        # 计算每个特征的缺陷比例
        for trait, rule in defect_rules.items():
            if trait not in df.columns:
                continue
                
            valid_records = df[trait].dropna()
            total = len(valid_records)
            if total == 0:
                continue
                
            # 修改这里：直接使用 isin
            defect_count = sum(valid_records.isin(rule['values']))
            defect_ratio = (defect_count / total * 100) if total > 0 else 0
            defect_ratios[trait] = defect_ratio
        
        # 创建水平条形图
        plt.figure(figsize=(15, 10))
        
        # 准备数据
        traits = list(defect_ratios.keys())
        ratios = list(defect_ratios.values())
        y_pos = np.arange(len(traits))
        
        # 绘制水平条形图
        bars = plt.barh(y_pos, ratios, color='royalblue', alpha=0.6)
        
        # 为每个条形添加数值标签
        for i, bar in enumerate(bars):
            width = bar.get_width()
            plt.text(width + 0.5, bar.get_y() + bar.get_height()/2,
                    f'{width:.1f}%',
                    ha='left', va='center',
                    fontsize=10)
        
        # 设置轴标签和标题
        plt.yticks(y_pos, traits, fontsize=12)
        plt.xlabel('缺陷比例 (%)', fontsize=12)
        plt.title('缺陷性状牛只占比', fontsize=14, pad=20)
        
        # 添加网格线
        plt.grid(True, alpha=0.3, axis='x')
        
        # 调整布局
        plt.tight_layout()
        
        # 保存图表
        output_file = os.path.join(output_dir, '缺陷性状牛只占比.png')
        plt.savefig(output_file, dpi=300, bbox_inches='tight')
        plt.close()
        
        # 同时保存数据到Excel文件
        results_df = pd.DataFrame({
            '性状': traits,
            '缺陷比例(%)': ratios
        })
        excel_file = os.path.join(output_dir, '缺陷性状牛只占比.xlsx')
        results_df.to_excel(excel_file, index=False)
        
        print("已生成缺陷性状占比图表和Excel文件")

        return defect_ratios

    def generate_defect_pie_charts(self, output_dir):
        # 设置中文字体
        plt.rcParams['font.sans-serif'] = ['SimHei']  
        plt.rcParams['axes.unicode_minus'] = False
        
        df = self.conformation_df
        pie_charts_dir = os.path.join(output_dir, "各性状缺陷性状占比-饼图")
        os.makedirs(pie_charts_dir, exist_ok=True)
        
        # 数据类型转换
        for col in df.columns:
            if col not in ['牧场', '牛号']:
                df[col] = pd.to_numeric(df[col], errors='coerce')

        # 定义三种分组方式
        groups = {
            # 第一类: 体高
            '体高': {
                'ranges': [(1, 3), (4, 8), (9, 9)],
                'labels': ['1-3', '4-8', '9'],
                'colors': ['#FFA07A', '#87CEEB', '#90EE90']  # 三种颜色
            },
            # 第二类: 大部分性状
            'type2': {
                'ranges': [(1, 3), (4, 6), (7, 9)],
                'labels': ['1-3', '4-6', '7-9'],
                'colors': ['#FFA07A', '#87CEEB', '#90EE90']
            },
            # 第三类: 乳房相关性状
            'type3': {
                'ranges': [(1, 2), (3, 7), (8, 9)],
                'labels': ['1-2', '3-7', '8-9'],
                'colors': ['#FFA07A', '#87CEEB', '#90EE90']
            }
        }

        # 定义每个性状的分组类型
        trait_types = {
            '体高': '体高',
            '乳房深度': 'type3', 
            '前乳头位置': 'type3',
            '后乳头位置': 'type3'
        }
        
        # 默认type2的性状
        type2_traits = ['胸宽', '体深', '腰强度', '尻宽', '蹄角度', '蹄踵深度', 
                    '骨质地', '后肢后视', '中央悬韧带', '前乳房附着', 
                    '后乳房附着高度', '后乳房附着宽度', '棱角性',
                    '尻角度', '后肢侧视', '前乳头长度']
        
        for trait in type2_traits:
            trait_types[trait] = 'type2'

        # 生成饼图
        for trait, group_type in trait_types.items():
            if trait not in df.columns:
                continue
                
            plt.figure(figsize=(10, 8))
            valid_records = df[trait].dropna()
            
            if len(valid_records) == 0:
                continue
                
            group_info = groups[group_type]
            proportions = []
            
            # 计算每个范围的比例
            for range_min, range_max in group_info['ranges']:
                count = len(valid_records[(valid_records >= range_min) & (valid_records <= range_max)])
                proportion = (count / len(valid_records)) * 100
                proportions.append(proportion)
            
            # 绘制饼图
            plt.pie(proportions,
                labels=group_info['labels'],
                colors=group_info['colors'],
                autopct='%1.0f%%',
                startangle=90)
            
            plt.title(trait, pad=20, fontsize=14)
            
            # 保存图片
            plt.savefig(os.path.join(pie_charts_dir, f'{trait}.png'),
                    bbox_inches='tight',
                    dpi=300)
            plt.close()
            
            print(f"已生成 {trait} 的饼图")


    # PPT生成-生成前检查是否已存在准备文件。
    def generate_ppt_report(self):
        # 检查是否选择了输出文件夹
        if not hasattr(self, 'output_folder') or not self.output_folder:
            QMessageBox.warning(self, "警告", "请先选择输出文件夹。")
            return

        # 检查是否生成了"结果-指数排名结果.xlsx"
        ranking_result_file = os.path.join(self.output_folder, "结果-指数排名结果.xlsx")
        if not os.path.exists(ranking_result_file):
            reply = QMessageBox.question(self, '确认', "未找到指数排名结果，是否现在生成？",
                                        QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No)
            if reply == QMessageBox.StandardButton.Yes:
                self.run_index_ranking()
            else:
                return

        # 检查是否生成了"结果-母牛关键性状指数.xlsx"
        key_traits_file = os.path.join(self.output_folder, "结果-母牛关键性状指数.xlsx")
        if not os.path.exists(key_traits_file):
            reply = QMessageBox.question(self, '确认', "未找到母牛关键性状指数，是否现在生成？",
                                        QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No)
            if reply == QMessageBox.StandardButton.Yes:
                self.calculate_key_traits_index()
            else:
                return

        # 获取牧场名称
        farm_name, ok = QInputDialog.getText(self, '输入牧场名称', '请输入牧场名称:')
        if ok and farm_name:
            self.farm_name = farm_name
        else:
            self.farm_name = "某某某"  # 默认牧场名称

        # 所有检查通过，开始生成PPT
        self.create_ppt_report()

    # PPT生成-主文件
    def create_ppt_report(self):
        try:
            # 检查是否存在PPT模板
            template_file = self.find_ppt_template()
            if template_file:
                # 直接使用模板创建新的演示文稿
                prs = Presentation(template_file)
                # 删除模板中的所有现有幻灯片
                for i in range(len(prs.slides) - 1, -1, -1):
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
            else:
                prs = Presentation()
                # 设置幻灯片大小为16:9
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)

            # 生成标题页-P1
            self.create_title_slide(prs)

            # 生成目录页-P2
            self.create_toc_slide(prs)

            # 生成牧场系谱记录分析页(标题分隔页）-P3
            self.create_pedigree_analysis_slide(prs)

            # 生成牧场系谱完整性分析页-P4
            self.create_pedigree_completeness_slide(prs)

            # 生成未识别牛只明细页-P5
            self.create_unidentified_cattle_slide(prs)

            # 生成系谱记录准确性页-P6
            self.create_pedigree_accuracy_slide(prs)
            
            # 生成牛群遗传数据评估页(标题分隔页）-P7
            self.create_key_traits_analysis_slide(prs)

            # 生成牧场关键育种性状分析页-P8
            self.create_key_breeding_traits_slide(prs)

            # 生成牧场净利润值NM$分布情况分析页-P9
            self.create_nm_distribution_slide(prs)

            # 生成牧场净利润值正态分布分析页-P10
            self.create_nm_distribution_analysis_slide(prs)

            # 生成牧场育种指数正态分布分析页-P11
            self.create_breeding_index_distribution_analysis_slide(prs)

            # 生成牧场关键育种性状进展分析页-P12+
            self.create_key_traits_progress_slides(prs)
            
            # 生成牧场关键育种性状进展分析页(标题分隔页）-NO_03        
            self.create_linear_score_title_slide(prs)

            # 生成牧场牛群体型外貌线性评分-3.1
            self.create_linear_score_slide(prs)

            # 生成牧场牛群体型外貌缺陷性状占比-3.2
            self.create_linear_defect_traits_slide(prs)

            # 生成牧场缺陷性状具体情况-3.3
            self.create_all_linear_defect_traits_slide(prs)            

            # 生成牧场牛群体型外貌进展情况-3.4
            self.create_linear_progress_slide(prs)

            # 生成牧场牧场公牛使用分析-04
            self.create_used_bull_slide(prs)

            # 生成牧场近几年公牛使用整体情况-根据配种记录-4.1
            self.create_bull_usage_slide(prs)

            # 生成牧场牧场近几年公牛使用明细-根据配种记录-4.2
            self.create_bull_usage_detail_slides(prs)

            # 生成牧场牧场近几年公牛使用进展-根据配种记录-4.3
            self.create_semen_usage_progress_slides(prs)

            # 生成牧场牧场近几年公牛使用单值图-根据配种记录-4.4
            self.create_semen_usage_scatter_plot_slide(prs)

            # 生成-根据配种记录-4.4
            self.create_bye_slide(prs)

            # TODO: 在这里添加生成其他页面的代码

            # 保存PPT
            output_path = os.path.join(self.output_folder, f"{self.farm_name}牧场遗传改良项目专项服务报告.pptx")
            prs.save(output_path)

            QMessageBox.information(self, "成功", f"PPT报告已生成：{output_path}")

        except Exception as e:
            logging.error(f"生成PPT报告时发生错误: {str(e)}")
            QMessageBox.critical(self, "错误", f"生成PPT报告时发生错误：{str(e)}")

    # PPT生成-寻找PPT模板
    def find_ppt_template(self):
        for file in os.listdir(self.output_folder):
            if file.lower().endswith(('.ppt', '.pptx')):
                return os.path.join(self.output_folder, file)
        return None
    
    # PPT生成-第一页-标题页
    def create_title_slide(self, prs):
        # 添加标题页
        slide_layout = prs.slide_layouts[0]  # 使用第一个布局（通常是标题页布局）
        slide = prs.slides.add_slide(slide_layout)

        # 设置标题
        title = f"{self.farm_name}牧场遗传改良项目专项服务"
        left = Inches(2)  # 调整左边距
        top = Inches(4.5)  # 调整上边距
        width = Inches(12)  # 调整宽度
        height = Inches(1.5)
        title_shape = slide.shapes.add_textbox(left, top, width, height)

        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.name = "微软雅黑"
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        title_frame.paragraphs[0].font.bold = True

        # 添加作者和日期
        left = Inches(9)  # 调整到右下角
        top = Inches(6)
        width = Inches(3.5)
        height = Inches(1.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = f"奶科院育种中心 {self.username}"
        p.font.name = "微软雅黑"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.RIGHT  # 右对齐

        p = tf.add_paragraph()
        current_date = datetime.datetime.now().strftime("%Y年%m月%d日")
        p.text = current_date
        p.font.name = "微软雅黑"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.RIGHT  # 右对齐

    # PPT生成-第二页-目录页
    def create_toc_slide(self, prs):
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加"目录"标题
        title_left = Inches(0.5)
        title_top = Inches(0.5)
        title_width = Inches(2)
        title_height = Inches(1)
        title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_shape.text_frame
        title_frame.text = "目录"
        title_frame.paragraphs[0].font.name = "微软雅黑"
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)  # 蓝色
        title_frame.paragraphs[0].font.bold = True

        # 定义网格布局
        cell_width = Inches(3)
        cell_height = Inches(2)  # 增加单元格高度
        grid_items = [
            ("01", "牧场系谱记录分析", "系谱识别率分析 未识别牛只明细 系谱记录准确性"),
            ("02", "牛群遗传数据评估", "牧场关键育种性状 净利润值分布情况 净利润值及育种指数正态分布 关键育种性状进展分析"),
            ("03", "牛群体型外貌评定", "体型外貌线性评分 体型外貌缺陷性状占比 缺陷性状具体情况 体型外貌进展情况"),
            ("04", "牧场公牛使用分析", "过去3年公牛使用明细 公牛使用时段与频率分析 隐性遗传缺陷基因分析"),
            ("05", "育种方向与指数确定", "育种方向确定 育种指数设定"),
            ("06", "牧场选配方案设计", "牧场选配方案推荐 不同选配方案对生产的预估影响")
        ]

        # 计算内容的总宽度和高度以实现居中
        columns = 3
        spacing = Inches(0.2)
        total_grid_width = columns * cell_width + (columns - 1) * spacing
        slide_width = prs.slide_width
        grid_left = (slide_width - total_grid_width) / 2  # 水平居中

        rows = (len(grid_items) + columns - 1) // columns
        total_grid_height = rows * cell_height + (rows - 1) * spacing
        slide_height = prs.slide_height
        grid_top = (slide_height - total_grid_height) / 2 + Inches(0.2)  # 垂直居中并稍微上移

        for i, (number, title, subtitle) in enumerate(grid_items):
            left = grid_left + (i % columns) * (cell_width + spacing)
            top = grid_top + (i // columns) * (cell_height + spacing)

            # 添加编号
            number_box = slide.shapes.add_textbox(left, top, Inches(0.5), Inches(0.5))
            number_frame = number_box.text_frame
            number_frame.text = number
            number_frame.paragraphs[0].font.name = "微软雅黑"
            number_frame.paragraphs[0].font.size = Pt(24)
            number_frame.paragraphs[0].font.color.rgb = RGBColor(0, 176, 240)  # 淡蓝色

            """# 添加横线
            line_top = top + Inches(0.4)
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                left,
                line_top,
                left + cell_width,
                line_top
            )
            line.line.color.rgb = RGBColor(0, 176, 240)  # 淡蓝色"""

            # 添加主标题
            title_box = slide.shapes.add_textbox(left, top + Inches(0.5), cell_width, Inches(0.4))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.name = "微软雅黑"
            title_frame.paragraphs[0].font.size = Pt(18)
            title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # 黑色
            title_frame.paragraphs[0].font.bold = True

            # 添加副标题，每个项另起一行并添加小点
            subtitle_box = slide.shapes.add_textbox(left, top + Inches(0.9), cell_width, Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.clear()  # 清除默认的空段落

            subtitle_items = subtitle.split()
            for item in subtitle_items:
                p = subtitle_frame.add_paragraph()
                p.text = "• " + item  # 直接在文本前添加项目符号
                p.font.name = "微软雅黑"
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(89, 89, 89)  # 深灰色
                p.level = 1  # 缩进级别

        # 添加右上角的logo
        logo_path = "path_to_your_logo.png"  # 替换为您的logo文件路径
        if os.path.exists(logo_path):
            left = Inches(9)
            top = Inches(0.5)
            width = Inches(1)
            height = Inches(1)
            slide.shapes.add_picture(logo_path, left, top, width, height)

    # PPT生成-第01-牧场系谱记录分析-标题
    def create_pedigree_analysis_slide(self, prs):
        slide_layout = prs.slide_layouts[1]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        """# 设置背景颜色为蓝色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 112, 192)  # 蓝色

        # 添加椭圆形状
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(6.5)
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白色
        shape.line.color.rgb = RGBColor(255, 255, 255)  # 白色边框"""

        # 添加序号 "01"
        number_box = slide.shapes.add_textbox(Inches(5.67), Inches(1.12), Inches(2), Inches(2))
        number_frame = number_box.text_frame
        number_frame.text = "01"
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = number_frame.paragraphs[0].font
        font.name = "Arial"
        font.size = Pt(166)
        font.color.rgb = RGBColor(255, 255, 255)  # 白色
        font.transparency = 69  # 设置透明度为69%

        # 添加标题 "牧场系谱记录分析"
        title_box = slide.shapes.add_textbox(Inches(2.67), Inches(4.13), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "牧场系谱记录分析"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = title_frame.paragraphs[0].font
        font.name = "微软雅黑"
        font.size = Pt(40)
        font.color.rgb = RGBColor(255, 255, 255)  # 白色

        # 添加页脚
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(1), Inches(0.3))
        footer_frame = footer_box.text_frame
        footer_frame.text = "3"
        footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = footer_frame.paragraphs[0].font
        font.name = "微软雅黑"
        font.size = Pt(12)
        font.color.rgb = RGBColor(255, 255, 255)  # 白色

    # PPT生成-第1.1-牧场系谱完整性（识别率）分析
    def create_pedigree_completeness_slide(self, prs):
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = f"1.1 {self.farm_name}牧场系谱完整性分析"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(1.9), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)

        # 添加表格
        table_rows, table_cols = 7, 9
        left, top, width, height = Inches(0.874), Inches(1.72), Inches(11.53), Inches(3.14)
        table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table

        # 检查是否存在系谱识别情况分析文件
        analysis_file = os.path.join(self.output_folder, "结果-系谱识别情况分析.xlsx")
        if not os.path.exists(analysis_file):
            reply = QMessageBox.question(self, '确认', "未找到系谱识别情况分析结果，点击Yes确认退出，点击No跳过",
                                        QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No)
            if reply == QMessageBox.StandardButton.Yes:
                # 用户选择退出
                return  # 这将退出当前函数
            else:
                # 用户选择跳过
                # 创建一个空的DataFrame来填充表格
                df = pd.DataFrame(columns=['是否在场', '出生年份', '头数', '父号可识别头数', '父号识别率', 
                                        '外祖父可识别头数', '外祖父识别率', '外曾外祖父可识别头数', '外曾外祖父识别率'])

        else:
            # 读取数据
            df = pd.read_excel(analysis_file)
            df = df[df['是否在场'] == '是']
            
            # 设置表头
            headers = df.columns.tolist()
            headers[headers.index('birth_year_group')] = '出生年份'
            
            # 设置表头
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = header
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
                para = cell.text_frame.paragraphs[0]
                para.font.name = "微软雅黑"
                para.font.size = Pt(14)
                para.font.bold = True
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 填充数据
            for row in range(1, min(7, len(df) + 1)):
                for col in range(9):
                    cell = table.cell(row, col)
                    value = str(df.iloc[row-1, col])
                    cell.text = value
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(233, 237, 244)
                    para = cell.text_frame.paragraphs[0]
                    para.font.name = "微软雅黑"
                    para.font.size = Pt(12)
                    para.font.color.rgb = RGBColor(0, 0, 0)
                    para.alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                    # 根据不同列的阈值设置字体样式
                    if headers[col] == '父号识别率':
                        try:
                            rate = float(value.strip('%'))
                            if rate < 90.00:
                                para.font.color.rgb = RGBColor(255, 0, 0)  # 红色
                                para.font.bold = True
                        except ValueError:
                            pass  # 如果无法转换为浮点数,保持原样式

                    elif headers[col] == '外祖父识别率':
                        try:
                            rate = float(value.strip('%'))
                            if rate < 70.00:
                                para.font.color.rgb = RGBColor(255, 0, 0)  # 红色
                                para.font.bold = True
                        except ValueError:
                            pass

                    elif headers[col] == '外曾外祖父识别率':
                        try:
                            rate = float(value.strip('%'))
                            if rate < 50.00:
                                para.font.color.rgb = RGBColor(255, 0, 0)  # 红色
                                para.font.bold = True
                        except ValueError:
                            pass

        # 添加文本框
        left, top, width, height = Inches(0.469), Inches(5.012), Inches(12.358), Inches(2.22)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = "未识别原因："
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        p = tf.add_paragraph()
        p.text = ""
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

        p = tf.add_paragraph()
        p.text = "父号未识别原因："
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

        p = tf.add_paragraph()
        p.text = ""
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

        p = tf.add_paragraph()
        p.text = "外祖父号未识别原因："
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

        p = tf.add_paragraph()
        p.text = ""
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

    # PPT生成-第1.2-未识别牛只明细
    def create_unidentified_cattle_slide(self, prs):
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = f"1.2 {self.farm_name}牧场未识别牛只明细"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(1.9), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)

        # 添加文本框
        left, top, width, height = Inches(0.583), Inches(1.5), Inches(11), Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = "未识别牛只明细：(暂不能自动化生成........需服务人员整理........）"
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = False
        p.alignment = PP_ALIGN.LEFT

    # PPT生成-第1.3-系谱记录准确性
    def create_pedigree_accuracy_slide(self, prs):
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = f"1.3 {self.farm_name}牧场系谱记录准确性分析"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(1.9), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)

        # 添加文本框
        left, top, width, height = Inches(0.583), Inches(1.5), Inches(11), Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = "系谱记录准确性：(暂不能自动化生成........需服务人员整理........根据基因组检测结果...........）"
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = False
        p.alignment = PP_ALIGN.LEFT

    # PPT生成-第02-牛群遗传数据评估-标题
    def create_key_traits_analysis_slide(self, prs):
        slide_layout = prs.slide_layouts[1]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加序号 "02"
        number_box = slide.shapes.add_textbox(Inches(5.67), Inches(1.12), Inches(2), Inches(2))
        number_frame = number_box.text_frame
        number_frame.text = "02"
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = number_frame.paragraphs[0].font
        font.name = "Arial"
        font.size = Pt(166)
        font.color.rgb = RGBColor(255, 255, 255)  # 白色
        font.transparency = 69  # 设置透明度为69%

        # 添加标题 "牧场系谱记录分析"
        title_box = slide.shapes.add_textbox(Inches(2.67), Inches(4.13), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "牛群遗传数据评估"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = title_frame.paragraphs[0].font
        font.name = "微软雅黑"
        font.size = Pt(40)
        font.color.rgb = RGBColor(255, 255, 255)  # 白色

    # PPT生成-第2.1-牛群遗传关键性状数据分析-按出生年份
    def create_key_breeding_traits_slide(self, prs):
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = f"2.1 {self.farm_name}牧场关键育种性状分析"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(1.9), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)

        # 检查是否存在牛群关键性状年度变化文件
        analysis_file = os.path.join(self.output_folder, "结果-牛群关键性状年度变化.xlsx")
        if not os.path.exists(analysis_file):
            reply = QMessageBox.question(self, '确认', "未找到牛群关键性状年度变化结果，点击Yes确认退出，点击No跳过",
                                        QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No)
            if reply == QMessageBox.StandardButton.Yes:
                return  # 退出当前函数
            else:
                # 创建一个空的DataFrame来填充表格
                df = pd.DataFrame(columns=['是否在场', '出生年份', '头数'])
        else:
            # 读取数据
            df = pd.read_excel(analysis_file)
            df = df[df['是否在场'] == '是']
            
            # 设置表头
            headers = df.columns.tolist()
            headers[headers.index('birth_year_group')] = '出生年份'

        # 创建翻译字典
        translation_dict = {
            'TPI': '育种综合指数', 'NM$': '净利润值', 'CM$': '奶酪净值', 'FM$': '液奶净值', 'GM$': '放牧净值',
            'MILK': '产奶量', 'FAT': '乳脂量', 'PROT': '乳蛋白量', 'FAT %': '乳脂率', 'PROT%': '乳蛋白率',
            'SCS': '体细胞指数', 'DPR': '女儿怀孕率', 'HCR': '青年牛受胎率', 'CCR': '成母牛受胎率', 'PL': '生产寿命',
            'SCE': '配种产犊难易度', 'DCE': '女儿产犊难易度', 'SSB': '配种死胎率', 'DSB': '女儿死淘率',
            'PTAT': '体型综合指数', 'UDC': '乳房综合指数', 'FLC': '肢蹄综合指数', 'BDC': '体重综合指数',
            'ST': '体高', 'SG': '强壮度', 'BD': '体深', 'DF': '乳用特征', 'RA': '尻角度', 'RW': '尻宽',
            'LS': '后肢侧视', 'LR': '后肢后视', 'FA': '蹄角度', 'FLS': '肢蹄评分', 'FU': '前方附着',
            'UH': '后房高度', 'UW': '后方宽度', 'UC': '悬韧带', 'UD': '乳房深度', 'FT': '前乳头位置',
            'RT': '后乳头位置', 'TL': '乳头长度', 'FE': '饲料效率指数', 'FI': '繁殖力指数', 'HI': '健康指数',
            'LIV': '存活率', 'GL': '妊娠长度', 'MAST': '乳房炎抗病性', 'MET': '子宫炎抗病性', 'RP': '胎衣不下抗病性',
            'KET': '酮病抗病性', 'DA': '变胃抗病性', 'MFV': '产后瘫抗病性', 'EFC': '首次产犊月龄',
            'HLiv': '后备牛存活率', 'FS': '饲料节约指数', 'RFI': '剩余饲料采食量', 'Milk Speed': '排乳速度',
            'Eval Date': '数据日期'
        }

        # 添加表格
        table_rows, table_cols = 8, len(headers)
        left, top, width, height = Inches(0.1535), Inches(1.5945), Inches(13.0276), Inches(3.76)
        table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table

        # 设置表头（两行）
        for i, header in enumerate(headers):
            # 第一行（中文）
            cell = table.cell(0, i)
            cell.text = translation_dict.get(header, header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
            para = cell.text_frame.paragraphs[0]
            para.font.name = "微软雅黑"
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 第二行（英文）
            cell = table.cell(1, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
            para = cell.text_frame.paragraphs[0]
            para.font.name = "Arial"
            para.font.size = Pt(10)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 填充数据
        for row in range(2, min(8, len(df) + 2)):
            for col in range(len(headers)):
                cell = table.cell(row, col)
                if row - 2 < len(df):
                    value = str(df.iloc[row-2, col])
                    cell.text = value
                else:
                    cell.text = ""
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(233, 237, 244)
                para = cell.text_frame.paragraphs[0]
                para.font.name = "微软雅黑"
                para.font.size = Pt(10)
                para.font.color.rgb = RGBColor(0, 0, 0)
                para.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 添加文本框
        left, top, width, height = Inches(0.469), Inches(5.6614), Inches(12.358), Inches(2.22)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = "需关注的关键性状："
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        p = tf.add_paragraph()
        p.text = ""
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

       
        p = tf.add_paragraph()
        p.text = "暂不能自动化生成........需服务人员整理........"
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

    # PPT生成-第2.2-牧场净利润值NM$分布情况分析
    def create_nm_distribution_slide(self, prs):
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = f"2.2 {self.farm_name}牧场净利润值NM$分布情况分析"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(9), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)

        # 读取数据
        excel_file = os.path.join(self.output_folder, "在群牛只净利润值分布.xlsx")
        df = pd.read_excel(excel_file, sheet_name="NM$分布")

        # 添加柱状图
        chart_left = Inches(0.5)
        chart_top = Inches(1.5)
        chart_width = Inches(6)
        chart_height = Inches(3.5)
        chart_data = CategoryChartData()
        chart_data.categories = df['NM$区间']
        chart_data.add_series('头数', df['头数'])

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, chart_left, chart_top, chart_width, chart_height, chart_data
        ).chart
        chart.has_title = False
        chart.has_legend = True

        # 设置柱状图字体大小和图例位置
        chart.font.size = Pt(9)
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.horz_offset = 0.1  # 调整这个值来将图例移动到更靠近顶部的位置
        chart.legend.include_in_layout = False

        # 添加数据标签
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(9)
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.number_format = '0'
        data_labels.show_value = True

        # 确保数据标签显示
        for point in plot.series[0].points:
            point.data_label.show_value = True

        # 设置 x 轴和 y 轴的字体大小
        chart.category_axis.tick_labels.font.size = Pt(9)
        chart.value_axis.tick_labels.font.size = Pt(9)

        # 添加饼图
        pie_left = Inches(7)
        pie_top = Inches(1.5)
        pie_width = Inches(6)
        pie_height = Inches(3.2244)
        pie_data = CategoryChartData()
        pie_data.categories = df['NM$区间']
        pie_data.add_series('占比', (df['占比'] ).tolist())

        pie = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, pie_left, pie_top, pie_width, pie_height, pie_data
        ).chart
        pie.has_title = False
        pie.has_legend = True

        # 设置饼图字体大小
        pie.font.size = Pt(9)

        # 添加数据标签（百分比形式）
        data_labels.show_value = False
        plot = pie.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(9)
        data_labels.number_format = '0.00%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

        # 添加表格
        table_rows, table_cols = len(df) + 1, 3
        left, top, width, height = Inches(7), Inches(4.7834), Inches(6), Inches(2.433)
        table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table

        # 设置表头
        headers = ['NM$区间', '头数', '占比']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
            para = cell.text_frame.paragraphs[0]
            para.font.name = "微软雅黑"
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 填充数据
        for row in range(1, table_rows):
            for col in range(table_cols):
                cell = table.cell(row, col)
                value = df.iloc[row-1, col]
                if col == 2:  # 占比列
                    cell.text = f"{value:.2%}"
                else:
                    cell.text = str(value)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(233, 237, 244)
                para = cell.text_frame.paragraphs[0]
                para.font.name = "微软雅黑"
                para.font.size = Pt(12)
                para.font.color.rgb = RGBColor(0, 0, 0)
                para.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 添加文本框
        left, top, width, height = Inches(0.5), Inches(5), Inches(6), Inches(2.4)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = "净利润值NM$分布情况："
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        p = tf.add_paragraph()
        p.text = "暂不能自动化生成........需服务人员整理........"
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

    # PPT生成-第2.3-牧场净利润值正态分布分析
    def create_nm_distribution_analysis_slide(self, prs):
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = f"2.3 {self.farm_name}牧场净利润值正态分布分析"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(9), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)

        # 添加图片
        img_path = os.path.join(self.output_folder, "NM$_年份正态分布.png")
        if os.path.exists(img_path):
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(12)
            height = Inches(4)
            slide.shapes.add_picture(img_path, left, top, width, height)
        else:
            reply = QMessageBox.question(self, '确认', "未找到NM$_年份正态分布.png，需进行'在群母牛关键性状计算'。点击Yes确认退出，点击No跳过",
                                        QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No)
            if reply == QMessageBox.StandardButton.Yes:
                return False  # 用户选择退出

        # 添加文本框
        left, top, width, height = Inches(0.469), Inches(5.6614), Inches(12.358), Inches(2.22)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = "牧场净利润值正态分布情况："
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        p = tf.add_paragraph()
        p.text = ""
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

        p = tf.add_paragraph()
        p.text = "暂不能自动化生成........需服务人员整理........"
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

        return True  # 成功创建幻灯片

    # PPT生成-第2.4-牧场育种指数正态分布分析
    def create_breeding_index_distribution_analysis_slide(self, prs):
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = f"2.4 {self.farm_name}牧场育种指数正态分布分析"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(9), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)

        # 添加图片
        img_path = os.path.join(self.output_folder, "育种指数得分_年份正态分布.png")
        if os.path.exists(img_path):
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(12)
            height = Inches(4)
            slide.shapes.add_picture(img_path, left, top, width, height)
        else:
            reply = QMessageBox.question(self, '确认', "未找到育种指数得分_年份正态分布.png，需进行'在群母牛指数计算及排名'。点击Yes确认退出，点击No跳过",
                                        QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No)
            if reply == QMessageBox.StandardButton.Yes:
                return False  # 用户选择退出

        # 添加文本框
        left, top, width, height = Inches(0.469), Inches(5.6614), Inches(12.358), Inches(2.22)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = "牧场育种指数正态分布情况："
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        p = tf.add_paragraph()
        p.text = ""
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

        p = tf.add_paragraph()
        p.text = "暂不能自动化生成........需服务人员整理........"
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

        return True  # 成功创建幻灯片

    # PPT生成-第2.5-牧场关键育种性状进展分析
    def create_key_traits_progress_slides(self, prs, selected_traits=None):
        if selected_traits is None:
            # 尝试从实例属性中获取
            if hasattr(self, 'selected_traits') and self.selected_traits:
                selected_traits = self.selected_traits
            else:
                # 从selected_traits_key_traits.txt文件中读取
                selected_traits_file = os.path.join(self.output_folder, "selected_traits_key_traits.txt")
                if os.path.exists(selected_traits_file):
                    with open(selected_traits_file, "r", encoding="utf-8") as f:
                        selected_traits = [line.strip() for line in f if line.strip()]
                else:
                    QMessageBox.warning(self, "警告", "未找到关键性状列表，请先计算关键性状。")
                    return False        




        charts_folder = os.path.join(self.output_folder, "结果-牛群关键性状进展图")

        if not os.path.exists(charts_folder):
            reply = QMessageBox.question(self, '确认', "未找到结果-牛群关键性状进展图文件夹，需进行'在群母牛关键性状计算'。点击Yes确认退出，点击No跳过",
                                        QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No)
            if reply == QMessageBox.StandardButton.Yes:
                return False  # 用户选择退出
            else:
                return True  # 用户选择跳过

        for trait in selected_traits:
            img_filename = f"{trait}进展情况.png"
            img_path = os.path.join(charts_folder, img_filename)
            if not os.path.exists(img_path):
                continue  # 如果图片不存在，跳过

            slide_layout = prs.slide_layouts[5]  # 使用空白布局
            slide = prs.slides.add_slide(slide_layout)

            # 添加标题
            title = f"2.5 {self.farm_name}牧场关键育种性状进展分析"
            left, top, width, height = Inches(0.583), Inches(0.732), Inches(9), Inches(0.571)
            title_shape = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_shape.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.alignment = PP_ALIGN.LEFT
            title_para.font.name = "微软雅黑"
            title_para.font.size = Pt(28)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(0, 0, 0)

            # 添加图片
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(12)
            slide.shapes.add_picture(img_path, left, top, width=width)

            # 添加文本框
            left, top, width, height = Inches(0.469), Inches(5.6614), Inches(12.358), Inches(2.22)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame

            translation_dict = {
                'TPI': '育种综合指数', 'NM$': '净利润值', 'CM$': '奶酪净值', 'FM$': '液奶净值', 'GM$': '放牧净值',
                'MILK': '产奶量', 'FAT': '乳脂量', 'PROT': '乳蛋白量', 'FAT %': '乳脂率', 'PROT%': '乳蛋白率',
                'SCS': '体细胞指数', 'DPR': '女儿怀孕率', 'HCR': '青年牛受胎率', 'CCR': '成母牛受胎率', 'PL': '生产寿命',
                'SCE': '配种产犊难易度', 'DCE': '女儿产犊难易度', 'SSB': '配种死胎率', 'DSB': '女儿死淘率',
                'PTAT': '体型综合指数', 'UDC': '乳房综合指数', 'FLC': '肢蹄综合指数', 'BDC': '体重综合指数',
                'ST': '体高', 'SG': '强壮度', 'BD': '体深', 'DF': '乳用特征', 'RA': '尻角度', 'RW': '尻宽',
                'LS': '后肢侧视', 'LR': '后肢后视', 'FA': '蹄角度', 'FLS': '肢蹄评分', 'FU': '前方附着',
                'UH': '后房高度', 'UW': '后方宽度', 'UC': '悬韧带', 'UD': '乳房深度', 'FT': '前乳头位置',
                'RT': '后乳头位置', 'TL': '乳头长度', 'FE': '饲料效率指数', 'FI': '繁殖力指数', 'HI': '健康指数',
                'LIV': '存活率', 'GL': '妊娠长度', 'MAST': '乳房炎抗病性', 'MET': '子宫炎抗病性', 'RP': '胎衣不下抗病性',
                'KET': '酮病抗病性', 'DA': '变胃抗病性', 'MFV': '产后瘫抗病性', 'EFC': '首次产犊月龄',
                'HLiv': '后备牛存活率', 'FS': '饲料节约指数', 'RFI': '剩余饲料采食量', 'Milk Speed': '排乳速度',
                'Eval Date': '数据日期'
            }


            # 设置标题和标签
            chinese_trait = translation_dict.get(trait, trait)  # 如果没有翻译，就使用原来的trait
            p = tf.add_paragraph()
            p.text = f"{chinese_trait} ({trait}) 性状进展情况："
            p.font.name = "微软雅黑"
            p.font.size = Pt(16)
            p.font.bold = True
            p.alignment = PP_ALIGN.LEFT

            p = tf.add_paragraph()
            p.text = ""
            p.font.name = "微软雅黑"
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.LEFT
            p.level = 1

            p = tf.add_paragraph()
            p.text = "暂不能自动化生成........需服务人员整理........"
            p.font.name = "微软雅黑"
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.LEFT
            p.level = 1

        return True  # 成功创建所有幻灯片

    # PPT生成-第03-牛群体型外貌评定-标题
    def create_linear_score_title_slide(self, prs):
        slide_layout = prs.slide_layouts[1]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加序号 "03"
        number_box = slide.shapes.add_textbox(Inches(5.67), Inches(1.12), Inches(2), Inches(2))
        number_frame = number_box.text_frame
        number_frame.text = "03"
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = number_frame.paragraphs[0].font
        font.name = "Arial"
        font.size = Pt(166)
        font.color.rgb = RGBColor(255, 255, 255)  # 白色
        font.transparency = 69  # 设置透明度为69%

        # 添加标题 "牧场系谱记录分析"
        title_box = slide.shapes.add_textbox(Inches(2.67), Inches(4.13), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "牛群体型外貌评定"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = title_frame.paragraphs[0].font
        font.name = "微软雅黑"
        font.size = Pt(40)
        font.color.rgb = RGBColor(255, 255, 255)  # 白色

    # PPT生成-第3.1-牛群体型外貌线性评分
    def create_linear_score_slide(self, prs):
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = f"3.1 {self.farm_name}牧场牛群体型外貌线性评分"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(1.9), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)

        # 添加文本框
        left, top, width, height = Inches(0.583), Inches(1.5), Inches(11), Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = "牛群体型外貌线性评分分析：(暂不能自动化生成........需服务人员整理........）"
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = False
        p.alignment = PP_ALIGN.LEFT

    # PPT生成-第3.2-牛群体型外貌缺陷性状占比
    def create_linear_defect_traits_slide(self, prs):
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = f"3.2 {self.farm_name}牧场牛群体型外貌缺陷性状占比"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(1.9), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)

        # 添加文本框
        left, top, width, height = Inches(0.583), Inches(1.5), Inches(11), Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = "牛群体型外貌缺陷性状占比分析：(暂不能自动化生成........需服务人员整理........）"
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = False
        p.alignment = PP_ALIGN.LEFT

    # PPT生成-第3.3-缺陷性状具体情况
    def create_all_linear_defect_traits_slide(self, prs):
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = f"3.3 {self.farm_name}牧场缺陷性状具体情况"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(1.9), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)

        # 添加文本框
        left, top, width, height = Inches(0.583), Inches(1.5), Inches(11), Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = "缺陷性状具体情况：(暂不能自动化生成........需服务人员整理........）"
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = False
        p.alignment = PP_ALIGN.LEFT

    # PPT生成-第3.4-牛群体型外貌进展情况
    def create_linear_progress_slide(self, prs):
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = f"3.4 {self.farm_name}牧场牛群体型外貌进展情况"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(1.9), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)

        # 添加文本框
        left, top, width, height = Inches(0.583), Inches(1.5), Inches(11), Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = "牛群体型外貌进展情况：(暂不能自动化生成........需服务人员整理........）"
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = False
        p.alignment = PP_ALIGN.LEFT

    # PPT生成-第04-牧场公牛使用分析-标题
    def create_used_bull_slide(self, prs):
        slide_layout = prs.slide_layouts[1]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加序号 "04"
        number_box = slide.shapes.add_textbox(Inches(5.67), Inches(1.12), Inches(2), Inches(2))
        number_frame = number_box.text_frame
        number_frame.text = "04"
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = number_frame.paragraphs[0].font
        font.name = "Arial"
        font.size = Pt(166)
        font.color.rgb = RGBColor(255, 255, 255)  # 白色
        font.transparency = 69  # 设置透明度为69%

        # 添加标题 "牧场系谱记录分析"
        title_box = slide.shapes.add_textbox(Inches(2.67), Inches(4.13), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "牧场公牛使用分析"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = title_frame.paragraphs[0].font
        font.name = "微软雅黑"
        font.size = Pt(40)
        font.color.rgb = RGBColor(255, 255, 255)  # 白色

    # PPT生成-第4.1-牧场近几年公牛使用整体情况-根据配种记录
    def create_bull_usage_slide(self, prs):
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = f"4.1 {self.farm_name}牧场近几年公牛使用整体情况"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(8), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)

        # 检查是否存在冻精使用情况报告文件
        report_file = os.path.join(self.output_folder, "冻精使用情况报告.xlsx")
        if not os.path.exists(report_file):
            reply = QMessageBox.question(self, '确认', "未找到冻精使用情况报告，点击Yes确认退出，点击No跳过",
                                        QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No)
            if reply == QMessageBox.StandardButton.Yes:
                return  # 退出当前函数
            else:
                # 创建一个空的DataFrame来填充表格
                df = pd.DataFrame()
        else:
            # 读取数据
            df = pd.read_excel(report_file)

        # 创建翻译字典
        translation_dict = {
            'TPI': '育种综合指数', 'NM$': '净利润值', 'CM$': '奶酪净值', 'FM$': '液奶净值', 'GM$': '放牧净值',
            'MILK': '产奶量', 'FAT': '乳脂量', 'PROT': '乳蛋白量', 'FAT %': '乳脂率', 'PROT%': '乳蛋白率',
            'SCS': '体细胞指数', 'DPR': '女儿怀孕率', 'HCR': '青年牛受胎率', 'CCR': '成母牛受胎率', 'PL': '生产寿命',
            'SCE': '配种产犊难易度', 'DCE': '女儿产犊难易度', 'SSB': '配种死胎率', 'DSB': '女儿死淘率',
            'PTAT': '体型综合指数', 'UDC': '乳房综合指数', 'FLC': '肢蹄综合指数', 'BDC': '体重综合指数',
            'ST': '体高', 'SG': '强壮度', 'BD': '体深', 'DF': '乳用特征', 'RA': '尻角度', 'RW': '尻宽',
            'LS': '后肢侧视', 'LR': '后肢后视', 'FA': '蹄角度', 'FLS': '肢蹄评分', 'FU': '前方附着',
            'UH': '后房高度', 'UW': '后方宽度', 'UC': '悬韧带', 'UD': '乳房深度', 'FT': '前乳头位置',
            'RT': '后乳头位置', 'TL': '乳头长度', 'FE': '饲料效率指数', 'FI': '繁殖力指数', 'HI': '健康指数',
            'LIV': '存活率', 'GL': '妊娠长度', 'MAST': '乳房炎抗病性', 'MET': '子宫炎抗病性', 'RP': '胎衣不下抗病性',
            'KET': '酮病抗病性', 'DA': '变胃抗病性', 'MFV': '产后瘫抗病性', 'EFC': '首次产犊月龄',
            'HLiv': '后备牛存活率', 'FS': '饲料节约指数', 'RFI': '剩余饲料采食量', 'Milk Speed': '排乳速度',
            'Eval Date': '数据日期'
        }


        # 添加表格
        headers = df.columns.tolist()
        table_rows, table_cols = len(df) + 2, len(headers)  # +2 for header rows
        left, top, width, height = Inches(0.1535), Inches(1.5945), Inches(13.0276), Inches(4.15)
        table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table

        # 设置表头（两行）
        for i, header in enumerate(headers):
            # 第一行（中文）
            cell = table.cell(0, i)
            cell.text = translation_dict.get(header, header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
            para = cell.text_frame.paragraphs[0]
            para.font.name = "微软雅黑"
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 第二行（英文）
            cell = table.cell(1, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
            para = cell.text_frame.paragraphs[0]
            para.font.name = "Arial"
            para.font.size = Pt(10)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 填充数据
        for row in range(2, table_rows):
            for col in range(table_cols):
                cell = table.cell(row, col)
                if row - 2 < len(df):
                    value = str(df.iloc[row-2, col])
                    cell.text = value
                else:
                    cell.text = ""
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(233, 237, 244)
                para = cell.text_frame.paragraphs[0]
                para.font.name = "微软雅黑"
                para.font.size = Pt(10)
                para.font.color.rgb = RGBColor(0, 0, 0)
                para.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 添加文本框
        left, top, width, height = Inches(0.469), Inches(5.91), Inches(12.358), Inches(1.22)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = "公牛使用情况分析："
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        p = tf.add_paragraph()
        p.text = "暂不能自动化生成........需服务人员整理........"
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

    # PPT生成-第4.2-牧场近几年公牛使用明细-根据配种记录
    def create_bull_usage_detail_slides(self, prs):
        report_file = os.path.join(self.output_folder, "冻精使用情况报告.xlsx")
        if not os.path.exists(report_file):
            QMessageBox.warning(self, '警告', "未找到冻精使用情况报告文件")
            return

        xls = pd.ExcelFile(report_file)
        sheet_names = xls.sheet_names

        for sheet_name in sheet_names:
            if '年冻精使用情况' not in sheet_name:
                continue

            df = pd.read_excel(report_file, sheet_name=sheet_name)
            df = df.drop(['BULL NAAB', 'BULL REG'], axis=1, errors='ignore')
            
            required_columns = ['年份', '冻精编号', '冻精类型', '使用次数']
            if not all(col in df.columns for col in required_columns):
                logging.warning(f"Sheet {sheet_name} 缺少必要的列，跳过此sheet")
                continue

            year = sheet_name.split('年')[0]

            # 分离均值行和其他数据
            mean_rows = df[df['冻精类型'] == '均值']
            other_rows = df[df['冻精类型'] != '均值']

            # 按冻精类型分组，但不包括均值
            grouped = other_rows.groupby('冻精类型')

            for semen_type, group_df in grouped:
                # 将当前类型的均值行添加到组数据的开头
                type_mean = mean_rows[mean_rows['冻精类型'] == semen_type]
                combined_df = pd.concat([type_mean, group_df]).reset_index(drop=True)
                
                # 创建幻灯片
                self.create_single_semen_type_slide(prs, year, semen_type, combined_df)

    def create_single_semen_type_slide(self, prs, year, semen_type, df):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = f"4.2 {self.farm_name}牧场{year}年公牛使用明细 - {semen_type}"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(8), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)

        # 创建翻译字典
        translation_dict = {
            '年份': 'Year', '冻精编号': 'Semen Code', '冻精类型': 'Semen Type', '使用次数': 'Usage Count',
            'TPI': '育种综合指数', 'NM$': '净利润值', 'CM$': '奶酪净值', 'FM$': '液奶净值', 'GM$': '放牧净值',
            'MILK': '产奶量', 'FAT': '乳脂量', 'PROT': '乳蛋白量', 'FAT %': '乳脂率', 'PROT%': '乳蛋白率',
            'SCS': '体细胞指数', 'DPR': '女儿怀孕率', 'HCR': '青年牛受胎率', 'CCR': '成母牛受胎率', 'PL': '生产寿命',
            'SCE': '配种产犊难易度', 'DCE': '女儿产犊难易度', 'SSB': '配种死胎率', 'DSB': '女儿死淘率',
            'PTAT': '体型综合指数', 'UDC': '乳房综合指数', 'FLC': '肢蹄综合指数', 'BDC': '体重综合指数',
            'ST': '体高', 'SG': '强壮度', 'BD': '体深', 'DF': '乳用特征', 'RA': '尻角度', 'RW': '尻宽',
            'LS': '后肢侧视', 'LR': '后肢后视', 'FA': '蹄角度', 'FLS': '肢蹄评分', 'FU': '前方附着',
            'UH': '后房高度', 'UW': '后方宽度', 'UC': '悬韧带', 'UD': '乳房深度', 'FT': '前乳头位置',
            'RT': '后乳头位置', 'TL': '乳头长度', 'FE': '饲料效率指数', 'FI': '繁殖力指数', 'HI': '健康指数',
            'LIV': '存活率', 'GL': '妊娠长度', 'MAST': '乳房炎抗病性', 'MET': '子宫炎抗病性', 'RP': '胎衣不下抗病性',
            'KET': '酮病抗病性', 'DA': '变胃抗病性', 'MFV': '产后瘫抗病性', 'EFC': '首次产犊月龄',
            'HLiv': '后备牛存活率', 'FS': '饲料节约指数', 'RFI': '剩余饲料采食量', 'Milk Speed': '排乳速度',
            'Eval Date': '数据日期'
        }

        # 添加表格
        headers = df.columns.tolist()
        table_rows = min(len(df) + 2, 20)  # 限制最大行数为20（包括2行表头）
        table_cols = len(headers)
        left, top, width, height = Inches(0.1535), Inches(1.5945), Inches(13.0276), Inches(5)
        table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table

        # 设置表头（两行）
        for i, header in enumerate(headers):
            # 第一行（中文）
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
            para = cell.text_frame.paragraphs[0]
            para.font.name = "微软雅黑"
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 第二行（英文）
            cell = table.cell(1, i)
            cell.text = translation_dict.get(header, header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
            para = cell.text_frame.paragraphs[0]
            para.font.name = "Arial"
            para.font.size = Pt(10)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 填充数据
        for row, (index, data) in enumerate(df.iterrows(), start=2):
            if row >= table_rows:
                break
            for col, value in enumerate(data):
                cell = table.cell(row, col)
                
                # 处理第一行的特殊情况
                if row == 2 and col == 1 and pd.isna(value):  # 假设冻精编号在第二列
                    cell.text = "均值"
                else:
                    cell.text = str(value)
                
                para = cell.text_frame.paragraphs[0]
                para.font.name = "微软雅黑"
                para.font.size = Pt(10)
                para.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                # 设置第一行的样式
                if row == 2:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(142, 180, 227)
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(255, 0, 0)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(233, 237, 244)

        # 调整列宽
        total_width = int(Inches(13.0276))  # 转换为EMU单位
        col_widths = []
        for i in range(table_cols):
            max_width = 0
            for row in range(table_rows):
                cell = table.cell(row, i)
                if cell.text_frame.text:
                    text_width = sum([len(p.text) for p in cell.text_frame.paragraphs])
                    max_width = max(max_width, text_width)
            col_widths.append(max_width)

        # 计算总文本宽度
        total_text_width = sum(col_widths)

        # 按比例分配实际宽度
        for i in range(table_cols):
            width_ratio = col_widths[i] / total_text_width
            col_width = int(width_ratio * total_width)
            table.columns[i].width = max(col_width, 1)  # 确保宽度至少为1

        # 添加文本框
        left, top, width, height = Inches(0.469), Inches(6.8614), Inches(12.358), Inches(1.22)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = f"{year}年{semen_type}使用分析："
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        p = tf.add_paragraph()
        p.text = "服务人员自行整理............."
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

    # PPT生成-第4.3-牧场近几年公牛使用进展-根据配种记录
    def create_semen_usage_progress_slides(self, prs):
        charts_folder = os.path.join(self.output_folder, "结果-冻精关键性状进展图")
        if not os.path.exists(charts_folder):
            QMessageBox.warning(self, '警告', "未找到冻精关键性状进展图文件夹，请先生成冻精使用报告。")
            return False

        translation_dict = {
            'TPI': '育种综合指数', 'NM$': '净利润值', 'CM$': '奶酪净值', 'FM$': '液奶净值', 'GM$': '放牧净值',
            'MILK': '产奶量', 'FAT': '乳脂量', 'PROT': '乳蛋白量', 'FAT %': '乳脂率', 'PROT%': '乳蛋白率',
            'SCS': '体细胞指数', 'DPR': '女儿怀孕率', 'HCR': '青年牛受胎率', 'CCR': '成母牛受胎率', 'PL': '生产寿命',
            'SCE': '配种产犊难易度', 'DCE': '女儿产犊难易度', 'SSB': '配种死胎率', 'DSB': '女儿死淘率',
            'PTAT': '体型综合指数', 'UDC': '乳房综合指数', 'FLC': '肢蹄综合指数', 'BDC': '体重综合指数',
            'ST': '体高', 'SG': '强壮度', 'BD': '体深', 'DF': '乳用特征', 'RA': '尻角度', 'RW': '尻宽',
            'LS': '后肢侧视', 'LR': '后肢后视', 'FA': '蹄角度', 'FLS': '肢蹄评分', 'FU': '前方附着',
            'UH': '后房高度', 'UW': '后方宽度', 'UC': '悬韧带', 'UD': '乳房深度', 'FT': '前乳头位置',
            'RT': '后乳头位置', 'TL': '乳头长度', 'FE': '饲料效率指数', 'FI': '繁殖力指数', 'HI': '健康指数',
            'LIV': '存活率', 'GL': '妊娠长度', 'MAST': '乳房炎抗病性', 'MET': '子宫炎抗病性', 'RP': '胎衣不下抗病性',
            'KET': '酮病抗病性', 'DA': '变胃抗病性', 'MFV': '产后瘫抗病性', 'EFC': '首次产犊月龄',
            'HLiv': '后备牛存活率', 'FS': '饲料节约指数', 'RFI': '剩余饲料采食量', 'Milk Speed': '排乳速度',
            'Eval Date': '数据日期'
        }

        for chart_file in os.listdir(charts_folder):
            if chart_file.endswith("冻精使用趋势.png"):
                trait = chart_file.split("冻精使用趋势.png")[0]
                chart_path = os.path.join(charts_folder, chart_file)

                slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)

                # 添加标题
                title = f"4.3 {self.farm_name}牧场近几年公牛使用进展 - {translation_dict.get(trait, trait)} ({trait})"
                left, top, width, height = Inches(0.583), Inches(0.732), Inches(9), Inches(0.571)
                title_shape = slide.shapes.add_textbox(left, top, width, height)
                title_frame = title_shape.text_frame
                title_frame.text = title
                title_para = title_frame.paragraphs[0]
                title_para.alignment = PP_ALIGN.LEFT
                title_para.font.name = "微软雅黑"
                title_para.font.size = Pt(28)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(0, 0, 0)

                # 添加图片
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(12)
                slide.shapes.add_picture(chart_path, left, top, width=width)

                # 添加文本框
                left, top, width, height = Inches(0.469), Inches(6.8614), Inches(12.358), Inches(1.22)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame

                p = tf.add_paragraph()
                p.text = f"{translation_dict.get(trait, trait)} ({trait}) 冻精使用进展分析："
                p.font.name = "微软雅黑"
                p.font.size = Pt(16)
                p.font.bold = True
                p.alignment = PP_ALIGN.LEFT

                p = tf.add_paragraph()
                p.text = "服务人员自行整理............."
                p.font.name = "微软雅黑"
                p.font.size = Pt(14)
                p.alignment = PP_ALIGN.LEFT
                p.level = 1

        return True

    # PPT生成-第4.3-牧场近几年使用公牛单值图-根据配种记录
    def create_semen_usage_scatter_plot_slide(self, prs):
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = f"4.4 {self.farm_name}牧场近几年使用公牛单值图"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(9), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)

        # 添加普通冻精图片
        regular_semen_path = os.path.join(self.output_folder, "普通冻精配种日期单值图.png")
        if os.path.exists(regular_semen_path):
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(6)  # 调整宽度以适应幻灯片
            height = Inches(3.5)  # 调整高度以保持纵横比
            slide.shapes.add_picture(regular_semen_path, left, top, width, height)

        # 添加性控冻精图片
        sexed_semen_path = os.path.join(self.output_folder, "性控冻精配种日期单值图.png")
        if os.path.exists(sexed_semen_path):
            left = Inches(7)  # 调整左边距以放置在右侧
            top = Inches(1.5)
            width = Inches(6)  # 调整宽度以适应幻灯片
            height = Inches(3.5)  # 调整高度以保持纵横比
            slide.shapes.add_picture(sexed_semen_path, left, top, width, height)

        # 添加文本框
        left, top, width, height = Inches(0.583), Inches(5.5), Inches(11), Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = "公牛使用单值图分析：(暂不能自动化生成........需服务人员整理........）"
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = False
        p.alignment = PP_ALIGN.LEFT

        return slide

    # PPT生成-第结尾页-再见
    def create_bye_slide(self, prs):
        slide_layout = prs.slide_layouts[6]  # 使用再见页
        slide = prs.slides.add_slide(slide_layout)

        




class RuleManagerDialog(QDialog):
    def __init__(self, username=''):
        super().__init__()
        self.setWindowTitle("管理公牛号转换规则")
        self.setGeometry(100, 100, 600, 300)
        self.username = username  # 将用户名存储为实例属性
        layout = QVBoxLayout()

        self.table = QTableWidget()
        self.table.setColumnCount(2)
        self.table.setHorizontalHeaderLabels(['常规NAAB', '性控NAAB'])
        layout.addWidget(self.table)

        input_layout = QHBoxLayout()
        self.regular_input = QLineEdit()
        self.regular_input.setPlaceholderText("输入常规NAAB")
        self.sexed_input = QLineEdit()
        self.sexed_input.setPlaceholderText("输入性控NAAB")
        add_button = QPushButton("添加")
        add_button.clicked.connect(self.add_rule)
        input_layout.addWidget(self.regular_input)
        input_layout.addWidget(self.sexed_input)
        input_layout.addWidget(add_button)
        layout.addLayout(input_layout)

        submit_button = QPushButton("提交")
        submit_button.clicked.connect(self.submit_rules)
        button_layout = QHBoxLayout()  # 在使用 button_layout 之前先创建它
        button_layout.addWidget(submit_button)
        layout.addLayout(button_layout)

        delete_button = QPushButton("删除选定规则")
        delete_button.clicked.connect(self.delete_selected_rule)
        button_layout = QHBoxLayout()

        button_layout.addWidget(delete_button)
        layout.addLayout(button_layout)

        self.setLayout(layout)

    def add_rule(self):
        regular = self.regular_input.text()
        sexed = self.sexed_input.text()
        if regular and sexed:
            row_position = self.table.rowCount()
            self.table.insertRow(row_position)
            self.table.setItem(row_position, 0, QTableWidgetItem(regular))
            self.table.setItem(row_position, 1, QTableWidgetItem(sexed))
            self.regular_input.clear()
            self.sexed_input.clear()

    def submit_rules(self):
        try:
            connection = pymysql.connect(
                host=cipher_suite.decrypt(encoded_host).decode(),
                port=int(cipher_suite.decrypt(encoded_port).decode()),
                user=cipher_suite.decrypt(encoded_user).decode(),
                password=cipher_suite.decrypt(encoded_password).decode(),
                database=cipher_suite.decrypt(encoded_db).decode(),
                charset='utf8mb4',
                cursorclass=pymysql.cursors.DictCursor
            )
            with connection.cursor() as cursor:
                for row in range(self.table.rowCount()):
                    regular = self.table.item(row, 0).text()
                    sexed = self.table.item(row, 1).text()
                    if self.username:  # 只提交当前用户添加的规则
                        sql = "INSERT INTO rules (`NAAB-conventional`, `NAAB-sexed`, `uploader`) VALUES (%s, %s, %s)"
                        cursor.execute(sql, (regular, sexed, self.username))

            connection.commit()
            connection.close()
            QMessageBox.information(self, "提交成功", "公牛号转换规则已成功提交到数据库。")
        except Exception as e:
            logging.error(f"提交规则时发生错误: {e}")
            QMessageBox.critical(self, "错误", f"提交规则时发生错误：{e},请联系管理员。")

    def delete_selected_rule(self):
        rows = list(set(index.row() for index in self.table.selectedIndexes()))
        try:
            connection = pymysql.connect(
                host=cipher_suite.decrypt(encoded_host).decode(),
                port=int(cipher_suite.decrypt(encoded_port).decode()),
                user=cipher_suite.decrypt(encoded_user).decode(),
                password=cipher_suite.decrypt(encoded_password).decode(),
                database=cipher_suite.decrypt(encoded_db).decode(),
                charset='utf8mb4',
                cursorclass=pymysql.cursors.DictCursor
            )
            with connection.cursor() as cursor:
                for row in sorted(rows, reverse=True):
                    regular = self.table.item(row, 0).text()
                    sexed = self.table.item(row, 1).text()
                    sql = "DELETE FROM rules WHERE `NAAB-conventional`=%s AND `NAAB-sexed`=%s"
                    cursor.execute(sql, (regular, sexed))
                    self.table.removeRow(row)
            connection.commit()
            connection.close()
        except Exception as e:
            logging.error(f"删除规则时发生错误: {e}")
            QMessageBox.critical(self, "错误", f"删除规则时发生错误：{e},请联系管理员。")

import pandas as pd
import os
import logging
from PyQt6.QtWidgets import (QDialog, QVBoxLayout, QHBoxLayout, QLabel, QCheckBox, 
                             QPushButton, QLineEdit, QTableWidget, QTableWidgetItem, 
                             QHeaderView, QMessageBox, QFileDialog)
from PyQt6.QtCore import Qt

import pandas as pd
import os
import logging
from PyQt6.QtWidgets import (QDialog, QVBoxLayout, QHBoxLayout, QLabel, QCheckBox, 
                             QPushButton, QLineEdit, QTableWidget, QTableWidgetItem, 
                             QHeaderView, QMessageBox, QFileDialog)
from PyQt6.QtCore import Qt

class GroupSelectionDialog(QDialog):
    def __init__(self, cow_df, bull_df, output_folder, parent=None):
        super().__init__(parent)
        self.cow_df = cow_df
        self.bull_df = bull_df
        self.output_folder = output_folder
        self.setWindowTitle("选择群体进行选配")
        self.setGeometry(100, 100, 800, 600)
        self.initUI()

    def initUI(self):
        layout = QVBoxLayout()

        # 选择分组(多选)
        group_label = QLabel("选择分组:")
        self.group_layout = QHBoxLayout()
        layout.addWidget(group_label)
        layout.addLayout(self.group_layout)

        # 选择胎次(多选)
        lac_label = QLabel("选择胎次:")
        self.lac_layout = QHBoxLayout()
        layout.addWidget(lac_label)
        layout.addLayout(self.lac_layout)

        # 选择是否在场(多选)
        in_herd_label = QLabel("选择是否在场:")
        self.in_herd_layout = QHBoxLayout()
        layout.addWidget(in_herd_label)
        layout.addLayout(self.in_herd_layout)

        # 选择性别(多选)
        sex_label = QLabel("选择性别:")
        self.sex_layout = QHBoxLayout()
        layout.addWidget(sex_label)
        layout.addLayout(self.sex_layout)

        # 显示群体信息按钮
        self.show_group_info_button = QPushButton("显示群体信息")
        self.show_group_info_button.clicked.connect(self.show_group_info)
        layout.addWidget(self.show_group_info_button)

        # 群体信息预览表格
        self.group_info_table = QTableWidget()
        self.group_info_table.setColumnCount(9)
        self.group_info_table.setHorizontalHeaderLabels(
            ['分组', '胎次', '头数', '1选完成数', '1选未完成数', '2选完成数', '2选未完成数', '3选完成数', '3选未完成数'])
        self.group_info_table.horizontalHeader().setSectionResizeMode(QHeaderView.ResizeMode.Stretch)
        layout.addWidget(self.group_info_table)

        # 选择冻精及分配比例
        semen_label = QLabel("选择冻精及分配比例:")
        layout.addWidget(semen_label)

        # 常规冻精
        regular_semen_label = QLabel("常规冻精:")
        layout.addWidget(regular_semen_label)
        regular_semen_layout = QHBoxLayout()
        self.regular_semen_ratios = []
        for bull_id in self.bull_df[self.bull_df['semen_type'] == '常规']['bull_id'].astype(str).unique():
            ratio_label = QLabel(bull_id)
            ratio_edit = QLineEdit()
            ratio_edit.setPlaceholderText("默认0%")
            regular_semen_layout.addWidget(ratio_label)
            regular_semen_layout.addWidget(ratio_edit)
            self.regular_semen_ratios.append(ratio_edit)
        layout.addLayout(regular_semen_layout)

        # 性控冻精
        sexed_semen_label = QLabel("性控冻精:")
        layout.addWidget(sexed_semen_label)
        sexed_semen_layout = QHBoxLayout()
        self.sexed_semen_ratios = []
        for bull_id in self.bull_df[self.bull_df['semen_type'] == '性控']['bull_id'].astype(str).unique():
            ratio_label = QLabel(bull_id)
            ratio_edit = QLineEdit()
            ratio_edit.setPlaceholderText("默认0%")
            sexed_semen_layout.addWidget(ratio_label)
            sexed_semen_layout.addWidget(ratio_edit)
            self.sexed_semen_ratios.append(ratio_edit)
        layout.addLayout(sexed_semen_layout)

        # 开始选配按钮
        self.start_breeding_button = QPushButton("开始选配")
        self.start_breeding_button.clicked.connect(self.start_breeding)
        layout.addWidget(self.start_breeding_button)

        self.setLayout(layout)
        
        # 初始化选项
        self.update_options()

    def update_options(self):
        # 更新分组选项
        self.clear_layout(self.group_layout)
        self.group_checkboxes = []
        groups = self.cow_df['group'].unique()
        for group in groups:
            checkbox = QCheckBox(str(group))
            self.group_checkboxes.append(checkbox)
            self.group_layout.addWidget(checkbox)

        # 更新胎次选项
        self.clear_layout(self.lac_layout)
        self.lac_checkboxes = []
        lacs = self.cow_df['lac'].dropna().unique().tolist()
        lacs = [lac for lac in lacs if pd.notna(lac)]
        lacs = [str(lac) for lac in lacs]
        lacs.sort()
        for lac in lacs:
            checkbox = QCheckBox(str(lac))
            self.lac_checkboxes.append(checkbox)
            self.lac_layout.addWidget(checkbox)

        # 更新是否在场选项
        self.clear_layout(self.in_herd_layout)
        self.in_herd_checkboxes = []
        in_herd_options = self.cow_df['是否在场'].unique()
        for option in in_herd_options:
            checkbox = QCheckBox(str(option))
            self.in_herd_checkboxes.append(checkbox)
            self.in_herd_layout.addWidget(checkbox)

        # 更新性别选项
        self.clear_layout(self.sex_layout)
        self.sex_checkboxes = []
        sex_options = self.cow_df['sex'].unique()
        for option in sex_options:
            checkbox = QCheckBox(str(option))
            self.sex_checkboxes.append(checkbox)
            self.sex_layout.addWidget(checkbox)

    def clear_layout(self, layout):
        while layout.count():
            child = layout.takeAt(0)
            if child.widget():
                child.widget().deleteLater()


    def show_group_info(self):
        try:
            selected_groups = [checkbox.text() for checkbox in self.group_checkboxes if checkbox.isChecked()]
            selected_lacs = [float(checkbox.text()) if checkbox.text() else checkbox.text() for checkbox in self.lac_checkboxes if checkbox.isChecked()]
            selected_in_herd = [checkbox.text() for checkbox in self.in_herd_checkboxes if checkbox.isChecked()]
            selected_sex = [checkbox.text() for checkbox in self.sex_checkboxes if checkbox.isChecked()]

            # 根据选择的条件筛选母牛数据
            filtered_cows = self.cow_df
            if selected_groups:
                filtered_cows = filtered_cows[filtered_cows['group'].isin(selected_groups)]
            if selected_lacs:
                filtered_cows = filtered_cows[filtered_cows['lac'].isin(selected_lacs)]
            if selected_in_herd:
                filtered_cows = filtered_cows[filtered_cows['是否在场'].isin(selected_in_herd)]
            if selected_sex:
                filtered_cows = filtered_cows[filtered_cows['sex'].isin(selected_sex)]

            # 统计各分组、胎次、是否在场和性别的头数
            group_info = []
            for group in selected_groups or filtered_cows['group'].unique():
                group_cows = filtered_cows[filtered_cows['group'] == group]
                for lac in selected_lacs or group_cows['lac'].unique():
                    lac_cows = group_cows[group_cows['lac'] == lac]
                    for in_herd in selected_in_herd or lac_cows['是否在场'].unique():
                        in_herd_cows = lac_cows[lac_cows['是否在场'] == in_herd]
                        for sex in selected_sex or in_herd_cows['sex'].unique():
                            sex_cows = in_herd_cows[in_herd_cows['sex'] == sex]
                            count = len(sex_cows)
                            group_info.append([group, lac, in_herd, sex, count])

            # 更新群体信息预览表格
            self.group_info_table.setColumnCount(5)
            self.group_info_table.setHorizontalHeaderLabels(['分组', '胎次', '是否在场', '性别', '头数'])
            self.group_info_table.setRowCount(len(group_info) + 1)  # 为总计行增加一行
            for row, info in enumerate(group_info):
                for col, value in enumerate(info):
                    item = QTableWidgetItem(str(value))
                    self.group_info_table.setItem(row, col, item)

            # 添加总计行
            total_count = sum(info[4] for info in group_info)
            self.group_info_table.setItem(len(group_info), 0, QTableWidgetItem('总计'))
            self.group_info_table.setItem(len(group_info), 4, QTableWidgetItem(str(total_count)))

        except Exception as e:
            logging.exception(f"显示群体信息时发生错误: {str(e)}")
            QMessageBox.critical(self, "错误", f"显示群体信息时发生错误,请查看日志文件了解详情。")
    
    def start_breeding(self):
        try:
            logging.info("开始选配...")
            
            # 初始化进度条
            progress_dialog = QProgressDialog("正在进行选配...", "取消", 0, 100, self)
            progress_dialog.setWindowModality(Qt.WindowModality.WindowModal)
            progress_dialog.setWindowFlags(progress_dialog.windowFlags() | Qt.WindowType.WindowStaysOnTopHint)
            progress_dialog.setValue(0)
            progress_dialog.show()
            progress_dialog.raise_()
            QApplication.processEvents()

            # 读取所需的数据文件
            cows, complete_report_df = self.read_data_files()
            if complete_report_df is None or complete_report_df.empty:
                logging.error("无法读取完整报告数据")
                QMessageBox.critical(self, "错误", "无法读取完整报告数据，请检查文件是否存在且格式正确。")
                return
            
            logging.info(f"读取的母牛数量: {len(cows)}")
            logging.info(f"完整报告数据形状: {complete_report_df.shape}")
            logging.info(f"完整报告数据列: {complete_report_df.columns.tolist()}")
            progress_dialog.setValue(1)
            QApplication.processEvents()

            # 读取指数排名结果
            ranking_result_file = os.path.join(self.output_folder, "结果-指数排名结果.xlsx")
            ranking_result_df = pd.read_excel(ranking_result_file)
            logging.info(f"排名结果数据形状: {ranking_result_df.shape}")
            progress_dialog.setValue(2)
            QApplication.processEvents()

            # 获取用户输入的参数
            sexed_semen, regular_semen, sexed_ratios, regular_ratios = self.get_user_inputs()
            
            # 验证用户输入的参数
            if not self.validate_user_inputs(sexed_semen, regular_semen, sexed_ratios, regular_ratios):
                return
            
            # 获取选择的group、胎次、是否在场和性别
            selected_groups = [checkbox.text() for checkbox in self.group_checkboxes if checkbox.isChecked()]
            selected_lacs = [float(checkbox.text()) if checkbox.text() else checkbox.text() for checkbox in self.lac_checkboxes if checkbox.isChecked()]
            selected_in_herd = [checkbox.text() for checkbox in self.in_herd_checkboxes if checkbox.isChecked()]
            selected_sex = [checkbox.text() for checkbox in self.sex_checkboxes if checkbox.isChecked()]
            
            logging.info(f"选择的分组: {selected_groups}")
            logging.info(f"选择的胎次: {selected_lacs}")
            logging.info(f"选择的是否在场: {selected_in_herd}")
            logging.info(f"选择的性别: {selected_sex}")
            
            # 根据选择的条件筛选母牛
            selected_cows = cows
            if selected_groups:
                selected_cows = selected_cows[selected_cows['group'].isin(selected_groups)]
            if selected_lacs:
                selected_cows = selected_cows[selected_cows['lac'].isin(selected_lacs)]
            if selected_in_herd:
                selected_cows = selected_cows[selected_cows['是否在场'].isin(selected_in_herd)]
            if selected_sex:
                selected_cows = selected_cows[selected_cows['sex'].isin(selected_sex)]
            
            logging.info(f"筛选后的母牛数量: {len(selected_cows)}")
            
            if len(selected_cows) == 0:
                logging.error("没有母牛被选中进行选配")
                QMessageBox.warning(self, "警告", "没有母牛被选中进行选配，请检查您的选择条件。")
                return
            
            progress_dialog.setValue(3)
            QApplication.processEvents()

            # 确保母牛ID为字符串类型
            selected_cows['cow_id'] = selected_cows['cow_id'].astype(str)
            complete_report_df['在群母牛号'] = complete_report_df['在群母牛号'].astype(str)
            ranking_result_df['cow_id'] = ranking_result_df['cow_id'].astype(str)

            # 初始化选配结果数据框和已分配公牛集合
            breeding_result_df, allocated_bulls = self.initialize_breeding_result(selected_cows)
            
            # 使用 merge 而不是 isin 来确保正确匹配
            sorted_cows = pd.merge(selected_cows[['cow_id']], ranking_result_df, on='cow_id', how='left')

            if 'Combine Index Score' in sorted_cows.columns:
                sorted_cows = sorted_cows.sort_values(by=['Combine Index Score'], ascending=False)
            else:
                logging.warning("'Combine Index Score' 列不存在，无法排序")

            logging.info(f"排序后的母牛数量: {len(sorted_cows)}")
            logging.info(f"排序后的母牛 'cow_id': {sorted_cows['cow_id'].tolist()}")

            if len(sorted_cows) == 0:
                logging.error("排序后没有母牛被选中")
                QMessageBox.warning(self, "警告", "没有符合条件的母牛被选中，请检查您的分组和胎次选择。")
                return
            
            # 获取近交系数阈值和隐性基因控制设置
            inbreeding_threshold, control_defect_genes = self.get_breeding_settings()
            logging.info(f"近交系数阈值: {inbreeding_threshold}, 控制隐性基因: {control_defect_genes}")
            progress_dialog.setValue(4)
            QApplication.processEvents()

            # 进行性控冻精和常规冻精的选配
            total_steps = len(sexed_semen) + len(regular_semen)  # 总步骤数，用于进度条计算
            current_step = 0

            for semen_type, selected_semen, ratios in [('性控', sexed_semen, sexed_ratios), ('常规', regular_semen, regular_ratios)]:
                if sum(ratios) == 0:
                    logging.info(f"跳过 {semen_type} 冻精选配，因为比例之和为0")
                    continue

                for round_number in range(1, 4):
                    self.allocate_semen(sorted_cows, semen_type, selected_semen, ratios, round_number, breeding_result_df, allocated_bulls, complete_report_df, inbreeding_threshold, control_defect_genes, progress_dialog, total_steps, current_step)
                    current_step += 1

            # 在保存结果之前
            logging.info(f"选配结果:\n{breeding_result_df}")
            logging.info(f"总行数: {len(breeding_result_df)}")
            logging.info(f"非空行数: {breeding_result_df.dropna().shape[0]}")

            # 检查选配结果
            self.check_breeding_results(breeding_result_df)

            # 保存选配结果到Excel文件
            self.save_breeding_result(breeding_result_df)
            progress_dialog.setValue(95)
            QApplication.processEvents()

            # 显示群体信息
            self.show_group_info()
            progress_dialog.setValue(100)
            QApplication.processEvents()

        except Exception as e:
            logging.exception(f"选配过程中出现未知错误: {str(e)}")
            QMessageBox.critical(self, "错误", f"选配过程中出现未知错误,请查看日志文件了解详情。")
        finally:
            progress_dialog.close()
 
    def check_breeding_results(self, breeding_result_df):
        logging.info("检查选配结果...")
        total_cows = len(breeding_result_df)
        cows_with_matches = breeding_result_df[breeding_result_df.iloc[:, 1:].notna().any(axis=1)]
        cows_without_matches = breeding_result_df[breeding_result_df.iloc[:, 1:].isna().all(axis=1)]
        
        logging.info(f"总母牛数: {total_cows}")
        logging.info(f"成功匹配的母牛数: {len(cows_with_matches)}")
        logging.info(f"未匹配的母牛数: {len(cows_without_matches)}")
        
        if not cows_without_matches.empty:
            logging.warning("以下母牛没有匹配到任何公牛:")
            for cow_id in cows_without_matches['母牛号']:
                logging.warning(f"  - {cow_id}")

    def allocate_semen(self, cows, semen_type, selected_semen, ratios, round_number, breeding_result_df, allocated_bulls, complete_report_df, inbreeding_threshold, control_defect_genes, progress_dialog, total_steps, current_step):
        logging.info(f"开始分配 {semen_type} 冻精，第 {round_number} 轮")
        logging.info(f"可用公牛数量: {len(selected_semen)}")
        logging.info(f"母牛数量: {len(cows)}")

        participating_semen = [semen for semen, ratio in zip(selected_semen, ratios) if ratio > 0]
        participating_ratios = [ratio for ratio in ratios if ratio > 0]
        semen_counts = [int(ratio / 100 * len(cows)) for ratio in participating_ratios]
        semen_counts[-1] += len(cows) - sum(semen_counts)

        logging.info(f"参与选配的公牛数量: {len(participating_semen)}")
        logging.info(f"参与选配的公牛比例: {participating_ratios}")
        logging.info(f"各公牛的分配数量: {semen_counts}")

        sorted_cows = cows.sort_values(by=['Combine Index Score'], ascending=False)

        for round in range(2):
            logging.info(f"开始第{round+1}次尝试 ({semen_type}, 第{round_number}选)...")

            for cow_id in sorted_cows['cow_id']:
                logging.info(f"处理母牛: {cow_id}")
                
                # 添加以下调试信息
                logging.info(f"完整报告数据中的母牛ID: {complete_report_df['在群母牛号'].unique().tolist()[:10]}")  # 显示前10个唯一的母牛ID
                logging.info(f"当前母牛ID在完整报告中: {cow_id in complete_report_df['在群母牛号'].values}")
                
                cow_report = complete_report_df[complete_report_df['在群母牛号'] == cow_id]
                if cow_report.empty:
                    logging.warning(f"未找到母牛 {cow_id} 的报告")
                    continue
                
                if len(cow_report) > 1:
                    logging.warning(f"母牛 {cow_id} 有多个报告，使用第一个")
                    cow_report = cow_report.iloc[0:1]

                try:
                    recommended_semen = [cow_report[f"推荐{semen_type}冻精{i}选"].iloc[0] for i in range(1, 7) if f"推荐{semen_type}冻精{i}选" in cow_report.columns]
                    recommended_scores = [cow_report[f"{semen_type}冻精{i}得分"].iloc[0] for i in range(1, 7) if f"{semen_type}冻精{i}得分" in cow_report.columns]
                    recommended_inbreeding = [float(cow_report[f"{semen_type}冻精{i}近交系数"].iloc[0].strip('%')) / 100 for i in range(1, 7) if f"{semen_type}冻精{i}近交系数" in cow_report.columns]
                    recommended_defect_genes = [cow_report[f"{semen_type}冻精{i}隐性基因情况"].iloc[0] for i in range(1, 7) if f"{semen_type}冻精{i}隐性基因情况" in cow_report.columns]
                    logging.info(f"母牛 {cow_id} 的推荐公牛: {recommended_semen}")
                except Exception as e:
                    logging.error(f"母牛ID: {cow_id}，获取推荐冻精信息时发生错误: {str(e)}")
                    continue

                prev_rounds = [f"{i}选{semen_type}" for i in range(1, round_number)]
                allocated_semen = breeding_result_df.loc[breeding_result_df['母牛号'] == cow_id, prev_rounds].values[0].tolist()
                allocated_semen = [s for s in allocated_semen if not pd.isna(s)]
                recommended_semen_scores = list(zip(recommended_semen, recommended_scores, recommended_inbreeding, recommended_defect_genes))
                filtered_semen_scores = [(s, score, inbreeding, defect_genes) for s, score, inbreeding, defect_genes in recommended_semen_scores if s not in allocated_semen]
                recommended_semen, recommended_scores, recommended_inbreeding, recommended_defect_genes = zip(*filtered_semen_scores) if filtered_semen_scores else ([], [], [], [])

                if breeding_result_df.loc[breeding_result_df['母牛号'] == cow_id, f"{round_number}选{semen_type}"].isnull().all():
                    logging.info(f"尝试为母牛 {cow_id} 分配 {semen_type} 冻精")
                    for bull_id, score, inbreeding, defect_genes in zip(recommended_semen, recommended_scores, recommended_inbreeding, recommended_defect_genes):
                        logging.info(f"考虑公牛 {bull_id}，近交系数: {inbreeding}，隐性基因: {defect_genes}")
                        if pd.isna(bull_id) or bull_id not in participating_semen:
                            logging.info(f"公牛 {bull_id} 不在参与名单中")
                            continue
                        if inbreeding > inbreeding_threshold:
                            logging.info(f"近交系数 {inbreeding} 超过阈值 {inbreeding_threshold}")
                            continue
                        if control_defect_genes and defect_genes == "NO safe":
                            logging.info("隐性基因不安全")
                            continue
                        if round == 0 and semen_counts[participating_semen.index(bull_id)] > 0:
                            breeding_result_df.loc[breeding_result_df['母牛号'] == cow_id, f"{round_number}选{semen_type}"] = bull_id
                            semen_counts[participating_semen.index(bull_id)] -= 1
                            allocated_bulls[cow_id].add(bull_id)
                            logging.info(f"成功为母牛 {cow_id} 分配公牛 {bull_id}")
                            break
                        elif round == 1:
                            breeding_result_df.loc[breeding_result_df['母牛号'] == cow_id, f"{round_number}选{semen_type}"] = bull_id
                            allocated_bulls[cow_id].add(bull_id)
                            logging.info(f"第二轮：成功为母牛 {cow_id} 分配公牛 {bull_id}")
                            break
                    else:
                        logging.warning(f"未能为母牛 {cow_id} 分配 {semen_type} 冻精")

        progress = int(((current_step + round_number / 3) / total_steps) * 100)
        progress_dialog.setValue(progress)
        QApplication.processEvents()
        if progress_dialog.wasCanceled():
            return

        logging.info(f"完成 {semen_type} 冻精分配，第 {round_number} 轮")
        logging.info(f"剩余可用公牛数量: {sum(semen_counts)}")

    def get_breeding_settings(self):
        settings_dialog = BreedingSettingsDialog(self)
        if settings_dialog.exec() == QDialog.DialogCode.Accepted:
            inbreeding_threshold, control_defect_genes = settings_dialog.get_settings()
            return inbreeding_threshold, control_defect_genes
        else:
            return 0.125, True

    def read_data_files(self):
        # 读取在群母牛信息和选配完整报告数据文件
        cows = self.cow_df
        complete_report_file = os.path.join(self.output_folder, "sty_14.enc")

        try:
            # 加载密钥
            with open(os.path.join(self.output_folder, 'natiommont-oasthandler.key'), 'rb') as key_file:
                key = key_file.read()

            # 生成加密器
            cipher = Fernet(key)

            # 读取加密文件
            with open(complete_report_file, 'rb') as file:
                encrypted_data = file.read()

            # 解密文件内容
            decrypted_data = cipher.decrypt(encrypted_data)

            # 读取解密的Excel内容
            complete_report_df = pd.read_excel(BytesIO(decrypted_data), engine='openpyxl')
            
            # 添加以下调试信息
            logging.info(f"完整报告数据形状: {complete_report_df.shape}")
            logging.info(f"完整报告数据列: {complete_report_df.columns.tolist()}")
            logging.info(f"完整报告数据前5行:\n{complete_report_df.head()}")
            logging.info(f"完整报告数据中的母牛ID: {complete_report_df['在群母牛号'].unique().tolist()[:10]}")  # 显示前10个唯一的母牛ID
            
            # 确保母牛ID为字符串类型
            complete_report_df['在群母牛号'] = complete_report_df['在群母牛号'].astype(str)
            cows['cow_id'] = cows['cow_id'].astype(str)
            
            return cows, complete_report_df
        
        except Exception as e:
            logging.error(f"读取完整报告数据时发生错误: {e}")
            return None, None  
            
    def get_user_inputs(self):
        # 获取性控冻精和常规冻精的比例
        sexed_ratios = []
        for ratio_edit in self.sexed_semen_ratios:
            try:
                ratio = int(ratio_edit.text())
                sexed_ratios.append(ratio)
            except ValueError:
                sexed_ratios.append(0)

        regular_ratios = []
        for ratio_edit in self.regular_semen_ratios:
            try:
                ratio = int(ratio_edit.text())
                regular_ratios.append(ratio)
            except ValueError:
                regular_ratios.append(0)

        # 获取性控冻精和常规冻精的编号
        sexed_semen = self.bull_df[self.bull_df['semen_type'] == '性控']['bull_id'].astype(str).tolist()
        regular_semen = self.bull_df[self.bull_df['semen_type'] == '常规']['bull_id'].astype(str).tolist()

        logging.info(f"性控冻精比例: {sexed_ratios}")
        logging.info(f"常规冻精比例: {regular_ratios}")
        logging.info(f"性控冻精数量: {len(sexed_semen)}")
        logging.info(f"常规冻精数量: {len(regular_semen)}")
        return sexed_semen, regular_semen, sexed_ratios, regular_ratios


    def validate_user_inputs(self, sexed_semen, regular_semen, sexed_ratios, regular_ratios):
        # 检查是否存在性控冻精
        has_sexed_semen = len(sexed_semen) > 0
        # 检查是否存在常规冻精
        has_regular_semen = len(regular_semen) > 0


        # 如果存在性控冻精,则验证性控冻精比例之和是否为100或0
        if has_sexed_semen and sum(sexed_ratios) not in [100, 0]:
            QMessageBox.warning(self, "警告", "性控冻精的比例之和必须等于100或0。")
            return False

        # 如果存在常规冻精,则验证常规冻精比例之和是否为100或0
        if has_regular_semen and sum(regular_ratios) not in [100, 0]:
            QMessageBox.warning(self, "警告", "常规冻精的比例之和必须等于100或0。")
            return False
        logging.info(f"性控冻精验证结果: {'通过' if has_sexed_semen and sum(sexed_ratios) in [100, 0] else '不通过'}")
        logging.info(f"常规冻精验证结果: {'通过' if has_regular_semen and sum(regular_ratios) in [100, 0] else '不通过'}")
        return True

    def initialize_breeding_result(self, cows):
        breeding_result_df = pd.DataFrame(columns=['母牛号'] + [f"{i}选{semen_type}" for semen_type in ['性控', '常规'] for i in range(1, 4)])
        breeding_result_df['母牛号'] = cows['cow_id'].astype(str)
        allocated_bulls = {cow_id: set() for cow_id in cows['cow_id'].astype(str)}
        return breeding_result_df, allocated_bulls

    def create_candidate_list(self, cows, semen_type, participating_semen, complete_report_df, inbreeding_threshold, control_defect_genes):
        # 为每头母牛创建候选冻精列表
        cow_candidate_semen = {cow_id: [] for cow_id in cows['cow_id']}
        for _, cow_row in cows.iterrows():
            cow_id = cow_row['cow_id']
            cow_report = complete_report_df[complete_report_df['在群母牛号'] == cow_id].iloc[0]
            for i in range(1, len(participating_semen) + 1):
                bull_id_col = f"推荐{semen_type}冻精{i}选"
                gene_col = f"{semen_type}冻精{i}隐性基因情况"
                inbreeding_col = f"{semen_type}冻精{i}近交系数"
                if bull_id_col in cow_report and pd.notna(cow_report[bull_id_col]):
                    bull_id = cow_report[bull_id_col]
                    gene_info = cow_report[gene_col]
                    inbreeding_coeff = cow_report[inbreeding_col]
                    if control_defect_genes and gene_info == "NO safe":
                        continue
                    if inbreeding_coeff != "N/A" and float(inbreeding_coeff.strip('%')) / 100 >= inbreeding_threshold:
                        continue
                    cow_candidate_semen[cow_id].append(bull_id)
        return cow_candidate_semen


    def save_breeding_result(self, breeding_result_df):
        # 读取指数排名结果
        ranking_result_file = os.path.join(self.output_folder, "结果-指数排名结果.xlsx")
        ranking_result_df = pd.read_excel(ranking_result_file)
        
        # 合并breeding_result_df和ranking_result_df
        breeding_result_df['母牛号'] = breeding_result_df['母牛号'].astype(str)
        ranking_result_df['cow_id'] = ranking_result_df['cow_id'].astype(str)
        merged_df = pd.merge(breeding_result_df, ranking_result_df[['cow_id', 'Combine Index Score']], 
                            left_on='母牛号', right_on='cow_id', how='left')

        
        # 添加其他列
        merged_df['建议选配方案'] = merged_df['母牛号'].map(self.cow_df.set_index('cow_id')['group'])
        merged_df['胎次'] = merged_df['母牛号'].map(self.cow_df.set_index('cow_id')['lac'])
        merged_df['本胎次奶厅高峰产量'] = merged_df['母牛号'].map(self.cow_df.set_index('cow_id')['peak_milk'])
        merged_df['305奶量'] = merged_df['母牛号'].map(self.cow_df.set_index('cow_id')['milk_305'])
        merged_df['泌乳天数'] = merged_df['母牛号'].map(self.cow_df.set_index('cow_id')['DIM'])
        merged_df['繁育状态'] = merged_df['母牛号'].map(self.cow_df.set_index('cow_id')['repro_status'])
        merged_df['指数得分'] = merged_df['Combine Index Score']
        merged_df['品种'] = merged_df['母牛号'].map(self.cow_df.set_index('cow_id')['breed'])
        merged_df['父亲'] = merged_df['母牛号'].map(self.cow_df.set_index('cow_id')['sire'])
        merged_df['外祖父'] = merged_df['母牛号'].map(self.cow_df.set_index('cow_id')['mgs'])
        merged_df['母亲'] = merged_df['母牛号'].map(self.cow_df.set_index('cow_id')['dam'])
        merged_df['外曾外祖父'] = merged_df['母牛号'].map(self.cow_df.set_index('cow_id')['mmgs'])
        merged_df['产犊日期'] = merged_df['母牛号'].map(self.cow_df.set_index('cow_id')['calving_date'])
        merged_df['出生日期'] = merged_df['母牛号'].map(self.cow_df.set_index('cow_id')['birth_date'])
        merged_df['数据提取日期'] = pd.Timestamp.now().date()
        merged_df['月龄'] = merged_df['母牛号'].map(self.cow_df.set_index('cow_id')['age']) 
        merged_df['配次'] = merged_df['母牛号'].map(self.cow_df.set_index('cow_id')['services_time'])

        # 重新排列列的顺序
        columns_order = ['母牛号', '建议选配方案', '1选性控', '2选性控', '3选性控', '1选常规', '2选常规', '3选常规', 
                        '胎次', '本胎次奶厅高峰产量', '305奶量', '泌乳天数', '繁育状态', '指数得分', '品种', 
                        '父亲', '外祖父', '母亲', '外曾外祖父', '产犊日期', '出生日期', '数据提取日期', '月龄', '配次']
        
        final_df = merged_df[columns_order]

        # 检查是否已经有 "选配结果.xlsx"
        breeding_result_file_path = os.path.join(self.output_folder, "结果-选配结果.xlsx")
        if os.path.exists(breeding_result_file_path):
            reply = QMessageBox(self)
            reply.setWindowTitle("选配结果已存在")
            reply.setText("已经存在一个选配结果文件,您想要怎么处理?\n\n\nYes:      覆盖\n\nNO:      继续在已有选配方案表中继续添加\n\nCancle:  取消")
            reply.setStandardButtons(QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No | QMessageBox.StandardButton.Cancel)
            reply.setDefaultButton(QMessageBox.StandardButton.No)
            reply.setWindowFlags(reply.windowFlags() | Qt.WindowType.WindowStaysOnTopHint)
            result = reply.exec()
            if result == QMessageBox.StandardButton.Yes:
                # 覆盖原来的选配结果
                final_df.to_excel(breeding_result_file_path, index=False)
                msg_box = QMessageBox(self)
                msg_box.setWindowTitle("选配完成")
                msg_box.setText(f"选配结果已覆盖原文件: {breeding_result_file_path}")
                msg_box.setWindowFlags(msg_box.windowFlags() | Qt.WindowType.WindowStaysOnTopHint)
                msg_box.exec()
            elif result == QMessageBox.StandardButton.No:
                # 将新生成的选配结果添加到现有的 "选配结果.xlsx"
                existing_df = pd.read_excel(breeding_result_file_path)
                combined_df = existing_df.copy()

                for i, row in final_df.iterrows():
                    cow_id = row['母牛号']
                    if cow_id in combined_df['母牛号'].values:
                        # 如果同一头母牛有新的选配结果,则覆盖原选配
                        for col in final_df.columns:
                            if not pd.isna(row[col]):
                                combined_df.loc[combined_df['母牛号'] == cow_id, col] = row[col]
                    else:
                        # 添加新的选配结果
                        combined_df = pd.concat([combined_df, final_df[final_df['母牛号'] == cow_id]], ignore_index=True)

                combined_df.to_excel(breeding_result_file_path, index=False)
                msg_box = QMessageBox(self)
                msg_box.setWindowTitle("选配完成")
                msg_box.setText(f"新的选配结果已添加到原文件: {breeding_result_file_path}")
                msg_box.setWindowFlags(msg_box.windowFlags() | Qt.WindowType.WindowStaysOnTopHint)
                msg_box.exec()
            else:
                # 用户取消操作
                msg_box = QMessageBox(self)
                msg_box.setWindowTitle("选配取消")
                msg_box.setText("选配结果未保存。")
                msg_box.setWindowFlags(msg_box.windowFlags() | Qt.WindowType.WindowStaysOnTopHint)
                msg_box.exec()
        else:
            # 直接生成 "选配结果.xlsx"
            final_df.to_excel(breeding_result_file_path, index=False)
            msg_box = QMessageBox(self)
            msg_box.setWindowTitle("选配完成")
            msg_box.setText(f"选配结果已保存至: {breeding_result_file_path}")
            msg_box.setWindowFlags(msg_box.windowFlags() | Qt.WindowType.WindowStaysOnTopHint)
            msg_box.exec()

class BreedingSettingsDialog(QDialog):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("选配设置")
        self.setGeometry(100, 100, 300, 200)
        self.initUI()

    def initUI(self):
        layout = QVBoxLayout()
        # 近交系数阈值选择
        inbreeding_label = QLabel("近交系数阈值:")
        self.inbreeding_combo = QComboBox()
        self.inbreeding_combo.addItems(["≤3.125%", "≤6.25%", "≤12.5%", "无视近交"])
        layout.addWidget(inbreeding_label)
        layout.addWidget(self.inbreeding_combo)
        # 隐性基因控制选择
        self.control_genes_checkbox = QCheckBox("控制隐性基因")
        self.control_genes_checkbox.setChecked(True)
        layout.addWidget(self.control_genes_checkbox)
        # 确定和取消按钮
        button_layout = QHBoxLayout()
        ok_button = QPushButton("确定")
        ok_button.clicked.connect(self.accept)
        cancel_button = QPushButton("取消")
        cancel_button.clicked.connect(self.reject)
        button_layout.addWidget(ok_button)
        button_layout.addWidget(cancel_button)
        layout.addLayout(button_layout)
        self.setLayout(layout)

    def get_settings(self):
        inbreeding_threshold = {
            "≤3.125%": 0.03125,
            "≤6.25%": 0.0625,
            "≤12.5%": 0.125,
            "无视近交": float('inf')
        }[self.inbreeding_combo.currentText()]
        control_defect_genes = self.control_genes_checkbox.isChecked()
        return inbreeding_threshold, control_defect_genes

class KeyTraitsDialog(QDialog):
    def __init__(self, default_traits, all_traits, parent=None):
        super().__init__(parent)
        self.setWindowTitle("选择关键性状")
        self.setGeometry(100, 100, 800, 600)
        
        self.default_traits = default_traits
        self.all_traits = all_traits
        self.required_trait = 'NM$'

        layout = QHBoxLayout()

        self.all_traits_list = QListWidget()
        self.all_traits_list.addItems(all_traits)
        self.all_traits_list.itemDoubleClicked.connect(self.add_trait)
        layout.addWidget(QLabel("全部性状"))
        layout.addWidget(self.all_traits_list)

        button_layout = QVBoxLayout()
        add_button = QPushButton("添加 >>")
        add_button.clicked.connect(self.add_trait)
        remove_button = QPushButton("<< 移除")
        remove_button.clicked.connect(self.remove_trait)
        select_all_button = QPushButton("全选")
        select_all_button.clicked.connect(self.select_all_traits)
        reset_button = QPushButton("恢复默认")
        reset_button.clicked.connect(self.reset_traits)
        button_layout.addWidget(add_button)
        button_layout.addWidget(remove_button)
        button_layout.addWidget(select_all_button)
        button_layout.addWidget(reset_button)
        layout.addLayout(button_layout)

        self.selected_traits_list = QListWidget()
        self.selected_traits_list.addItems(default_traits)
        self.selected_traits_list.setDragDropMode(QAbstractItemView.DragDropMode.InternalMove)
        self.selected_traits_list.itemDoubleClicked.connect(self.remove_trait)
        layout.addWidget(QLabel("已选择性状\n（可拖拽调整顺序）"))
        layout.addWidget(self.selected_traits_list)

        button_box = QDialogButtonBox(QDialogButtonBox.StandardButton.Ok | QDialogButtonBox.StandardButton.Cancel)
        button_box.accepted.connect(self.accept)
        button_box.rejected.connect(self.reject)
        layout.addWidget(button_box)

        self.setLayout(layout)

    def add_trait(self, item=None):
        if not item:
            item = self.all_traits_list.currentItem()
        if item and not self.is_trait_selected(item.text()):
            self.selected_traits_list.addItem(item.text())

    def remove_trait(self, item=None):
        if not item:
            item = self.selected_traits_list.currentItem()
        if item and item.text() != self.required_trait:
            self.selected_traits_list.takeItem(self.selected_traits_list.row(item))

    def is_trait_selected(self, trait):
        for i in range(self.selected_traits_list.count()):
            if self.selected_traits_list.item(i).text() == trait:
                return True
        return False

    def select_all_traits(self):
        self.selected_traits_list.clear()
        self.selected_traits_list.addItems(self.all_traits)

    def reset_traits(self):
        self.selected_traits_list.clear()
        self.selected_traits_list.addItems(self.default_traits)

    def get_selected_traits(self):
        return [self.selected_traits_list.item(i).text() for i in range(self.selected_traits_list.count())]


if __name__ == "__main__":
    import sys
    app = QApplication(sys.argv)
    login_dialog = LoginDialog()
    if login_dialog.exec() == QDialog.DialogCode.Accepted:
        main_window = MainWindow(login_dialog.username)
        main_window.show()
        sys.exit(app.exec())
    else:
        sys.exit(0)



