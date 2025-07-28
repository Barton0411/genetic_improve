"""
PPT生成模块
"""

import os
import logging
from pathlib import Path
from typing import Optional, Tuple
import pandas as pd
import numpy as np
from datetime import datetime

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# 设置日志
logger = logging.getLogger(__name__)

class PPTGenerator:
    """PPT报告生成器"""
    
    def __init__(self, project_path: Path, farm_name: str = "牧场"):
        """
        初始化PPT生成器
        
        Args:
            project_path: 项目路径
            farm_name: 牧场名称
        """
        self.project_path = Path(project_path)
        self.farm_name = farm_name
        self.output_folder = self.project_path / "analysis_results"
        self.output_folder.mkdir(parents=True, exist_ok=True)
        
        # 获取模板路径（使用程序内置的模板）
        self.template_path = Path(__file__).parent.parent.parent / "PPT模版.pptx"
        
    def check_required_files(self) -> Tuple[bool, str]:
        """
        检查生成PPT所需的文件是否存在
        
        Returns:
            Tuple[bool, str]: (是否满足条件, 错误信息)
        """
        # 系谱分析是可自动生成的，所以单独检查
        optional_files = {
            "系谱识别情况分析": self.output_folder / "结果-系谱识别情况分析.xlsx",
        }
        
        # 核心必须文件（需要用户手动执行分析）
        required_files = {
            "母牛关键性状指数": self.output_folder / "processed_cow_data_key_traits_scores_genomic.xlsx",
            "母牛育种指数": self.output_folder / "processed_index_cow_index_scores.xlsx",
        }
        
        # 增强的可选文件（如果存在则加入PPT）
        self.optional_analysis_files = {
            "公牛关键性状": self.output_folder / "processed_bull_data_key_traits.xlsx",
            "公牛育种指数": self.output_folder / "processed_index_bull_scores.xlsx",
            "近交系数分析": self.output_folder / "inbreeding_coefficients.xlsx",
            "隐性基因筛查统计_备选公牛": self.output_folder / "隐性基因筛查统计_备选公牛.xlsx",
            "隐性基因筛查统计_已配公牛": self.output_folder / "隐性基因筛查统计_已配公牛.xlsx",
            "年度性状趋势": self.output_folder / "processed_cow_data_key_traits_mean_by_year.xlsx",
            "选配推荐报告": self.project_path / "选配推荐报告" / "常规冻精选配推荐.xlsx",
        }
        
        missing_files = []
        
        # 检查可选文件
        for name, file_path in optional_files.items():
            if not file_path.exists():
                # 检查是否有母牛数据（可以自动生成系谱分析）
                cow_data_file = self.project_path / "standardized_data" / "processed_cow_data.xlsx"
                if not cow_data_file.exists():
                    missing_files.append(name + "（需要先上传母牛数据）")
                else:
                    missing_files.append(name)
        
        # 检查必须文件
        for name, file_path in required_files.items():
            if not file_path.exists():
                missing_files.append(name)
                
        if missing_files:
            return False, f"缺少以下必要文件：{', '.join(missing_files)}"
            
        return True, ""
    
    def generate_ppt(self, progress_callback=None) -> bool:
        """
        生成PPT报告
        
        Args:
            progress_callback: 进度回调函数
            
        Returns:
            bool: 是否成功生成
        """
        try:
            # 创建演示文稿
            if self.template_path.exists():
                # 使用模板创建
                prs = Presentation(str(self.template_path))
                # 删除模板中的所有现有幻灯片
                for i in range(len(prs.slides) - 1, -1, -1):
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
                logger.info("使用PPT模板创建演示文稿")
            else:
                # 创建空白演示文稿
                prs = Presentation()
                # 设置幻灯片大小为16:9
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)
                logger.info("创建空白演示文稿")
                
            if progress_callback:
                progress_callback("正在生成标题页...", 10)
                
            # 生成各个页面
            self._create_title_slide(prs)
            
            if progress_callback:
                progress_callback("正在生成目录页...", 10)
            self._create_toc_slide(prs)
            
            # 系谱分析部分
            if progress_callback:
                progress_callback("正在生成系谱分析页...", 15)
            self._create_pedigree_analysis_slide(prs)
            
            if progress_callback:
                progress_callback("正在生成系谱完整性分析页...", 20)
            self._create_pedigree_completeness_slide(prs)
            
            # 遗传评估部分
            if progress_callback:
                progress_callback("正在生成关键性状分析页...", 25)
            self._create_key_traits_analysis_slide(prs)
            
            if progress_callback:
                progress_callback("正在生成母牛关键性状详情页...", 30)
            self._create_cow_key_traits_detail_slide(prs)
            
            if progress_callback:
                progress_callback("正在生成育种指数分布页...", 35)
            self._create_breeding_index_distribution_slide(prs)
            
            if progress_callback:
                progress_callback("正在生成公牛排名分析页...", 40)
            self._create_bull_ranking_slide(prs)
            
            if progress_callback:
                progress_callback("正在生成年度遗传趋势页...", 45)
            self._create_genetic_trend_slide(prs)
            
            # 选配方案部分
            if progress_callback:
                progress_callback("正在生成选配推荐页...", 50)
            self._create_mating_recommendation_slide(prs)
            
            if progress_callback:
                progress_callback("正在生成近交系数分析页...", 55)
            self._create_inbreeding_analysis_slide(prs)
            
            if progress_callback:
                progress_callback("正在生成隐性基因筛查页...", 60)
            self._create_recessive_gene_slide(prs)
            
            if progress_callback:
                progress_callback("正在生成选配推荐详情页...", 65)
            self._create_mating_detail_slide(prs)
            
            # 总结建议部分
            if progress_callback:
                progress_callback("正在生成项目总结页...", 70)
            self._create_summary_slide(prs)
            
            if progress_callback:
                progress_callback("正在生成改良建议页...", 75)
            self._create_improvement_suggestions_slide(prs)
            
            if progress_callback:
                progress_callback("正在生成结束页...", 80)
            self._create_bye_slide(prs)
            
            # 保存PPT
            output_path = self.output_folder / f"{self.farm_name}牧场遗传改良项目专项服务报告.pptx"
            prs.save(str(output_path))
            
            if progress_callback:
                progress_callback("PPT报告生成完成！", 100)
                
            logger.info(f"PPT报告已保存至: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"生成PPT报告时发生错误: {str(e)}")
            raise
            
    def _create_title_slide(self, prs):
        """创建标题页"""
        slide_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title = f"{self.farm_name}牧场遗传改良项目专项服务"
        left = Inches(2)
        top = Inches(4.5)
        width = Inches(12)
        height = Inches(1.5)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        font = title_para.font
        font.name = "微软雅黑"
        font.size = Pt(48)
        font.bold = True
        font.color.rgb = RGBColor(0, 0, 0)
        
        # 添加日期
        date_str = datetime.now().strftime("%Y年%m月")
        date_left = Inches(2)
        date_top = Inches(6.5)
        date_width = Inches(12)
        date_height = Inches(0.5)
        date_shape = slide.shapes.add_textbox(date_left, date_top, date_width, date_height)
        
        date_frame = date_shape.text_frame
        date_frame.text = date_str
        date_para = date_frame.paragraphs[0]
        date_para.alignment = PP_ALIGN.CENTER
        date_font = date_para.font
        date_font.name = "微软雅黑"
        date_font.size = Pt(24)
        date_font.color.rgb = RGBColor(89, 89, 89)
        
    def _create_toc_slide(self, prs):
        """创建目录页"""
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title = "目录"
        left, top, width, height = Inches(0.5), Inches(0.5), Inches(15), Inches(1)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        
        # 定义目录项
        toc_items = [
            ("01", "牧场系谱记录分析", "系谱完整性 系谱准确性"),
            ("02", "牛群遗传数据评估", "关键性状 育种指数 分布分析"),
            ("03", "选配推荐方案", "个体选配 近交控制 隐性基因"),
            ("04", "项目总结与建议", "改良建议 后续计划")
        ]
        
        # 创建网格布局
        columns = 2
        spacing = Inches(0.5)
        cell_width = Inches(6.5)
        cell_height = Inches(2.5)
        
        total_grid_width = columns * cell_width + (columns - 1) * spacing
        slide_width = prs.slide_width
        grid_left = (slide_width - total_grid_width) / 2
        
        rows = (len(toc_items) + columns - 1) // columns
        total_grid_height = rows * cell_height + (rows - 1) * spacing
        slide_height = prs.slide_height
        grid_top = (slide_height - total_grid_height) / 2 + Inches(0.5)
        
        for i, (number, title, subtitle) in enumerate(toc_items):
            left = grid_left + (i % columns) * (cell_width + spacing)
            top = grid_top + (i // columns) * (cell_height + spacing)
            
            # 添加编号
            number_box = slide.shapes.add_textbox(left, top, Inches(0.5), Inches(0.5))
            number_frame = number_box.text_frame
            number_frame.text = number
            number_frame.paragraphs[0].font.name = "微软雅黑"
            number_frame.paragraphs[0].font.size = Pt(24)
            number_frame.paragraphs[0].font.color.rgb = RGBColor(0, 176, 240)
            
            # 添加主标题
            title_box = slide.shapes.add_textbox(left, top + Inches(0.5), cell_width, Inches(0.4))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.name = "微软雅黑"
            title_frame.paragraphs[0].font.size = Pt(18)
            title_frame.paragraphs[0].font.bold = True
            
            # 添加副标题
            subtitle_box = slide.shapes.add_textbox(left, top + Inches(0.9), cell_width, Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.clear()
            
            subtitle_items = subtitle.split()
            for item in subtitle_items:
                p = subtitle_frame.add_paragraph()
                p.text = "• " + item
                p.font.name = "微软雅黑"
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(89, 89, 89)
                p.level = 1
                
    def _create_pedigree_analysis_slide(self, prs):
        """创建系谱分析标题页"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加序号
        number_box = slide.shapes.add_textbox(Inches(5.67), Inches(1.12), Inches(2), Inches(2))
        number_frame = number_box.text_frame
        number_frame.text = "01"
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = number_frame.paragraphs[0].font
        font.name = "Arial"
        font.size = Pt(166)
        font.color.rgb = RGBColor(255, 255, 255)
        font.color.transparency = 0.69
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(2.67), Inches(4.13), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "牧场系谱记录分析"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = title_frame.paragraphs[0].font
        font.name = "微软雅黑"
        font.size = Pt(40)
        font.color.rgb = RGBColor(255, 255, 255)
        
    def _create_pedigree_completeness_slide(self, prs):
        """创建系谱完整性分析页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title = f"1.1 {self.farm_name}牧场系谱完整性分析"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(10), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        
        # 尝试加载数据并创建表格
        try:
            analysis_file = self.output_folder / "结果-系谱识别情况分析.xlsx"
            if analysis_file.exists():
                df = pd.read_excel(analysis_file)
                df = df[df['是否在场'] == '是'].head(6)  # 只显示在场的前6行
                
                # 创建表格
                table_rows, table_cols = min(len(df) + 1, 7), 9
                left, top, width, height = Inches(0.874), Inches(1.72), Inches(14.5), Inches(5)
                table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table
                
                # 设置表头
                headers = ['是否在场', '出生年份', '头数', '父号可识别头数', '父号识别率', 
                          '外祖父可识别头数', '外祖父识别率', '外曾外祖父可识别头数', '外曾外祖父识别率']
                
                for i, header in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = header
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
                    para = cell.text_frame.paragraphs[0]
                    para.font.name = "微软雅黑"
                    para.font.size = Pt(14)
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(255, 255, 255)
                    para.alignment = PP_ALIGN.CENTER
                    
                # 填充数据
                for row in range(1, min(table_rows, len(df) + 1)):
                    for col in range(table_cols):
                        cell = table.cell(row, col)
                        if col == 1:  # 出生年份列
                            value = str(df.iloc[row-1]['birth_year_group'])
                        else:
                            value = str(df.iloc[row-1, col])
                        cell.text = value
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(233, 237, 244)
                        para = cell.text_frame.paragraphs[0]
                        para.font.name = "微软雅黑"
                        para.font.size = Pt(12)
                        para.alignment = PP_ALIGN.CENTER
                        
        except Exception as e:
            logger.warning(f"创建系谱完整性表格时出错: {e}")
            
    def _create_key_traits_analysis_slide(self, prs):
        """创建关键性状分析页"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加序号
        number_box = slide.shapes.add_textbox(Inches(5.67), Inches(1.12), Inches(2), Inches(2))
        number_frame = number_box.text_frame
        number_frame.text = "02"
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = number_frame.paragraphs[0].font
        font.name = "Arial"
        font.size = Pt(166)
        font.color.rgb = RGBColor(255, 255, 255)
        font.color.transparency = 0.69
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(2.67), Inches(4.13), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "牛群遗传数据评估"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = title_frame.paragraphs[0].font
        font.name = "微软雅黑"
        font.size = Pt(40)
        font.color.rgb = RGBColor(255, 255, 255)
        
    def _create_breeding_index_distribution_slide(self, prs):
        """创建育种指数分布页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title = f"2.1 {self.farm_name}牧场育种指数分布情况"
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(10), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        
        # 添加说明文字
        info_text = "育种指数反映了母牛的综合遗传潜力，分布情况显示了牛群的整体遗传水平。"
        info_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(0.5))
        info_frame = info_box.text_frame
        info_frame.text = info_text
        info_para = info_frame.paragraphs[0]
        info_para.font.name = "微软雅黑"
        info_para.font.size = Pt(16)
        info_para.font.color.rgb = RGBColor(89, 89, 89)
        
    def _create_mating_recommendation_slide(self, prs):
        """创建选配推荐汇总页"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加序号
        number_box = slide.shapes.add_textbox(Inches(5.67), Inches(1.12), Inches(2), Inches(2))
        number_frame = number_box.text_frame
        number_frame.text = "03"
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = number_frame.paragraphs[0].font
        font.name = "Arial"
        font.size = Pt(166)
        font.color.rgb = RGBColor(255, 255, 255)
        font.color.transparency = 0.69
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(2.67), Inches(4.13), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "选配推荐方案"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = title_frame.paragraphs[0].font
        font.name = "微软雅黑"
        font.size = Pt(40)
        font.color.rgb = RGBColor(255, 255, 255)
        
    def _create_bye_slide(self, prs):
        """创建结束页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加感谢文字
        thanks_text = "谢谢"
        left = Inches(2)
        top = Inches(3.5)
        width = Inches(12)
        height = Inches(2)
        thanks_shape = slide.shapes.add_textbox(left, top, width, height)
        
        thanks_frame = thanks_shape.text_frame
        thanks_frame.text = thanks_text
        thanks_para = thanks_frame.paragraphs[0]
        thanks_para.alignment = PP_ALIGN.CENTER
        font = thanks_para.font
        font.name = "微软雅黑"
        font.size = Pt(72)
        font.bold = True
        font.color.rgb = RGBColor(0, 0, 0)
    
    def _create_cow_key_traits_detail_slide(self, prs):
        """创建母牛关键性状详情页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title = f"2.2 {self.farm_name}牧场母牛关键性状分析"
        self._add_slide_title(slide, title)
        
        try:
            # 加载关键性状数据
            traits_file = self.output_folder / "processed_cow_data_key_traits_scores_genomic.xlsx"
            if traits_file.exists():
                df = pd.read_excel(traits_file)
                
                # 计算关键性状的统计信息
                traits_stats = []
                traits_cols = ['产奶量', '乳脂率', '乳蛋白率', '体细胞', '女儿妊娠率', '生产寿命']
                
                for trait in traits_cols:
                    if trait in df.columns:
                        stats = {
                            '性状': trait,
                            '平均值': df[trait].mean(),
                            '标准差': df[trait].std(),
                            '最小值': df[trait].min(),
                            '最大值': df[trait].max(),
                            'CV%': (df[trait].std() / df[trait].mean() * 100) if df[trait].mean() != 0 else 0
                        }
                        traits_stats.append(stats)
                
                if traits_stats:
                    # 创建表格
                    rows = len(traits_stats) + 1
                    cols = 6
                    left, top, width, height = Inches(1), Inches(2), Inches(14), Inches(4)
                    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                    
                    # 设置表头
                    headers = ['性状', '平均值', '标准差', '最小值', '最大值', 'CV%']
                    for i, header in enumerate(headers):
                        cell = table.cell(0, i)
                        cell.text = header
                        self._format_cell(cell, is_header=True)
                    
                    # 填充数据
                    for row_idx, stats in enumerate(traits_stats):
                        table.cell(row_idx + 1, 0).text = stats['性状']
                        table.cell(row_idx + 1, 1).text = f"{stats['平均值']:.2f}"
                        table.cell(row_idx + 1, 2).text = f"{stats['标准差']:.2f}"
                        table.cell(row_idx + 1, 3).text = f"{stats['最小值']:.2f}"
                        table.cell(row_idx + 1, 4).text = f"{stats['最大值']:.2f}"
                        table.cell(row_idx + 1, 5).text = f"{stats['CV%']:.1f}%"
                        
                        for col in range(cols):
                            self._format_cell(table.cell(row_idx + 1, col))
                
        except Exception as e:
            logger.warning(f"创建母牛关键性状详情页时出错: {e}")
    
    def _create_bull_ranking_slide(self, prs):
        """创建公牛排名分析页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title = f"2.3 {self.farm_name}牧场公牛育种指数排名"
        self._add_slide_title(slide, title)
        
        try:
            # 加载公牛育种指数数据
            bull_file = self.output_folder / "processed_index_bull_scores.xlsx"
            if bull_file.exists():
                df = pd.read_excel(bull_file)
                # 按育种指数排序，取前10名
                df_top = df.nlargest(10, 'Combine Index Score')[['bull_id', 'name_cn', 'Combine Index Score']]
                
                # 创建表格
                rows = len(df_top) + 1
                cols = 3
                left, top, width, height = Inches(3), Inches(2), Inches(10), Inches(5)
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # 设置表头
                headers = ['排名', '公牛号/中文名', '育种指数']
                for i, header in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = header
                    self._format_cell(cell, is_header=True)
                
                # 填充数据
                for idx, row in df_top.iterrows():
                    rank = idx + 1
                    table.cell(rank, 0).text = str(rank)
                    table.cell(rank, 1).text = f"{row['bull_id']} / {row.get('name_cn', '')}"
                    table.cell(rank, 2).text = f"{row['Combine Index Score']:.0f}"
                    
                    for col in range(cols):
                        self._format_cell(table.cell(rank, col))
                        
        except Exception as e:
            logger.warning(f"创建公牛排名分析页时出错: {e}")
    
    def _create_genetic_trend_slide(self, prs):
        """创建年度遗传趋势页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title = f"2.4 {self.farm_name}牧场年度遗传进展趋势"
        self._add_slide_title(slide, title)
        
        try:
            # 加载年度趋势数据
            trend_file = self.output_folder / "processed_cow_data_key_traits_mean_by_year.xlsx"
            if trend_file.exists():
                df = pd.read_excel(trend_file)
                
                # 添加说明文字
                info_text = "通过分析不同出生年份母牛的遗传性状均值，可以评估牧场的遗传改良进展。"
                info_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(0.5))
                info_frame = info_box.text_frame
                info_frame.text = info_text
                info_para = info_frame.paragraphs[0]
                info_para.font.name = "微软雅黑"
                info_para.font.size = Pt(14)
                info_para.font.color.rgb = RGBColor(89, 89, 89)
                
                # 这里可以添加图表展示年度趋势
                # 暂时添加文字说明
                summary_text = f"数据显示了{len(df)}个年份的遗传趋势数据。"
                summary_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(3))
                summary_frame = summary_box.text_frame
                summary_frame.text = summary_text
                summary_para = summary_frame.paragraphs[0]
                summary_para.font.name = "微软雅黑"
                summary_para.font.size = Pt(16)
                
        except Exception as e:
            logger.warning(f"创建年度遗传趋势页时出错: {e}")
    
    def _create_inbreeding_analysis_slide(self, prs):
        """创建近交系数分析页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title = f"3.1 {self.farm_name}牧场近交系数控制分析"
        self._add_slide_title(slide, title)
        
        try:
            # 加载近交系数数据
            inbreeding_file = self.output_folder / "inbreeding_coefficients.xlsx"
            if inbreeding_file.exists():
                df = pd.read_excel(inbreeding_file)
                
                # 统计近交系数分布
                total_count = len(df)
                low_risk = len(df[df['inbreeding_coefficient'] < 0.03125])
                medium_risk = len(df[(df['inbreeding_coefficient'] >= 0.03125) & (df['inbreeding_coefficient'] < 0.0625)])
                high_risk = len(df[df['inbreeding_coefficient'] >= 0.0625])
                
                # 创建统计表格
                rows, cols = 5, 3
                left, top, width, height = Inches(3), Inches(2), Inches(10), Inches(4)
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # 设置表头
                headers = ['风险等级', '近交系数范围', '头数（占比）']
                for i, header in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = header
                    self._format_cell(cell, is_header=True)
                
                # 填充数据
                data = [
                    ('低风险', '< 3.125%', f"{low_risk} ({low_risk/total_count*100:.1f}%)"),
                    ('中风险', '3.125% - 6.25%', f"{medium_risk} ({medium_risk/total_count*100:.1f}%)"),
                    ('高风险', '≥ 6.25%', f"{high_risk} ({high_risk/total_count*100:.1f}%)"),
                    ('合计', '-', f"{total_count} (100.0%)")
                ]
                
                for row_idx, (level, range_, count) in enumerate(data):
                    table.cell(row_idx + 1, 0).text = level
                    table.cell(row_idx + 1, 1).text = range_
                    table.cell(row_idx + 1, 2).text = count
                    
                    for col in range(cols):
                        self._format_cell(table.cell(row_idx + 1, col))
                        
        except Exception as e:
            logger.warning(f"创建近交系数分析页时出错: {e}")
    
    def _create_recessive_gene_slide(self, prs):
        """创建隐性基因筛查页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title = f"3.2 {self.farm_name}牧场隐性基因筛查结果"
        self._add_slide_title(slide, title)
        
        try:
            # 尝试加载备选公牛或已配公牛的隐性基因筛查结果
            candidate_file = self.output_folder / "隐性基因筛查统计_备选公牛.xlsx"
            mated_file = self.output_folder / "隐性基因筛查统计_已配公牛.xlsx"
            
            df = None
            if candidate_file.exists():
                df = pd.read_excel(candidate_file)
                subtitle = "备选公牛隐性基因筛查统计"
            elif mated_file.exists():
                df = pd.read_excel(mated_file)
                subtitle = "已配公牛隐性基因筛查统计"
            
            if df is not None:
                # 添加副标题
                subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(0.5))
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.text = subtitle
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.name = "微软雅黑"
                subtitle_para.font.size = Pt(16)
                subtitle_para.font.bold = True
                
                # 创建表格展示统计结果
                rows = min(len(df) + 1, 10)  # 最多显示10行
                cols = min(len(df.columns), 6)  # 最多显示6列
                left, top, width, height = Inches(1), Inches(2.5), Inches(14), Inches(4)
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # 设置表头
                for i, col_name in enumerate(df.columns[:cols]):
                    cell = table.cell(0, i)
                    cell.text = str(col_name)
                    self._format_cell(cell, is_header=True)
                
                # 填充数据
                for row_idx in range(min(len(df), rows - 1)):
                    for col_idx in range(cols):
                        cell = table.cell(row_idx + 1, col_idx)
                        cell.text = str(df.iloc[row_idx, col_idx])
                        self._format_cell(cell)
                        
        except Exception as e:
            logger.warning(f"创建隐性基因筛查页时出错: {e}")
    
    def _create_mating_detail_slide(self, prs):
        """创建选配推荐详情页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title = f"3.3 {self.farm_name}牧场个体选配推荐示例"
        self._add_slide_title(slide, title)
        
        try:
            # 查找选配推荐报告
            mating_file = self.project_path / "选配推荐报告" / "常规冻精选配推荐.xlsx"
            if mating_file.exists():
                df = pd.read_excel(mating_file)
                # 选择前5头牛的推荐结果作为示例
                df_sample = df.head(5)
                
                # 创建表格
                rows = len(df_sample) + 1
                cols = 5  # 母牛号、推荐公牛1、近交系数1、推荐公牛2、近交系数2
                left, top, width, height = Inches(1), Inches(2), Inches(14), Inches(4)
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # 设置表头
                headers = ['母牛号', '首选公牛', '近交系数', '次选公牛', '近交系数']
                for i, header in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = header
                    self._format_cell(cell, is_header=True)
                
                # 填充数据
                for idx, row in df_sample.iterrows():
                    row_idx = idx + 1
                    table.cell(row_idx, 0).text = str(row.get('cow_id', ''))
                    table.cell(row_idx, 1).text = str(row.get('推荐常规冻精1选', ''))
                    table.cell(row_idx, 2).text = f"{row.get('常规冻精1近交系数', 0):.2%}"
                    table.cell(row_idx, 3).text = str(row.get('推荐常规冻精2选', ''))
                    table.cell(row_idx, 4).text = f"{row.get('常规冻精2近交系数', 0):.2%}"
                    
                    for col in range(cols):
                        self._format_cell(table.cell(row_idx, col))
            else:
                # 添加提示信息
                info_text = "暂无选配推荐数据。请先执行个体选配分析。"
                info_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(12), Inches(2))
                info_frame = info_box.text_frame
                info_frame.text = info_text
                info_para = info_frame.paragraphs[0]
                info_para.font.name = "微软雅黑"
                info_para.font.size = Pt(18)
                info_para.font.color.rgb = RGBColor(150, 150, 150)
                info_para.alignment = PP_ALIGN.CENTER
                
        except Exception as e:
            logger.warning(f"创建选配推荐详情页时出错: {e}")
    
    def _create_summary_slide(self, prs):
        """创建项目总结页"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加序号
        number_box = slide.shapes.add_textbox(Inches(5.67), Inches(1.12), Inches(2), Inches(2))
        number_frame = number_box.text_frame
        number_frame.text = "04"
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = number_frame.paragraphs[0].font
        font.name = "Arial"
        font.size = Pt(166)
        font.color.rgb = RGBColor(255, 255, 255)
        font.color.transparency = 0.69
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(2.67), Inches(4.13), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "项目总结与建议"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = title_frame.paragraphs[0].font
        font.name = "微软雅黑"
        font.size = Pt(40)
        font.color.rgb = RGBColor(255, 255, 255)
    
    def _create_improvement_suggestions_slide(self, prs):
        """创建改良建议页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title = f"4.1 {self.farm_name}牧场遗传改良建议"
        self._add_slide_title(slide, title)
        
        # 添加建议内容
        suggestions = [
            "1. 系谱管理建议",
            "   • 加强系谱记录的完整性，特别是外祖父和外曾外祖父信息",
            "   • 建立标准化的系谱录入和验证流程",
            "",
            "2. 选配策略优化",
            "   • 优先使用高育种指数的公牛进行配种",
            "   • 严格控制近交系数，避免近亲配种",
            "   • 关注隐性基因携带情况，降低遗传缺陷风险",
            "",
            "3. 后备牛培育",
            "   • 重点培育高遗传潜力的后备母牛",
            "   • 制定科学的后备牛选育标准",
            "",
            "4. 数据管理提升",
            "   • 定期更新基因组检测数据",
            "   • 建立完善的育种数据管理系统"
        ]
        
        # 创建文本框
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(6))
        text_frame = text_box.text_frame
        
        for suggestion in suggestions:
            p = text_frame.add_paragraph()
            p.text = suggestion
            p.font.name = "微软雅黑"
            p.font.size = Pt(16) if suggestion.startswith("   ") else Pt(18)
            p.font.bold = not suggestion.startswith("   ") and suggestion != ""
            if suggestion.startswith("   "):
                p.level = 1
    
    def _add_slide_title(self, slide, title):
        """添加幻灯片标题的辅助方法"""
        left, top, width, height = Inches(0.583), Inches(0.732), Inches(15), Inches(0.571)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
    
    def _format_cell(self, cell, is_header=False):
        """格式化表格单元格的辅助方法"""
        cell.fill.solid()
        if is_header:
            cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
            para = cell.text_frame.paragraphs[0]
            para.font.name = "微软雅黑"
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
        else:
            cell.fill.fore_color.rgb = RGBColor(233, 237, 244)
            para = cell.text_frame.paragraphs[0]
            para.font.name = "微软雅黑"
            para.font.size = Pt(12)
        para.alignment = PP_ALIGN.CENTER 