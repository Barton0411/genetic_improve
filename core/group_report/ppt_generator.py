"""根据牧场组汇总 Excel 生成精简管理汇报 PPT。"""

from __future__ import annotations

import logging
import os
import re
import tempfile
import zipfile
from collections import defaultdict
from datetime import datetime
from pathlib import Path
from typing import Callable, Dict, Iterable, List, Optional, Tuple

from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from utils.file_manager import FileManager


logger = logging.getLogger(__name__)

BLUE = RGBColor(0, 139, 206)
DARK_BLUE = RGBColor(0, 66, 112)
GREEN = RGBColor(138, 184, 74)
DARK = RGBColor(64, 64, 64)
GRAY = RGBColor(128, 128, 128)
LIGHT = RGBColor(242, 247, 250)
WHITE = RGBColor(255, 255, 255)


def _number(value, default=0.0):
    try:
        if value is None or value == "":
            return default
        return float(value)
    except (TypeError, ValueError):
        return default


class GroupPPTReportGenerator:
    """生成固定篇幅、不会按牧场数倍增的牧场组 PPT。"""

    def __init__(
        self,
        project_path: Path,
        reporter_name: str = "",
        progress_callback: Optional[Callable[[str, int], None]] = None,
    ):
        self.project_path = Path(project_path)
        self.reporter_name = reporter_name or ""
        self.progress_callback = progress_callback
        self.last_output_path: Optional[Path] = None
        self.metadata = FileManager.load_project_metadata(self.project_path)

    def _progress(self, message: str, value: int) -> None:
        if self.progress_callback:
            self.progress_callback(message, value)

    def _find_excel(self) -> Optional[Path]:
        reports = list(
            (self.project_path / "reports").glob("牧场组育种分析汇总报告_*.xlsx")
        )
        return max(reports, key=lambda path: path.stat().st_mtime) if reports else None

    @staticmethod
    def _read_table(path: Path, sheet_name: str) -> List[Dict]:
        workbook = load_workbook(path, read_only=True, data_only=True)
        try:
            if sheet_name not in workbook.sheetnames:
                return []
            sheet = workbook[sheet_name]
            rows = list(sheet.iter_rows(min_row=2, values_only=True))
            if not rows:
                return []
            headers = [str(value or "") for value in rows[0]]
            return [
                dict(zip(headers, row))
                for row in rows[1:]
                if any(value is not None for value in row)
            ]
        finally:
            workbook.close()

    @staticmethod
    def _set_text(shape, text: str, size: int, *, bold=False, color=DARK, align=PP_ALIGN.LEFT):
        shape.text = str(text)
        frame = shape.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = str(text)
        paragraph.alignment = align
        run = paragraph.runs[0]
        run.font.name = "微软雅黑"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame.word_wrap = True

    def _base_slide(self, prs: Presentation, title: str, page: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), prs.slide_height
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = GREEN
        bar.line.fill.background()
        title_box = slide.shapes.add_textbox(
            Inches(0.65), Inches(0.35), Inches(11.8), Inches(0.55)
        )
        self._set_text(title_box, title, 28, bold=True, color=DARK_BLUE)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.03), Inches(12.0), Inches(0.025)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = BLUE
        line.line.fill.background()
        footer = slide.shapes.add_textbox(
            Inches(11.8), Inches(7.12), Inches(0.7), Inches(0.22)
        )
        self._set_text(footer, str(page), 10, color=GRAY, align=PP_ALIGN.RIGHT)
        return slide

    def _add_cover(self, prs: Presentation, farm_count: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE
        block = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(3.0), prs.slide_height
        )
        block.fill.solid()
        block.fill.fore_color.rgb = DARK_BLUE
        block.line.fill.background()
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(2.82), 0, Inches(0.18), prs.slide_height
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = GREEN
        accent.line.fill.background()
        title = slide.shapes.add_textbox(
            Inches(3.65), Inches(2.05), Inches(8.5), Inches(1.35)
        )
        self._set_text(title, "牧场组育种分析汇报", 38, bold=True, color=DARK_BLUE)
        subtitle = slide.shapes.add_textbox(
            Inches(3.68), Inches(3.45), Inches(7.6), Inches(0.6)
        )
        self._set_text(subtitle, f"{farm_count} 个牧场 · 最终汇总", 22, color=BLUE)
        meta = slide.shapes.add_textbox(
            Inches(3.68), Inches(5.35), Inches(7.8), Inches(0.8)
        )
        text = datetime.now().strftime("%Y年%m月%d日")
        if self.reporter_name:
            text += f"    汇报人：{self.reporter_name}"
        self._set_text(meta, text, 16, color=GRAY)
        return slide

    def _add_kpis(self, slide, items: Iterable[Tuple[str, str]], top=1.5):
        items = list(items)
        width = 11.9 / max(len(items), 1)
        for index, (label, value) in enumerate(items):
            left = 0.7 + index * width
            value_box = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width - 0.15), Inches(0.75)
            )
            self._set_text(value_box, value, 27, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
            label_box = slide.shapes.add_textbox(
                Inches(left), Inches(top + 0.72), Inches(width - 0.15), Inches(0.42)
            )
            self._set_text(label_box, label, 14, color=GRAY, align=PP_ALIGN.CENTER)

    def _add_bar_chart(
        self,
        slide,
        categories: List[str],
        series: List[Tuple[str, List[float]]],
        *,
        left=0.75,
        top=1.45,
        width=12.0,
        height=5.25,
        horizontal=False,
        legend=True,
    ):
        if not categories or not series:
            box = slide.shapes.add_textbox(
                Inches(left), Inches(top + 1.5), Inches(width), Inches(0.8)
            )
            self._set_text(box, "当前没有可展示的数据", 20, color=GRAY, align=PP_ALIGN.CENTER)
            return
        data = CategoryChartData()
        data.categories = categories
        for name, values in series:
            data.add_series(name, values)
        chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
        chart = slide.shapes.add_chart(
            chart_type, Inches(left), Inches(top), Inches(width), Inches(height), data
        ).chart
        chart.has_legend = legend and len(series) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.font.size = Pt(11)
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_labels.font.size = Pt(10)
        chart.category_axis.tick_labels.font.size = Pt(10)
        chart.chart_style = 10

    def _add_table(self, slide, headers: List[str], rows: List[List], *, top=1.45):
        rows = rows[:12]
        table_shape = slide.shapes.add_table(
            len(rows) + 1,
            len(headers),
            Inches(0.75),
            Inches(top),
            Inches(12.0),
            Inches(5.45),
        )
        table = table_shape.table
        for column, header in enumerate(headers):
            cell = table.cell(0, column)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(13)
                    run.font.bold = True
                    run.font.color.rgb = WHITE
        for row_index, values in enumerate(rows, start=1):
            for column, value in enumerate(values):
                cell = table.cell(row_index, column)
                cell.text = str(value if value is not None else "")
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if row_index % 2 == 0 else WHITE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.name = "微软雅黑"
                        run.font.size = Pt(11)
                        run.font.color.rgb = DARK

    @staticmethod
    def _top_farms(overview: List[Dict], limit: int = 12) -> List[Dict]:
        return sorted(
            overview,
            key=lambda row: _number(row.get("在群头数")),
            reverse=True,
        )[:limit]

    @staticmethod
    def _normalize_chart_axis_ids(output_path: Path) -> None:
        """将 python-pptx 的负轴 ID 转为等价无符号值，兼容更多解析器。"""
        output_path = Path(output_path)
        temp_fd, temp_name = tempfile.mkstemp(
            suffix=".pptx", dir=str(output_path.parent)
        )
        os.close(temp_fd)
        temp_path = Path(temp_name)
        pattern = re.compile(
            rb'(<c:(?:axId|crossAx)\b[^>]*\bval=")(-\d+)(")'
        )

        def replace(match):
            signed = int(match.group(2))
            unsigned = signed % (1 << 32)
            return match.group(1) + str(unsigned).encode("ascii") + match.group(3)

        try:
            with zipfile.ZipFile(output_path, "r") as source, zipfile.ZipFile(
                temp_path, "w", compression=zipfile.ZIP_DEFLATED
            ) as target:
                for item in source.infolist():
                    data = source.read(item.filename)
                    if (
                        item.filename.startswith("ppt/charts/chart")
                        and item.filename.endswith(".xml")
                    ):
                        data = pattern.sub(replace, data)
                    target.writestr(item, data)
            os.replace(temp_path, output_path)
        finally:
            if temp_path.exists():
                temp_path.unlink()

    def generate(self) -> Tuple[bool, str]:
        if self.metadata.get("project_type") != "multi_farm_group":
            return False, "当前项目不是牧场组项目"
        if not self.metadata.get("all_tasks_complete"):
            return False, "全部牧场子任务完成后才能生成牧场组PPT"
        if self.metadata.get("group_results", {}).get("status") == "stale":
            return False, "牧场任务或汇总范围已变化，请先重新生成牧场组汇总Excel"
        excel_path = self._find_excel()
        if not excel_path:
            return False, "请先生成牧场组汇总Excel"

        try:
            self._progress("正在读取牧场组汇总Excel...", 5)
            overview = self._read_table(excel_path, "牧场对比总览")
            pedigree = self._read_table(excel_path, "系谱识别对比")
            traits = self._read_table(excel_path, "关键性状对比")
            yearly = self._read_table(excel_path, "年份遗传进展")
            ranking = self._read_table(excel_path, "跨牧场牛只排名")
            availability = self._read_table(excel_path, "数据可用性")
            farms = self._top_farms(overview)

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            self._add_cover(prs, len(overview))

            self._progress("正在生成数据范围与牛群概况...", 15)
            slide = self._base_slide(prs, "数据范围与任务完成情况", 2)
            current_total = sum(int(_number(row.get("在群头数"))) for row in overview)
            all_total = sum(int(_number(row.get("全部头数"))) for row in overview)
            pedigree_covered = sum(1 for row in pedigree if _number(row.get("在群头数")) > 0)
            self._add_kpis(slide, [
                ("完成牧场", f"{len(overview)} 个"),
                ("在群母牛", f"{current_total:,} 头"),
                ("全部母牛", f"{all_total:,} 头"),
                ("系谱覆盖牧场", f"{pedigree_covered} 个"),
            ])
            note = slide.shapes.add_textbox(Inches(1.0), Inches(3.45), Inches(11.3), Inches(1.7))
            self._set_text(note, "所有牧场均作为独立子任务完成计算；牧场组均值按有效头数加权，比例按汇总分子和分母重新计算。", 20, color=DARK, align=PP_ALIGN.CENTER)

            slide = self._base_slide(prs, "各牧场在群母牛规模", 3)
            self._add_bar_chart(slide, [row.get("牧场名称", "") for row in farms], [("在群头数", [_number(row.get("在群头数")) for row in farms])], horizontal=True, legend=False)

            self._progress("正在生成系谱和核心指数对比...", 35)
            pedigree_map = {str(row.get("牧场编号")): row for row in pedigree}
            pedigree_farms = [row for row in farms if str(row.get("牧场编号")) in pedigree_map]
            slide = self._base_slide(prs, "在群母牛系谱识别率", 4)
            self._add_bar_chart(
                slide,
                [row.get("牧场名称", "") for row in pedigree_farms],
                [
                    ("父号", [_number(pedigree_map[str(row.get("牧场编号"))].get("父号识别率")) * 100 for row in pedigree_farms]),
                    ("外祖父", [_number(pedigree_map[str(row.get("牧场编号"))].get("外祖父识别率")) * 100 for row in pedigree_farms]),
                ],
            )

            slide = self._base_slide(prs, "核心育种指数牧场对比", 5)
            self._add_bar_chart(
                slide,
                [row.get("牧场名称", "") for row in farms],
                [
                    ("平均NM$", [_number(row.get("平均NM$")) for row in farms]),
                    ("平均TPI/10", [_number(row.get("平均TPI")) / 10 for row in farms]),
                ],
            )

            trait_map = defaultdict(dict)
            for row in traits:
                trait_map[str(row.get("牧场编号"))][str(row.get("性状"))] = _number(row.get("在群平均值"))

            def add_trait_slide(title, trait_names, page):
                slide = self._base_slide(prs, title, page)
                series = []
                for trait in trait_names:
                    values = [trait_map[str(row.get("牧场编号"))].get(trait, 0) for row in farms]
                    series.append((trait, values))
                self._add_bar_chart(slide, [row.get("牧场名称", "") for row in farms], series)

            self._progress("正在生成关键性状分组对比...", 55)
            add_trait_slide("经济与产量性状对比", ("MILK", "FAT", "PROT"), 6)
            add_trait_slide("健康与繁殖性状对比", ("SCS", "PL", "DPR"), 7)
            add_trait_slide("体型与饲料效率性状对比", ("PTAT", "UDC", "FLC", "RFI"), 8)

            self._progress("正在生成年份趋势与牛只排名...", 72)
            year_groups = defaultdict(lambda: {"count": 0, "nm_sum": 0.0, "tpi_sum": 0.0})
            for row in yearly:
                label = str(row.get("出生年份") or "")
                headcount = int(_number(row.get("头数")))
                year_groups[label]["count"] += headcount
                year_groups[label]["nm_sum"] += _number(row.get("平均NM$")) * headcount
                year_groups[label]["tpi_sum"] += _number(row.get("平均TPI")) * headcount
            year_labels = list(year_groups)
            slide = self._base_slide(prs, "牧场组年度遗传进展", 9)
            self._add_bar_chart(
                slide,
                year_labels,
                [
                    ("加权平均NM$", [year_groups[label]["nm_sum"] / year_groups[label]["count"] if year_groups[label]["count"] else 0 for label in year_labels]),
                    ("加权平均TPI/10", [(year_groups[label]["tpi_sum"] / year_groups[label]["count"] / 10) if year_groups[label]["count"] else 0 for label in year_labels]),
                ],
            )

            top_ranking = [row for row in ranking if row.get("排名类型") == "前列"][:10]
            slide = self._base_slide(prs, "跨牧场优秀母牛排名", 10)
            self._add_table(slide, ["牧场组排名", "牧场名称", "原始牛号", "综合指数"], [[row.get("牧场组排名"), row.get("牧场名称"), row.get("原始牛号"), round(_number(row.get("综合指数")), 2)] for row in top_ranking])

            self._progress("正在生成数据可用性与重点结论...", 87)
            slide = self._base_slide(prs, "分析模块数据可用性", 11)
            self._add_table(slide, ["牧场名称", "配种记录", "备选公牛", "已配公牛分析", "近交及隐性基因"], [[row.get("牧场名称"), row.get("配种记录"), row.get("备选公牛"), row.get("已配公牛分析"), row.get("近交及隐性基因")] for row in availability])

            slide = self._base_slide(prs, "重点结论与后续使用", 12)
            best_nm = max(overview, key=lambda row: _number(row.get("平均NM$")), default={})
            best_tpi = max(overview, key=lambda row: _number(row.get("平均TPI")), default={})
            lowest_pedigree = min(pedigree, key=lambda row: _number(row.get("父号识别率"), 1), default={})
            missing_breeding = sum(1 for row in availability if row.get("配种记录") != "有")
            bullets = [
                f"平均NM$最高：{best_nm.get('牧场名称', '—')}（{_number(best_nm.get('平均NM$')):.2f}）",
                f"平均TPI最高：{best_tpi.get('牧场名称', '—')}（{_number(best_tpi.get('平均TPI')):.2f}）",
                f"父号识别率最低：{lowest_pedigree.get('牧场名称', '—')}（{_number(lowest_pedigree.get('父号识别率')):.1%}）",
                f"缺少配种记录的牧场：{missing_breeding} 个；相关章节按数据可用性展示",
                "单牧场技术明细和完整报告保留在各子项目目录中。",
            ]
            box = slide.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11.2), Inches(4.8))
            frame = box.text_frame
            frame.clear()
            frame.word_wrap = True
            for index, text in enumerate(bullets):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = f"• {text}"
                paragraph.space_after = Pt(18)
                for run in paragraph.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(20)
                    run.font.color.rgb = DARK

            self._progress("正在保存牧场组PPT...", 96)
            output_path = self.project_path / "reports" / f"牧场组育种分析汇报_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            prs.save(output_path)
            self._normalize_chart_axis_ids(output_path)
            self.last_output_path = output_path
            FileManager.update_group_result(self.project_path, ppt_path=output_path)
            self._progress("牧场组PPT生成完成", 100)
            return True, str(output_path)
        except Exception as exc:
            logger.exception("生成牧场组PPT失败")
            return False, f"生成牧场组PPT失败：{exc}"
