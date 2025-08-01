"""
PPT生成器主模块

负责生成遗传改良项目服务报告PPT
"""

import os
import json
import logging
import pandas as pd
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PyQt6.QtWidgets import QMessageBox, QInputDialog, QProgressDialog
from PyQt6.QtCore import Qt

from .data_preparation import DataPreparation
from .data_adapter import DataAdapter
from .data_validator import DataValidator
from .slide_generators import (
    TitleSlideGenerator,
    TOCSlideGenerator,
    PedigreeAnalysisGenerator,
    GeneticEvaluationGenerator,
    LinearScoreGenerator,
    BullUsageGenerator
)

logger = logging.getLogger(__name__)


class PPTGenerator:
    """PPT生成器类"""
    
    def __init__(self, output_folder: str, username: str = "用户"):
        """
        初始化PPT生成器
        
        Args:
            output_folder: 输出文件夹路径
            username: 用户名
        """
        self.output_folder = Path(output_folder)
        self.username = username
        self.farm_name = "某某某"  # 默认牧场名称
        
        # 加载配置
        self.load_configurations()
        
        # 初始化数据准备器和验证器
        self.data_prep = DataPreparation(output_folder)
        self.data_validator = DataValidator()
        
        # 初始化各部分生成器
        self.slide_generators = {
            'title': TitleSlideGenerator(self.trait_translations),
            'toc': TOCSlideGenerator(self.trait_translations),
            'pedigree': PedigreeAnalysisGenerator(self.trait_translations),
            'genetic': GeneticEvaluationGenerator(self.trait_translations),
            'linear': LinearScoreGenerator(self.trait_translations),
            'bull': BullUsageGenerator(self.trait_translations)
        }
        
    def load_configurations(self):
        """加载配置文件"""
        # 加载性状翻译配置
        config_path = Path(__file__).parent.parent.parent / 'config' / 'trait_translations.json'
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                config = json.load(f)
                self.trait_translations = config.get('traits', {})
                self.column_translations = config.get('columns', {})
                self.thresholds = config.get('thresholds', {})
        except Exception as e:
            logger.error(f"加载配置文件失败: {str(e)}")
            self.trait_translations = {}
            self.column_translations = {}
            self.thresholds = {}
    
    def generate_report(self, parent_widget=None) -> bool:
        """
        生成PPT报告的主函数
        
        Args:
            parent_widget: 父窗口部件
            
        Returns:
            是否成功生成
        """
        try:
            # 1. 准备数据
            data_ready = self.prepare_data(parent_widget)
            if not data_ready:
                return False
            
            # 2. 获取牧场名称
            if parent_widget:
                farm_name, ok = QInputDialog.getText(
                    parent_widget, 
                    '输入牧场名称', 
                    '请输入牧场名称:'
                )
                if ok and farm_name:
                    self.farm_name = farm_name
            
            # 3. 创建进度对话框
            if parent_widget:
                progress = QProgressDialog("正在生成PPT报告...", "取消", 0, 100, parent_widget)
                progress.setWindowModality(Qt.WindowModality.WindowModal)
                progress.show()
            
            # 4. 创建PPT
            prs = self.create_presentation()
            
            # 5. 生成各个部分
            slide_count = 0
            total_slides = 23  # 估计的总页数
            
            # 标题页和目录
            if parent_widget:
                progress.setValue(int(slide_count / total_slides * 100))
                progress.setLabelText("生成标题页...")
            self.slide_generators['title'].create_title_slide(prs, self.farm_name, self.username)
            slide_count += 1
            
            if parent_widget:
                progress.setValue(int(slide_count / total_slides * 100))
                progress.setLabelText("生成目录页...")
            self.slide_generators['toc'].create_toc_slide(prs)
            slide_count += 1
            
            # 系谱记录分析部分
            if parent_widget:
                progress.setLabelText("生成系谱记录分析...")
            slide_count += self.slide_generators['pedigree'].generate_all_slides(
                prs, self.output_folder, self.farm_name
            )
            if parent_widget:
                progress.setValue(int(slide_count / total_slides * 100))
            
            # 牛群遗传数据评估部分
            if parent_widget:
                progress.setLabelText("生成遗传数据评估...")
            slide_count += self.slide_generators['genetic'].generate_all_slides(
                prs, self.output_folder, self.farm_name
            )
            if parent_widget:
                progress.setValue(int(slide_count / total_slides * 100))
            
            # 体型外貌评定部分（手动制作部分）
            if parent_widget:
                progress.setLabelText("生成体型外貌评定...")
            slide_count += self.slide_generators['linear'].generate_all_slides(
                prs, self.output_folder, self.farm_name
            )
            if parent_widget:
                progress.setValue(int(slide_count / total_slides * 100))
            
            # 公牛使用分析部分
            if parent_widget:
                progress.setLabelText("生成公牛使用分析...")
            slide_count += self.slide_generators['bull'].generate_all_slides(
                prs, self.output_folder, self.farm_name
            )
            if parent_widget:
                progress.setValue(int(slide_count / total_slides * 100))
            
            # 6. 保存PPT
            output_path = self.output_folder / f"{self.farm_name}牧场遗传改良项目专项服务报告.pptx"
            prs.save(str(output_path))
            
            if parent_widget:
                progress.setValue(100)
                progress.close()
                QMessageBox.information(parent_widget, "成功", f"PPT报告已生成：\n{output_path}")
            
            logger.info(f"PPT报告生成成功: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"生成PPT报告时发生错误: {str(e)}")
            if parent_widget:
                if 'progress' in locals():
                    progress.close()
                QMessageBox.critical(parent_widget, "错误", f"生成PPT报告时发生错误：\n{str(e)}")
            return False
    
    def create_presentation(self) -> Presentation:
        """
        创建PPT演示文稿
        
        Returns:
            Presentation对象
        """
        # 查找模板文件
        template_file = self.find_template()
        
        if template_file:
            # 使用模板创建
            prs = Presentation(str(template_file))
            # 删除模板中的所有现有幻灯片
            for i in range(len(prs.slides) - 1, -1, -1):
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i]
        else:
            # 创建新的演示文稿
            prs = Presentation()
            # 设置幻灯片大小为16:9
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            
        return prs
    
    def find_template(self) -> Optional[Path]:
        """
        查找PPT模板文件
        
        Returns:
            模板文件路径，如果没找到返回None
        """
        # 支持多种模板文件名
        template_keywords = ['template', '模板', 'model', '模版']
        
        # 首先在输出文件夹查找
        for file in self.output_folder.iterdir():
            if file.suffix.lower() in ['.ppt', '.pptx']:
                for keyword in template_keywords:
                    if keyword in file.name.lower():
                        logger.info(f"找到PPT模板: {file}")
                        return file
        
        # 在项目根目录的templates文件夹查找
        templates_dir = Path(__file__).parent.parent.parent / 'templates'
        if templates_dir.exists():
            for file in templates_dir.iterdir():
                if file.suffix.lower() in ['.ppt', '.pptx']:
                    logger.info(f"找到PPT模板: {file}")
                    return file
                    
        logger.info("未找到PPT模板，将使用空白模板")
        return None
    
    def prepare_data(self, parent_widget=None) -> bool:
        """
        准备PPT生成所需的所有数据
        
        Args:
            parent_widget: 父窗口部件
            
        Returns:
            是否成功准备所有数据
        """
        try:
            # 1. 检查现有数据文件
            all_ready, missing_files = self.data_prep.check_all_files()
            
            if not all_ready:
                if parent_widget:
                    reply = QMessageBox.question(
                        parent_widget,
                        '数据缺失',
                        f"以下文件缺失：\n{chr(10).join(missing_files[:10])}\n\n是否尝试生成这些数据？",
                        QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
                    )
                    if reply != QMessageBox.StandardButton.Yes:
                        return False
                        
                # 2. 尝试加载必要的数据
                data_loaded = self.load_source_data()
                if not data_loaded:
                    if parent_widget:
                        QMessageBox.warning(
                            parent_widget, 
                            "警告", 
                            "无法加载源数据文件。请确保已完成基础数据分析。"
                        )
                    return False
                    
                # 2.5. 验证加载的数据
                all_valid, validation_report = self.validate_all_data()
                if not all_valid:
                    logger.warning(f"数据验证发现问题:\n{validation_report}")
                    if parent_widget:
                        reply = QMessageBox.question(
                            parent_widget,
                            '数据验证警告',
                            f"数据验证发现一些问题，是否继续生成PPT？\n\n{validation_report[:500]}...",
                            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
                        )
                        if reply != QMessageBox.StandardButton.Yes:
                            return False
                    
                # 3. 调用数据准备流程
                logger.info("开始准备PPT所需数据...")
                success = self.data_prep.prepare_all_data(
                    parent_widget=parent_widget,
                    cow_df=self.cow_df,
                    bull_traits=self.bull_traits,
                    selected_traits=self.selected_traits,
                    table_c=self.table_c,
                    merged_cow_traits=self.merged_cow_traits,
                    breeding_df=self.breeding_df,
                    cow_traits=self.cow_traits,
                    breeding_index_scores=self.breeding_index_scores
                )
                
                if not success:
                    logger.error("数据准备失败")
                    return False
                    
            return True
            
        except Exception as e:
            logger.error(f"准备数据时发生错误: {str(e)}")
            if parent_widget:
                QMessageBox.critical(parent_widget, "错误", f"准备数据时发生错误：{str(e)}")
            return False
    
    def load_source_data(self) -> bool:
        """
        加载源数据文件
        
        Returns:
            是否成功加载所有必要的源数据
        """
        try:
            import pandas as pd
            
            # 1. 加载母牛基本信息
            # 尝试多个可能的文件名
            cow_file_names = [
                "processed_cow_data.xlsx",
                "processed_cow_data_key_traits_scores_genomic.xlsx",
                "processed_cow_data_key_traits_scores.xlsx",
                "牛只信息.xlsx",
                "cow_data.xlsx",
                "母牛数据.xlsx"
            ]
            
            self.cow_df = None
            for filename in cow_file_names:
                cow_file = self.output_folder / filename
                if cow_file.exists():
                    df = pd.read_excel(cow_file)
                    # 使用数据适配器标准化列名
                    self.cow_df = DataAdapter.standardize_columns(
                        df, 
                        required_columns=['cow_id', 'birth_date', 'lac', '是否在场']
                    )
                    # 转换日期列
                    self.cow_df = DataAdapter.convert_date_columns(self.cow_df)
                    # 添加缺失列
                    self.cow_df = DataAdapter.add_missing_columns(
                        self.cow_df,
                        {'sex': '母', '是否在场': '是'}
                    )
                    logger.info(f"加载母牛数据: {filename}")
                    # 验证数据
                    is_valid, errors = self.data_validator.validate_dataframe(self.cow_df, 'cow_data')
                    if not is_valid:
                        logger.warning(f"母牛数据验证失败: {errors}")
                    break
                    
            if self.cow_df is None:
                logger.warning("未找到母牛数据文件")
                
            # 2. 加载公牛性状数据
            # 尝试多个可能的文件名
            bull_file_names = [
                "processed_bull_data.xlsx",
                "processed_bull_data_key_traits.xlsx",
                "公牛性状数据.xlsx",
                "bull_data.xlsx"
            ]
            
            self.bull_traits = None
            for filename in bull_file_names:
                bull_file = self.output_folder / filename
                if bull_file.exists():
                    df = pd.read_excel(bull_file)
                    # 使用数据适配器标准化列名
                    self.bull_traits = DataAdapter.standardize_columns(
                        df,
                        required_columns=['NAAB']
                    )
                    logger.info(f"加载公牛数据: {filename}")
                    # 验证数据
                    is_valid, errors = self.data_validator.validate_dataframe(self.bull_traits, 'bull_traits')
                    if not is_valid:
                        logger.warning(f"公牛数据验证失败: {errors}")
                    break
                    
            if self.bull_traits is None:
                logger.warning("未找到公牛性状数据文件")
                
            # 3. 加载选中的性状
            traits_file = self.output_folder / "selected_traits_key_traits.txt"
            if traits_file.exists():
                with open(traits_file, 'r', encoding='utf-8') as f:
                    self.selected_traits = [line.strip() for line in f if line.strip()]
            else:
                # 使用默认性状列表
                self.selected_traits = ['TPI', 'NM$', 'MILK', 'FAT', 'PROT', 'SCS', 'PL', 'DPR']
                logger.info(f"使用默认性状列表: {self.selected_traits}")
                
            # 4. 加载系谱信息（table_c）
            # 如果没有单独的系谱文件，尝试从母牛数据中提取
            pedigree_file_names = [
                "系谱信息.xlsx",
                "pedigree_data.xlsx"
            ]
            
            self.table_c = None
            for filename in pedigree_file_names:
                pedigree_file = self.output_folder / filename
                if pedigree_file.exists():
                    self.table_c = pd.read_excel(pedigree_file)
                    logger.info(f"加载系谱数据: {filename}")
                    break
                    
            # 如果没有单独的系谱文件，尝试从母牛数据中构建
            if self.table_c is None and self.cow_df is not None:
                self.table_c = self.build_pedigree_data(self.cow_df)
                
            if self.table_c is None:
                logger.warning("未找到系谱信息文件")
                
            # 5. 加载合并的母牛性状数据
            merged_file_names = [
                "processed_cow_data_key_traits_scores_genomic.xlsx",
                "processed_cow_data_key_traits_mean_by_year.xlsx",
                "母牛关键性状合并.xlsx",
                "merged_cow_traits.xlsx"
            ]
            
            self.merged_cow_traits = None
            for filename in merged_file_names:
                merged_file = self.output_folder / filename
                if merged_file.exists():
                    df = pd.read_excel(merged_file)
                    # 使用数据适配器标准化列名
                    self.merged_cow_traits = DataAdapter.standardize_columns(df)
                    # 转换日期列
                    self.merged_cow_traits = DataAdapter.convert_date_columns(self.merged_cow_traits)
                    logger.info(f"加载合并母牛性状数据: {filename}")
                    break
                    
            # 如果没有合并文件，使用原始母牛数据
            if self.merged_cow_traits is None and self.cow_df is not None:
                self.merged_cow_traits = self.cow_df
                logger.info("使用原始母牛数据作为合并数据")
                
            if self.merged_cow_traits is None:
                logger.warning("未找到合并的母牛性状数据")
                
            # 6. 加载配种记录
            breeding_file_names = [
                "配种记录.xlsx",
                "breeding_records.xlsx",
                "配种数据.xlsx"
            ]
            
            self.breeding_df = None
            for filename in breeding_file_names:
                breeding_file = self.output_folder / filename
                if breeding_file.exists():
                    self.breeding_df = pd.read_excel(breeding_file)
                    logger.info(f"加载配种记录: {filename}")
                    # 验证数据
                    is_valid, errors = self.data_validator.validate_dataframe(self.breeding_df, 'breeding_records')
                    if not is_valid:
                        logger.warning(f"配种记录验证失败: {errors}")
                    break
                    
            if self.breeding_df is None:
                logger.warning("未找到配种记录文件，创建空的配种记录")
                # 创建空的配种记录DataFrame以避免后续错误
                self.breeding_df = pd.DataFrame({
                    'BULL NAAB': [],
                    '配种日期': [],
                    '冻精类型': [],
                    '母牛号': []
                })
                
            # 7. 加载母牛性状原始数据
            # 如果已经有cow_df，直接使用
            if self.cow_df is not None:
                self.cow_traits = self.cow_df
                logger.info("使用cow_df作为母牛性状数据")
            else:
                cow_traits_file_names = [
                    "母牛性状数据.xlsx",
                    "cow_traits.xlsx"
                ]
                
                self.cow_traits = None
                for filename in cow_traits_file_names:
                    cow_traits_file = self.output_folder / filename
                    if cow_traits_file.exists():
                        self.cow_traits = pd.read_excel(cow_traits_file)
                        logger.info(f"加载母牛性状数据: {filename}")
                        break
                        
                if self.cow_traits is None:
                    logger.warning("未找到母牛性状数据文件")
                
            # 8. 加载育种指数得分
            index_file_names = [
                "processed_index_cow_index_scores.xlsx",
                "TPI_index_cow_index_scores.xlsx",
                "processed_index_cow_scores.xlsx",
                "育种指数得分.xlsx",
                "breeding_index_scores.xlsx"
            ]
            
            self.breeding_index_scores = None
            for filename in index_file_names:
                index_file = self.output_folder / filename
                if index_file.exists():
                    df = pd.read_excel(index_file)
                    # 使用数据适配器标准化列名
                    self.breeding_index_scores = DataAdapter.standardize_columns(
                        df,
                        required_columns=['cow_id', '育种指数得分']
                    )
                    logger.info(f"加载育种指数得分: {filename}")
                    # 验证数据
                    is_valid, errors = self.data_validator.validate_dataframe(self.breeding_index_scores, 'breeding_index')
                    if not is_valid:
                        logger.warning(f"育种指数数据验证失败: {errors}")
                    break
                    
            if self.breeding_index_scores is None:
                logger.warning("未找到育种指数得分文件")
                
            # 检查关键数据是否存在
            if self.cow_df is None:
                # 如果没有单独的母牛数据，但有合并数据，使用合并数据
                if self.merged_cow_traits is not None:
                    self.cow_df = self.merged_cow_traits
                    logger.info("使用合并母牛性状数据作为主数据")
                else:
                    logger.error("缺少关键数据：牛只信息")
                    return False
                
            return True
            
        except Exception as e:
            logger.error(f"加载源数据时发生错误: {str(e)}")
            return False
    
    def build_pedigree_data(self, cow_df: pd.DataFrame) -> pd.DataFrame:
        """
        从母牛数据中构建系谱数据
        
        Args:
            cow_df: 母牛数据
            
        Returns:
            系谱数据框
        """
        try:
            # 复制母牛数据
            table_c = cow_df.copy()
            
            # 添加必要的系谱列
            if 'birth_date' in table_c.columns:
                table_c['birth_year'] = pd.to_datetime(table_c['birth_date']).dt.year
            
            # 添加系谱识别列（如果不存在）
            if 'sire_id' not in table_c.columns:
                table_c['sire_id'] = ''
            if 'dam_id' not in table_c.columns:
                table_c['dam_id'] = ''
                
            # 添加识别状态列
            table_c['sire_identified'] = table_c['sire_id'].apply(
                lambda x: '已识别' if pd.notna(x) and str(x).strip() else '未识别'
            )
            table_c['mgs_identified'] = '未识别'  # 默认未识别
            table_c['mmgs_identified'] = '未识别'  # 默认未识别
            
            # 添加性别列（如果不存在）
            if 'sex' not in table_c.columns:
                table_c['sex'] = '母'
                
            # 添加在场状态（如果不存在）
            if '是否在场' not in table_c.columns:
                table_c['是否在场'] = '是'
                
            logger.info("从母牛数据构建系谱信息成功")
            return table_c
            
        except Exception as e:
            logger.error(f"构建系谱数据失败: {str(e)}")
            return None
    
    def validate_all_data(self) -> Tuple[bool, str]:
        """
        验证所有加载的数据
        
        Returns:
            (是否全部有效, 验证报告)
        """
        validation_results = {}
        
        # 验证各种数据
        if self.cow_df is not None:
            validation_results['母牛数据'] = self.data_validator.validate_dataframe(
                self.cow_df, 'cow_data'
            )
        else:
            validation_results['母牛数据'] = (False, ['数据未加载'])
            
        if self.bull_traits is not None:
            validation_results['公牛性状'] = self.data_validator.validate_dataframe(
                self.bull_traits, 'bull_traits'
            )
        else:
            validation_results['公牛性状'] = (False, ['数据未加载'])
            
        if self.breeding_df is not None:
            validation_results['配种记录'] = self.data_validator.validate_dataframe(
                self.breeding_df, 'breeding_records'
            )
        else:
            validation_results['配种记录'] = (False, ['数据未加载'])
            
        if self.breeding_index_scores is not None:
            validation_results['育种指数'] = self.data_validator.validate_dataframe(
                self.breeding_index_scores, 'breeding_index'
            )
        else:
            validation_results['育种指数'] = (False, ['数据未加载'])
            
        if self.table_c is not None:
            validation_results['系谱数据'] = self.data_validator.validate_dataframe(
                self.table_c, 'pedigree'
            )
        else:
            validation_results['系谱数据'] = (False, ['数据未加载'])
            
        # 生成验证报告
        report = self.data_validator.generate_validation_report(validation_results)
        
        # 判断是否全部有效
        all_valid = all(result[0] for result in validation_results.values())
        
        return all_valid, report