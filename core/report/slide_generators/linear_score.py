"""
牛群体型外貌评定部分生成器（手动制作部分）
"""

from pathlib import Path
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import logging

from .base import BaseSlideGenerator
from ..analysis_text_generator import AnalysisTextGenerator

logger = logging.getLogger(__name__)


class LinearScoreGenerator(BaseSlideGenerator):
    """体型外貌评定幻灯片生成器"""
    
    def __init__(self, trait_translations):
        """初始化生成器"""
        super().__init__(trait_translations)
        self.text_generator = AnalysisTextGenerator(trait_translations)
    
    def generate_all_slides(self, prs, output_folder: Path, farm_name: str) -> int:
        """
        生成所有体型外貌评定相关幻灯片
        
        注意：这部分内容需要手动制作，这里只生成占位页面
        
        Args:
            prs: Presentation对象
            output_folder: 输出文件夹路径
            farm_name: 牧场名称
            
        Returns:
            生成的幻灯片数量
        """
        slide_count = 0
        
        # 1. 标题分隔页
        self.create_section_title_slide(prs)
        slide_count += 1
        
        # 2. 体型外貌线性评分
        self.create_linear_score_slide(prs, farm_name)
        slide_count += 1
        
        # 3. 体型外貌缺陷性状占比
        self.create_defect_traits_slide(prs, farm_name)
        slide_count += 1
        
        # 4. 缺陷性状具体情况
        self.create_defect_details_slide(prs, farm_name)
        slide_count += 1
        
        # 5. 体型外貌进展情况
        self.create_linear_progress_slide(prs, farm_name)
        slide_count += 1
        
        return slide_count
    
    def create_section_title_slide(self, prs):
        """创建标题分隔页"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # 添加序号 "03"
        number_box = slide.shapes.add_textbox(Inches(5.67), Inches(1.12), Inches(2), Inches(2))
        number_frame = number_box.text_frame
        number_frame.text = "03"
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = number_frame.paragraphs[0].font
        font.name = "Arial"
        font.size = Pt(166)
        font.color.rgb = RGBColor(255, 255, 255)

        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(2.67), Inches(4.13), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "牛群体型外貌评定"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = title_frame.paragraphs[0].font
        font.name = "微软雅黑"
        font.size = Pt(40)
        font.color.rgb = RGBColor(255, 255, 255)
        
    def create_linear_score_slide(self, prs, farm_name: str):
        """创建体型外貌线性评分页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        self.add_title(slide, f"3.1 {farm_name}牧场牛群体型外貌线性评分")
        
        # 添加占位文本
        self.add_text(
            slide,
            "体型外貌线性评分内容：\n\n(暂不能自动化生成........需服务人员整理........）",
            0.583, 1.5, 11, 5,
            font_size=16
        )
        
    def create_defect_traits_slide(self, prs, farm_name: str):
        """创建体型外貌缺陷性状占比页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        self.add_title(slide, f"3.2 {farm_name}牧场牛群体型外貌缺陷性状占比")
        
        # 添加占位文本
        self.add_text(
            slide,
            "体型外貌缺陷性状占比内容：\n\n(暂不能自动化生成........需服务人员整理........）",
            0.583, 1.5, 11, 5,
            font_size=16
        )
        
    def create_defect_details_slide(self, prs, farm_name: str):
        """创建缺陷性状具体情况页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        self.add_title(slide, f"3.3 {farm_name}牧场缺陷性状具体情况")
        
        # 添加占位文本
        self.add_text(
            slide,
            "缺陷性状具体情况内容：\n\n(暂不能自动化生成........需服务人员整理........）",
            0.583, 1.5, 11, 5,
            font_size=16
        )
        
    def create_linear_progress_slide(self, prs, farm_name: str):
        """创建体型外貌进展情况页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        self.add_title(slide, f"3.4 {farm_name}牧场牛群体型外貌进展情况")
        
        # 添加占位文本
        self.add_text(
            slide,
            "体型外貌进展情况内容：\n\n(暂不能自动化生成........需服务人员整理........）",
            0.583, 1.5, 11, 5,
            font_size=16
        )