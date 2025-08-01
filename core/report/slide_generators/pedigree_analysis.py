"""
系谱记录分析部分生成器
"""

import pandas as pd
from pathlib import Path
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import logging

from .base import BaseSlideGenerator
from ..analysis_text_generator import AnalysisTextGenerator

logger = logging.getLogger(__name__)


class PedigreeAnalysisGenerator(BaseSlideGenerator):
    """系谱记录分析幻灯片生成器"""
    
    def __init__(self, trait_translations):
        """初始化生成器"""
        super().__init__(trait_translations)
        self.text_generator = AnalysisTextGenerator(trait_translations)
    
    def generate_all_slides(self, prs, output_folder: Path, farm_name: str) -> int:
        """
        生成所有系谱分析相关幻灯片
        
        Args:
            prs: Presentation对象
            output_folder: 输出文件夹路径
            farm_name: 牧场名称
            
        Returns:
            生成的幻灯片数量
        """
        slide_count = 0
        
        # 1. 标题分隔页
        self.create_section_title_slide(prs)
        slide_count += 1
        
        # 2. 系谱完整性分析
        self.create_pedigree_completeness_slide(prs, output_folder, farm_name)
        slide_count += 1
        
        # 3. 未识别牛只明细（手动制作）
        self.create_unidentified_cattle_slide(prs, farm_name)
        slide_count += 1
        
        # 4. 系谱记录准确性（手动制作）
        self.create_pedigree_accuracy_slide(prs, farm_name)
        slide_count += 1
        
        return slide_count
    
    def create_section_title_slide(self, prs):
        """创建标题分隔页"""
        slide_layout = prs.slide_layouts[1]  
        slide = prs.slides.add_slide(slide_layout)

        # 添加序号 "01"
        number_box = slide.shapes.add_textbox(Inches(5.67), Inches(1.12), Inches(2), Inches(2))
        number_frame = number_box.text_frame
        number_frame.text = "01"
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = number_frame.paragraphs[0].font
        font.name = "Arial"
        font.size = Pt(166)
        font.color.rgb = RGBColor(255, 255, 255)  
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(2.67), Inches(4.13), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "牧场系谱记录分析"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = title_frame.paragraphs[0].font
        font.name = "微软雅黑"
        font.size = Pt(40)
        font.color.rgb = RGBColor(255, 255, 255)  
        
        # 添加页脚
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(1), Inches(0.3))
        footer_frame = footer_box.text_frame
        footer_frame.text = "3"
        footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = footer_frame.paragraphs[0].font
        font.name = "微软雅黑"
        font.size = Pt(12)
        font.color.rgb = RGBColor(255, 255, 255)
        
    def create_pedigree_completeness_slide(self, prs, output_folder: Path, farm_name: str):
        """创建系谱完整性分析页"""
        slide_layout = prs.slide_layouts[5]  
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        self.add_title(slide, f"1.1 {farm_name}牧场系谱完整性分析")

        # 读取数据
        analysis_file = output_folder / "结果-系谱识别情况分析.xlsx"
        if analysis_file.exists():
            df = pd.read_excel(analysis_file)
            df = df[df['是否在场'] == '是']
            
            # 生成分析文本
            analysis_texts = self.text_generator.generate_pedigree_analysis_text(df)
            
            # 添加分析文本框
            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(1))
            text_frame = text_box.text_frame
            text_frame.text = analysis_texts.get('overview', '')
            text_frame.word_wrap = True
            for paragraph in text_frame.paragraphs:
                paragraph.font.name = "微软雅黑"
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)
            
            # 添加表格
            table = self.add_table(slide, 7, 9, 0.874, 2.5, 11.53, 3.14)
            
            # 设置表头
            headers = list(df.columns)
            if 'birth_year_group' in headers:
                headers[headers.index('birth_year_group')] = '出生年份'
            
            self.format_table_header(table, headers)
            
            # 填充数据
            for row in range(1, min(7, len(df) + 1)):
                for col in range(9):
                    if row - 1 < len(df):
                        value = str(df.iloc[row-1, col])
                        cell = table.cell(row, col)
                        
                        # 根据不同列的阈值设置样式
                        text_color = RGBColor(0, 0, 0)
                        bold = False
                        
                        if headers[col] == '父号识别率':
                            try:
                                rate = float(value.strip('%'))
                                if rate < 90.00:
                                    text_color = RGBColor(255, 0, 0)
                                    bold = True
                            except ValueError:
                                pass
                        elif headers[col] == '外祖父识别率':
                            try:
                                rate = float(value.strip('%'))
                                if rate < 70.00:
                                    text_color = RGBColor(255, 0, 0)
                                    bold = True
                            except ValueError:
                                pass
                        elif headers[col] == '外曾外祖父识别率':
                            try:
                                rate = float(value.strip('%'))
                                if rate < 50.00:
                                    text_color = RGBColor(255, 0, 0)
                                    bold = True
                            except ValueError:
                                pass
                        
                        self.format_table_cell(cell, value, text_color=text_color, bold=bold)
        
        # 添加说明文本
        textbox = slide.shapes.add_textbox(
            Inches(0.469), Inches(5.012), 
            Inches(12.358), Inches(2.22)
        )
        tf = textbox.text_frame
        
        # 未识别原因标题
        p = tf.add_paragraph()
        p.text = "未识别原因："
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        # 添加占位文本
        for section in ["父号未识别原因：", "外祖父号未识别原因："]:
            p = tf.add_paragraph()
            p.text = ""
            p.font.name = "微软雅黑"
            p.font.size = Pt(14)
            p.level = 1
            
            p = tf.add_paragraph()
            p.text = section
            p.font.name = "微软雅黑"
            p.font.size = Pt(14)
            p.font.bold = True
            p.level = 1
            
            p = tf.add_paragraph()
            p.text = ""
            p.font.name = "微软雅黑"
            p.font.size = Pt(14)
            p.level = 1
            
    def create_unidentified_cattle_slide(self, prs, farm_name: str):
        """创建未识别牛只明细页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        self.add_title(slide, f"1.2 {farm_name}牧场未识别牛只明细")
        
        # 添加说明文本
        self.add_text(
            slide, 
            "未识别牛只明细：(暂不能自动化生成........需服务人员整理........）",
            0.583, 1.5, 11, 0.5,
            font_size=16
        )
        
    def create_pedigree_accuracy_slide(self, prs, farm_name: str):
        """创建系谱记录准确性页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        self.add_title(slide, f"1.3 {farm_name}牧场系谱记录准确性分析")
        
        # 添加说明文本
        self.add_text(
            slide, 
            "系谱记录准确性：(暂不能自动化生成........需服务人员整理........根据基因组检测结果...........）",
            0.583, 1.5, 11, 0.5,
            font_size=16
        )