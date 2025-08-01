"""
牛群遗传数据评估部分生成器
"""

import os
import pandas as pd
from pathlib import Path
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
import logging

from .base import BaseSlideGenerator
from ..analysis_text_generator import AnalysisTextGenerator

logger = logging.getLogger(__name__)


class GeneticEvaluationGenerator(BaseSlideGenerator):
    """遗传数据评估幻灯片生成器"""
    
    def __init__(self, trait_translations):
        """初始化生成器"""
        super().__init__(trait_translations)
        self.text_generator = AnalysisTextGenerator(trait_translations)
    
    def generate_all_slides(self, prs, output_folder: Path, farm_name: str) -> int:
        """
        生成所有遗传评估相关幻灯片
        
        Args:
            prs: Presentation对象
            output_folder: 输出文件夹路径
            farm_name: 牧场名称
            
        Returns:
            生成的幻灯片数量
        """
        slide_count = 0
        
        # 1. 标题分隔页
        self.create_section_title_slide(prs)
        slide_count += 1
        
        # 2. 关键育种性状分析
        self.create_key_breeding_traits_slide(prs, output_folder, farm_name)
        slide_count += 1
        
        # 3. NM$分布情况分析
        self.create_nm_distribution_slide(prs, output_folder, farm_name)
        slide_count += 1
        
        # 4. NM$正态分布分析
        if self.create_nm_normal_distribution_slide(prs, output_folder, farm_name):
            slide_count += 1
            
        # 5. 育种指数正态分布分析
        if self.create_index_normal_distribution_slide(prs, output_folder, farm_name):
            slide_count += 1
            
        # 6. 关键育种性状进展分析
        progress_count = self.create_traits_progress_slides(prs, output_folder, farm_name)
        slide_count += progress_count
        
        return slide_count
    
    def create_section_title_slide(self, prs):
        """创建标题分隔页"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # 添加序号 "02"
        number_box = slide.shapes.add_textbox(Inches(5.67), Inches(1.12), Inches(2), Inches(2))
        number_frame = number_box.text_frame
        number_frame.text = "02"
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = number_frame.paragraphs[0].font
        font.name = "Arial"
        font.size = Pt(166)
        font.color.rgb = RGBColor(255, 255, 255)

        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(2.67), Inches(4.13), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "牛群遗传数据评估"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = title_frame.paragraphs[0].font
        font.name = "微软雅黑"
        font.size = Pt(40)
        font.color.rgb = RGBColor(255, 255, 255)
        
    def create_key_breeding_traits_slide(self, prs, output_folder: Path, farm_name: str):
        """创建关键育种性状分析页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        self.add_title(slide, f"2.1 {farm_name}牧场关键育种性状分析")
        
        # 读取数据
        analysis_file = output_folder / "结果-牛群关键性状年度变化.xlsx"
        if analysis_file.exists():
            df = pd.read_excel(analysis_file)
            df = df[df['是否在场'] == '是']
            
            # 准备表头
            headers = list(df.columns)
            if 'birth_year_group' in headers:
                headers[headers.index('birth_year_group')] = '出生年份'
            
            # 添加表格（双行表头）
            table = self.add_table(slide, 8, len(headers), 0.1535, 1.5945, 13.0276, 3.76)
            
            # 设置第一行表头（中文）
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                chinese_name = self.translate_trait(header) if header in self.trait_translations else header
                self.format_table_cell(
                    cell, chinese_name,
                    font_size=12,
                    bg_color=RGBColor(79, 129, 189),
                    text_color=RGBColor(255, 255, 255),
                    bold=True
                )
                
            # 设置第二行表头（英文）
            for i, header in enumerate(headers):
                cell = table.cell(1, i)
                self.format_table_cell(
                    cell, header,
                    font_size=10,
                    bg_color=RGBColor(79, 129, 189),
                    text_color=RGBColor(255, 255, 255),
                    bold=True
                )
                
            # 填充数据
            for row in range(2, min(8, len(df) + 2)):
                for col in range(len(headers)):
                    cell = table.cell(row, col)
                    if row - 2 < len(df):
                        value = str(df.iloc[row-2, col])
                        self.format_table_cell(cell, value, font_size=10)
                    else:
                        self.format_table_cell(cell, "", font_size=10)
                        
        # 添加说明文本
        textbox = slide.shapes.add_textbox(
            Inches(0.469), Inches(5.6614), 
            Inches(12.358), Inches(2.22)
        )
        tf = textbox.text_frame
        
        p = tf.add_paragraph()
        p.text = "需关注的关键性状："
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        p = tf.add_paragraph()
        p.text = "暂不能自动化生成........需服务人员整理........"
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1
        
    def create_nm_distribution_slide(self, prs, output_folder: Path, farm_name: str):
        """创建NM$分布情况分析页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        self.add_title(slide, f"2.2 {farm_name}牧场净利润值NM$分布情况分析")
        
        # 读取数据
        excel_file = output_folder / "在群牛只净利润值分布.xlsx"
        if excel_file.exists():
            df = pd.read_excel(excel_file, sheet_name="NM$分布")
            
            # 添加柱状图
            chart_data = CategoryChartData()
            chart_data.categories = df['NM$区间']
            chart_data.add_series('头数', df['头数'])
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, 
                Inches(0.5), Inches(1.5), 
                Inches(6), Inches(3.5), 
                chart_data
            ).chart
            
            chart.has_title = False
            chart.has_legend = True
            chart.font.size = Pt(9)
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False
            
            # 添加数据标签
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = Pt(9)
            data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            
            # 添加饼图
            pie_data = CategoryChartData()
            pie_data.categories = df['NM$区间']
            pie_data.add_series('占比', df['占比'].tolist())
            
            pie = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, 
                Inches(7), Inches(1.5), 
                Inches(6), Inches(3.2244), 
                pie_data
            ).chart
            
            pie.has_title = False
            pie.has_legend = True
            pie.font.size = Pt(9)
            
            # 添加表格
            table = self.add_table(slide, len(df) + 1, 3, 7, 4.7834, 6, 2.433)
            
            # 设置表头
            headers = ['NM$区间', '头数', '占比']
            self.format_table_header(table, headers)
            
            # 填充数据
            for row in range(1, len(df) + 1):
                for col in range(3):
                    cell = table.cell(row, col)
                    value = df.iloc[row-1, col]
                    if col == 2:  # 占比列
                        value = f"{value:.2%}"
                    self.format_table_cell(cell, str(value))
                    
        # 添加说明文本
        self.add_text(
            slide,
            "净利润值NM$分布情况：\n\n暂不能自动化生成........需服务人员整理........",
            0.5, 5, 6, 2.4,
            font_size=14
        )
        
    def create_nm_normal_distribution_slide(self, prs, output_folder: Path, farm_name: str) -> bool:
        """创建NM$正态分布分析页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        self.add_title(slide, f"2.3 {farm_name}牧场净利润值正态分布分析")
        
        # 添加图片
        img_path = output_folder / "NM$_年份正态分布.png"
        if img_path.exists():
            self.add_image(slide, str(img_path), 0.5, 1.5, 12, 4)
        else:
            return False
            
        # 添加说明文本
        textbox = slide.shapes.add_textbox(
            Inches(0.469), Inches(5.6614), 
            Inches(12.358), Inches(2.22)
        )
        tf = textbox.text_frame
        
        p = tf.add_paragraph()
        p.text = "牧场净利润值正态分布情况："
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "暂不能自动化生成........需服务人员整理........"
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.level = 1
        
        return True
        
    def create_index_normal_distribution_slide(self, prs, output_folder: Path, farm_name: str) -> bool:
        """创建育种指数正态分布分析页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        self.add_title(slide, f"2.4 {farm_name}牧场育种指数正态分布分析")
        
        # 添加图片
        img_path = output_folder / "育种指数得分_年份正态分布.png"
        if img_path.exists():
            self.add_image(slide, str(img_path), 0.5, 1.5, 12, 4)
        else:
            return False
            
        # 添加说明文本
        textbox = slide.shapes.add_textbox(
            Inches(0.469), Inches(5.6614), 
            Inches(12.358), Inches(2.22)
        )
        tf = textbox.text_frame
        
        p = tf.add_paragraph()
        p.text = "牧场育种指数正态分布情况："
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "暂不能自动化生成........需服务人员整理........"
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.level = 1
        
        return True
        
    def create_traits_progress_slides(self, prs, output_folder: Path, farm_name: str) -> int:
        """创建关键性状进展分析页"""
        slide_count = 0
        
        # 获取选中的性状
        traits_file = output_folder / "selected_traits_key_traits.txt"
        if traits_file.exists():
            with open(traits_file, 'r', encoding='utf-8') as f:
                selected_traits = [line.strip() for line in f if line.strip()]
        else:
            return 0
            
        # 进展图文件夹
        charts_folder = output_folder / "结果-牛群关键性状进展图"
        if not charts_folder.exists():
            return 0
            
        # 为每个性状创建一页
        for trait in selected_traits:
            img_filename = f"{trait}进展情况.png"
            img_path = charts_folder / img_filename
            if not img_path.exists():
                continue
                
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            
            # 添加标题
            self.add_title(slide, f"2.5 {farm_name}牧场关键育种性状进展分析")
            
            # 添加图片
            self.add_image(slide, str(img_path), 0.5, 1.5, width=12)
            
            # 添加说明文本
            chinese_trait = self.translate_trait(trait)
            textbox = slide.shapes.add_textbox(
                Inches(0.469), Inches(5.6614), 
                Inches(12.358), Inches(2.22)
            )
            tf = textbox.text_frame
            
            p = tf.add_paragraph()
            p.text = f"{chinese_trait} ({trait}) 性状进展情况："
            p.font.name = "微软雅黑"
            p.font.size = Pt(16)
            p.font.bold = True
            
            p = tf.add_paragraph()
            p.text = "暂不能自动化生成........需服务人员整理........"
            p.font.name = "微软雅黑"
            p.font.size = Pt(14)
            p.level = 1
            
            slide_count += 1
            
        return slide_count