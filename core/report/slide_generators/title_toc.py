"""
标题页和目录页生成器
"""

from datetime import datetime
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from .base import BaseSlideGenerator


class TitleSlideGenerator(BaseSlideGenerator):
    """标题页生成器"""
    
    def create_title_slide(self, prs, farm_name: str, username: str):
        """
        创建标题页
        
        Args:
            prs: Presentation对象
            farm_name: 牧场名称
            username: 用户名
        """
        slide_layout = prs.slide_layouts[0]  # 使用第一个布局（通常是标题页布局）
        slide = prs.slides.add_slide(slide_layout)

        # 设置标题
        title = f"{farm_name}牧场遗传改良项目专项服务"
        left = Inches(2)  
        top = Inches(4.5)  
        width = Inches(12)  
        height = Inches(1.5)
        title_shape = slide.shapes.add_textbox(left, top, width, height)

        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.name = "微软雅黑"
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        title_frame.paragraphs[0].font.bold = True

        # 添加作者和日期
        left = Inches(9)  
        top = Inches(6)
        width = Inches(3.5)
        height = Inches(1.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.add_paragraph()
        p.text = f"奶科院育种中心 {username}"
        p.font.name = "微软雅黑"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.RIGHT  

        p = tf.add_paragraph()
        current_date = datetime.now().strftime("%Y年%m月%d日")
        p.text = current_date
        p.font.name = "微软雅黑"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.RIGHT


class TOCSlideGenerator(BaseSlideGenerator):
    """目录页生成器"""
    
    def create_toc_slide(self, prs):
        """创建目录页"""
        slide_layout = prs.slide_layouts[5]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        self.add_title(slide, "目录", left=2.5, width=8)

        # 创建网格布局的目录项
        grid_items = [
            ("01", "牧场系谱记录分析", "系谱识别率 系谱准确性"),
            ("02", "牛群遗传数据评估", "育种指数 性状分析 遗传进展"),
            ("03", "牛群体型外貌评定", "线性评分 缺陷性状 进展分析"),
            ("04", "牧场公牛使用分析", "使用情况 使用进展"),
        ]

        # 设置网格布局参数
        columns = 2
        cell_width = Inches(5.5)
        cell_height = Inches(2)
        spacing = Inches(0.5)

        # 计算网格的起始位置（居中）
        total_grid_width = columns * cell_width + (columns - 1) * spacing
        slide_width = prs.slide_width
        grid_left = (slide_width - total_grid_width) / 2  

        rows = (len(grid_items) + columns - 1) // columns
        total_grid_height = rows * cell_height + (rows - 1) * spacing
        slide_height = prs.slide_height
        grid_top = (slide_height - total_grid_height) / 2 + Inches(0.2)  

        for i, (number, title, subtitle) in enumerate(grid_items):
            left = grid_left + (i % columns) * (cell_width + spacing)
            top = grid_top + (i // columns) * (cell_height + spacing)

            # 添加编号
            number_box = slide.shapes.add_textbox(left, top, Inches(0.5), Inches(0.5))
            number_frame = number_box.text_frame
            number_frame.text = number
            number_frame.paragraphs[0].font.name = "微软雅黑"
            number_frame.paragraphs[0].font.size = Pt(24)
            number_frame.paragraphs[0].font.color.rgb = RGBColor(0, 176, 240)  # 淡蓝色

            # 添加主标题
            title_box = slide.shapes.add_textbox(left, top + Inches(0.5), cell_width, Inches(0.4))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.name = "微软雅黑"
            title_frame.paragraphs[0].font.size = Pt(18)
            title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  
            title_frame.paragraphs[0].font.bold = True

            # 添加副标题
            subtitle_box = slide.shapes.add_textbox(left, top + Inches(0.9), cell_width, Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.clear()  

            subtitle_items = subtitle.split()
            for item in subtitle_items:
                p = subtitle_frame.add_paragraph()
                p.text = "• " + item  
                p.font.name = "微软雅黑"
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(89, 89, 89)  
                p.level = 1