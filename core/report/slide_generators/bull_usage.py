"""
牧场公牛使用分析部分生成器
"""

import pandas as pd
from pathlib import Path
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import logging

from .base import BaseSlideGenerator
from ..analysis_text_generator import AnalysisTextGenerator

logger = logging.getLogger(__name__)


class BullUsageGenerator(BaseSlideGenerator):
    """公牛使用分析幻灯片生成器"""
    
    def __init__(self, trait_translations):
        """初始化生成器"""
        super().__init__(trait_translations)
        self.text_generator = AnalysisTextGenerator(trait_translations)
    
    def generate_all_slides(self, prs, output_folder: Path, farm_name: str) -> int:
        """
        生成所有公牛使用分析相关幻灯片
        
        Args:
            prs: Presentation对象
            output_folder: 输出文件夹路径
            farm_name: 牧场名称
            
        Returns:
            生成的幻灯片数量
        """
        slide_count = 0
        
        # 1. 标题分隔页
        self.create_section_title_slide(prs)
        slide_count += 1
        
        # 2. 公牛使用整体情况
        self.create_bull_usage_overview_slide(prs, farm_name)
        slide_count += 1
        
        # 3. 公牛使用明细
        detail_count = self.create_bull_usage_detail_slides(prs, output_folder, farm_name)
        slide_count += detail_count
        
        # 4. 公牛使用进展
        self.create_semen_usage_progress_slide(prs, output_folder, farm_name)
        slide_count += 1
        
        # 5. 公牛使用散点图
        self.create_semen_usage_scatter_slide(prs, farm_name)
        slide_count += 1
        
        # 6. 结束页
        self.create_bye_slide(prs)
        slide_count += 1
        
        return slide_count
    
    def create_section_title_slide(self, prs):
        """创建标题分隔页"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # 添加序号 "04"
        number_box = slide.shapes.add_textbox(Inches(5.67), Inches(1.12), Inches(2), Inches(2))
        number_frame = number_box.text_frame
        number_frame.text = "04"
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = number_frame.paragraphs[0].font
        font.name = "Arial"
        font.size = Pt(166)
        font.color.rgb = RGBColor(255, 255, 255)

        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(2.67), Inches(4.13), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "牧场公牛使用分析"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = title_frame.paragraphs[0].font
        font.name = "微软雅黑"
        font.size = Pt(40)
        font.color.rgb = RGBColor(255, 255, 255)
        
    def create_bull_usage_overview_slide(self, prs, farm_name: str):
        """创建公牛使用整体情况页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        self.add_title(slide, f"4.1 {farm_name}牧场近几年公牛使用整体情况")
        
        # 添加说明文本
        textbox = slide.shapes.add_textbox(
            Inches(0.583), Inches(1.5), 
            Inches(12), Inches(5.5)
        )
        tf = textbox.text_frame
        
        p = tf.add_paragraph()
        p.text = "公牛使用整体情况分析："
        p.font.name = "微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "\n基于配种记录数据的分析内容："
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        
        p = tf.add_paragraph()
        p.text = "• 年度公牛使用数量统计"
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• 冻精类型分布分析"
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• 主要公牛使用频次"
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "\n(需根据配种记录数据生成具体内容）"
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        p.font.italic = True
        
    def create_bull_usage_detail_slides(self, prs, output_folder: Path, farm_name: str) -> int:
        """创建公牛使用明细页"""
        # 这里简化处理，只创建一页示例
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        self.add_title(slide, f"4.2 {farm_name}牧场近几年公牛使用明细")
        
        # 添加说明文本
        self.add_text(
            slide,
            "公牛使用明细内容：\n\n• 按年份分组展示\n• 按冻精类型分类\n• 详细使用记录表\n\n(需根据配种记录数据生成具体内容）",
            0.583, 1.5, 11, 5,
            font_size=14
        )
        
        return 1  # 返回生成的页数
        
    def create_semen_usage_progress_slide(self, prs, output_folder: Path, farm_name: str):
        """创建公牛使用进展页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        self.add_title(slide, f"4.3 {farm_name}牧场近几年公牛使用进展")
        
        # 检查数据文件
        trend_file = output_folder / "结果-冻精使用趋势图.xlsx"
        if trend_file.exists():
            # TODO: 读取数据并生成图表
            self.add_text(
                slide,
                "冻精使用趋势分析：\n\n(基于 结果-冻精使用趋势图.xlsx 数据生成图表）",
                0.583, 1.5, 11, 5,
                font_size=14
            )
        else:
            self.add_text(
                slide,
                "冻精使用趋势分析：\n\n(数据文件未找到，需要生成冻精使用趋势数据）",
                0.583, 1.5, 11, 5,
                font_size=14
            )
            
    def create_semen_usage_scatter_slide(self, prs, farm_name: str):
        """创建公牛使用散点图页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        self.add_title(slide, f"4.4 {farm_name}牧场近几年公牛使用散点图")
        
        # 添加说明文本
        self.add_text(
            slide,
            "公牛使用散点图分析：\n\n• 横轴：公牛性状指标\n• 纵轴：使用频次\n• 点大小：代表使用量\n\n(需根据配种记录和公牛性状数据生成散点图）",
            0.583, 1.5, 11, 5,
            font_size=14
        )
        
    def create_bye_slide(self, prs):
        """创建结束页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加感谢文字
        thank_box = slide.shapes.add_textbox(
            Inches(4), Inches(3.5), 
            Inches(8), Inches(2)
        )
        tf = thank_box.text_frame
        
        p = tf.paragraphs[0]
        p.text = "谢谢！"
        p.font.name = "微软雅黑"
        p.font.size = Pt(60)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(0, 112, 192)  # 蓝色
        
        # 添加联系信息
        contact_box = slide.shapes.add_textbox(
            Inches(4), Inches(5.5), 
            Inches(8), Inches(1)
        )
        cf = contact_box.text_frame
        
        p = cf.paragraphs[0]
        p.text = "奶科院育种中心"
        p.font.name = "微软雅黑"
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(89, 89, 89)