"""
幻灯片生成器基类
"""

import logging
from pathlib import Path
from typing import Dict, Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


class BaseSlideGenerator:
    """幻灯片生成器基类"""
    
    def __init__(self, trait_translations: Dict[str, str]):
        """
        初始化生成器
        
        Args:
            trait_translations: 性状翻译字典
        """
        self.trait_translations = trait_translations
        
    def add_title(self, slide, title: str, 
                  left: float = 0.583, top: float = 0.732, 
                  width: float = 9, height: float = 0.571):
        """
        为幻灯片添加标题
        
        Args:
            slide: 幻灯片对象
            title: 标题文本
            left, top, width, height: 位置和大小（英寸）
        """
        title_shape = slide.shapes.add_textbox(
            Inches(left), Inches(top), 
            Inches(width), Inches(height)
        )
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "微软雅黑"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)
        
    def add_text(self, slide, text: str, 
                 left: float, top: float, width: float, height: float,
                 font_size: int = 14, bold: bool = False, 
                 alignment: PP_ALIGN = PP_ALIGN.LEFT):
        """
        添加文本框
        
        Args:
            slide: 幻灯片对象
            text: 文本内容
            left, top, width, height: 位置和大小（英寸）
            font_size: 字体大小
            bold: 是否加粗
            alignment: 对齐方式
        """
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), 
            Inches(width), Inches(height)
        )
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "微软雅黑"
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.alignment = alignment
        
    def add_table(self, slide, rows: int, cols: int,
                  left: float, top: float, width: float, height: float):
        """
        添加表格
        
        Args:
            slide: 幻灯片对象
            rows: 行数
            cols: 列数
            left, top, width, height: 位置和大小（英寸）
            
        Returns:
            表格对象
        """
        table = slide.shapes.add_table(
            rows, cols,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        ).table
        return table
        
    def format_table_header(self, table, headers: list, 
                           bg_color: RGBColor = RGBColor(79, 129, 189)):
        """
        格式化表格表头
        
        Args:
            table: 表格对象
            headers: 表头列表
            bg_color: 背景颜色
        """
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            para = cell.text_frame.paragraphs[0]
            para.font.name = "微软雅黑"
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
    def format_table_cell(self, cell, value: str, 
                         font_size: int = 12,
                         bg_color: RGBColor = RGBColor(233, 237, 244),
                         text_color: RGBColor = RGBColor(0, 0, 0),
                         bold: bool = False,
                         alignment: PP_ALIGN = PP_ALIGN.CENTER):
        """
        格式化表格单元格
        
        Args:
            cell: 单元格对象
            value: 单元格值
            font_size: 字体大小
            bg_color: 背景颜色
            text_color: 文字颜色
            bold: 是否加粗
            alignment: 对齐方式
        """
        cell.text = str(value)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
        para = cell.text_frame.paragraphs[0]
        para.font.name = "微软雅黑"
        para.font.size = Pt(font_size)
        para.font.color.rgb = text_color
        para.font.bold = bold
        para.alignment = alignment
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
    def add_image(self, slide, image_path: str,
                  left: float, top: float, 
                  width: float = None, height: float = None):
        """
        添加图片
        
        Args:
            slide: 幻灯片对象
            image_path: 图片路径
            left, top: 位置（英寸）
            width, height: 大小（英寸），如果不指定则使用原始大小
        """
        if width and height:
            slide.shapes.add_picture(
                image_path, 
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )
        elif width:
            slide.shapes.add_picture(
                image_path, 
                Inches(left), Inches(top),
                width=Inches(width)
            )
        elif height:
            slide.shapes.add_picture(
                image_path, 
                Inches(left), Inches(top),
                height=Inches(height)
            )
        else:
            slide.shapes.add_picture(
                image_path, 
                Inches(left), Inches(top)
            )
            
    def translate_trait(self, trait: str) -> str:
        """
        翻译性状名称
        
        Args:
            trait: 英文性状名
            
        Returns:
            中文性状名
        """
        return self.trait_translations.get(trait, trait)