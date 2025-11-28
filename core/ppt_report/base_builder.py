"""
基础幻灯片构建器 - 提供统一的格式化方法
"""

import logging
from pathlib import Path
from typing import Optional, List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.slide import Slide

from .config import (
    LAYOUT_SECTION, LAYOUT_CONTENT_1,
    COLOR_PRIMARY, COLOR_TEXT_TITLE, COLOR_TEXT_BODY, COLOR_TEXT_SECONDARY, COLOR_TEXT_COMMENT,
    COLOR_TABLE_HEADER, COLOR_TABLE_ODD, COLOR_TABLE_EVEN,
    FONT_NAME_CN, FONT_NAME_EN,
    FONT_SIZE_TITLE, FONT_SIZE_SUBTITLE, FONT_SIZE_BODY, FONT_SIZE_COMMENT,
    FONT_SIZE_SECTION_TITLE, FONT_SIZE_SECTION_NUMBER,
    FONT_SIZE_TABLE_HEADER, FONT_SIZE_TABLE_DATA,
    LINE_SPACING_TITLE, LINE_SPACING_BODY,
    TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
    CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT,
    COMMENT_LEFT, COMMENT_TOP, COMMENT_WIDTH, COMMENT_HEIGHT,
    TABLE_ROW_HEIGHT
)

logger = logging.getLogger(__name__)


class BaseSlideBuilder:
    """基础幻灯片构建器"""

    def __init__(self, prs: Presentation, chart_creator):
        """
        初始化构建器

        Args:
            prs: PowerPoint演示文稿对象
            chart_creator: 图表创建器实例
        """
        self.prs = prs
        self.chart_creator = chart_creator

    def add_title(self, slide: Slide, title: str, subtitle: Optional[str] = None):
        """
        添加页面标题（优先使用占位符）

        Args:
            slide: 幻灯片对象
            title: 标题文本
            subtitle: 副标题文本（可选）
        """
        # 尝试使用布局中的 TITLE 占位符
        title_placeholder = None
        for shape in slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type == 1:  # TITLE = 1
                    title_placeholder = shape
                    break

        if title_placeholder:
            # 使用占位符（继承模板样式）
            title_frame = title_placeholder.text_frame
            title_frame.text = title
            # 占位符会自动应用模板样式，但可以微调
            if title_frame.paragraphs:
                para = title_frame.paragraphs[0]
                para.font.name = FONT_NAME_CN
                para.font.bold = True
        else:
            # 如果没有占位符，手动创建文本框
            title_box = slide.shapes.add_textbox(
                TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.word_wrap = True

            # 格式化标题
            para = title_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            para.font.name = FONT_NAME_CN
            para.font.size = FONT_SIZE_TITLE  # 32pt
            para.font.bold = True
            para.font.color.rgb = COLOR_TEXT_TITLE
            para.line_spacing = 1.1  # 标题行距1.1倍

        # 副标题（如果有）
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                TITLE_LEFT,
                TITLE_TOP + TITLE_HEIGHT,
                TITLE_WIDTH,
                Inches(0.4)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle

            sub_para = subtitle_frame.paragraphs[0]
            sub_para.alignment = PP_ALIGN.LEFT
            sub_para.font.name = FONT_NAME_CN
            sub_para.font.size = FONT_SIZE_SUBTITLE  # 20pt
            sub_para.font.color.rgb = COLOR_TEXT_SECONDARY  # 中灰色（层次感）
            sub_para.line_spacing = 1.2  # 副标题行距1.2倍

    def add_comment(self, slide: Slide, comment: str):
        """
        添加评论框（页面底部） - 应用专业设计标准

        Args:
            slide: 幻灯片对象
            comment: 评论文本
        """
        comment_box = slide.shapes.add_textbox(
            COMMENT_LEFT, COMMENT_TOP, COMMENT_WIDTH, COMMENT_HEIGHT
        )
        comment_frame = comment_box.text_frame
        comment_frame.text = f"💡 分析：{comment}"
        comment_frame.word_wrap = True

        # 格式化评论 - 应用专业设计标准
        para = comment_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        para.font.name = FONT_NAME_CN
        para.font.size = FONT_SIZE_COMMENT  # 16pt（可读性好）
        para.font.color.rgb = COLOR_TEXT_COMMENT  # 浅灰色
        para.font.italic = True
        para.line_spacing = 1.3  # 行距1.3倍

        # 添加浅色背景（留白原则：背景与内容区分）
        fill = comment_box.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(247, 247, 247)  # 更浅的灰色背景

    def add_table(
        self,
        slide: Slide,
        data: List[List[str]],
        has_header: bool = True,
        position: Optional[Tuple[float, float, float, float]] = None
    ):
        """
        添加表格

        Args:
            slide: 幻灯片对象
            data: 表格数据（二维列表）
            has_header: 第一行是否为表头
            position: 表格位置 (left, top, width, height)，默认使用内容区域
        """
        if not data or not data[0]:
            logger.warning("表格数据为空")
            return None

        rows = len(data)
        cols = len(data[0])

        # 使用默认位置或自定义位置
        if position:
            left, top, width, height = position
        else:
            left = CONTENT_LEFT
            top = CONTENT_TOP
            width = CONTENT_WIDTH
            height = min(CONTENT_HEIGHT, TABLE_ROW_HEIGHT * (rows + 1))

        # 创建表格
        table = slide.shapes.add_table(
            rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
        ).table

        # 填充数据并格式化
        for row_idx, row_data in enumerate(data):
            for col_idx, cell_value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_value) if cell_value is not None else ""

                # 格式化单元格
                para = cell.text_frame.paragraphs[0]
                para.alignment = PP_ALIGN.CENTER
                para.font.name = FONT_NAME_CN

                if row_idx == 0 and has_header:
                    # 表头格式
                    para.font.size = FONT_SIZE_TABLE_HEADER
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(255, 255, 255)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_TABLE_HEADER
                else:
                    # 数据行格式
                    para.font.size = FONT_SIZE_TABLE_DATA
                    cell.fill.solid()
                    # 交替行颜色
                    if row_idx % 2 == 1:
                        cell.fill.fore_color.rgb = COLOR_TABLE_ODD
                    else:
                        cell.fill.fore_color.rgb = COLOR_TABLE_EVEN

        logger.info(f"添加表格: {rows}行 x {cols}列")
        return table

    def add_image(
        self,
        slide: Slide,
        image_path: Path,
        position: Optional[Tuple[float, float, float, float]] = None
    ):
        """
        添加图片

        Args:
            slide: 幻灯片对象
            image_path: 图片路径
            position: 图片位置 (left, top, width, height)，默认使用内容区域
        """
        if not image_path.exists():
            logger.error(f"图片不存在: {image_path}")
            return None

        # 使用默认位置或自定义位置
        if position:
            left, top, width, height = position
        else:
            left = CONTENT_LEFT
            top = CONTENT_TOP
            width = CONTENT_WIDTH
            height = CONTENT_HEIGHT

        try:
            picture = slide.shapes.add_picture(
                str(image_path),
                Inches(left),
                Inches(top),
                width=Inches(width),
                height=Inches(height)
            )
            logger.info(f"添加图片: {image_path.name}")
            return picture
        except Exception as e:
            logger.error(f"添加图片失败: {e}")
            return None

    def add_section_slide(self, section_number: str, section_title: str):
        """
        添加章节标题页（使用间隔页布局）

        Args:
            section_number: 章节编号 (如 "01", "02")
            section_title: 章节标题
        """
        # 使用布局1：间隔页
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[LAYOUT_SECTION])

        # 计算居中位置
        slide_width = self.prs.slide_width.inches
        slide_height = self.prs.slide_height.inches

        # 添加大号章节编号（水印效果，半透明）
        number_width = 3
        number_height = 2.5
        number_box = slide.shapes.add_textbox(
            Inches((slide_width - number_width) / 2),
            Inches(1.5),
            Inches(number_width),
            Inches(number_height)
        )
        number_frame = number_box.text_frame
        number_frame.text = section_number
        number_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        para = number_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.font.name = FONT_NAME_EN
        para.font.size = FONT_SIZE_SECTION_NUMBER  # 180pt水印
        para.font.color.rgb = RGBColor(230, 230, 230)  # 更浅的灰色水印效果
        para.font.bold = True

        # 添加章节标题（居中） - 专业设计标准
        title_width = 10
        title_box = slide.shapes.add_textbox(
            Inches((slide_width - title_width) / 2),
            Inches(slide_height / 2 - 0.5),
            Inches(title_width),
            Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.name = FONT_NAME_CN
        title_para.font.size = FONT_SIZE_SECTION_TITLE  # 48pt（设计标准40-48pt）
        title_para.font.bold = True
        title_para.font.color.rgb = COLOR_PRIMARY  # 统一深蓝色
        title_para.line_spacing = LINE_SPACING_TITLE  # 1.1倍行距

        logger.info(f"添加章节页: {section_number} - {section_title}")
        return slide

    def add_content_slide(self, title: str, layout_type: int = LAYOUT_CONTENT_1):
        """
        添加内容页（正文页）

        Args:
            title: 页面标题
            layout_type: 布局类型（LAYOUT_CONTENT_1 或 LAYOUT_CONTENT_2）

        Returns:
            幻灯片对象
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_type])
        self.add_title(slide, title)
        logger.info(f"添加内容页: {title}")
        return slide

    def add_bullet_points(
        self,
        slide: Slide,
        points: List[str],
        position: Optional[Tuple[float, float, float, float]] = None
    ):
        """
        添加项目符号列表

        Args:
            slide: 幻灯片对象
            points: 要点列表
            position: 文本框位置，默认使用内容区域
        """
        if position:
            left, top, width, height = position
        else:
            left = CONTENT_LEFT
            top = CONTENT_TOP
            width = CONTENT_WIDTH
            height = CONTENT_HEIGHT

        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for i, point in enumerate(points):
            if i == 0:
                para = text_frame.paragraphs[0]
            else:
                para = text_frame.add_paragraph()

            para.text = point
            para.level = 0
            para.font.name = FONT_NAME_CN
            para.font.size = FONT_SIZE_BODY
            para.space_before = Pt(6)
            para.space_after = Pt(6)

            # 添加项目符号
            para.bullet = True

        logger.info(f"添加项目符号列表: {len(points)}项")
        return text_box

    def add_statistics_cards(
        self,
        slide: Slide,
        stats: List[Tuple[str, str, str]],
        position: Optional[Tuple[float, float, float, float]] = None
    ):
        """
        添加统计卡片（大数字展示）

        Args:
            slide: 幻灯片对象
            stats: 统计数据列表 [(标签, 数值, 单位), ...]
            position: 起始位置，默认使用内容区域
        """
        if position:
            left, top, _, _ = position
        else:
            left = CONTENT_LEFT
            top = CONTENT_TOP

        card_width = 2.5
        card_height = 1.5
        spacing = 0.3

        for i, (label, value, unit) in enumerate(stats):
            card_left = left + i * (card_width + spacing)

            # 卡片背景
            card_box = slide.shapes.add_textbox(
                Inches(card_left), Inches(top), Inches(card_width), Inches(card_height)
            )

            # 背景色
            fill = card_box.fill
            fill.solid()
            fill.fore_color.rgb = COLOR_PRIMARY

            # 数值
            value_box = slide.shapes.add_textbox(
                Inches(card_left), Inches(top + 0.3), Inches(card_width), Inches(0.6)
            )
            value_frame = value_box.text_frame
            value_frame.text = f"{value}{unit}"

            value_para = value_frame.paragraphs[0]
            value_para.alignment = PP_ALIGN.CENTER
            value_para.font.name = FONT_NAME_EN
            value_para.font.size = Pt(36)
            value_para.font.bold = True
            value_para.font.color.rgb = RGBColor(255, 255, 255)

            # 标签
            label_box = slide.shapes.add_textbox(
                Inches(card_left), Inches(top + 1.0), Inches(card_width), Inches(0.4)
            )
            label_frame = label_box.text_frame
            label_frame.text = label

            label_para = label_frame.paragraphs[0]
            label_para.alignment = PP_ALIGN.CENTER
            label_para.font.name = FONT_NAME_CN
            label_para.font.size = Pt(14)
            label_para.font.color.rgb = RGBColor(255, 255, 255)

        logger.info(f"添加统计卡片: {len(stats)}个")

    def find_slide_by_text(self, search_text: str, start_index: int = 0) -> Optional[int]:
        """
        通过搜索幻灯片中的文本内容来查找幻灯片索引

        Args:
            search_text: 要搜索的文本（支持部分匹配）
            start_index: 从哪个索引开始搜索（默认从0开始）

        Returns:
            找到的幻灯片索引（0-based），未找到返回None
        """
        for slide_idx in range(start_index, len(self.prs.slides)):
            slide = self.prs.slides[slide_idx]

            # 搜索所有文本框和表格中的文本
            for shape in slide.shapes:
                try:
                    # 检查文本框
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        if search_text in shape.text_frame.text:
                            logger.info(f"在第{slide_idx + 1}页（索引{slide_idx}）找到文本: {search_text}")
                            return slide_idx

                    # 检查表格 (使用 shape_type 避免访问非表格shape的table属性)
                    if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                if search_text in cell.text:
                                    logger.info(f"在第{slide_idx + 1}页（索引{slide_idx}）的表格中找到文本: {search_text}")
                                    return slide_idx
                except:
                    continue

        logger.warning(f"未找到包含文本的幻灯片: {search_text}")
        return None

    def find_slides_by_text(self, search_text: str, start_index: int = 0, max_count: int = None) -> List[int]:
        """
        查找所有包含指定文本的幻灯片索引

        Args:
            search_text: 要搜索的文本
            start_index: 从哪个索引开始搜索
            max_count: 最多返回多少个结果（None表示不限制）

        Returns:
            找到的幻灯片索引列表
        """
        found_indices = []
        logger.debug(f"开始搜索文本: '{search_text}', 从第{start_index+1}页开始")

        for slide_idx in range(start_index, len(self.prs.slides)):
            if max_count and len(found_indices) >= max_count:
                break

            slide = self.prs.slides[slide_idx]
            slide_found = False

            # 搜索所有文本框和表格中的文本
            for shape in slide.shapes:
                try:
                    # 检查文本框
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text_content = shape.text_frame.text
                        if search_text in text_content:
                            if not slide_found:
                                found_indices.append(slide_idx)
                                slide_found = True
                                logger.info(f"✓ 在第{slide_idx + 1}页（索引{slide_idx}）的文本框[{shape.name}]中找到: {search_text}")
                                logger.debug(f"  完整内容: {text_content[:100]}")
                            break

                    # 检查表格 (使用 shape_type 避免访问非表格shape的table属性)
                    if shape.shape_type == 19 and not slide_found:  # MSO_SHAPE_TYPE.TABLE
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                if search_text in cell.text:
                                    if not slide_found:
                                        found_indices.append(slide_idx)
                                        slide_found = True
                                        logger.info(f"✓ 在第{slide_idx + 1}页（索引{slide_idx}）的表格中找到: {search_text}")
                                    break
                            if slide_found:
                                break
                except Exception as e:
                    logger.debug(f"检查形状 {shape.name} 时出错: {e}")
                    continue

                if slide_found:
                    break

        if not found_indices:
            logger.warning(f"❌ 未找到包含文本 '{search_text}' 的幻灯片")
        else:
            logger.info(f"✓ 共找到 {len(found_indices)} 个匹配的幻灯片")

        return found_indices

    # ------------------------------------------------------------------ #
    # P4优化：统一的表格操作方法
    # ------------------------------------------------------------------ #

    def _find_table(self, slide, name: str):
        """
        在幻灯片中按名称查找表格

        Args:
            slide: 幻灯片对象
            name: 表格名称（如"表格 1"）

        Returns:
            表格对象或None
        """
        for shape in slide.shapes:
            if shape.name == name and getattr(shape, "has_table", False):
                return shape.table
        return None

    @staticmethod
    def _set_cell_text(cell, text: str, font_size: int = None, bold: bool = None):
        """
        更新单元格文字，并可选设置字体样式

        Args:
            cell: 单元格对象
            text: 文本内容
            font_size: 字体大小（可选）
            bold: 是否加粗（可选）
        """
        tf = cell.text_frame
        if not tf.paragraphs:
            para = tf.add_paragraph()
        else:
            para = tf.paragraphs[0]
        if para.runs:
            run = para.runs[0]
        else:
            run = para.add_run()
        run.text = text
        if font_size is not None or bold is not None:
            font = run.font
            font.name = FONT_NAME_CN
            if font_size is not None:
                font.size = Pt(font_size)
            if bold is not None:
                font.bold = bold
        # 清理其它 run/段落中的旧文本
        for extra_run in para.runs[1:]:
            extra_run.text = ""
        for extra_para in tf.paragraphs[1:]:
            for extra_run in extra_para.runs:
                extra_run.text = ""

    def _add_table_row(self, table):
        """
        添加表格行（复制最后一行的结构）

        Args:
            table: 表格对象
        """
        try:
            from copy import deepcopy
            tbl = table._tbl
            # 复制最后一行的结构
            last_row_idx = len(table.rows) - 1
            last_tr = table.rows[last_row_idx]._tr
            new_tr = deepcopy(last_tr)
            # 添加到表格
            tbl.append(new_tr)

            # 清空新行的文本
            new_row_idx = len(table.rows) - 1
            for col_idx in range(len(table.columns)):
                try:
                    cell = table.cell(new_row_idx, col_idx)
                    cell.text = ""
                except:
                    pass
        except Exception as e:
            logger.warning(f"添加表格行失败: {e}")

    def _delete_table_row(self, table, row_idx: int):
        """
        删除表格中的指定行

        Args:
            table: 表格对象
            row_idx: 要删除的行索引
        """
        try:
            tbl = table._tbl
            tr = table.rows[row_idx]._tr
            tbl.remove(tr)
        except Exception as e:
            logger.warning(f"删除表格行{row_idx}失败: {e}")
