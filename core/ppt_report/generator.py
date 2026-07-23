"""
主PPT生成器 - 基于Excel综合报告
"""

import logging
import json
import re
import zipfile
from pathlib import Path
from typing import Optional, Tuple, Callable
from datetime import datetime
import time

from pptx import Presentation
from pptx.util import Inches
from openpyxl import load_workbook

from .utils import find_excel_report
from .data_collector import DataCollector
from .chart_creator import ChartCreator
from .config import *

logger = logging.getLogger(__name__)


class ExcelBasedPPTGenerator:
    """
    基于Excel综合报告的PPT生成器

    与现有PPTGenerator接口兼容，但数据源改为Excel综合报告
    """

    def __init__(self, project_path: Path, farm_name: str = "牧场", reporter_name: str = "用户"):
        """
        初始化PPT生成器

        Args:
            project_path: 项目路径
            farm_name: 牧场名称
            reporter_name: 汇报人姓名
        """
        self.project_path = Path(project_path)
        self.farm_name = farm_name
        self.reporter_name = reporter_name
        self.reports_folder = self.project_path / "reports"
        self.output_folder = self.project_path / "analysis_results"
        self.output_folder.mkdir(parents=True, exist_ok=True)

        # 查找Excel报告
        self.excel_report_path = None

        # 初始化组件（延迟初始化）
        self.data_collector = None
        self.chart_creator = None
        self.prs = None
        self.farm_info = {}
        self.last_output_path: Optional[Path] = None

        # 缓存的openpyxl workbook（避免重复加载，每次加载需要约21秒）
        self._cached_workbook = None
        self._cached_workbook_data_only = None

        logger.info(f"初始化PPT生成器: {farm_name}, 汇报人: {reporter_name}")

    def check_required_files(self) -> Tuple[bool, str]:
        """
        检查必需文件是否存在

        Returns:
            (是否满足条件, 错误信息)
        """
        # 检查Excel综合报告
        self.excel_report_path = find_excel_report(self.reports_folder)

        if self.excel_report_path is None:
            return False, "未找到Excel综合报告，请先生成Excel报告（点击'生成Excel报告'按钮）"

        logger.info(f"找到Excel报告: {self.excel_report_path.name}")
        return True, ""

    def generate_ppt(self, progress_callback: Optional[Callable] = None) -> bool:
        """
        生成PPT报告（主入口）

        Args:
            progress_callback: 进度回调函数 callback(message: str, progress: int)

        Returns:
            是否成功生成
        """
        try:
            logger.info("=" * 60)
            logger.info("开始生成PPT报告 v2.0（基于Excel）")
            logger.info("=" * 60)

            # 1-3. 并行初始化：workbook加载 + 数据收集 + PPT创建
            self._report_progress(progress_callback, "正在初始化（并行加载）...", 0)
            self._initialize_and_load_parallel(progress_callback)
            self._report_progress(progress_callback, "✓ 初始化完成", 20)

            # 4. 生成各部分幻灯片 (20-90%)
            data = self._collected_data
            # 按实际耗时分配权重：Part4(遗传评估~12页)和Part6(公牛分析)最耗时
            part_weights = [3, 5, 5, 35, 5, 25, 5]  # 总共83 → 归一化到70%
            total_weight = sum(part_weights)

            # Part 1: 封面与目录
            self._report_progress(progress_callback, "[1/7] 生成封面与目录...", 20)
            self._build_part1_cover_and_toc()
            p1_end = 20 + int(70 * part_weights[0] / total_weight)
            self._report_progress(progress_callback, "✓ Part 1 完成", p1_end)

            # Part 2: 牧场概况
            self._report_progress(progress_callback, "[2/7] 生成牧场概况...", p1_end)
            self._build_part2_farm_overview(data)
            p2_end = 20 + int(70 * sum(part_weights[:2]) / total_weight)
            self._report_progress(progress_callback, "✓ Part 2 完成", p2_end)

            # Part 3: 系谱分析
            self._report_progress(progress_callback, "[3/7] 生成系谱分析...", p2_end)
            self._build_part3_pedigree(data)
            p3_end = 20 + int(70 * sum(part_weights[:3]) / total_weight)
            self._report_progress(progress_callback, "✓ Part 3 完成", p3_end)

            # Part 4: 遗传评估（最耗时：~12页正态分布+表格+图表）
            self._report_progress(progress_callback, "[4/7] 生成遗传评估（共~12页，请耐心等待）...", p3_end)
            self._build_part4_genetics(data)
            p4_end = 20 + int(70 * sum(part_weights[:4]) / total_weight)
            self._report_progress(progress_callback, "✓ Part 4 完成", p4_end)

            # Part 5: 配种记录
            self._report_progress(progress_callback, "[5/7] 生成配种分析...", p4_end)
            self._build_part5_breeding(data)
            p5_end = 20 + int(70 * sum(part_weights[:5]) / total_weight)
            self._report_progress(progress_callback, "✓ Part 5 完成", p5_end)

            # Part 6: 公牛使用（较耗时：折线图数据提取+更新）
            self._report_progress(progress_callback, "[6/7] 生成公牛分析...", p5_end)
            self._build_part6_bulls(data)
            p6_end = 20 + int(70 * sum(part_weights[:6]) / total_weight)
            self._report_progress(progress_callback, "✓ Part 6 完成", p6_end)

            # Part 7: 选配推荐
            self._report_progress(progress_callback, "[7/7] 生成选配推荐...", p6_end)
            self._build_part7_mating(data)
            self._report_progress(progress_callback, "✓ Part 7 完成", 90)

            # 4.5 清理空数据页面
            self._report_progress(progress_callback, "正在清理空数据页面...", 90)
            self._cleanup_empty_slides()
            self._cleanup_hmy_optional_sections()

            # 5. 保存PPT (90-100%)
            self._report_progress(progress_callback, "正在保存PPT文件...", 90)
            output_path = self._save_presentation()
            self._report_progress(progress_callback, "✓ PPT保存完成", 95)

            # 6. 清理临时文件 (95-100%)
            self._report_progress(progress_callback, "正在清理临时文件...", 95)
            self.chart_creator.cleanup()
            # 关闭缓存的workbook
            if self._cached_workbook:
                self._cached_workbook.close()
                self._cached_workbook = None
            if self._cached_workbook_data_only:
                self._cached_workbook_data_only.close()
                self._cached_workbook_data_only = None
            self._report_progress(progress_callback, "✓ 清理完成", 100)

            logger.info("=" * 60)
            logger.info(f"✅ PPT报告生成成功: {output_path}")
            logger.info("=" * 60)

            return True

        except Exception as e:
            logger.error(f"生成PPT失败: {e}", exc_info=True)
            if progress_callback:
                progress_callback(f"❌ 生成失败: {e}", 0)
            return False

    def _initialize_components(self):
        """初始化各组件（不含workbook预加载）"""
        # 查找Excel报告（如果还没找到）
        if self.excel_report_path is None:
            self.excel_report_path = find_excel_report(self.reports_folder)
            if self.excel_report_path is None:
                raise FileNotFoundError("未找到Excel综合报告，请先生成Excel报告")
            logger.info(f"找到Excel报告: {self.excel_report_path.name}")

        # 数据收集器
        self.data_collector = DataCollector(self.excel_report_path)

        # 图表生成器
        self.chart_creator = ChartCreator()

        # 跟踪所有 builder 实例（用于收集待删除的空数据页面）
        self._builders = []

        # data_only=True/False 的workbook缓存（由并行加载填充）
        self._cached_workbook_data_only = None
        self._cached_workbook = None

        logger.info("组件初始化完成")

    def _initialize_and_load_parallel(self, progress_callback=None):
        """
        并行初始化：两个workbook加载 与 数据收集+PPT创建 同时进行

        三个任务并行：
        - 线程1: load_workbook(data_only=True)  ~15s（图表/表格数据读取用）
        - 线程2: load_workbook(data_only=False) ~20s（时间线图片提取用）
        - 主线程: 数据收集 ~8s + PPT创建 ~1s

        总耗时 ≈ max(20s, 9s) ≈ 20s，避免Part 6延迟加载20秒
        """
        from concurrent.futures import ThreadPoolExecutor

        # 基础初始化（快速，不含workbook加载）
        self._initialize_components()

        excel_path_str = str(self.excel_report_path)
        wb_data_only_result = [None]
        wb_full_result = [None]

        def _load_wb_data_only():
            """后台线程1：加载data_only workbook（图表数据读取）"""
            logger.info("预加载Excel workbook (data_only=True)...")
            t0 = time.perf_counter()
            wb = load_workbook(excel_path_str, data_only=True)
            logger.info(f"✓ data_only workbook加载完成，耗时: {time.perf_counter()-t0:.2f}秒")
            wb_data_only_result[0] = wb

        def _load_wb_full():
            """后台线程2：加载完整workbook（时间线图片提取）"""
            logger.info("预加载Excel workbook (data_only=False)...")
            t0 = time.perf_counter()
            wb = load_workbook(excel_path_str, data_only=False)
            logger.info(f"✓ 完整workbook加载完成，耗时: {time.perf_counter()-t0:.2f}秒")
            wb_full_result[0] = wb

        # 三路并行：2个workbook加载 + 主线程数据收集
        with ThreadPoolExecutor(max_workers=2) as executor:
            wb_future1 = executor.submit(_load_wb_data_only)
            wb_future2 = executor.submit(_load_wb_full)

            # 主线程同时执行：数据收集 + PPT创建
            self._report_progress(progress_callback, "正在读取Excel报告数据...", 3)
            data = self.data_collector.collect_all_data()
            self.farm_info = data.get('farm_info_dict', {}) or {}
            if self.farm_name and self.farm_name != "牧场":
                self.farm_info["farm_name"] = self.farm_name
                data["farm_info_dict"] = self.farm_info
            self._report_progress(progress_callback, "✓ 数据读取完成", 12)

            self._report_progress(progress_callback, "正在创建PPT...", 12)
            self._create_presentation()
            self._report_progress(progress_callback, "✓ PPT创建完成", 15)

            # 等待两个workbook加载完成
            self._report_progress(progress_callback, "等待workbook加载完成...", 15)
            wb_future1.result()
            wb_future2.result()

        self._cached_workbook_data_only = wb_data_only_result[0]
        self._cached_workbook = wb_full_result[0]
        self._report_progress(progress_callback, "✓ 全部初始化完成", 18)

        # 确保data被存储以供后续使用
        self._collected_data = data

    def _create_presentation(self):
        """创建PPT演示文稿"""
        # 查找模板（从程序根目录）
        program_root = Path(__file__).parent.parent.parent
        preferred_template = program_root / "牧场牧场育种分析报告-模版.pptx"
        fallback_template = program_root / "PPT模版.pptx"

        template_path = preferred_template if preferred_template.exists() else fallback_template

        if template_path.exists():
            self.prs = Presentation(str(template_path))
            logger.info(f"使用PPT模板创建: {template_path}")
        else:
            self.prs = Presentation()
            self.prs.slide_width = Inches(13.33)
            self.prs.slide_height = Inches(7.5)
            logger.warning("未找到模板，使用空白PPT")

    def _save_presentation(self) -> Path:
        """保存PPT文件"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        # 调用方传入的牧场名称比项目目录名更可靠。自动报告的临时项目
        # 目录可能是随机名称，不应覆盖真实牧场名。
        farm_name = (self.farm_name or self.farm_info.get("farm_name") or "牧场").strip()
        report_farm_name = farm_name if farm_name.endswith("牧场") else f"{farm_name}牧场"
        filename = f"{report_farm_name}育种分析报告_{timestamp}.pptx"

        # 报告输出到 reports 目录；analysis_results 仅用于中间结果
        self.reports_folder.mkdir(parents=True, exist_ok=True)
        output_path = self.reports_folder / filename

        self.prs.save(str(output_path))
        self._normalize_chart_axis_ids(output_path)
        self.last_output_path = output_path
        logger.info(f"PPT已保存: {output_path}")

        return output_path

    @staticmethod
    def _normalize_chart_axis_ids(pptx_path: Path) -> None:
        """
        将模板图表中的负轴 ID 转为合法的 UInt32。

        WPS 生成的部分模板会把轴 ID 保存成带符号的 32 位整数。PowerPoint
        有时会自动修复，但严格的 OpenXML 读取器会直接拒绝该文件。轴之间
        通过同一 ID 关联，因此按二进制等价的无符号值统一转换即可。
        """
        axis_id_pattern = re.compile(
            rb'(<c:(?:axId|crossAx)\s+val=")(-\d+)(")'
        )
        temp_path = pptx_path.with_suffix(".normalized.pptx")
        replacement_count = 0

        try:
            with zipfile.ZipFile(pptx_path, "r") as source_zip:
                with zipfile.ZipFile(temp_path, "w") as target_zip:
                    for member in source_zip.infolist():
                        content = source_zip.read(member.filename)
                        if member.filename.startswith("ppt/charts/") and member.filename.endswith(".xml"):
                            def replace_axis_id(match):
                                nonlocal replacement_count
                                replacement_count += 1
                                unsigned_value = int(match.group(2)) & 0xFFFFFFFF
                                return match.group(1) + str(unsigned_value).encode("ascii") + match.group(3)

                            content = axis_id_pattern.sub(replace_axis_id, content)
                        target_zip.writestr(member, content)

            temp_path.replace(pptx_path)
        except Exception:
            temp_path.unlink(missing_ok=True)
            raise

        if replacement_count:
            logger.info(f"已规范化 {replacement_count} 个图表轴ID")

    def _report_progress(self, callback: Optional[Callable], message: str, progress: int):
        """报告进度"""
        if callback:
            callback(message, progress)
        logger.info(f"[{progress}%] {message}")

    def _cleanup_empty_slides(self):
        """
        统一删除所有标记为空数据的页面

        收集所有 builder 标记的待删除页面，从后往前删除避免索引偏移
        """
        # 收集所有待删除的页面索引
        all_indices = set()
        for builder in self._builders:
            if hasattr(builder, 'get_slides_to_delete'):
                indices = builder.get_slides_to_delete()
                if indices:
                    all_indices.update(indices)
                    logger.debug(f"{builder.__class__.__name__} 标记了 {len(indices)} 个页面待删除")

        if not all_indices:
            logger.info("没有空数据页面需要删除")
            return

        logger.info(f"开始清理空数据页面，共 {len(all_indices)} 个")

        # 从后往前删除，避免索引偏移问题
        for idx in sorted(all_indices, reverse=True):
            try:
                if idx < 0 or idx >= len(self.prs.slides):
                    logger.warning(f"页面索引 {idx} 超出范围，跳过")
                    continue
                rId = self.prs.slides._sldIdLst[idx].rId
                self.prs.part.drop_rel(rId)
                del self.prs.slides._sldIdLst[idx]
                logger.info(f"✓ 已删除空数据页面：第 {idx + 1} 页（索引{idx}）")
            except Exception as e:
                logger.warning(f"删除页面 {idx} 失败: {e}")

        logger.info(f"✓ 空数据页面清理完成")

    @staticmethod
    def _set_shape_text(shape, value: str):
        """替换文本并尽量保留模板中首个文本 run 的格式。"""
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None or not text_frame.paragraphs:
            return

        first_paragraph = text_frame.paragraphs[0]
        if first_paragraph.runs:
            first_paragraph.runs[0].text = value
            for run in first_paragraph.runs[1:]:
                run.text = ""
        else:
            first_paragraph.text = value
        for paragraph in text_frame.paragraphs[1:]:
            for run in paragraph.runs:
                run.text = ""

    def _cleanup_hmy_optional_sections(self):
        """慧牧云报告仅保留实际有数据的章节，并同步目录编号。"""
        metadata_path = self.project_path / "project_metadata.json"
        try:
            metadata = json.loads(metadata_path.read_text(encoding="utf-8"))
        except (OSError, ValueError):
            return
        if metadata.get("data_source") != "慧牧云":
            return

        workbook = self._cached_workbook_data_only or self._cached_workbook
        sheet_names = set(workbook.sheetnames) if workbook is not None else set()
        section_titles = [
            "牧场概况",
            "系谱记录分析",
            "牛群遗传评估",
            "配种记录分析",
            "公牛使用分析",
            "选配推荐方案",
            "项目总结建议",
        ]
        availability = {
            "牧场概况": "牧场基础信息" in sheet_names,
            "系谱记录分析": "系谱识别分析" in sheet_names,
            "牛群遗传评估": bool(
                sheet_names
                & {
                    "年份汇总与性状进展",
                    "育种性状明细",
                    "育种指数分布分析",
                    "母牛指数排名明细",
                }
            ),
            "配种记录分析": bool(
                sheet_names
                & {
                    "配种记录-隐性基因分析",
                    "配种记录-近交系数分析",
                    "配种记录-隐性基因及近交系数明细",
                }
            ),
            "公牛使用分析": bool(
                sheet_names
                & {
                    "已用公牛性状汇总",
                    "已用公牛性状明细",
                    "备选公牛排名",
                    "备选公牛-隐性基因分析",
                    "备选公牛-近交系数分析",
                    "备选公牛-明细表",
                }
            ),
            "选配推荐方案": "个体选配推荐结果" in sheet_names,
            # 当前总结页仍是模板占位内容，避免输出无数据的静态章节。
            "项目总结建议": False,
        }

        def slide_texts(slide):
            return {
                shape.text.strip()
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            }

        cover_positions = {}
        for index, slide in enumerate(self.prs.slides):
            if index < 2:
                continue
            texts = slide_texts(slide)
            for title in section_titles:
                if title in texts:
                    cover_positions[title] = index

        indices_to_remove = set()
        ordered_positions = sorted(cover_positions.values())
        for title, start in cover_positions.items():
            if availability.get(title, False):
                continue
            end = next(
                (position for position in ordered_positions if position > start),
                len(self.prs.slides),
            )
            indices_to_remove.update(range(start, end))

        for index in sorted(indices_to_remove, reverse=True):
            r_id = self.prs.slides._sldIdLst[index].rId
            self.prs.part.drop_rel(r_id)
            del self.prs.slides._sldIdLst[index]

        surviving_titles = []
        for title in section_titles:
            if (
                availability.get(title, False)
                and title in cover_positions
                and any(
                    title in slide_texts(slide)
                    for slide in list(self.prs.slides)[2:]
                )
            ):
                surviving_titles.append(title)
        number_by_title = {
            title: f"{index:02d}"
            for index, title in enumerate(surviving_titles, start=1)
        }

        # 更新目录：缺失章节清空，保留章节重新连续编号。
        if len(self.prs.slides) > 1:
            toc_slide = self.prs.slides[1]
            text_shapes = [
                shape
                for shape in toc_slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            ]
            for title in section_titles:
                title_index = next(
                    (
                        index
                        for index, shape in enumerate(text_shapes)
                        if shape.text.strip() == title
                    ),
                    None,
                )
                if title_index is None:
                    continue
                number_shape = text_shapes[title_index - 1] if title_index > 0 else None
                description_shape = (
                    text_shapes[title_index + 1]
                    if title_index + 1 < len(text_shapes)
                    else None
                )
                if title in number_by_title:
                    if number_shape is not None:
                        self._set_shape_text(number_shape, number_by_title[title])
                else:
                    if number_shape is not None:
                        self._set_shape_text(number_shape, "")
                    self._set_shape_text(text_shapes[title_index], "")
                    if description_shape is not None:
                        self._set_shape_text(description_shape, "")

        # 更新各章节封面的序号。
        for slide in self.prs.slides:
            texts = slide_texts(slide)
            title = next((item for item in surviving_titles if item in texts), None)
            if not title:
                continue
            for shape in slide.shapes:
                text = getattr(shape, "text", "").strip()
                if len(text) == 2 and text.isdigit():
                    self._set_shape_text(shape, number_by_title[title])
                    break

        logger.info(
            "慧牧云报告章节已按数据可用性清理: %s",
            ", ".join(surviving_titles),
        )

    # ==================== 各部分构建方法（占位） ====================

    def _build_part1_cover_and_toc(self):
        """构建Part 1: 封面与目录"""
        from .slide_builders import Part1CoverBuilder

        # 使用初始化时传入的汇报人姓名
        builder = Part1CoverBuilder(
            self.prs,
            self.chart_creator,
            farm_info=self.farm_info,
            fallback_farm_name=self.farm_name,
            fallback_reporter=self.reporter_name
        )
        builder.build()

    def _build_part2_farm_overview(self, data: dict):
        """构建Part 2: 牧场概况"""
        from .slide_builders import Part2FarmOverviewBuilder

        logger.info("构建Part 2: 牧场概况")
        builder = Part2FarmOverviewBuilder(self.prs, self.chart_creator, self.farm_name)
        builder.build(data)
        self._builders.append(builder)

    def _build_part3_pedigree(self, data: dict):
        """构建Part 3: 系谱分析"""
        from .slide_builders.part3_pedigree import Part3PedigreeBuilder

        logger.info("构建Part 3: 系谱分析")
        builder = Part3PedigreeBuilder(self.prs, self.chart_creator, self.farm_name)
        builder.build(data)
        self._builders.append(builder)

    def _build_part4_genetics(self, data: dict):
        """构建Part 4: 遗传评估"""
        from .slide_builders.part4_genetics import Part4GeneticsBuilder

        logger.info("构建Part 4: 遗传评估")
        # 传递缓存的workbook给需要的Builder
        data['_cached_workbook'] = self._cached_workbook
        data['_cached_workbook_data_only'] = self._cached_workbook_data_only

        builder = Part4GeneticsBuilder(self.prs, self.chart_creator, self.farm_name)
        builder.build(data)
        self._builders.append(builder)

    def _build_part5_breeding(self, data: dict):
        """构建Part 5: 配种记录分析"""
        from .slide_builders.part5_breeding_records import Part5BreedingRecordsBuilder
        from .slide_builders.part5_inbreeding import Part5InbreedingBuilder

        logger.info("构建Part 5: 配种记录分析")

        # 隐性基因分析（2页）
        genes_builder = Part5BreedingRecordsBuilder(self.prs, self.chart_creator, self.farm_name)
        genes_builder.build(data)
        self._builders.append(genes_builder)

        # 近交分析（2页）
        inbreeding_builder = Part5InbreedingBuilder(self.prs, self.chart_creator, self.farm_name)
        inbreeding_builder.build(data)
        self._builders.append(inbreeding_builder)

    def _build_part6_bulls(self, data: dict):
        """构建Part 6: 公牛使用"""
        from .slide_builders.part6_bulls_usage import Part6BullsUsageBuilder
        from .slide_builders.part6_bulls_detail import Part6BullsDetailBuilder
        from .slide_builders.part6_traits_trends import Part6TraitsTrendsBuilder
        from .slide_builders.part6_timeline import Part6TimelineBuilder

        logger.info("构建Part 6: 公牛使用")

        # 传递缓存的workbook给需要的Builder
        data['_cached_workbook_data_only'] = self._cached_workbook_data_only

        # 已用公牛性状汇总表（1页）
        bulls_builder = Part6BullsUsageBuilder(self.prs, self.chart_creator, self.farm_name)
        bulls_builder.build(data)
        self._builders.append(bulls_builder)

        # 已用公牛明细（多页，每年一页）
        bulls_detail_builder = Part6BullsDetailBuilder(self.prs, self.chart_creator, self.farm_name)
        bulls_detail_builder.build(data)
        self._builders.append(bulls_detail_builder)

        # 性状进展折线图（多页，每页2个图表）
        traits_trends_builder = Part6TraitsTrendsBuilder(self.prs, self.chart_creator, self.farm_name)
        traits_trends_builder.build(data)
        self._builders.append(traits_trends_builder)

        # 配种记录时间线（2页：全部记录、近一年）
        # 优化：优先使用预导出的图片，避免加载data_only=False workbook（节省约20秒）
        timeline_builder = Part6TimelineBuilder(self.prs, self.chart_creator, self.farm_name)

        # 先检查是否有预导出的时间线图片
        pregenerated_images = timeline_builder._find_pregenerated_images(str(self.excel_report_path))

        if not pregenerated_images or len(pregenerated_images) < 2:
            # 没有预导出图片，需要从Excel提取，必须加载data_only=False workbook
            if self._cached_workbook is None:
                logger.info("延迟加载Excel workbook (openpyxl data_only=False) 用于提取图片...")
                t0 = time.perf_counter()
                self._cached_workbook = load_workbook(
                    str(self.excel_report_path), data_only=False
                )
                t1 = time.perf_counter()
                logger.info(f"✓ Excel workbook加载完成，耗时: {t1-t0:.2f}秒")
            data['_cached_workbook'] = self._cached_workbook
        else:
            logger.info("✓ 检测到预导出的时间线图片，跳过加载data_only=False workbook（节省约20秒）")

        timeline_builder.build(data)
        self._builders.append(timeline_builder)

    def _build_part7_mating(self, data: dict):
        """构建Part 7: 选配推荐"""
        from .slide_builders.part7_candidate_bulls_ranking import Part7CandidateBullsRankingBuilder
        from ..excel_report.data_collectors import collect_bull_ranking_data
        from .slide_builders.template_slide_copier import copy_template_slide

        logger.info("构建Part 7: 选配推荐")

        # 注意：不复制模板页面，直接使用模板中已有的第172-174页
        # 这些页面在模板加载时就已经存在于PPT中

        # 添加Excel路径到数据字典
        if self.excel_report_path:
            data['excel_path'] = str(self.excel_report_path)
            logger.info(f"✓ Excel报告路径: {self.excel_report_path.name}")
        else:
            logger.warning("⚠️  Excel报告路径未找到")

        # 传递缓存的workbook给需要的Builder（避免重复加载）
        data['_cached_workbook_data_only'] = self._cached_workbook_data_only

        # 收集备选公牛排名数据（从analysis_results）
        try:
            bull_ranking_data = collect_bull_ranking_data(self.output_folder)
            if bull_ranking_data:
                data['bull_ranking'] = bull_ranking_data
                logger.info(f"✓ 收集备选公牛排名数据: {len(bull_ranking_data.get('ranking_df', []))}头公牛")
            else:
                logger.warning("⚠️  备选公牛排名数据收集失败")
        except Exception as e:
            logger.warning(f"收集备选公牛排名数据时出错: {e}")

        # 备选公牛排名分析（1页）
        ranking_builder = Part7CandidateBullsRankingBuilder(self.prs, self.farm_name)
        ranking_builder.build(data)
        self._builders.append(ranking_builder)

        # 备选公牛-隐性基因分析（动态页数）
        from .slide_builders.part7_candidate_bulls_genes import Part7CandidateBullsGenesBuilder
        genes_builder = Part7CandidateBullsGenesBuilder(self.prs, self.farm_name)
        genes_builder.build(data)
        self._builders.append(genes_builder)

        # 备选公牛-近交系数分析（动态页数）
        from .slide_builders.part7_candidate_bulls_inbreeding import Part7CandidateBullsInbreedingBuilder
        inbreeding_builder = Part7CandidateBullsInbreedingBuilder(self.prs, self.farm_name)
        inbreeding_builder.build(data)
        self._builders.append(inbreeding_builder)

        # 选配推荐方案章节页（页172）
        from .slide_builders.part7_mating_section import Part7MatingSectionBuilder
        section_builder = Part7MatingSectionBuilder(self.prs, self.chart_creator, self.farm_name)
        section_builder.build(data)
        self._builders.append(section_builder)

        # 个体选配推荐结果/选配统计摘要（页173）
        from .slide_builders.part7_mating_recommendation import Part7MatingRecommendationBuilder
        mating_rec_builder = Part7MatingRecommendationBuilder(self.prs, self.chart_creator, self.farm_name)
        mating_rec_builder.build(data)
        self._builders.append(mating_rec_builder)

        # 项目总结建议（页174）
        from .slide_builders.part7_summary_suggestions import Part7SummaryBuilder
        summary_builder = Part7SummaryBuilder(self.prs, self.chart_creator, self.farm_name)
        summary_builder.build(data)
        self._builders.append(summary_builder)
