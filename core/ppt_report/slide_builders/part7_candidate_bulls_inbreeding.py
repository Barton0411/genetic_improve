"""
Part 7: 备选公牛-近交系数分析 幻灯片构建器
"""

import logging
from pathlib import Path
from typing import List, Dict, Any, Optional
from pptx.chart.data import CategoryChartData
from pptx.util import Pt
from pptx.dml.color import RGBColor
from openpyxl import load_workbook

from ..base_builder import BaseSlideBuilder
from ..config import FONT_NAME_CN

logger = logging.getLogger(__name__)


class Part7CandidateBullsInbreedingBuilder(BaseSlideBuilder):
    """
    Part 7: 备选公牛-近交系数分析构建器

    功能：
    - 动态管理26页模板页面
    - 填充8×8表格（4个近交系数区间 + 高风险小计）
    - 更新风险占比柱状图
    """

    # 近交系数区间定义
    INBREEDING_INTERVALS = [
        "<3.125%",
        "3.125%-6.25%",
        "6.25%-12.5%",
        ">12.5%"
    ]

    def build(self, data: Dict[str, Any]):
        """
        构建备选公牛近交系数分析页面

        Args:
            data: 包含excel_path的数据字典
        """
        try:
            logger.info("=" * 60)
            logger.info("开始构建 Part 7 - 备选公牛近交系数分析")
            logger.info("=" * 60)

            # 1. 读取Excel数据
            excel_path = data.get('excel_path')
            if not excel_path or not Path(excel_path).exists():
                logger.error("Excel报告路径未提供或文件不存在")
                return

            logger.info(f"从Excel读取数据: {Path(excel_path).name}")
            bulls_data = self._read_bulls_inbreeding_data(excel_path)

            if not bulls_data:
                logger.warning("未读取到备选公牛近交系数分析数据")
                return

            bull_count = len(bulls_data)
            logger.info(f"读取到 {bull_count} 头备选公牛的近交系数数据")

            # 2. 定位模板页面（查找"备选公牛-近交系数分析"标题）
            pages = self.find_slides_by_text("备选公牛-近交系数分析")

            if not pages:
                logger.warning("未找到备选公牛-近交系数分析模板页面")
                return

            start_slide_index = pages[0]
            logger.info(f"找到起始页面索引: {start_slide_index} (第{start_slide_index + 1}页)")

            # 3. 动态管理页面数量（26页模板）
            self._manage_pages(start_slide_index, bull_count)

            # 4. 填充每头公牛的近交系数分析页面
            for i, bull_data in enumerate(bulls_data):
                slide_index = start_slide_index + i
                slide = self.prs.slides[slide_index]

                logger.info(f"填充第 {i+1}/{bull_count} 头公牛 (第{slide_index+1}页)")
                self._fill_bull_inbreeding_page(slide, bull_data, i + 1)

            logger.info("=" * 60)
            logger.info(f"✅ Part 7 - 备选公牛近交系数分析构建完成 ({bull_count}页)")
            logger.info("=" * 60)

        except Exception as e:
            logger.error(f"构建备选公牛近交系数分析失败: {e}", exc_info=True)

    def _read_bulls_inbreeding_data(self, excel_path: str) -> List[Dict[str, Any]]:
        """
        从Excel读取备选公牛近交系数分析数据

        数据结构：每头公牛占12行
        - 行1: 公牛标识
        - 行2: 空行
        - 行3: 表头
        - 行4-7: 4个近交系数区间数据
        - 行8: 空行
        - 行9: 表头
        - 行10: 高风险小计（>6.25%）
        - 行11-12: 空行/备注

        Args:
            excel_path: Excel文件路径

        Returns:
            公牛近交系数数据列表
        """
        try:
            wb = load_workbook(excel_path, data_only=True)

            # 查找"备选公牛-近交系数分析" sheet
            sheet_name = None
            for name in wb.sheetnames:
                if "备选公牛-近交系数分析" in name:
                    sheet_name = name
                    break

            if not sheet_name:
                logger.warning("未找到'备选公牛-近交系数分析'工作表")
                return []

            ws = wb[sheet_name]
            logger.info(f"读取工作表: {sheet_name}")

            # 解析数据（每头公牛12行）
            bulls_data = []
            data_start_row = 2  # 假设从第2行开始
            row_idx = data_start_row

            while row_idx <= ws.max_row:
                # 读取公牛标识（第1行）
                bull_id_cell = ws.cell(row=row_idx, column=1)
                bull_id = str(bull_id_cell.value or "").strip()

                # 如果是空行或不是公牛标识，跳过
                if not bull_id or "【公牛" not in bull_id:
                    row_idx += 1
                    continue

                logger.info(f"解析公牛: {bull_id[:50]}")

                # 读取4个近交系数区间数据（行4-7）
                intervals_data = []
                for i in range(4):
                    interval_row = row_idx + 3 + i  # 行4、5、6、7

                    interval_data = {
                        'interval': self.INBREEDING_INTERVALS[i],
                        'adult_count': ws.cell(row=interval_row, column=2).value or 0,
                        'adult_ratio': ws.cell(row=interval_row, column=3).value or "0.0%",
                        'heifer_count': ws.cell(row=interval_row, column=4).value or 0,
                        'heifer_ratio': ws.cell(row=interval_row, column=5).value or "0.0%",
                        'total_count': ws.cell(row=interval_row, column=6).value or 0,
                        'total_ratio': ws.cell(row=interval_row, column=7).value or "0.0%"
                    }
                    intervals_data.append(interval_data)

                # 读取高风险小计（>6.25%）（行9）
                high_risk_row = row_idx + 8
                high_risk_summary = {
                    'label': '小计（高风险：>6.25%）',
                    'adult_count': ws.cell(row=high_risk_row, column=2).value or 0,
                    'adult_ratio': ws.cell(row=high_risk_row, column=3).value or "0.0%",
                    'heifer_count': ws.cell(row=high_risk_row, column=4).value or 0,
                    'heifer_ratio': ws.cell(row=high_risk_row, column=5).value or "0.0%",
                    'total_count': ws.cell(row=high_risk_row, column=6).value or 0,
                    'total_ratio': ws.cell(row=high_risk_row, column=7).value or "0.0%"
                }

                bull_data = {
                    'bull_id': bull_id,
                    'intervals': intervals_data,
                    'high_risk_summary': high_risk_summary
                }

                bulls_data.append(bull_data)
                logger.info(f"  ✓ 解析完成: {len(intervals_data)}个区间 + 高风险小计")

                # 跳到下一头公牛（12行为一个数据块）
                row_idx += 12

            logger.info(f"共解析 {len(bulls_data)} 头公牛的近交系数数据")
            return bulls_data

        except Exception as e:
            logger.error(f"读取备选公牛近交系数数据失败: {e}", exc_info=True)
            return []

    def _manage_pages(self, start_index: int, bull_count: int):
        """
        动态管理页面数量

        模板提供26页，根据实际公牛数量：
        - 如果公牛数 < 26：删除多余页面
        - 如果公牛数 > 26：复制模板页面

        Args:
            start_index: 起始页面索引
            bull_count: 公牛数量
        """
        template_pages = 23  # 模板页面数（第149-171页，共23页）

        if bull_count == template_pages:
            logger.info(f"公牛数量({bull_count})与模板页面数({template_pages})一致，无需调整")
            return

        elif bull_count < template_pages:
            # 删除多余页面（从后往前删）
            pages_to_delete = template_pages - bull_count
            logger.info(f"公牛数量({bull_count}) < 模板页面数({template_pages})，需删除 {pages_to_delete} 页")

            for i in range(pages_to_delete):
                # 始终删除最后一页模板页
                delete_index = start_index + template_pages - 1 - i
                if delete_index < len(self.prs.slides):
                    self._delete_slide(delete_index)
                    logger.info(f"  删除第 {delete_index + 1} 页")

            logger.info(f"✓ 删除完成，保留 {bull_count} 页")

        else:
            # 复制模板页面
            pages_to_add = bull_count - template_pages
            logger.info(f"公牛数量({bull_count}) > 模板页面数({template_pages})，需复制 {pages_to_add} 页")

            # 使用第一页作为模板进行复制
            source_index = start_index

            for i in range(pages_to_add):
                # 插入位置：现有页面的末尾
                insert_at = start_index + template_pages + i
                self._copy_slide(source_index, insert_at)
                logger.info(f"  复制模板页到第 {insert_at + 1} 页")

            logger.info(f"✓ 复制完成，共 {bull_count} 页")

    def _fill_bull_inbreeding_page(self, slide, bull_data: Dict[str, Any], bull_number: int):
        """
        填充单个公牛的近交系数分析页面

        Args:
            slide: 幻灯片对象
            bull_data: 公牛数据
            bull_number: 公牛编号（用于日志）
        """
        try:
            bull_id = bull_data.get('bull_id', '')
            logger.info(f"  填充公牛 #{bull_number}: {bull_id[:50]}")

            # 1. 更新公牛标识文本
            self._update_bull_id(slide, bull_id)

            # 2. 填充近交系数表格（8×8）
            table = self._find_inbreeding_table(slide)
            if table:
                self._fill_inbreeding_table(table, bull_data)
                logger.info(f"    ✓ 表格填充完成")
            else:
                logger.warning(f"    未找到8×8表格")

            # 3. 更新风险占比柱状图
            self._update_risk_chart(slide, bull_data)

            # 4. 填充分析文本
            self._fill_analysis_text(slide, bull_data)

        except Exception as e:
            logger.error(f"填充公牛近交系数页面失败: {e}", exc_info=True)

    def _update_bull_id(self, slide, bull_id: str):
        """更新公牛标识文本"""
        try:
            # 查找包含"【公牛"的文本框或表格单元格
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if "【公牛" in shape.text or "公牛标识" in shape.text:
                        shape.text = bull_id
                        logger.info(f"    ✓ 更新公牛标识: {bull_id[:30]}")
                        return
                elif shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if "【公牛" in cell.text or "公牛标识" in cell.text:
                                cell.text = bull_id
                                # 设置字体格式：14号，白色
                                from pptx.enum.text import PP_ALIGN
                                self._format_cell(cell, font_size=14, alignment=PP_ALIGN.LEFT,
                                                font_color=(255, 255, 255))
                                logger.info(f"    ✓ 更新公牛标识(表格): {bull_id[:30]}")
                                return
        except Exception as e:
            logger.warning(f"更新公牛标识失败: {e}")

    def _find_inbreeding_table(self, slide):
        """
        查找近交系数表格（8行×8列）

        预期结构：
        - 行1-2: 合并标题行
        - 行3: 表头
        - 行4-7: 4个近交系数区间数据
        - 行8: 高风险小计

        Returns:
            表格对象或None
        """
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # 查找8行×8列的表格
                if len(table.rows) == 8 and len(table.columns) == 8:
                    return table
        return None

    def _fill_inbreeding_table(self, table, bull_data: Dict[str, Any]):
        """
        填充8×8近交系数表格

        表格结构：
        - 第3行: 表头（近交系数区间, 成母牛数量, 成母牛占比, 后备牛数量, 后备牛占比, 总数量, 总占比）
        - 第4-7行: 4个区间数据
        - 第8行: 高风险小计

        Args:
            table: 表格对象
            bull_data: 公牛数据
        """
        try:
            intervals_data = bull_data.get('intervals', [])
            high_risk_summary = bull_data.get('high_risk_summary', {})

            # 填充4个区间数据（第4-7行）
            from pptx.enum.text import PP_ALIGN

            for i, interval_data in enumerate(intervals_data[:4]):
                row_idx = 3 + i  # 行4、5、6、7（索引3、4、5、6）

                # 列1: 近交系数区间
                table.cell(row_idx, 1).text = str(interval_data.get('interval', ''))
                self._format_cell(table.cell(row_idx, 1), font_size=12, alignment=PP_ALIGN.CENTER)

                # 列2-3: 成母牛
                table.cell(row_idx, 2).text = str(interval_data.get('adult_count', 0))
                self._format_cell(table.cell(row_idx, 2), font_size=12, alignment=PP_ALIGN.CENTER)
                table.cell(row_idx, 3).text = str(interval_data.get('adult_ratio', '0.0%'))
                self._format_cell(table.cell(row_idx, 3), font_size=12, alignment=PP_ALIGN.CENTER)

                # 列4-5: 后备牛
                table.cell(row_idx, 4).text = str(interval_data.get('heifer_count', 0))
                self._format_cell(table.cell(row_idx, 4), font_size=12, alignment=PP_ALIGN.CENTER)
                table.cell(row_idx, 5).text = str(interval_data.get('heifer_ratio', '0.0%'))
                self._format_cell(table.cell(row_idx, 5), font_size=12, alignment=PP_ALIGN.CENTER)

                # 列6-7: 全群
                table.cell(row_idx, 6).text = str(interval_data.get('total_count', 0))
                self._format_cell(table.cell(row_idx, 6), font_size=12, alignment=PP_ALIGN.CENTER)
                table.cell(row_idx, 7).text = str(interval_data.get('total_ratio', '0.0%'))
                self._format_cell(table.cell(row_idx, 7), font_size=12, alignment=PP_ALIGN.CENTER)

            # 填充高风险小计（第8行，索引7）
            table.cell(7, 1).text = high_risk_summary.get('label', '小计（高风险：>6.25%）')
            self._format_cell(table.cell(7, 1), font_size=12, alignment=PP_ALIGN.CENTER, bold=True)

            table.cell(7, 2).text = str(high_risk_summary.get('adult_count', 0))
            self._format_cell(table.cell(7, 2), font_size=12, alignment=PP_ALIGN.CENTER, bold=True)

            table.cell(7, 3).text = str(high_risk_summary.get('adult_ratio', '0.0%'))
            self._format_cell(table.cell(7, 3), font_size=12, alignment=PP_ALIGN.CENTER, bold=True)

            table.cell(7, 4).text = str(high_risk_summary.get('heifer_count', 0))
            self._format_cell(table.cell(7, 4), font_size=12, alignment=PP_ALIGN.CENTER, bold=True)

            table.cell(7, 5).text = str(high_risk_summary.get('heifer_ratio', '0.0%'))
            self._format_cell(table.cell(7, 5), font_size=12, alignment=PP_ALIGN.CENTER, bold=True)

            table.cell(7, 6).text = str(high_risk_summary.get('total_count', 0))
            self._format_cell(table.cell(7, 6), font_size=12, alignment=PP_ALIGN.CENTER, bold=True)

            table.cell(7, 7).text = str(high_risk_summary.get('total_ratio', '0.0%'))
            self._format_cell(table.cell(7, 7), font_size=12, alignment=PP_ALIGN.CENTER, bold=True)

            logger.info(f"      表格数据填充: 4个区间 + 高风险小计")

        except Exception as e:
            logger.error(f"填充近交系数表格失败: {e}", exc_info=True)

    def _update_risk_chart(self, slide, bull_data: Dict[str, Any]):
        """
        更新高风险占比柱状图

        图表标题: 高风险近交（>6.25%）影响占比
        数据: 成母牛、后备牛、全群的高风险占比

        Args:
            slide: 幻灯片对象
            bull_data: 公牛数据
        """
        try:
            # 查找图表（通过标题识别）
            chart = None
            for shape in slide.shapes:
                if shape.has_chart:
                    c = shape.chart
                    if c.has_title and "高风险" in c.chart_title.text_frame.text and "占比" in c.chart_title.text_frame.text:
                        chart = c
                        break

            if not chart:
                logger.warning(f"    未找到高风险占比图表")
                return

            # 提取高风险小计数据
            high_risk_summary = bull_data.get('high_risk_summary', {})

            # 转换百分比为小数
            mature_ratio = self._parse_percentage(high_risk_summary.get('adult_ratio', '0.0%'))
            heifer_ratio = self._parse_percentage(high_risk_summary.get('heifer_ratio', '0.0%'))
            total_ratio = self._parse_percentage(high_risk_summary.get('total_ratio', '0.0%'))

            logger.info(f"    图表数据: 成母牛={mature_ratio:.4f}, 后备牛={heifer_ratio:.4f}, 全群={total_ratio:.4f}")

            # 创建图表数据
            chart_data = CategoryChartData()
            chart_data.categories = ['成母牛', '后备牛', '全群']
            chart_data.add_series('高风险占比', [mature_ratio, heifer_ratio, total_ratio])

            # 更新图表
            chart.replace_data(chart_data)
            logger.info(f"    ✓ 图表更新完成")

        except Exception as e:
            logger.error(f"更新高风险占比图表失败: {e}", exc_info=True)

    def _fill_analysis_text(self, slide, bull_data: Dict[str, Any]):
        """
        填充近交系数分析文本

        Args:
            slide: 幻灯片对象
            bull_data: 公牛数据
        """
        try:
            # 查找分析文本框（优先位置靠下的文本框，避免匹配标题）
            textbox = None
            candidate_textboxes = []

            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    name = shape.name if shape.name else ""
                    shape_text = shape.text if hasattr(shape, "text") else ""
                    # 跳过标题和公牛标识
                    if "备选公牛" in shape_text or "近交系数分析" in shape_text:
                        continue
                    if "【公牛" in shape_text:
                        continue
                    if "文本框" in name or "TextBox" in name:
                        candidate_textboxes.append(shape)

            # 按位置排序，选择最靠下的
            if candidate_textboxes:
                candidate_textboxes.sort(key=lambda s: getattr(s, 'top', 0), reverse=True)
                textbox = candidate_textboxes[0]

            if not textbox:
                logger.debug("    未找到分析文本框")
                return

            # 生成分析文本
            analysis = self._generate_inbreeding_analysis(bull_data)

            # 清空文本框并设置新内容
            tf = textbox.text_frame
            tf.clear()
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            run = p.add_run()
            run.text = analysis

            # 设置字体：微软雅黑15号非加粗
            run.font.name = FONT_NAME_CN
            run.font.size = Pt(15)
            run.font.bold = False

            logger.info(f"    ✓ 填充近交系数分析文本")

        except Exception as e:
            logger.warning(f"填充近交系数分析文本失败: {e}")

    def _generate_inbreeding_analysis(self, bull_data: Dict[str, Any]) -> str:
        """
        生成近交系数分析文本

        Args:
            bull_data: 公牛数据

        Returns:
            分析文本字符串
        """
        try:
            intervals = bull_data.get('intervals', [])
            high_risk = bull_data.get('high_risk_summary', {})

            # 获取高风险比例
            total_high_risk_ratio = self._parse_percentage(high_risk.get('total_ratio', '0.0%'))
            adult_high_risk_ratio = self._parse_percentage(high_risk.get('adult_ratio', '0.0%'))
            heifer_high_risk_ratio = self._parse_percentage(high_risk.get('heifer_ratio', '0.0%'))

            # 获取各区间占比
            low_risk_ratio = 0.0  # <3.125%
            medium_risk_ratio = 0.0  # 3.125%-6.25%
            if intervals:
                low_risk_ratio = self._parse_percentage(intervals[0].get('total_ratio', '0.0%'))
                if len(intervals) > 1:
                    medium_risk_ratio = self._parse_percentage(intervals[1].get('total_ratio', '0.0%'))

            # 生成分析文本
            parts = ["分析："]

            # 评估整体风险水平
            if total_high_risk_ratio < 0.05:  # <5%
                parts.append(f"该公牛与牧场母牛群的近交风险较低，高风险（>6.25%）近交占比仅{total_high_risk_ratio*100:.1f}%，")
                parts.append("适合大范围使用。")
            elif total_high_risk_ratio < 0.15:  # 5%-15%
                parts.append(f"该公牛与牧场母牛群存在一定近交风险，高风险（>6.25%）近交占比为{total_high_risk_ratio*100:.1f}%，")
                parts.append("建议在使用时注意避开高风险母牛。")
            elif total_high_risk_ratio < 0.30:  # 15%-30%
                parts.append(f"该公牛与牧场母牛群近交风险较高，高风险（>6.25%）近交占比达{total_high_risk_ratio*100:.1f}%，")
                parts.append("需谨慎使用，建议优先选配低风险母牛。")
            else:  # >30%
                parts.append(f"该公牛与牧场母牛群近交风险很高，高风险（>6.25%）近交占比达{total_high_risk_ratio*100:.1f}%，")
                parts.append("不建议大范围使用，仅推荐用于特定低近交风险母牛。")

            # 补充成母牛和后备牛的差异
            if abs(adult_high_risk_ratio - heifer_high_risk_ratio) > 0.05:
                if adult_high_risk_ratio > heifer_high_risk_ratio:
                    parts.append(f"成母牛群高风险占比({adult_high_risk_ratio*100:.1f}%)高于后备牛群({heifer_high_risk_ratio*100:.1f}%)，")
                    parts.append("建议优先用于后备牛配种。")
                else:
                    parts.append(f"后备牛群高风险占比({heifer_high_risk_ratio*100:.1f}%)高于成母牛群({adult_high_risk_ratio*100:.1f}%)，")
                    parts.append("建议优先用于成母牛配种。")

            return "".join(parts)

        except Exception as e:
            logger.warning(f"生成近交系数分析失败: {e}")
            return "分析：该公牛近交系数分析数据请参考上方表格和图表。"

    def _parse_percentage(self, percent_value) -> float:
        """
        将百分比字符串转换为小数

        Args:
            percent_value: 百分比值（字符串或数字）

        Returns:
            小数值（例如："0.3%" → 0.003）
        """
        try:
            if percent_value is None:
                return 0.0

            # 转换为字符串并清理
            percent_str = str(percent_value).strip().replace('%', '').strip()

            if not percent_str or percent_str == '':
                return 0.0

            # 转换为浮点数并除以100
            return float(percent_str) / 100.0

        except Exception as e:
            logger.warning(f"百分比转换失败: {percent_value} -> {e}")
            return 0.0

    def _delete_slide(self, slide_index: int):
        """
        删除指定索引的幻灯片

        Args:
            slide_index: 幻灯片索引
        """
        try:
            rId = self.prs.slides._sldIdLst[slide_index].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[slide_index]
        except Exception as e:
            logger.error(f"删除第{slide_index + 1}页失败: {e}")

    def _copy_slide(self, source_index: int, insert_at: int):
        """
        复制幻灯片到指定位置

        Args:
            source_index: 源幻灯片索引
            insert_at: 插入位置索引
        """
        try:
            import copy

            source = self.prs.slides[source_index]

            # 复制幻灯片XML
            blank_slide_layout = self.prs.slide_layouts[6]  # 空白布局
            dest = self.prs.slides.add_slide(blank_slide_layout)

            # 移动到指定位置
            xml_slides = self.prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[-1])
            xml_slides.insert(insert_at, slides[-1])

            # 复制所有形状
            for shape in source.shapes:
                el = shape.element
                newel = copy.deepcopy(el)
                dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

        except Exception as e:
            logger.error(f"复制第{source_index + 1}页到第{insert_at + 1}页失败: {e}")

    def _format_cell(self, cell, font_size: int = 12, alignment=None, bold: bool = False,
                     font_color: tuple = None, font_name: str = "微软雅黑"):
        """
        格式化单元格

        Args:
            cell: 单元格对象
            font_size: 字号，默认12
            alignment: 对齐方式
            bold: 是否加粗
            font_color: 字体颜色RGB元组，如 (255, 255, 255) 代表白色
            font_name: 字体名称，默认"微软雅黑"
        """
        try:
            from pptx.enum.text import PP_ALIGN

            if not cell.text_frame or not cell.text_frame.paragraphs:
                return

            for paragraph in cell.text_frame.paragraphs:
                if alignment is not None:
                    paragraph.alignment = alignment
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    if bold:
                        run.font.bold = True
                    if font_color:
                        run.font.color.rgb = RGBColor(*font_color)
                    if font_name:
                        run.font.name = font_name
        except Exception as e:
            logger.debug(f"    格式化单元格失败: {e}")
