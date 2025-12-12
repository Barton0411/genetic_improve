"""
Part 3: 系谱分析构建器
"""

import logging
from typing import Dict, List, Optional

import pandas as pd
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.util import Pt

from ..base_builder import BaseSlideBuilder
from ..config import FONT_NAME_CN

logger = logging.getLogger(__name__)


class Part3PedigreeBuilder(BaseSlideBuilder):
    """Part 3: 系谱分析构建器"""

    SECTION_SLIDE_INDEX = 6  # Slide 7 (0-based)
    SUMMARY_SLIDE_INDEX = 7  # Slide 8 (0-based)

    def __init__(self, prs, chart_creator, farm_name: str):
        super().__init__(prs, chart_creator)
        self.farm_name = farm_name

    def build(self, data: Dict, versions: Optional[List[str]] = None):
        """
        构建 Part 3: 系谱分析（填充模板中的第 7-8 页）
        """
        logger.info("=" * 60)
        logger.info("开始构建Part 3: 系谱分析")
        logger.info("=" * 60)

        # 检查数据是否为空
        df = data.get("pedigree_analysis")
        if df is None or df.empty:
            logger.warning("pedigree_analysis 数据为空，标记Part3页面待删除")
            slides_to_delete = [
                self.SECTION_SLIDE_INDEX,
                self.SUMMARY_SLIDE_INDEX
            ]
            self.mark_slides_for_deletion(slides_to_delete)
            return

        # 章节页（Slide 7）模板已就绪，这里暂不做动态修改
        self._fill_summary_slide(data)
        logger.info("✓ Part 3 模板页更新完成")

    # ------------------------------------------------------------------ #
    def _find_table(self, slide, name: str):
        for shape in slide.shapes:
            if shape.name == name and getattr(shape, "has_table", False):
                return shape.table
        return None

    @staticmethod
    def _delete_row(table, row_idx: int):
        """
        从导出的 PPT 中物理删除一行（不影响模板文件）。

        python-pptx 没有公开的删除行 API，这里通过底层 XML 删除 <tr> 节点。
        """
        try:
            row = table.rows[row_idx]
        except IndexError:
            return
        tbl = table._tbl  # 底层 <a:tbl>
        tr = row._tr      # 底层 <a:tr>
        tbl.remove(tr)

    # ------------------------------------------------------------------ #
    def _fill_summary_slide(self, data: Dict):
        """填充 Slide 8 - 系谱识别率可视化分析"""
        try:
            slide = self.prs.slides[self.SUMMARY_SLIDE_INDEX]
        except IndexError:
            logger.warning("模板缺少Part3汇总页（Slide 8）")
            return

        df = data.get("pedigree_analysis")
        if df is None or df.empty:
            logger.warning("pedigree_analysis 数据为空，跳过 Slide 8 填充")
            return

        # Excel中该区域的结构示例（列共 8 列）:
        # 0: 出生年份, 1: 总头数, 2: 可识别父号牛数, 3: 可识别父号占比,
        # 4: 可识别外祖父牛数, 5: 可识别外祖父占比, 6: 可识别外曾外祖父牛数, 7: 可识别外曾外祖父占比

        # 找到“出生年份”表头行
        header_row_idx = None
        for i in range(len(df)):
            if str(df.iloc[i, 0]).strip() == "出生年份":
                header_row_idx = i
                break
        if header_row_idx is None:
            logger.warning("未在系谱识别分析中找到“出生年份”表头，跳过填充")
            return

        # 从表头下一行开始，收集到“合计”行为止（包含“合计”），中间空行跳过
        rows = []
        for i in range(header_row_idx + 1, len(df)):
            first_cell = df.iloc[i, 0]
            if pd.isna(first_cell):
                # 空行直接跳过，避免被截断
                continue

            text = str(first_cell).strip()
            row_values = [
                df.iloc[i, j] if j < df.shape[1] else "" for j in range(8)
            ]
            rows.append([str(v) if not pd.isna(v) else "" for v in row_values])

            if text == "合计":
                break

        if not rows:
            logger.warning("未能从系谱识别分析中解析出任何汇总行")
            return

        # 将数据写入模板中的“表格 2”
        table = self._find_table(slide, "表格 2")
        if not table:
            logger.warning("未找到系谱汇总表（表格 2）")
            return

        # 1) 识别模板中的“合计”行：优先查找首列文本包含“合计”的行
        total_template_rows = len(table.rows)
        total_row_idx = None
        for r in range(1, total_template_rows):
            cell_text = table.cell(r, 0).text_frame.text.strip()
            if "合计" in cell_text:
                total_row_idx = r
                break
        if total_row_idx is None:
            # 找不到就退回用最后一行作为“合计”
            total_row_idx = total_template_rows - 1

        # 2) 根据模板结构计算可容纳的年份行行数
        year_slots = max(total_row_idx - 1, 0)

        # rows 中最后一行应为“合计”，其余为年份行
        if rows[-1][0] == "" or "合计" not in rows[-1][0]:
            logger.warning("Excel 最后一行不是“合计”，按顺序最后一行作为合计处理")
        year_rows = rows[:-1]
        total_row = rows[-1]

        # 如果年份行多于模板可容纳的行数，只保留最近的若干年
        if len(year_rows) > year_slots:
            logger.warning(
                "年份行数(%s)超过模板年份行数(%s)，仅保留最近 %s 年",
                len(year_rows),
                year_slots,
                year_slots,
            )
            year_rows = year_rows[-year_slots:]

        # 3) 写入年份行，从第 1 行到 total_row_idx-1
        for slot in range(year_slots):
            row_idx = 1 + slot
            if slot < len(year_rows):
                # 有对应年份数据
                row_values = year_rows[slot]
                for col_idx in range(len(table.columns)):
                    value = row_values[col_idx] if col_idx < len(row_values) else ""
                    self._set_cell_text(table.cell(row_idx, col_idx), value)

        # 4) 写入合计行到模板中的“合计”行（保持该行的浅蓝底样式）
        for col_idx in range(len(table.columns)):
            value = total_row[col_idx] if col_idx < len(total_row) else ""
            self._set_cell_text(table.cell(total_row_idx, col_idx), value)

        # 5) 删除“多余的模板行”，让导出的 PPT 在结构上也只保留
        #    （header + 实际年份行 + 合计行）这些需要的行。
        desired_year_rows = len(year_rows)
        data_start = 1
        current_total_rows = len(table.rows)
        rows_to_keep = {0, total_row_idx}
        # 保留前 N 个年份行（如果模板年份槽比实际年份多，多出的会被删除）
        for i in range(desired_year_rows):
            rows_to_keep.add(data_start + i)
        rows_to_delete = [
            r
            for r in range(1, current_total_rows)
            if r not in rows_to_keep
        ]
        # 从下往上删除，避免索引错位
        for r in sorted(rows_to_delete, reverse=True):
            self._delete_row(table, r)
        total_template_rows = len(table.rows)

        # 6) 根据“占比<80%标粉色”的规则设置底色（包括合计行）
        # 占比列索引：可识别父号占比(3)、可识别外祖父占比(5)、可识别外曾外祖父占比(7)
        percent_columns = [3, 5, 7]

        # 尝试从模板中找一个已有的粉色单元格作为颜色来源
        pink_rgb = None
        try:
            for r in range(1, total_template_rows):
                for c in range(len(table.columns)):
                    cell = table.cell(r, c)
                    fill = cell.fill
                    if fill.type == MSO_FILL.SOLID and fill.fore_color.rgb is not None:
                        # 找到第一个有填充色的正文单元格，作为粉色模板
                        pink_rgb = fill.fore_color.rgb
                        break
                if pink_rgb is not None:
                    break
        except Exception:
            pink_rgb = None

        if pink_rgb is None:
            # 兜底使用一个柔和的粉色
            pink_rgb = RGBColor(255, 230, 230)

        # 从第 1 行到最后一行（包括合计行），对占比列应用 <80% 标粉色规则
        for r in range(1, total_template_rows):
            for c in percent_columns:
                if c >= len(table.columns):
                    continue
                cell = table.cell(r, c)
                text = cell.text_frame.text.strip()
                if not text:
                    continue
                # 去掉百分号并解析数值
                clean = text.replace("%", "").replace("％", "").strip()
                try:
                    value = float(clean)
                except ValueError:
                    continue

                if value < 80.0:
                    # 小于80%标记为粉色
                    fill = cell.fill
                    fill.solid()
                    fill.fore_color.rgb = pink_rgb

        # 7) 创建三个系谱识别率饼图
        self._add_pedigree_pie_charts(slide, df, header_row_idx)

        # 8) 更新分析文本框
        self._update_pedigree_analysis(slide, df, header_row_idx, rows)

    # ------------------------------------------------------------------ #
    def _update_pedigree_analysis(self, slide, df: pd.DataFrame, header_row_idx: int, rows: list):
        """
        第8页：系谱识别率分析

        Args:
            slide: 幻灯片对象
            df: 系谱识别分析DataFrame
            header_row_idx: 表头行索引
            rows: 解析后的数据行列表
        """
        # 查找分析文本框
        shape = None
        for s in slide.shapes:
            if s.name == "文本框 13":
                shape = s
                break

        if not shape or not shape.has_text_frame:
            logger.warning("未找到系谱识别分析文本框（文本框 13）")
            return

        if not rows:
            self._set_analysis_text(shape, "分析：暂无系谱识别率数据。")
            return

        # 从合计行获取整体识别率
        total_row = rows[-1] if rows[-1][0] == "合计" else None
        if not total_row:
            self._set_analysis_text(shape, "分析：暂无系谱识别率数据。")
            return

        # 解析三代识别率（列3、5、7为占比）
        def parse_pct(val):
            try:
                return float(str(val).replace('%', '').replace('％', '').strip())
            except (ValueError, TypeError):
                return 0.0

        sire_pct = parse_pct(total_row[3]) if len(total_row) > 3 else 0
        mgs_pct = parse_pct(total_row[5]) if len(total_row) > 5 else 0
        mggs_pct = parse_pct(total_row[7]) if len(total_row) > 7 else 0

        # 识别问题指标（<80%）
        problems = []
        if sire_pct < 80:
            problems.append(f"父号识别率{sire_pct:.1f}%")
        if mgs_pct < 80:
            problems.append(f"外祖父识别率{mgs_pct:.1f}%")
        if mggs_pct < 80:
            problems.append(f"外曾外祖父识别率{mggs_pct:.1f}%")

        # 整体评价
        avg_pct = (sire_pct + mgs_pct + mggs_pct) / 3
        if avg_pct >= 90:
            overall = "系谱档案完整度优秀"
        elif avg_pct >= 80:
            overall = "系谱档案完整度良好"
        elif avg_pct >= 60:
            overall = "系谱档案有待完善"
        else:
            overall = "系谱档案亟需补充"

        # 年度趋势分析（找出识别率最低的年份）
        trend_info = ""
        year_rows = [r for r in rows if r[0] != "合计"]
        if year_rows:
            # 找父号识别率最低的年份
            min_row = min(year_rows, key=lambda r: parse_pct(r[3]) if len(r) > 3 else 100)
            worst_year = min_row[0]
            worst_pct = parse_pct(min_row[3]) if len(min_row) > 3 else 0
            if worst_pct < 80:
                trend_info = f"{worst_year}年父号识别率最低（{worst_pct:.1f}%），建议重点补充该年份系谱信息。"

        # 建议
        if problems:
            suggestion = f"当前{' / '.join(problems)}偏低，建议通过基因组检测或系谱追溯提升识别率。"
        else:
            suggestion = "各代系谱识别率均达标，为精准选配和近交控制提供了可靠基础。"

        text = (
            f"分析：牧场{overall}，父号识别率{sire_pct:.1f}%、外祖父识别率{mgs_pct:.1f}%、"
            f"外曾外祖父识别率{mggs_pct:.1f}%。{trend_info}{suggestion}"
        )
        self._set_analysis_text(shape, text)
        logger.info("✓ 系谱识别率分析文本更新完成")

    @staticmethod
    def _set_analysis_text(shape, text: str):
        """设置分析文本框内容，设置微软雅黑15号非加粗"""
        if not shape or not shape.has_text_frame:
            return
        tf = shape.text_frame
        # 清空并重新设置
        tf.clear()
        if not tf.paragraphs:
            para = tf.add_paragraph()
        else:
            para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text

        # 设置字体：微软雅黑15号非加粗
        run.font.name = FONT_NAME_CN
        run.font.size = Pt(15)
        run.font.bold = False

        # 清理多余段落
        for extra_para in tf.paragraphs[1:]:
            for extra_run in extra_para.runs:
                extra_run.text = ""

    # ------------------------------------------------------------------ #
    def _add_pedigree_pie_charts(self, slide, df: pd.DataFrame, header_row_idx: int):
        """
        更新三个系谱识别率饼图的数据（保持模板格式不变）

        Args:
            slide: 幻灯片对象
            df: 系谱识别分析DataFrame
            header_row_idx: 表头行索引
        """
        try:
            from pptx.chart.data import CategoryChartData

            # 查找"二、系谱识别率可视化分析"区域的数据
            viz_start_idx = None
            for i in range(header_row_idx, len(df)):
                cell = str(df.iloc[i, 0]).strip() if i < len(df) else ""
                if "可视化分析" in cell:
                    viz_start_idx = i
                    break

            if viz_start_idx is None:
                logger.warning("未找到'系谱识别率可视化分析'区域，跳过饼图更新")
                return

            # 读取已识别/未识别数据（通常在viz_start_idx+2和+3行）
            identified_row_idx = viz_start_idx + 2
            unidentified_row_idx = viz_start_idx + 3

            if identified_row_idx >= len(df) or unidentified_row_idx >= len(df):
                logger.warning("可视化分析数据行不足，跳过饼图更新")
                return

            # 提取数据
            # 列结构：col0=标签, col1=父号数值, col2=空, col3=标签, col4=外祖父数值,
            #          col5=空, col6=标签, col7=外曾外祖父数值
            sire_identified = self._parse_number(df.iloc[identified_row_idx, 1])
            sire_unidentified = self._parse_number(df.iloc[unidentified_row_idx, 1])

            mgs_identified = self._parse_number(df.iloc[identified_row_idx, 4])
            mgs_unidentified = self._parse_number(df.iloc[unidentified_row_idx, 4])

            mggs_identified = self._parse_number(df.iloc[identified_row_idx, 7])
            mggs_unidentified = self._parse_number(df.iloc[unidentified_row_idx, 7])

            logger.info(f"系谱识别数据: 父号({sire_identified}/{sire_unidentified}), "
                       f"外祖父({mgs_identified}/{mgs_unidentified}), "
                       f"外曾外祖父({mggs_identified}/{mggs_unidentified})")

            # 收集模板中的饼图（按位置从左到右排序）
            chart_shapes = []
            for shape in slide.shapes:
                if shape.has_chart:
                    chart_shapes.append(shape)

            # 按左边位置排序（从左到右对应：父号、外祖父、外曾外祖父）
            chart_shapes.sort(key=lambda s: s.left)

            if len(chart_shapes) < 3:
                logger.warning(f"模板中只找到 {len(chart_shapes)} 个饼图，预期3个")

            # 三个饼图的数据配置
            chart_data_list = [
                {'identified': sire_identified, 'unidentified': sire_unidentified},
                {'identified': mgs_identified, 'unidentified': mgs_unidentified},
                {'identified': mggs_identified, 'unidentified': mggs_unidentified},
            ]

            # 更新每个饼图的数据（保持模板格式不变）
            for i, shape in enumerate(chart_shapes[:3]):
                if i >= len(chart_data_list):
                    break

                data = chart_data_list[i]
                chart = shape.chart

                # 准备新数据
                chart_data = CategoryChartData()
                chart_data.categories = ['已识别', '未识别']
                chart_data.add_series('数量', [data['identified'], data['unidentified']])

                # 替换图表数据（保持格式不变）
                chart.replace_data(chart_data)

            logger.info("✓ 成功更新3个系谱识别率饼图数据")

        except Exception as e:
            logger.error(f"更新系谱识别率饼图失败: {e}", exc_info=True)

    @staticmethod
    def _parse_number(value) -> int:
        """解析数值，处理空值和非数字"""
        if pd.isna(value):
            return 0
        try:
            return int(float(value))
        except (ValueError, TypeError):
            return 0

    # ------------------------------------------------------------------ #
    @staticmethod
    def _set_cell_text(cell, text: str):
        """只更新单元格文字，保持模板样式"""
        tf = cell.text_frame
        if not tf.paragraphs:
            para = tf.add_paragraph()
        else:
            para = tf.paragraphs[0]
        if para.runs:
            run = para.runs[0]
        else:
            run = para.add_run()
        run.text = text
        # 清理其它 run/段落中的旧文本
        for extra_run in para.runs[1:]:
            extra_run.text = ""
        for extra_para in tf.paragraphs[1:]:
            for extra_run in extra_para.runs:
                extra_run.text = ""
