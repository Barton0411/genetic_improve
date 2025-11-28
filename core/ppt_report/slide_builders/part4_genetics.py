"""
Part 4: 牛群遗传评估（在群牛关键性状年均进展）

实现内容：
- 第10页：在群牛关键性状年均进展（表格）——直接写入Excel“年份汇总与性状进展”的在群母牛年份汇总区域
- 从第11页开始：在群牛关键性状年均进展（折线图），根据性状列数量动态生成页数，不设上限
"""

import logging
import re
from copy import deepcopy
from typing import Dict, Optional, List, Tuple, Callable

import pandas as pd
from openpyxl import load_workbook
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor

from ..base_builder import BaseSlideBuilder
from ..config import CONTENT_LEFT, CONTENT_TOP, CHART_WIDTH, CHART_HEIGHT, FONT_NAME_CN

logger = logging.getLogger(__name__)


class Part4GeneticsBuilder(BaseSlideBuilder):
    """Part 4: 牛群遗传评估（在群牛关键性状年均进展）"""

    SECTION_SLIDE_INDEX = 8  # Slide 9 (0-based) - 章节页，当前保持模板静态
    YEARLY_SLIDE_INDEX = 9   # Slide 10 (0-based) - 在群牛关键性状年均进展
    HEADER_KEYWORDS = (
        "分组",
        "排名分组",
        "排名",
        "样本量",
        "头数",
        "均值",
        "标准差",
        "变异系数",
        "中位数",
        "Q1-Q3",
        "最大值",
        "最小值",
        "平均得分",
    )

    # 性状定义字典（来源：CDCB官网 uscdcb.com）
    TRAIT_DEFINITIONS = {
        # ========== 选择指数 ==========
        "NM$": (
            "终身净收益指数",
            "美国主要遗传选择指数，综合约45个经济性状按其对牧场盈利能力的影响加权，"
            "预测每头动物传递给后代的终身利润差异（美元），数值越高表示综合经济效益越好"
        ),
        "TPI": (
            "总性能指数",
            "美国荷斯坦协会(Holstein USA)的综合选择指数，整合产量、健康和体型性状。"
            "与NM$类似但权重分配不同，更侧重体型性状，帮助牧场追踪遗传改良并做出育种决策"
        ),
        "CM$": ("奶酪收益指数", "针对奶酪生产牧场的选择指数，对产奶量赋予负权重，更侧重乳蛋白产量"),
        "FM$": ("液态奶收益指数", "针对液态奶市场的选择指数，显著增加产奶量权重，不考虑乳蛋白"),
        "GM$": ("放牧收益指数", "针对放牧型牧场的选择指数，繁殖性状权重是其他指数的2.5倍，强调季节性产犊"),

        # ========== 产量性状 ==========
        "MILK": (
            "产奶量PTA",
            "预测传递力(Predicted Transmitting Ability)，单位磅(lbs)。表示后代成年女儿相对于群体平均的产奶量差异，"
            "正值表示增产。在NM$中权重约为-1%（因高产奶量会稀释乳成分）"
        ),
        "FAT": (
            "乳脂产量PTA",
            "预测后代成年女儿的乳脂产量（磅），正值表示乳脂增产。"
            "在NM$中权重最高（约32%），是最重要的经济性状之一"
        ),
        "FAT%": ("乳脂率PTA", "遗传传递乳脂浓度的能力，以百分点表示。乳脂率=乳脂产量/产奶量×100%"),
        "PROT": (
            "乳蛋白产量PTA",
            "预测后代成年女儿的乳蛋白产量（磅），正值表示乳蛋白增产。"
            "在NM$中权重约13%，是第二重要的产量性状"
        ),
        "PROT%": ("乳蛋白率PTA", "遗传传递乳蛋白浓度的能力，以百分点表示。乳蛋白率=乳蛋白产量/产奶量×100%"),

        # ========== 饲料效率 ==========
        "RFI": (
            "剩余采食量",
            "实际采食量与基于体型、产量预期采食量的差值。负值表示在相同产量下采食更少，饲料效率更高。"
            "在NM$中权重约-4%，有利于降低饲料成本"
        ),
        "FSAV": ("节省饲料量", "每个泌乳期预期减少的干物质采食量（磅），正值表示更节省饲料"),
        "BWC": ("体重综合指数", "反映体型大小，在NM$中权重约-9%，因大体型牛只维持消耗更高"),

        # ========== 健康与长寿 ==========
        "SCS": (
            "体细胞评分",
            "通过体细胞数(SCC)反映乳房健康的遗传易感性，采用对数转换：SCS=log2(SCC/100)+3。"
            "数值越低乳房越健康，乳房炎发病率越低。在NM$中权重约-3%"
        ),
        "PL": (
            "生产寿命",
            "预测后代母牛在泌乳群中被淘汰前的留群时间（月）。数值越高表示使用年限越长，"
            "减少更新成本。在NM$中权重约15%，是重要的长寿性状"
        ),
        "LIV": (
            "母牛存活力",
            "预测后代母牛在泌乳群中存活（非因死亡淘汰）的差异，以百分点表示。"
            "数值越高表示非自愿淘汰率越低"
        ),
        "HLV": ("青年牛存活力", "预测后备牛从出生到初产期间的存活能力"),
        "MAST": (
            "乳房炎抗性",
            "预测后代对临床乳房炎的抵抗能力，以百分点表示。数值越高抗性越强，发病率越低"
        ),
        "DA": ("真胃移位抗性", "预测后代对真胃移位(Displaced Abomasum)的抵抗能力，以百分点表示"),
        "KETO": ("酮病抗性", "预测后代对酮病(Ketosis)的抵抗能力，以百分点表示"),
        "KET": ("酮病抗性", "预测后代对酮病(Ketosis)的抵抗能力，以百分点表示"),
        "MF": ("低血钙症抗性", "预测后代对产后瘫痪/低血钙症(Milk Fever/Hypocalcemia)的抵抗能力"),
        "METR": ("子宫炎抗性", "预测后代对子宫炎(Metritis)的抵抗能力，以百分点表示"),
        "MET": ("子宫炎抗性", "预测后代对子宫炎(Metritis)的抵抗能力，以百分点表示"),
        "RP": ("胎衣不下抗性", "预测后代对胎衣不下(Retained Placenta)的抵抗能力，以百分点表示"),
        "HTH$": ("健康综合指数", "综合六个健康性状(乳房炎、酮病、真胃移位、胎衣不下、子宫炎、低血钙症)的经济指数"),

        # ========== 繁殖性状 ==========
        "DPR": (
            "女儿妊娠率",
            "每21天发情周期内非妊娠母牛的预期受孕百分比。每增加1%意味着空怀天数减少约4天。"
            "在NM$中权重约5%，是主要的繁殖性状"
        ),
        "CCR": (
            "经产牛受孕率",
            "泌乳母牛的受孕能力，定义为预期受孕的百分比。在NM$中权重约1%"
        ),
        "HCR": (
            "青年牛受孕率",
            "未产青年牛的受孕能力，定义为预期受孕的百分比。在NM$中权重约1%"
        ),
        "GL": ("妊娠期", "预测后代对妊娠天数的影响，较短妊娠期可减少难产风险"),
        "EFC": ("早期初产", "反映后备牛早期投产能力，在NM$中权重约1%"),
        "CA$": ("产犊能力指数", "综合产犊难易度、死胎率等产犊相关性状的经济指数，在NM$中权重约5%"),

        # ========== 体型综合性状 ==========
        "PTAT": (
            "体型综合PTA",
            "遗传传递品种鉴定最终评分的能力。基于体型线性鉴定计算，正值表示体型优于平均"
        ),
        "UDC": (
            "乳房综合指数",
            "综合前乳房附着、后乳房高度、后乳房宽度、乳房深度、乳头位置等乳房性状，"
            "按经济效益加权。在NM$中权重约7%，良好的乳房结构有利于机器挤奶和乳房健康"
        ),
        "FLC": (
            "肢蹄综合指数",
            "综合蹄角度、后肢侧视、后肢后视等肢蹄性状。在NM$中权重约3%，"
            "良好的肢蹄结构有利于奶牛行走和使用寿命"
        ),
        "FS": ("最终得分", "体型鉴定的最终分类评分，反映整体外貌质量"),

        # ========== 体型线性性状 ==========
        "ST": ("体高", "从腰角到地面的垂直高度，反映奶牛的身高"),
        "SR": ("体躯强度", "胸部宽度，反映胸宽和体躯结实程度"),
        "BD": ("体深", "从背线到腹底的深度，反映奶牛的体躯容积"),
        "DF": ("乳用特征", "反映奶牛的乳用型外貌特征，骨骼棱角分明、皮薄毛细的程度"),
        "RA": ("尻角度", "从腰角到坐骨的倾斜度，适中角度有利于产犊和排出胎衣"),
        "RW": ("尻宽", "两坐骨间的宽度，较宽有利于产犊"),
        "RLS": ("后肢侧视", "从侧面观察飞节角度，适中角度有利于运动和承重"),
        "RLR": ("后肢后视", "从后面观察后肢的角度，理想为两腿平行"),
        "FA": ("蹄角度", "前蹄蹄壁与地面的角度，45度左右为理想"),
        "FUA": ("前乳房附着", "前乳房与腹壁的连接强度和位置"),
        "RUH": ("后乳房高度", "后乳房附着点相对于阴门的高度"),
        "RUW": ("后乳房宽度", "后乳房在后视时的宽度"),
        "UC": ("乳房深度", "乳房底部相对于飞节的位置，理想为高于飞节"),
        "TP": ("乳头位置", "前乳头在乳区的位置，理想为位于乳区中央"),
        "TL": ("乳头长度", "乳头的长度，中等长度便于机器挤奶"),
    }

    def __init__(self, prs, chart_creator, farm_name: str):
        super().__init__(prs, chart_creator)
        self.farm_name = farm_name
        self.excel_path: Optional[str] = None
        self._excel_wb = None
        self._distribution_cache: Dict[str, Dict[str, List[Dict]]] = {}

    def build(self, data: Dict, versions: Optional[list] = None):
        """
        构建 Part 4: 牛群遗传评估

        当前仅填充模板中的第 10 页（在群牛关键性状年均进展表格）。
        """
        logger.info("=" * 60)
        logger.info("开始构建Part 4: 牛群遗传评估")
        logger.info("=" * 60)

        self.excel_path = data.get("excel_path")
        if not self.excel_path:
            logger.warning("未获取到Excel路径，无法从源Excel复制统计表")

        # 优先使用缓存的workbook（避免重复加载，每次加载需要约21秒）
        cached_wb = data.get("_cached_workbook_data_only")
        if cached_wb is not None:
            self._excel_wb = cached_wb
            logger.debug("使用缓存的workbook (data_only)")

        self._fill_yearly_traits_slide(data)

        # 从第 11 页开始：在群牛关键性状年均进展折线图（动态）
        self._fill_yearly_traits_charts(data)

        # 在群牛遗传分布（模板中第 81 页）
        self._fill_genetic_distribution_slide(data)

        # 在群牛NM$正态分布相关页面（整体 / 成母牛vs后备牛 / 不同阶段 / 不同出生年份）
        self._fill_nm_normal_distribution_slides(data)

        # 在群牛TPI分布及TPI正态分布相关页面
        self._fill_tpi_distribution_slide(data)
        self._fill_tpi_normal_distribution_slides(data)

        # 在群牛育种指数分布及指数正态分布相关页面
        self._fill_index_distribution_slide(data)
        self._fill_index_normal_distribution_slides(data)

        logger.info("✓ Part 4 模板页更新完成")

    # ------------------------------------------------------------------ #
    def _find_table(self, slide, name: str):
        for shape in slide.shapes:
            if shape.name == name and getattr(shape, "has_table", False):
                return shape.table
        return None

    # ------------------------------------------------------------------ #
    def _fill_yearly_traits_slide(self, data: Dict):
        """
        填充 Slide 10 - 在群牛关键性状年均进展（表格部分）

        使用 Excel《年份汇总与性状进展》Sheet 中“在群母牛年份汇总”区域：
        - 前 9 行（包含表头、年份行、合计行、对比行）
        - 前 16 列（出生年份、头数、NM$、TPI、MILK、FAT、FAT%、PROT、PROT%、SCS、PL、DPR、PTAT、UDC、FLC、RFI）
        """
        try:
            slide = self.prs.slides[self.YEARLY_SLIDE_INDEX]
        except IndexError:
            logger.warning("模板缺少Part4 年度性状页（Slide 10）")
            return

        df = data.get("traits_yearly")
        if df is None or df.empty:
            logger.warning("traits_yearly 数据为空，跳过 Slide 10 填充")
            return

        table = self._find_table(slide, "表格 2")
        if not table:
            logger.warning("未找到在群牛关键性状年均进展表（表格 2）")
            return

        # 限定为前 9 行、前 16 列（与模板行列数一致）
        max_rows = min(len(table.rows), len(df))
        max_cols = min(len(table.columns), 16, df.shape[1])

        logger.info(
            "开始填充 Slide 10 表格：使用前 %s 行 × %s 列数据",
            max_rows,
            max_cols,
        )

        # P1优化：批量转换为numpy数组，减少iloc访问开销
        data_array = df.iloc[:max_rows, :max_cols].fillna("").astype(str).to_numpy()

        for r in range(max_rows):
            for c in range(len(table.columns)):
                cell_text = data_array[r, c] if c < max_cols else ""
                self._set_cell_text(table.cell(r, c), cell_text)

        # 若模板行数多于数据行数，清空剩余行内容
        for r in range(max_rows, len(table.rows)):
            for c in range(len(table.columns)):
                self._set_cell_text(table.cell(r, c), "")

        # 生成并填充分析文本
        analysis_text = self._generate_yearly_traits_analysis(df)
        if analysis_text:
            self._update_yearly_traits_analysis(slide, analysis_text)

    # ------------------------------------------------------------------ #
    def _generate_yearly_traits_analysis(self, df: pd.DataFrame) -> str:
        """
        生成在群母牛年份汇总表的分析文本

        Args:
            df: 年份汇总数据表（包含出生年份、头数及各性状均值）

        Returns:
            分析文本字符串
        """
        if df is None or df.empty:
            return ""

        try:
            # 表格结构：第一列=出生年份，第二列=头数，后续列=各性状
            # 最后几行通常是"合计"、"总计"、"对比"等汇总行
            # 年份格式可能是"2021年"、"2020年及以前"等
            year_col = df.columns[0]

            def is_year_row(val):
                """判断是否为年份数据行（支持"2021年"、"2020年及以前"等格式）"""
                if pd.isna(val):
                    return False
                val_str = str(val).strip()
                # 排除汇总行
                if any(kw in val_str for kw in ["总计", "合计", "对比", "小计"]):
                    return False
                # 包含"年"且有4位数字的行认为是年份行
                import re
                if "年" in val_str and re.search(r'\d{4}', val_str):
                    return True
                # 纯4位数字也是年份
                if val_str.isdigit() and len(val_str) == 4:
                    return True
                return False

            data_rows = df[df[year_col].apply(is_year_row)].copy()

            if data_rows.empty:
                return ""

            # 提取年份数字
            import re
            def extract_year(val):
                match = re.search(r'(\d{4})', str(val))
                return int(match.group(1)) if match else 0

            years = [extract_year(v) for v in data_rows[year_col]]
            years = [y for y in years if y > 0]
            if not years:
                return ""
            min_year = min(years)
            max_year = max(years)

            # 统计各年份头数
            count_col = df.columns[1]
            total_count = data_rows[count_col].sum()

            # 关键性状分析（选取NM$、TPI、MILK、FAT、PROT、SCS等）
            key_traits = ["NM$", "TPI", "MILK", "FAT", "PROT", "SCS", "PL", "DPR"]
            trait_trends = []

            for trait in key_traits:
                if trait in df.columns:
                    trait_idx = df.columns.get_loc(trait)
                    # 计算近年趋势（最近3-5年）
                    recent_data = data_rows.iloc[-min(5, len(data_rows)):][df.columns[trait_idx]]
                    recent_data = pd.to_numeric(recent_data, errors='coerce').dropna()
                    if len(recent_data) >= 2:
                        first_val = recent_data.iloc[0]
                        last_val = recent_data.iloc[-1]
                        if first_val != 0:
                            change_pct = ((last_val - first_val) / abs(first_val)) * 100
                            if abs(change_pct) > 5:
                                direction = "提升" if change_pct > 0 else "下降"
                                # SCS特殊处理：下降是好事
                                if trait == "SCS" and change_pct < 0:
                                    direction = "改善"
                                trait_trends.append((trait, direction, abs(change_pct)))

            # 构建分析文本
            parts = []
            parts.append(f"在群母牛出生年份分布为{min_year}-{max_year}年，共计{int(total_count)}头。")

            # 整体趋势判断
            if trait_trends:
                improving = [t for t in trait_trends if t[1] in ("提升", "改善")]
                declining = [t for t in trait_trends if t[1] == "下降"]

                if len(improving) > len(declining):
                    parts.append("近年遗传进展整体向好，")
                    if improving:
                        top_improve = sorted(improving, key=lambda x: x[2], reverse=True)[:2]
                        traits_str = "、".join([f"{t[0]}" for t in top_improve])
                        parts.append(f"{traits_str}等性状表现突出。")
                elif declining:
                    parts.append("部分性状有下降趋势，建议关注育种方向调整。")
                else:
                    parts.append("各性状整体保持稳定。")
            else:
                parts.append("各性状年度均值较为稳定。")

            return "".join(parts)

        except Exception as e:
            logger.warning(f"生成年份汇总分析文本时出错: {e}")
            return ""

    def _update_yearly_traits_analysis(self, slide, analysis_text: str):
        """
        更新Slide 10的分析文本框

        查找包含"分析："前缀的文本框（通常在页面底部）
        """
        if not analysis_text:
            return

        analysis_box = None

        # 方法1: 查找包含"分析："或"分析:"的文本框（这是最可靠的标识）
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip() if shape.text else ""
                # 查找以"分析"开头的文本框
                if text.startswith("分析：") or text.startswith("分析:"):
                    analysis_box = shape
                    logger.debug(f"找到分析文本框（通过'分析：'前缀）: {shape.name}")
                    break

        # 方法2: 查找名称为"文本框 14"的形状（常见的分析框名称）
        if not analysis_box:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.name == "文本框 14":
                    analysis_box = shape
                    logger.debug(f"找到分析文本框（通过名称'文本框 14'）")
                    break

        # 方法3: 查找页面下半部分的文本框（排除标题和表格）
        if not analysis_box:
            # 收集所有候选文本框
            candidates = []
            for shape in slide.shapes:
                if shape.has_text_frame and not getattr(shape, "has_table", False):
                    # 排除标题（通常在顶部，top值很小）
                    if hasattr(shape, "top") and shape.top > 4000000:  # 约4cm以下
                        text = shape.text.strip() if shape.text else ""
                        # 排除明显的标题（字数很少且不包含"分析"）
                        if "分析" in text or len(text) > 20:
                            candidates.append((shape.top, shape))

            # 选择位置最靠下的文本框
            if candidates:
                candidates.sort(key=lambda x: x[0], reverse=True)
                analysis_box = candidates[0][1]
                logger.debug(f"找到分析文本框（通过位置）: {analysis_box.name}")

        if analysis_box:
            tf = analysis_box.text_frame
            # 保留第一段的格式
            if tf.paragraphs:
                para = tf.paragraphs[0]
                # 保存原始格式
                font_size = None
                font_name = None
                font_color = None
                if para.runs:
                    run = para.runs[0]
                    font_size = run.font.size
                    font_name = run.font.name
                    # 安全获取颜色（处理_NoneColor等情况）
                    try:
                        if run.font.color and run.font.color.type is not None:
                            font_color = run.font.color.rgb
                    except (AttributeError, TypeError):
                        font_color = None

                # 清空并设置新文本（避免重复"分析："前缀）
                para.clear()
                run = para.add_run()
                if analysis_text.startswith("分析：") or analysis_text.startswith("分析:"):
                    run.text = analysis_text
                else:
                    run.text = f"分析：{analysis_text}"

                # 设置字体：微软雅黑15号非加粗
                run.font.name = FONT_NAME_CN
                run.font.size = Pt(15)
                run.font.bold = False

            logger.info(f"✓ 第10页年份汇总分析文本已更新")
        else:
            logger.warning("未找到第10页的分析文本框")

    # ------------------------------------------------------------------ #
    def _fill_genetic_distribution_slide(self, data: Dict):
        """
        填充“在群牛遗传分布”页（模板中第81页，index=80）

        数据来源：Excel《NM$分布分析》Sheet
        - 左侧表格：在群母牛NM$分布（分布区间 + 头数 + 占比）
        - 上方饼图：在群母牛各分布区间占比
        - 上方柱状图：在群母牛 vs 全部母牛的头数分布
        """
        # 由于前面可能删除了多余的关键性状页，这里的索引会变化，
        # 因此通过文本内容动态查找目标页，而不是使用固定index。
        slide = None

        # 1）优先根据小标题“在群母牛NM$分布”查找
        for s in self.prs.slides:
            has_subtitle = any(
                hasattr(shape, "text") and "在群母牛NM$分布" in str(shape.text)
                for shape in s.shapes
            )
            if has_subtitle:
                slide = s
                break

        # 2）如果找不到小标题，再退回到大标题“在群牛遗传分布”
        if slide is None:
            for s in self.prs.slides:
                has_title = any(
                    hasattr(shape, "text") and "在群牛遗传分布" in str(shape.text)
                    for shape in s.shapes
                )
                if has_title:
                    slide = s
                    break

        if slide is None:
            logger.warning("未在模板中找到包含“在群母牛NM$分布”或“在群牛遗传分布”的幻灯片，跳过填充")
            return

        nm_df = data.get("traits_nm")
        if nm_df is None or nm_df.empty:
            logger.warning("NM$分布分析数据为空，跳过在群牛遗传分布页填充")
            return

        # 解析在群母牛和全部母牛的分布表
        # 结构示例：
        # 列0: 在群母牛NM$分布   -> 分布区间 / （表头）
        # 列1: Unnamed:1        -> 头数
        # 列2: Unnamed:2        -> 占比
        # 列9: 全部母牛NM$分布    -> 分布区间
        # 列10: Unnamed:10      -> 头数
        # 列11: Unnamed:11      -> 占比
        try:
            present_col_interval = nm_df.columns[0]
            present_col_count = nm_df.columns[1]
            present_col_percent = nm_df.columns[2]

            all_col_interval = nm_df.columns[9]
            all_col_count = nm_df.columns[10]
        except Exception as e:
            logger.warning("NM$分布分析Sheet列结构异常，跳过在群牛遗传分布页填充: %s", e)
            return

        # 去掉表头行，只保留实际数据行
        present_data = nm_df.iloc[1:].copy()
        present_data = present_data[
            present_data[present_col_interval].notna()
            & present_data[present_col_count].notna()
        ]

        # 左侧表格：在群母牛NM$分布
        table = self._find_table(slide, "表格 1")
        if table:
            max_rows = min(len(table.rows), len(present_data) + 1)  # 包含表头
            max_cols = min(len(table.columns), 3)

            # 表头
            headers = [str(v) for v in nm_df.iloc[0, :3].tolist()]
            for c in range(max_cols):
                self._set_cell_text(table.cell(0, c), headers[c] if c < len(headers) else "")

            # 数据行
            for r in range(1, max_rows):
                row = present_data.iloc[r - 1]
                values = [
                    row.get(present_col_interval, ""),
                    row.get(present_col_count, ""),
                    row.get(present_col_percent, ""),
                ]
                for c in range(max_cols):
                    self._set_cell_text(table.cell(r, c), "" if pd.isna(values[c]) else str(values[c]))

            # 清空剩余行
            for r in range(max_rows, len(table.rows)):
                for c in range(len(table.columns)):
                    self._set_cell_text(table.cell(r, c), "")
        else:
            logger.warning("在群牛遗传分布页未找到表格 1，跳过表格填充")

        # 处理图表数据
        # 分类区间统一使用在群母牛的分布区间
        categories = present_data[present_col_interval].astype(str).tolist()

        # 在群母牛头数
        def _to_float_list(series):
            out = []
            for v in series:
                if pd.isna(v):
                    out.append(None)
                else:
                    try:
                        out.append(float(str(v).replace(",", "")))
                    except Exception:
                        out.append(None)
            return out

        present_counts = _to_float_list(present_data[present_col_count])

        # 在群母牛占比（用于饼图）
        present_percents = []
        for v in present_data[present_col_percent]:
            if pd.isna(v):
                present_percents.append(None)
            else:
                s = str(v).replace("%", "").strip()
                try:
                    present_percents.append(float(s))
                except Exception:
                    present_percents.append(None)

        # 全部母牛头数（如果有）
        if all_col_count in nm_df.columns:
            all_data = nm_df.iloc[1:].copy()
            all_data = all_data[all_data[all_col_interval].notna() & all_data[all_col_count].notna()]
            all_counts = _to_float_list(all_data[all_col_count])
        else:
            all_counts = None

        # 填充饼图（Chart 1）：在群母牛占比
        pie_shape = next((s for s in slide.shapes if getattr(s, "has_chart", False)
                          and s.chart.chart_type == XL_CHART_TYPE.PIE_EXPLODED), None)
        if pie_shape:
            pie_chart = pie_shape.chart
            pie_data = CategoryChartData()
            pie_data.categories = categories
            pie_data.add_series("在群母牛", [v if v is not None else 0 for v in present_percents])
            pie_chart.replace_data(pie_data)
        else:
            logger.warning("在群牛遗传分布页未找到饼图（Chart 1）")

        # 填充柱状图（Chart 2）：在群母牛头数分布（不再叠加全部母牛）
        bar_shape = next((s for s in slide.shapes if getattr(s, "has_chart", False)
                          and s.chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED), None)
        if bar_shape:
            bar_chart = bar_shape.chart
            bar_data = CategoryChartData()
            bar_data.categories = categories
            bar_data.add_series("在群母牛", [v if v is not None else 0 for v in present_counts])
            bar_chart.replace_data(bar_data)
        else:
            logger.warning("在群牛遗传分布页未找到柱状图（Chart 2）")

        # 生成并更新分析文本
        distribution_data = [
            {"interval": cat, "count": cnt, "percent": pct}
            for cat, cnt, pct in zip(categories, present_counts, present_percents)
            if cnt is not None
        ]
        total_count = sum(c for c in present_counts if c is not None)
        analysis_text = self._generate_distribution_analysis("NM$", distribution_data, int(total_count))
        self._update_distribution_analysis_textbox(slide, analysis_text)
        logger.info("✓ NM$分布页分析文本更新完成")

    # ------------------------------------------------------------------ #
    def _fill_tpi_distribution_slide(self, data: Dict):
        """
        填充“在群母牛TPI分布”页：
        - 来源：Excel《TPI分布分析》Sheet（traits_tpi）
        - 左侧表格：在群母牛TPI分布（表格 1）
        - 饼图：在群母牛各区间占比（Chart 1）
        - 柱状图：在群母牛头数分布（Chart 2）
        """
        slide = None

        # 优先根据小标题“在群母牛TPI分布”查找
        for s in self.prs.slides:
            has_subtitle = any(
                hasattr(shape, "text") and "在群母牛TPI分布" in str(shape.text)
                for shape in s.shapes
            )
            if has_subtitle:
                slide = s
                break

        if slide is None:
            logger.warning("未在模板中找到包含“在群母牛TPI分布”的幻灯片，跳过TPI分布页填充")
            return

        tpi_df = data.get("traits_tpi")
        if tpi_df is None or tpi_df.empty:
            logger.warning("TPI分布分析数据为空，跳过在群母牛TPI分布页填充")
            return

        try:
            present_col_interval = tpi_df.columns[0]
            present_col_count = tpi_df.columns[1]
            present_col_percent = tpi_df.columns[2]
        except Exception as e:
            logger.warning("TPI分布分析Sheet列结构异常，跳过在群母牛TPI分布页填充: %s", e)
            return

        present_data = tpi_df.iloc[1:].copy()
        present_data = present_data[
            present_data[present_col_interval].notna()
            & present_data[present_col_count].notna()
        ]

        table = self._find_table(slide, "表格 1")
        if table:
            max_rows = min(len(table.rows), len(present_data) + 1)
            max_cols = min(len(table.columns), 3)

            headers = [str(v) for v in tpi_df.iloc[0, :3].tolist()]
            for c in range(max_cols):
                self._set_cell_text(table.cell(0, c), headers[c] if c < len(headers) else "")

            for r in range(1, max_rows):
                row = present_data.iloc[r - 1]
                values = [
                    row.get(present_col_interval, ""),
                    row.get(present_col_count, ""),
                    row.get(present_col_percent, ""),
                ]
                for c in range(max_cols):
                    val = values[c]
                    self._set_cell_text(
                        table.cell(r, c),
                        "" if pd.isna(val) else str(val),
                    )

            for r in range(max_rows, len(table.rows)):
                for c in range(len(table.columns)):
                    self._set_cell_text(table.cell(r, c), "")
        else:
            logger.warning("在群母牛TPI分布页未找到表格 1，跳过表格填充")

        categories = present_data[present_col_interval].astype(str).tolist()

        def _to_float_list(series):
            out = []
            for v in series:
                if pd.isna(v):
                    out.append(None)
                else:
                    try:
                        out.append(float(str(v).replace(",", "")))
                    except Exception:
                        out.append(None)
            return out

        present_counts = _to_float_list(present_data[present_col_count])

        present_percents = []
        for v in present_data[present_col_percent]:
            if pd.isna(v):
                present_percents.append(None)
            else:
                s = str(v).replace("%", "").strip()
                try:
                    present_percents.append(float(s))
                except Exception:
                    present_percents.append(None)

        pie_shape = next(
            (
                s
                for s in slide.shapes
                if getattr(s, "has_chart", False)
                and s.chart.chart_type == XL_CHART_TYPE.PIE_EXPLODED
            ),
            None,
        )
        if pie_shape:
            pie_chart = pie_shape.chart
            pie_data = CategoryChartData()
            pie_data.categories = categories
            pie_data.add_series(
                "在群母牛", [v if v is not None else 0 for v in present_percents]
            )
            pie_chart.replace_data(pie_data)
        else:
            logger.warning("在群母牛TPI分布页未找到饼图（Chart 1）")

        bar_shape = next(
            (
                s
                for s in slide.shapes
                if getattr(s, "has_chart", False)
                and s.chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
            ),
            None,
        )
        if bar_shape:
            bar_chart = bar_shape.chart
            bar_data = CategoryChartData()
            bar_data.categories = categories
            bar_data.add_series(
                "在群母牛", [v if v is not None else 0 for v in present_counts]
            )
            bar_chart.replace_data(bar_data)
        else:
            logger.warning("在群母牛TPI分布页未找到柱状图（Chart 2）")

        # 生成并更新分析文本
        distribution_data = [
            {"interval": cat, "count": cnt, "percent": pct}
            for cat, cnt, pct in zip(categories, present_counts, present_percents)
            if cnt is not None
        ]
        total_count = sum(c for c in present_counts if c is not None)
        analysis_text = self._generate_distribution_analysis("TPI", distribution_data, int(total_count))
        self._update_distribution_analysis_textbox(slide, analysis_text)
        logger.info("✓ TPI分布页分析文本更新完成")

    # ------------------------------------------------------------------ #
    def _fill_index_distribution_slide(self, data: Dict):
        """
        填充“在群母牛育种指数分布”页：
        - 来源：Excel《育种指数分布分析》Sheet（traits_index_distribution）
        - 表格 5 + 两个图表
        """
        slide = None

        for s in self.prs.slides:
            has_subtitle = any(
                hasattr(shape, "text") and "在群母牛育种指数分布" in str(shape.text)
                for shape in s.shapes
            )
            if has_subtitle:
                slide = s
                break

        if slide is None:
            logger.warning("未在模板中找到“在群母牛育种指数分布”页，跳过指数分布填充")
            return

        idx_df = data.get("traits_index_distribution")
        if idx_df is None or idx_df.empty:
            logger.warning("育种指数分布分析数据为空，跳过在群母牛育种指数分布页填充")
            return

        try:
            present_col_interval = idx_df.columns[0]
            present_col_count = idx_df.columns[1]
            present_col_percent = idx_df.columns[2]
        except Exception as e:
            logger.warning("育种指数分布分析Sheet列结构异常，跳过在群母牛育种指数分布页填充: %s", e)
            return

        present_data = idx_df.iloc[1:].copy()
        present_data = present_data[
            present_data[present_col_interval].notna()
            & present_data[present_col_count].notna()
        ]

        table = self._find_table(slide, "表格 5")
        if table:
            max_rows = min(len(table.rows), len(present_data) + 1)
            max_cols = min(len(table.columns), 3)

            headers = [str(v) for v in idx_df.iloc[0, :3].tolist()]
            for c in range(max_cols):
                self._set_cell_text(table.cell(0, c), headers[c] if c < len(headers) else "")

            for r in range(1, max_rows):
                row = present_data.iloc[r - 1]
                values = [
                    row.get(present_col_interval, ""),
                    row.get(present_col_count, ""),
                    row.get(present_col_percent, ""),
                ]
                for c in range(max_cols):
                    val = values[c]
                    self._set_cell_text(
                        table.cell(r, c),
                        "" if pd.isna(val) else str(val),
                    )

            for r in range(max_rows, len(table.rows)):
                for c in range(len(table.columns)):
                    self._set_cell_text(table.cell(r, c), "")
        else:
            logger.warning("在群母牛育种指数分布页未找到表格 5，跳过表格填充")

        categories = present_data[present_col_interval].astype(str).tolist()

        def _to_float_list(series):
            out = []
            for v in series:
                if pd.isna(v):
                    out.append(None)
                else:
                    try:
                        out.append(float(str(v).replace(",", "")))
                    except Exception:
                        out.append(None)
            return out

        present_counts = _to_float_list(present_data[present_col_count])

        present_percents = []
        for v in present_data[present_col_percent]:
            if pd.isna(v):
                present_percents.append(None)
            else:
                s = str(v).replace("%", "").strip()
                try:
                    present_percents.append(float(s))
                except Exception:
                    present_percents.append(None)

        pie_shape = next(
            (
                s
                for s in slide.shapes
                if getattr(s, "has_chart", False)
                and s.chart.chart_type == XL_CHART_TYPE.PIE_EXPLODED
            ),
            None,
        )
        if pie_shape:
            pie_chart = pie_shape.chart
            pie_data = CategoryChartData()
            pie_data.categories = categories
            pie_data.add_series(
                "在群母牛", [v if v is not None else 0 for v in present_percents]
            )
            pie_chart.replace_data(pie_data)
        else:
            logger.warning("在群母牛育种指数分布页未找到饼图（Chart 1）")

        bar_shape = next(
            (
                s
                for s in slide.shapes
                if getattr(s, "has_chart", False)
                and s.chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
            ),
            None,
        )
        if bar_shape:
            bar_chart = bar_shape.chart
            bar_data = CategoryChartData()
            bar_data.categories = categories
            bar_data.add_series(
                "在群母牛", [v if v is not None else 0 for v in present_counts]
            )
            bar_chart.replace_data(bar_data)
        else:
            logger.warning("在群母牛育种指数分布页未找到柱状图（Chart 2）")

        # 生成并更新分析文本
        distribution_data = [
            {"interval": cat, "count": cnt, "percent": pct}
            for cat, cnt, pct in zip(categories, present_counts, present_percents)
            if cnt is not None
        ]
        total_count = sum(c for c in present_counts if c is not None)
        analysis_text = self._generate_distribution_analysis("育种指数", distribution_data, int(total_count))
        self._update_distribution_analysis_textbox(slide, analysis_text)
        logger.info("✓ 育种指数分布页分析文本更新完成")

    # ------------------------------------------------------------------ #
    def _collect_slide_text(self, slide) -> str:
        """收集某页所有可见文本内容，用于关键字匹配"""
        texts: List[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                try:
                    texts.append(shape.text)
                except Exception:
                    continue
        return "\n".join(texts)

    def _find_slide_by_keywords(
        self,
        include_keywords: List[str],
        exclude_keywords: Optional[List[str]] = None,
    ):
        """
        按文本关键字查找幻灯片

        include_keywords 中的每个关键字都必须出现在该页文本中；
        exclude_keywords 中的关键字若出现则跳过。
        """
        exclude_keywords = exclude_keywords or []

        for slide in self.prs.slides:
            text = self._collect_slide_text(slide)
            if not text:
                continue

            if all(k in text for k in include_keywords) and not any(
                k in text for k in exclude_keywords
            ):
                return slide
        return None

    def _update_slide_picture(self, slide, image_path: str, context_desc: str):
        """
        使用生成的图片覆盖幻灯片中的主图片位置

        行为：删除原有图片占位符，然后在相同位置插入新图片（等效“替换”）。
        """
        try:
            picture_shapes = [
                s
                for s in slide.shapes
                if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
            ]
            if not picture_shapes:
                logger.warning(
                    "未在“%s”相关页面找到图片占位符，跳过图片更新", context_desc
                )
                return

            # 仅使用第一个图片作为尺寸和位置参考
            placeholder = picture_shapes[0]
            left = placeholder.left
            top = placeholder.top
            width = placeholder.width
            height = placeholder.height

            # 删除原来的图片占位符
            placeholder_element = placeholder.element
            placeholder_element.getparent().remove(placeholder_element)

            slide.shapes.add_picture(image_path, left, top, width=width, height=height)
            logger.info("✓ 已更新“%s”页面的正态分布图: %s", context_desc, image_path)
        except Exception as e:
            logger.error("更新“%s”页面正态分布图失败: %s", context_desc, e)

    # ------------------------------------------------------------------ #
    def _update_year_text_boxes(self, slide, actual_year_labels: list):
        """
        更新年份页的文本框标题，将模板占位符替换为实际年份

        Args:
            slide: PPT幻灯片对象
            actual_year_labels: 实际年份标签列表，如 ["在群母牛-2024年出生", "在群母牛-2023年出生", ...]
        """
        if not slide or not actual_year_labels:
            return

        # 查找所有包含"年出生"的文本框
        year_text_boxes = []
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                try:
                    text = shape.text.strip()
                    if '年出生' in text and '在群母牛-' in text:
                        year_text_boxes.append(shape)
                except Exception:
                    continue

        # 按位置排序（从左到右，从上到下）
        year_text_boxes.sort(key=lambda s: (s.top, s.left))

        # 用实际年份标签替换文本框内容
        for i, (text_box, actual_label) in enumerate(zip(year_text_boxes, actual_year_labels)):
            try:
                if text_box.text_frame and text_box.text_frame.paragraphs:
                    para = text_box.text_frame.paragraphs[0]
                    # 清除所有现有runs，避免文本拼接
                    for run in list(para.runs):
                        run.text = ""
                    # 设置新文本
                    if para.runs:
                        para.runs[0].text = actual_label
                    else:
                        para.text = actual_label
                    logger.debug(f"更新年份文本框 {i+1}: \"{actual_label}\"")
            except Exception as e:
                logger.warning(f"更新年份文本框失败: {e}")

    def _get_excel_sheet(self, sheet_name: str):
        if not self.excel_path:
            logger.warning("未提供Excel路径，无法读取Sheet: %s", sheet_name)
            return None
        try:
            if self._excel_wb is None:
                self._excel_wb = load_workbook(self.excel_path, data_only=True)
            return self._excel_wb[sheet_name]
        except Exception as e:
            logger.warning("读取Excel Sheet失败 %s: %s", sheet_name, e)
            return None

    def _is_empty_value(self, value) -> bool:
        if value is None:
            return True
        if isinstance(value, str) and not value.strip():
            return True
        return False

    def _row_is_empty(self, row: List[str]) -> bool:
        return all(self._is_empty_value(cell) for cell in row)

    def _normalize_label(self, text: Optional[str]) -> str:
        if text is None:
            return ""
        return re.sub(r"\s+", "", str(text))

    def _row_looks_like_header(self, row: List[str]) -> bool:
        for cell in row:
            if self._is_empty_value(cell):
                continue
            text = str(cell)
            for keyword in self.HEADER_KEYWORDS:
                if keyword in text:
                    return True
        return False

    def _find_index_column(self, df: pd.DataFrame) -> Optional[str]:
        """在育种性状明细表中查找育种指数列，兼容多种命名。"""
        if df is None or df.empty:
            return None

        # 统一列名格式
        normalized_pairs: List[Tuple[str, str]] = []
        for col in df.columns:
            if not isinstance(col, str):
                continue
            col_str = col.strip()
            normalized_pairs.append((col, col_str))

        def _normalize(text: str) -> str:
            return re.sub(r"\s+", "", text).lower()

        candidate_exact = [
            "index_score",
            "indexescore",
            "combineindexscore",
            "index得分",
            "指数得分",
            "母牛指数得分",
            "育种指数得分",
            "育种指数",
        ]

        normalized_exact = {_normalize(name): name for name in candidate_exact}

        for col, text in normalized_pairs:
            normalized = _normalize(text)
            if normalized in normalized_exact:
                return col

        candidate_keywords = ["index", "育种指数", "指数"]
        exclude_keywords = ["nm", "tpi", "milk", "fat", "prot", "scc", "ptat", "udc", "flc"]

        for col, text in normalized_pairs:
            normalized = _normalize(text)
            if any(keyword in normalized for keyword in candidate_keywords):
                if any(exclude in normalized for exclude in exclude_keywords):
                    continue
                return col

        for col in df.columns:
            if not isinstance(col, str):
                continue
            lower_col = col.lower()
            if lower_col.endswith("_score") and "index" in lower_col:
                return col

        return None

    def _prepare_block_rows(
        self,
        rows: List[List[str]],
        label: Optional[str],
        remove_label: bool = False,
        drop_header: bool = False,
    ) -> List[List[str]]:
        if not rows:
            return []

        cleaned = [row[:] for row in rows if row is not None]

        while cleaned and self._row_is_empty(cleaned[0]):
            cleaned.pop(0)

        if remove_label and cleaned and label:
            first_cell = cleaned[0][0] if cleaned[0] else ""
            if self._normalize_label(first_cell) == self._normalize_label(label):
                cleaned.pop(0)

        while cleaned and self._row_is_empty(cleaned[0]):
            cleaned.pop(0)
        while cleaned and self._row_is_empty(cleaned[-1]):
            cleaned.pop()

        if drop_header and cleaned and self._row_looks_like_header(cleaned[0]):
            cleaned.pop(0)

        return cleaned

    def _clear_table(self, slide, table_name: str):
        table = self._find_table(slide, table_name)
        if not table:
            return
        for r in range(len(table.rows)):
            for c in range(len(table.columns)):
                self._set_cell_text(table.cell(r, c), "")

    def _delete_table(self, slide, table_name: str):
        """删除指定名称的表格"""
        if not slide:
            return

        # 查找表格对应的shape
        target_shape = None
        for shape in slide.shapes:
            if hasattr(shape, 'table') and shape.table is not None:
                shape_name = getattr(shape, 'name', '')
                if table_name in shape_name or shape_name == table_name:
                    target_shape = shape
                    break

        if target_shape:
            # 从slide中删除shape
            sp = target_shape.element
            sp.getparent().remove(sp)
            logger.debug(f"已删除表格: {table_name}")

    def _set_shape_text(self, shape, text: str):
        if not shape or not getattr(shape, "has_text_frame", False):
            return
        text_frame = shape.text_frame
        if not text_frame.paragraphs:
            para = text_frame.add_paragraph()
        else:
            para = text_frame.paragraphs[0]
        if para.runs:
            para.runs[0].text = text
        else:
            para.text = text
        for extra_run in para.runs[1:]:
            extra_run.text = ""
        for extra_para in text_frame.paragraphs[1:]:
            for extra_run in extra_para.runs:
                extra_run.text = ""

    def _update_birth_textboxes(self, slide, titles: List[str]):
        if not slide:
            return
        text_boxes = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if "在群母牛-" in text and ("年出生" in text or "****" in text):
                    text_boxes.append(shape)
        text_boxes.sort(key=lambda s: (s.top, s.left))
        for idx, shape in enumerate(text_boxes):
            title = titles[idx] if idx < len(titles) else "在群母牛-****年出生"
            self._set_shape_text(shape, title)

    def _update_table_title(self, slide, table_name: str, title: str):
        """更新表格第一行或相关文本框的标题。"""
        if not slide:
            return

        table = self._find_table(slide, table_name)
        if table and table.rows and table.columns:
            cell = table.cell(0, 0)
            tf = cell.text_frame
            if tf and tf.paragraphs:
                para = tf.paragraphs[0]
                if para.runs:
                    para.runs[0].text = title
                else:
                    para.text = title
                for extra_run in para.runs[1:]:
                    extra_run.text = ""
                for extra_para in tf.paragraphs[1:]:
                    for extra_run in extra_para.runs:
                        extra_run.text = ""
            else:
                cell.text = title

        # 如果存在命名包含表格名的文本框，则同步更新
        try:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                shape_name = getattr(shape, "name", "")
                if table_name not in shape_name:
                    continue
                self._set_shape_text(shape, title)
        except Exception:
            pass

    def _extract_block_from_position(self, sheet, start_row: int, start_col: int, width: int = 7, include_title: bool = False):
        """
        从指定位置提取表格块

        Args:
            sheet: Excel worksheet
            start_row: 标签所在行号
            start_col: 起始列号
            width: 列宽
            include_title: 是否包含标题行（标签所在行）
        """
        block_rows: List[List[str]] = []

        # 如果需要包含标题行，则从标签行开始
        row = start_row if include_title else start_row + 1

        while row <= sheet.max_row:
            row_values = []
            for offset in range(width):
                cell_val = sheet.cell(row=row, column=start_col + offset).value
                row_values.append("" if self._is_empty_value(cell_val) else str(cell_val))
            if all(self._is_empty_value(val) for val in row_values):
                break
            block_rows.append(row_values)
            row += 1
        return block_rows

    def _get_distribution_blocks(self, sheet_name: str):
        if sheet_name in self._distribution_cache:
            return self._distribution_cache[sheet_name]

        sheet = self._get_excel_sheet(sheet_name)
        if sheet is None:
            return {}

        label_map: Dict[str, List[Dict]] = {}

        for row in range(1, sheet.max_row + 1):
            for col in range(1, sheet.max_column + 1):
                value = sheet.cell(row=row, column=col).value
                if not isinstance(value, str):
                    continue
                text = value.strip()
                if not text:
                    continue

                # 匹配统计表（在S列，第19列）或五等份表标签（在C列，第3列）
                if ("在群母牛" in text) or text.endswith("统计") or "五等份分析" in text:
                    # 判断是否为五等份表（包括整体页、阶段页、年份页的五等份表）
                    is_quintile = "五等份分析" in text or (
                        "在群母牛-" in text and
                        not text.endswith("统计") and
                        # 匹配所有五等份表：含"组"、"出生"或以"牛"结尾
                        ("组" in text or "出生" in text or text.endswith("牛"))
                    )

                    if is_quintile:
                        # 五等份表：向左扩展1列（包含行标签），宽度8列，不包含标题行
                        start_col = max(1, col - 1)  # 向左1列到B列
                        width = 8  # B-H列共8列
                        include_title = False  # 不包含标题行
                    elif text.endswith("统计"):
                        # 统计表：保持原列，宽度7列，包含标题行
                        start_col = col
                        width = 7
                        include_title = True  # 包含"NM$统计"标题行
                    else:
                        # 其他表：保持原列，宽度7列，不包含标题
                        start_col = col
                        width = 7
                        include_title = False

                    block = self._extract_block_from_position(sheet, row, start_col, width, include_title)
                    if block:
                        label_map.setdefault(text, []).append(
                            {"rows": block, "row_index": row, "col_index": col}
                        )

        self._distribution_cache[sheet_name] = label_map
        return label_map

    def _get_block(self, sheet_name: str, label: str, predicate: Optional[Callable[[List[List[str]]], bool]] = None):
        blocks = self._get_distribution_blocks(sheet_name).get(label, [])
        blocks = sorted(blocks, key=lambda b: b["row_index"])
        for block in blocks:
            rows = block["rows"]
            if not rows:
                continue
            cleaned_rows = self._prepare_block_rows(rows, label, remove_label=False)
            if not cleaned_rows:
                continue

            if predicate is None:
                return cleaned_rows

            predicate_rows = self._prepare_block_rows(
                rows, label, remove_label=True, drop_header=True
            )

            try:
                if predicate(predicate_rows):
                    return cleaned_rows
            except Exception:
                continue
        return None

    def _apply_block_to_table(self, table, block_rows: List[List[str]], font_size: int, start_row: int = 0):
        if not block_rows:
            return

        max_rows = min(len(block_rows), len(table.rows) - start_row)
        max_cols = len(table.columns)
        for r in range(max_rows):
            row_values = block_rows[r]
            for c in range(max_cols):
                value = row_values[c] if c < len(row_values) else ""
                # bold只用于block的第一行（相对于block，不是table）
                self._set_cell_text(table.cell(start_row + r, c), value, font_size=font_size, bold=(r == 0))

        # 清空剩余行（从start_row + max_rows开始）
        for r in range(start_row + max_rows, len(table.rows)):
            for c in range(len(table.columns)):
                self._set_cell_text(table.cell(r, c), "")

    def _fill_table_from_excel(
        self,
        slide,
        table_name: str,
        sheet_name: str,
        label: str,
        font_size: int = 9,
        predicate: Optional[Callable[[List[List[str]]], bool]] = None,
        max_rows: Optional[int] = None,
        prepend_title: bool = False,
        start_row: int = 0,
    ):
        """
        从Excel填充PPT表格

        Args:
            slide: PPT幻灯片对象
            table_name: PPT表格名称
            sheet_name: Excel Sheet名称
            label: Excel标签名称
            font_size: 字体大小
            predicate: 用于选择block的谓词
            max_rows: 最大行数限制（用于截断数据）
            prepend_title: 是否在数据前插入标题行（用于年份表格）
            start_row: 从表格的第几行开始填充（用于保留标题行）
        """
        if slide is None:
            return
        table = self._find_table(slide, table_name)
        if not table:
            logger.warning("未在目标页面找到表格 %s", table_name)
            return
        block = self._get_block(sheet_name, label, predicate)
        if not block:
            logger.warning("Sheet %s 中未找到标签 %s 的数据块", sheet_name, label)
            return

        # 如果指定了max_rows，截断block
        if max_rows is not None and len(block) > max_rows:
            block = block[:max_rows]

        # 如果需要在数据前插入标题行（年份表格）
        if prepend_title and block:
            # 创建标题行：第一列为空，第二列为标题，其余列为空
            title_row = ['', label] + [''] * (len(block[0]) - 2 if block else 6)
            block = [title_row] + block

        self._apply_block_to_table(table, block, font_size=font_size, start_row=start_row)

    def _export_excel_chart_to_image(self, sheet_name: str, image_index: int) -> Optional[str]:
        """
        从Excel导出图表为图片

        Args:
            sheet_name: Excel sheet名称
            image_index: 图表在sheet中的索引（从0开始，对应Excel中的Image 5就是索引4）

        Returns:
            导出的图片路径，失败返回None
        """
        import tempfile
        from PIL import Image

        try:
            if not self.excel_path or not self._excel_wb:
                logger.warning(f"Excel文件未加载，无法导出图表")
                return None

            sheet = self._excel_wb[sheet_name]

            # 获取sheet中的所有图片对象
            if not hasattr(sheet, '_images') or not sheet._images:
                logger.warning(f"Sheet '{sheet_name}' 中没有找到图片对象")
                return None

            if image_index >= len(sheet._images):
                logger.warning(f"Sheet '{sheet_name}' 中图片索引 {image_index} 超出范围")
                return None

            # 获取图片对象
            image_obj = sheet._images[image_index]

            # 导出为临时文件
            temp_dir = tempfile.gettempdir()
            image_path = f"{temp_dir}/excel_chart_{sheet_name}_{image_index}.png"

            # 从图片对象提取数据并保存
            if hasattr(image_obj, '_data'):
                with open(image_path, 'wb') as f:
                    f.write(image_obj._data())
                logger.info(f"成功导出图表: {sheet_name} Image {image_index+1} -> {image_path}")
                return image_path
            else:
                logger.warning(f"无法从图片对象中提取数据")
                return None

        except Exception as e:
            logger.error(f"导出Excel图表失败: {e}", exc_info=True)
            return None

    # ------------------------------------------------------------------ #
    def _fill_nm_normal_distribution_slides(self, data: Dict):
        """
        填充"在群牛遗传分布"下的NM$正态分布相关页面：
        - 在群牛NM$整体正态分布
        - 成母牛/后备牛NM$正态分布
        - 不同阶段牛群NM$正态分布
        - 不同出生年份牛群NM$正态分布

        数据来源：从Excel的"NM$分布分析"sheet读取表格和图表
        """
        sheet_name = "NM$分布分析"
        logger.info(f"开始填充NM$正态分布页面（从Excel: {sheet_name}）")

        # 找到4个幻灯片
        overall_slide = self._find_slide_by_keywords(
            ["在群牛遗传分布", "整体正态分布"], exclude_keywords=["指数"]
        )
        maturity_slide = self._find_slide_by_keywords(
            ["在群牛遗传分布", "成母牛/后备牛", "NM$"], exclude_keywords=["指数"]
        )
        stage_slide = self._find_slide_by_keywords(
            ["在群牛遗传分布", "不同阶段牛群", "NM$"], exclude_keywords=["指数"]
        )
        birth_year_slide = self._find_slide_by_keywords(
            ["在群牛遗传分布", "不同出生年份牛群", "NM$"], exclude_keywords=["指数"]
        )

        # 1) 在群牛NM$整体正态分布 - 索引0
        if overall_slide:
            # 导出并插入Excel图表
            image_path = self._export_excel_chart_to_image(sheet_name, 0)
            if image_path:
                self._update_slide_picture(overall_slide, image_path, "在群牛NM$整体正态分布")

            # 填充表格数据
            self._fill_table_from_excel(
                overall_slide, "表格 2", sheet_name, "NM$统计", font_size=9, max_rows=3
            )
            self._fill_table_from_excel(
                overall_slide, "表格 7", sheet_name, "在群母牛-NM$五等份分析", font_size=12
            )

            # 填充分析文本
            analysis = "分析：NM$（净效益）反映牛只的综合经济价值。整体正态分布图展示了牧场在群母牛NM$的分布特征，分布越集中说明牛群遗传水平越一致，均值越高说明整体经济价值越好。"
            self._update_distribution_analysis_textbox(overall_slide, analysis)

        # 2) 成母牛/后备牛NM$正态分布 - 索引2
        if maturity_slide:
            # 导出并插入Excel图表
            image_path = self._export_excel_chart_to_image(sheet_name, 2)
            if image_path:
                self._update_slide_picture(maturity_slide, image_path, "成母牛/后备牛NM$正态分布")

            # 填充表格数据
            maturity_predicate = lambda rows: rows and "成母牛" in str(rows[0][0])
            self._fill_table_from_excel(
                maturity_slide, "表格 3", sheet_name, "NM$统计",
                font_size=9, predicate=maturity_predicate
            )
            self._fill_table_from_excel(
                maturity_slide, "表格 5", sheet_name, "在群母牛-成母牛组", font_size=12
            )
            self._fill_table_from_excel(
                maturity_slide, "表格 10", sheet_name, "在群母牛-后备牛组", font_size=12
            )

            # 填充分析文本
            analysis = "分析：成母牛与后备牛的NM$分布对比，可以反映牧场育种选配的效果。通常后备牛的NM$均值应高于成母牛，说明遗传进展良好；若后备牛均值偏低，需关注选种选配策略。"
            self._update_distribution_analysis_textbox(maturity_slide, analysis)

        # 3) 不同阶段牛群NM$正态分布 - 索引4
        if stage_slide:
            # 导出并插入Excel图表
            image_path = self._export_excel_chart_to_image(sheet_name, 4)
            if image_path:
                self._update_slide_picture(stage_slide, image_path, "不同阶段牛群NM$正态分布")

            # 填充表格数据
            stage_predicate = lambda rows: rows and "胎" in str(rows[0][0])
            self._fill_table_from_excel(
                stage_slide, "表格 3", sheet_name, "NM$统计",
                font_size=9, predicate=stage_predicate
            )
            self._fill_table_from_excel(
                stage_slide, "表格 14", sheet_name, "在群母牛-2胎及以上组", font_size=12
            )
            self._fill_table_from_excel(
                stage_slide, "表格 16", sheet_name, "在群母牛-1胎组", font_size=12
            )
            self._fill_table_from_excel(
                stage_slide, "表格 18", sheet_name, "在群母牛-12月龄以上0胎牛", font_size=12
            )
            self._fill_table_from_excel(
                stage_slide, "表格 20", sheet_name, "在群母牛-12月龄以下0胎牛", font_size=12
            )

            # 填充分析文本
            analysis = "分析：不同阶段（胎次）牛群的NM$分布反映了各年龄段牛只的遗传水平。通过对比可评估淘汰策略的效果，理想状态下低胎次牛的NM$均值应高于高胎次牛。"
            self._update_distribution_analysis_textbox(stage_slide, analysis)

        # 4) 不同出生年份牛群NM$正态分布 - 索引6
        if birth_year_slide:
            # 导出并插入Excel图表
            image_path = self._export_excel_chart_to_image(sheet_name, 6)
            if image_path:
                self._update_slide_picture(birth_year_slide, image_path, "不同出生年份牛群NM$正态分布")

            # 填充统计表格（使用predicate匹配包含年份信息的数据块）
            nm_year_predicate = lambda rows: rows and any(
                "年出生" in str(row[0]) or (str(row[0]).replace("年", "").isdigit() and len(str(row[0]).replace("年", "")) == 4)
                for row in rows if row and row[0]
            )
            self._fill_table_from_excel(
                birth_year_slide, "表格 9", sheet_name, "NM$统计",
                font_size=9, predicate=nm_year_predicate
            )

            # 填充年份表格
            table_names = ["表格 2", "表格 3", "表格 15"]
            all_blocks = self._get_distribution_blocks(sheet_name)
            year_labels = [label for label in all_blocks.keys() if '在群母牛-' in label and '年出生' in label]
            year_labels_sorted = sorted(year_labels, reverse=True)
            titles = year_labels_sorted[:len(table_names)]

            while len(titles) < len(table_names):
                titles.append("在群母牛-****年出生")

            for table_name, title in zip(table_names, titles):
                if title and "****" not in title:
                    logger.info(f"NM$：填充 {table_name} - {title}")

                    # 替换标题
                    table = self._find_table(birth_year_slide, table_name)
                    if table and len(table.rows) > 0:
                        first_cell = table.cell(0, 0)
                        if "****" in first_cell.text:
                            tf = first_cell.text_frame
                            tf.text = ""
                            para = tf.paragraphs[0]
                            para.text = title
                            para.alignment = PP_ALIGN.CENTER
                            run = para.runs[0]
                            font = run.font
                            font.name = "微软雅黑"
                            font.size = Pt(14)
                            font.bold = True
                            font.color.rgb = RGBColor(50, 140, 207)
                            logger.info(f"  替换标题: {title}")

                    # 从第二行开始填充数据
                    self._fill_table_from_excel(
                        birth_year_slide, table_name, sheet_name, title,
                        font_size=12, start_row=1
                    )
                else:
                    self._delete_table(birth_year_slide, table_name)

            # 填充分析文本
            analysis = "分析：不同出生年份牛群的NM$分布反映了牧场遗传进展趋势。通常近年出生牛的NM$均值应高于早期出生牛，说明选种选配效果良好。"
            self._update_distribution_analysis_textbox(birth_year_slide, analysis)

        logger.info("NM$正态分布页面填充完成")

    # ------------------------------------------------------------------ #
    def _fill_tpi_normal_distribution_slides(self, data: Dict):
        """
        填充TPI相关正态分布页面：
        - 在群牛TPI整体正态分布
        - 成母牛/后备牛TPI正态分布
        - 不同阶段牛群TPI正态分布
        - 不同出生年份牛群TPI正态分布

        数据来源：从Excel的"TPI分布分析"sheet读取表格和图表
        """
        sheet_name = "TPI分布分析"
        logger.info(f"开始填充TPI正态分布页面（从Excel: {sheet_name}）")

        # 找到4个幻灯片
        overall_slide = self._find_slide_by_keywords(
            ["在群牛TPI整体正态分布"], exclude_keywords=["指数"]
        )
        maturity_slide = self._find_slide_by_keywords(
            ["成母牛/后备牛TPI正态分布"], exclude_keywords=["指数"]
        )
        stage_slide = self._find_slide_by_keywords(
            ["不同阶段牛群TPI正态分布"], exclude_keywords=["指数"]
        )
        birth_year_slide = self._find_slide_by_keywords(
            ["不同出生年份牛群TPI正态分布"], exclude_keywords=["指数"]
        )

        # 1) 在群牛TPI整体正态分布 - 索引0
        if overall_slide:
            # 导出并插入Excel图表
            image_path = self._export_excel_chart_to_image(sheet_name, 0)
            if image_path:
                self._update_slide_picture(overall_slide, image_path, "在群牛TPI整体正态分布")

            # 填充表格数据
            self._fill_table_from_excel(
                overall_slide, "表格 2", sheet_name, "TPI统计", font_size=9, max_rows=3
            )
            self._fill_table_from_excel(
                overall_slide, "表格 7", sheet_name, "在群母牛-TPI五等份分析", font_size=12
            )

            # 填充分析文本
            analysis = "分析：TPI（总性能指数）综合反映牛只的生产性能和功能性状。整体正态分布展示了牧场在群母牛TPI的分布特征，均值越高说明牛群整体性能越优秀。"
            self._update_distribution_analysis_textbox(overall_slide, analysis)

        # 2) 成母牛/后备牛TPI正态分布 - 索引2
        if maturity_slide:
            # 导出并插入Excel图表
            image_path = self._export_excel_chart_to_image(sheet_name, 2)
            if image_path:
                self._update_slide_picture(maturity_slide, image_path, "成母牛/后备牛TPI正态分布")

            # 填充表格数据
            maturity_predicate = lambda rows: rows and "成母牛" in str(rows[0][0])
            self._fill_table_from_excel(
                maturity_slide, "表格 3", sheet_name, "在群母牛TPI统计",
                font_size=9, predicate=maturity_predicate
            )
            self._fill_table_from_excel(
                maturity_slide, "表格 5", sheet_name, "在群母牛-成母牛组", font_size=12
            )
            self._fill_table_from_excel(
                maturity_slide, "表格 10", sheet_name, "在群母牛-后备牛组", font_size=12
            )

            # 填充分析文本
            analysis = "分析：成母牛与后备牛的TPI分布对比反映了牧场育种效果。后备牛TPI均值应高于成母牛，说明遗传进展良好；差距越大，遗传进展越明显。"
            self._update_distribution_analysis_textbox(maturity_slide, analysis)

        # 3) 不同阶段牛群TPI正态分布 - 索引4
        if stage_slide:
            # 导出并插入Excel图表
            image_path = self._export_excel_chart_to_image(sheet_name, 4)
            if image_path:
                self._update_slide_picture(stage_slide, image_path, "不同阶段牛群TPI正态分布")

            # 填充表格数据
            stage_predicate = lambda rows: rows and "胎" in str(rows[0][0])
            self._fill_table_from_excel(
                stage_slide, "表格 3", sheet_name, "在群母牛TPI统计",
                font_size=9, predicate=stage_predicate
            )
            self._fill_table_from_excel(
                stage_slide, "表格 14", sheet_name, "在群母牛-2胎及以上组", font_size=12
            )
            self._fill_table_from_excel(
                stage_slide, "表格 16", sheet_name, "在群母牛-1胎组", font_size=12
            )
            self._fill_table_from_excel(
                stage_slide, "表格 18", sheet_name, "在群母牛-12月龄以上0胎牛", font_size=12
            )
            self._fill_table_from_excel(
                stage_slide, "表格 20", sheet_name, "在群母牛-12月龄以下0胎牛", font_size=12
            )

            # 填充分析文本
            analysis = "分析：不同阶段（胎次）牛群的TPI分布反映了各年龄段牛只的综合性能水平。低胎次牛TPI均值高于高胎次牛，说明淘汰策略有效，遗传水平逐步提升。"
            self._update_distribution_analysis_textbox(stage_slide, analysis)

        # 4) 不同出生年份牛群TPI正态分布 - 索引6
        if birth_year_slide:
            # 导出并插入Excel图表
            image_path = self._export_excel_chart_to_image(sheet_name, 6)
            if image_path:
                self._update_slide_picture(birth_year_slide, image_path, "不同出生年份牛群TPI正态分布")

            # 填充统计表格（使用predicate匹配包含年份信息的数据块）
            tpi_year_predicate = lambda rows: rows and any(
                "年出生" in str(row[0]) or (str(row[0]).replace("年", "").isdigit() and len(str(row[0]).replace("年", "")) == 4)
                for row in rows if row and row[0]
            )
            self._fill_table_from_excel(
                birth_year_slide, "表格 9", sheet_name, "在群母牛TPI统计",
                font_size=9, predicate=tpi_year_predicate
            )

            # 填充年份表格
            table_names = ["表格 3", "表格 12", "表格 15"]
            all_blocks = self._get_distribution_blocks(sheet_name)
            year_labels = [label for label in all_blocks.keys() if '在群母牛-' in label and '年出生' in label]
            year_labels_sorted = sorted(year_labels, reverse=True)
            titles = year_labels_sorted[:len(table_names)]

            while len(titles) < len(table_names):
                titles.append("在群母牛-****年出生")

            for table_name, title in zip(table_names, titles):
                if title and "****" not in title:
                    logger.info(f"TPI：填充 {table_name} - {title}")

                    # 替换标题
                    table = self._find_table(birth_year_slide, table_name)
                    if table and len(table.rows) > 0:
                        first_cell = table.cell(0, 0)
                        if "****" in first_cell.text:
                            tf = first_cell.text_frame
                            tf.text = ""
                            para = tf.paragraphs[0]
                            para.text = title
                            para.alignment = PP_ALIGN.CENTER
                            run = para.runs[0]
                            font = run.font
                            font.name = "微软雅黑"
                            font.size = Pt(14)
                            font.bold = True
                            font.color.rgb = RGBColor(50, 140, 207)
                            logger.info(f"  替换标题: {title}")

                    # 从第二行开始填充数据
                    self._fill_table_from_excel(
                        birth_year_slide, table_name, sheet_name, title,
                        font_size=12, start_row=1
                    )
                else:
                    self._delete_table(birth_year_slide, table_name)

            # 填充分析文本
            analysis = "分析：不同出生年份牛群的TPI分布反映了牧场遗传进展趋势。近年出生牛的TPI均值应高于早期出生牛，体现选种选配效果。"
            self._update_distribution_analysis_textbox(birth_year_slide, analysis)

        logger.info("TPI正态分布页面填充完成")

    # ------------------------------------------------------------------ #
    def _fill_index_normal_distribution_slides(self, data: Dict):
        """
        填充育种指数相关正态分布页面：
        - 在群牛育种指数整体正态分布
        - 成母牛/后备牛育种指数正态分布
        - 不同阶段牛群育种指数正态分布
        - 不同出生年份牛群育种指数正态分布

        数据来源：从Excel的"育种指数分布分析"sheet读取表格和图表
        """
        sheet_name = "育种指数分布分析"
        logger.info(f"开始填充育种指数正态分布页面（从Excel: {sheet_name}）")

        # 找到4个幻灯片
        overall_slide = self._find_slide_by_keywords(
            ["在群牛育种指数整体正态分布"], exclude_keywords=[]
        )
        maturity_slide = self._find_slide_by_keywords(
            ["成母牛/后备牛育种指数正态分布"], exclude_keywords=[]
        )
        stage_slide = self._find_slide_by_keywords(
            ["不同阶段牛群育种指数正态分布"], exclude_keywords=[]
        )
        birth_year_slide = self._find_slide_by_keywords(
            ["不同出生年份牛群育种指数正态分布"], exclude_keywords=[]
        )

        # 1) 在群牛育种指数整体正态分布 - 索引0
        if overall_slide:
            # 导出并插入Excel图表
            image_path = self._export_excel_chart_to_image(sheet_name, 0)
            if image_path:
                self._update_slide_picture(overall_slide, image_path, "在群牛育种指数整体正态分布")

            # 填充表格数据
            self._fill_table_from_excel(
                overall_slide, "表格 2", sheet_name, "育种指数统计", font_size=9, max_rows=3
            )
            self._fill_table_from_excel(
                overall_slide, "表格 7", sheet_name, "在群母牛-育种指数五等份分析", font_size=12
            )

            # 填充分析文本
            analysis = "分析：育种指数是牧场自定义的综合选择指数，反映牛只在牧场育种目标下的综合遗传价值。整体正态分布展示了牧场在群母牛育种指数的分布特征。"
            self._update_distribution_analysis_textbox(overall_slide, analysis)

        # 2) 成母牛/后备牛育种指数正态分布 - 索引2
        if maturity_slide:
            # 导出并插入Excel图表
            image_path = self._export_excel_chart_to_image(sheet_name, 2)
            if image_path:
                self._update_slide_picture(maturity_slide, image_path, "成母牛/后备牛育种指数正态分布")

            # 填充表格数据
            index_maturity_predicate = lambda rows: rows and "成母牛" in str(rows[0][0])
            self._fill_table_from_excel(
                maturity_slide, "表格 3", sheet_name, "在群母牛index统计",
                font_size=9, predicate=index_maturity_predicate
            )
            self._fill_table_from_excel(
                maturity_slide, "表格 5", sheet_name, "在群母牛-成母牛组", font_size=12
            )
            self._fill_table_from_excel(
                maturity_slide, "表格 10", sheet_name, "在群母牛-后备牛组", font_size=12
            )

            # 填充分析文本
            analysis = "分析：成母牛与后备牛的育种指数分布对比反映了牧场选种选配效果。后备牛育种指数均值应高于成母牛，说明遗传改良方向正确。"
            self._update_distribution_analysis_textbox(maturity_slide, analysis)

        # 3) 不同阶段牛群育种指数正态分布 - 索引4
        if stage_slide:
            # 导出并插入Excel图表
            image_path = self._export_excel_chart_to_image(sheet_name, 4)
            if image_path:
                self._update_slide_picture(stage_slide, image_path, "不同阶段牛群育种指数正态分布")

            # 填充表格数据
            index_stage_predicate = lambda rows: rows and "胎" in str(rows[0][0])
            self._fill_table_from_excel(
                stage_slide, "表格 5", sheet_name, "在群母牛index统计",
                font_size=9, predicate=index_stage_predicate
            )
            for table_name, label in [
                ("表格 14", "在群母牛-2胎及以上组"),
                ("表格 16", "在群母牛-1胎组"),
                ("表格 18", "在群母牛-12月龄以上0胎牛"),
                ("表格 20", "在群母牛-12月龄以下0胎牛"),
            ]:
                self._fill_table_from_excel(
                    stage_slide, table_name, sheet_name, label, font_size=12
                )

            # 填充分析文本
            analysis = "分析：不同阶段（胎次）牛群的育种指数分布反映了各年龄段牛只在牧场育种目标下的价值。低胎次牛育种指数高于高胎次牛，说明牧场遗传改良有效。"
            self._update_distribution_analysis_textbox(stage_slide, analysis)

        # 4) 不同出生年份牛群育种指数正态分布 - 索引6
        if birth_year_slide:
            # 导出并插入Excel图表
            image_path = self._export_excel_chart_to_image(sheet_name, 6)
            if image_path:
                self._update_slide_picture(birth_year_slide, image_path, "不同出生年份牛群育种指数正态分布")

            # 填充统计表格（使用predicate匹配包含年份信息的数据块）
            index_year_predicate = lambda rows: rows and any(
                "年出生" in str(row[0]) or (str(row[0]).replace("年", "").isdigit() and len(str(row[0]).replace("年", "")) == 4)
                for row in rows if row and row[0]
            )
            self._fill_table_from_excel(
                birth_year_slide, "表格 9", sheet_name, "在群母牛index统计",
                font_size=9, predicate=index_year_predicate
            )

            # 填充年份表格
            table_names = ["表格 11", "表格 12", "表格 15"]
            all_blocks = self._get_distribution_blocks(sheet_name)
            logger.info(f"育种指数分布分析：找到 {len(all_blocks)} 个标签")

            year_labels = [label for label in all_blocks.keys() if '在群母牛-' in label and '年出生' in label]
            logger.info(f"育种指数分布分析：找到 {len(year_labels)} 个年份标签: {year_labels[:3]}...")

            year_labels_sorted = sorted(year_labels, reverse=True)
            logger.info(f"育种指数分布分析：排序后前3个: {year_labels_sorted[:3]}")

            titles = year_labels_sorted[:len(table_names)]

            while len(titles) < len(table_names):
                titles.append("在群母牛-****年出生")

            for table_name, title in zip(table_names, titles):
                if title and "****" not in title:
                    logger.info(f"育种指数：填充 {table_name} - {title}")

                    # 替换标题
                    table = self._find_table(birth_year_slide, table_name)
                    if table and len(table.rows) > 0:
                        first_cell = table.cell(0, 0)
                        if "****" in first_cell.text:
                            tf = first_cell.text_frame
                            tf.text = ""
                            para = tf.paragraphs[0]
                            para.text = title
                            para.alignment = PP_ALIGN.CENTER
                            run = para.runs[0]
                            font = run.font
                            font.name = "微软雅黑"
                            font.size = Pt(14)
                            font.bold = True
                            font.color.rgb = RGBColor(50, 140, 207)
                            logger.info(f"  替换标题: {title}")

                    # 从第二行开始填充数据
                    self._fill_table_from_excel(
                        birth_year_slide, table_name, sheet_name, title,
                        font_size=12, start_row=1
                    )
                else:
                    logger.info(f"育种指数：删除占位表格 {table_name} (title={title})")
                    self._delete_table(birth_year_slide, table_name)

            # 填充分析文本
            analysis = "分析：不同出生年份牛群的育种指数分布反映了牧场遗传进展趋势。近年出生牛的育种指数均值应高于早期出生牛，体现育种目标的持续实现。"
            self._update_distribution_analysis_textbox(birth_year_slide, analysis)

        logger.info("育种指数正态分布页面填充完成")

    # ------------------------------------------------------------------ #
    def _fill_yearly_traits_charts(self, data: Dict):
        """
        填充 Slide 11 及后续“在群牛关键性状年均进展”折线图（动态）

        - 使用《年份汇总与性状进展》Sheet 中“在群母牛年份汇总”区域作为当前牧场曲线
        - 使用同一Sheet右侧“对比数据（用于图表）”区域作为对比牧场/外部参考曲线
        - 根据“在群母牛年份汇总”区域中的所有“平均*”性状列动态生成折线图页，性状数量不设上限
        """
        df = data.get("traits_yearly")
        if df is None or df.empty:
            logger.warning("traits_yearly 数据为空，跳过年度性状折线图填充")
            return

        try:
            columns = list(df.columns)
            present_end_idx = columns.index("全部母牛年份汇总")
        except ValueError:
            logger.warning("traits_yearly 中未找到“全部母牛年份汇总”列，无法解析在群母牛区域")
            return

        # 在群母牛年份汇总区域（左侧块）
        present_block = df.iloc[:, :present_end_idx]
        if present_block.empty:
            logger.warning("在群母牛年份汇总区域为空，跳过折线图填充")
            return

        # 第 1 行是表头，后续行为数据
        header = present_block.iloc[0].tolist()
        present = present_block.iloc[1:].copy()
        present.columns = header

        if "出生年份" not in present.columns:
            logger.warning("在群母牛年份汇总区域缺少“出生年份”列，无法绘制年度曲线")
            return

        # 去掉空行
        present = present[present["出生年份"].notna()].copy()

        # 性状列：所有以“平均”开头的列，顺序与Excel保持一致
        trait_columns: List[str] = [
            col for col in present.columns if isinstance(col, str) and col.startswith("平均")
        ]
        if not trait_columns:
            logger.warning("在群母牛年份汇总中未找到任何“平均*”性状列，跳过折线图填充")
            return

        # 右侧“对比数据（用于图表）”区域（如果存在）
        comparison_data = None
        try:
            comp_start_idx = columns.index("对比数据（用于图表）")
        except ValueError:
            logger.info("traits_yearly 中未找到“对比数据（用于图表）”列，仅绘制当前牧场曲线")
        else:
            comp_block = df.iloc[:, comp_start_idx:]
            if not comp_block.empty:
                comp_header = comp_block.iloc[0].tolist()
                comparison_data = comp_block.iloc[1:].copy()
                comparison_data.columns = comp_header

        # ==== 步骤1：收集模板中的“在群牛关键性状年均进展”折线图页 ====
        trait_count = len(trait_columns)

        chart_slots: List[int] = []
        for slide_idx in range(self.YEARLY_SLIDE_INDEX + 1, len(self.prs.slides)):
            slide = self.prs.slides[slide_idx]

            has_title = any(
                getattr(shape, "text", "") and "在群牛关键性状年均进展" in shape.text
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
            has_chart = any(getattr(shape, "has_chart", False) for shape in slide.shapes)

            if has_title and has_chart:
                chart_slots.append(slide_idx)

        if not chart_slots:
            logger.warning("模板中未找到“在群牛关键性状年均进展”折线图页，跳过折线图填充")
            return

        logger.info(
            "年度性状折线图：发现 %s 个性状，模板中现有 %s 个折线图页",
            trait_count,
            len(chart_slots),
        )

        # 如果模板页数不足，只填充前面可用的页，不再自动复制新页
        slots_to_use = min(trait_count, len(chart_slots))
        if trait_count > len(chart_slots):
            logger.warning(
                "性状数量(%s)超过模板页数(%s)，仅填充前 %s 个性状",
                trait_count,
                len(chart_slots),
                slots_to_use,
            )

        # ==== 步骤2：按顺序将性状数据填入每一页 ====
        for i in range(slots_to_use):
            trait_col = trait_columns[i]
            trait_name = trait_col.replace("平均", "")
            slide_idx = chart_slots[i]
            slide = self.prs.slides[slide_idx]

            # 找到该页中的图表
            chart_shape = next(
                (shape for shape in slide.shapes if getattr(shape, "has_chart", False)),
                None,
            )
            if not chart_shape:
                logger.warning("Slide %s 未找到图表，跳过该页", slide_idx + 1)
                continue

            chart = chart_shape.chart

            # 更新Y轴标题为当前性状名
            try:
                if chart.value_axis is not None and chart.value_axis.has_title:
                    tf = chart.value_axis.axis_title.text_frame
                    if tf.paragraphs:
                        para = tf.paragraphs[0]
                        if para.runs:
                            para.runs[0].text = trait_name
                            for r in para.runs[1:]:
                                r.text = ""
                        else:
                            run = para.add_run()
                            run.text = trait_name
                    else:
                        para = tf.add_paragraph()
                        run = para.add_run()
                        run.text = trait_name
            except Exception as e:
                logger.warning("更新 Slide %s Y轴标题失败: %s", slide_idx + 1, e)

            # 替换图表数据和小标题
            self._update_trait_chart(chart, present, comparison_data, trait_col, trait_name)

            # 生成并更新分析文本
            analysis_text = self._generate_trait_analysis(trait_name, present, trait_col)
            self._update_trait_analysis_textbox(slide, analysis_text)

        # ==== 步骤3：如果模板有多余的页（>性状数），整页删除 ====
        if len(chart_slots) > trait_count:
            extra_indices = chart_slots[trait_count:]
            for slide_idx in sorted(extra_indices, reverse=True):
                self._delete_slide(slide_idx)
                logger.info("已删除多余的年度性状折线图页: Slide %s", slide_idx + 1)

    def _update_trait_chart(
        self,
        chart,
        present_df: pd.DataFrame,
        comparison_df: Optional[pd.DataFrame],
        trait_col: str,
        trait_name: str,
    ):
        """
        根据在群母牛 + 对比数据更新单个折线图
        """
        # 1. 当前牧场数据：排除“总计”和“对比”行
        valid_present = present_df[
            ~present_df["出生年份"].astype(str).str.contains("总计|对比", na=False)
        ].copy()
        if valid_present.empty:
            logger.warning("性状 %s 没有有效年份数据，跳过图表更新", trait_name)
            return

        # 构造清洗后的年份（与Excel中“对比数据（用于图表）”保持一致）
        years_raw = valid_present["出生年份"].astype(str).tolist()
        clean_years = []
        for y in years_raw:
            # 例如 “2020年及以前” -> “2020年”
            cleaned = y
            if "及以前" in y:
                import re

                match = re.search(r"(\\d{4})", y)
                if match:
                    cleaned = f"{match.group(1)}年"
            clean_years.append(cleaned)

        # 当前牧场该性状的数值
        current_values = pd.to_numeric(valid_present[trait_col], errors="coerce").tolist()

        # 2. 对比系列（来自“对比数据（用于图表）”区域）
        comparison_series = []
        if comparison_df is not None and "年份" in comparison_df.columns:
            comp_valid = comparison_df[comparison_df["年份"].notna()].copy()
            if not comp_valid.empty:
                comp_valid["年份"] = comp_valid["年份"].astype(str)
                comp_valid = comp_valid.set_index("年份")

                for col in comp_valid.columns:
                    if not col or col == "年份":
                        continue
                    if not col.endswith(trait_col):
                        continue

                    # 列名形如“对比牧场二_平均NM$”，截取系列名
                    series_name = col[: -len(trait_col)].rstrip("_")

                    values = []
                    has_data = False
                    for year in clean_years:
                        if year in comp_valid.index:
                            v = comp_valid.at[year, col]
                        else:
                            v = None
                        if pd.notna(v):
                            has_data = True
                            values.append(float(v))
                        else:
                            values.append(None)

                    if has_data:
                        comparison_series.append((series_name, values))

        # 3. 组装ChartData并替换模板中的占位数据
        chart_data = CategoryChartData()
        chart_data.categories = clean_years

        # 主系列：当前牧场
        chart_data.add_series("在群母牛", current_values)

        # 对比系列：对比牧场 / 外部参考
        for series_name, values in comparison_series:
            chart_data.add_series(series_name, values)

        # 用新的数据替换图表，样式（颜色/线型/标记）沿用PPT模板默认设置
        chart.replace_data(chart_data)

        # 更新图表标题为“{性状}性状进展”
        try:
            if not chart.has_title:
                chart.has_title = True

            tf = chart.chart_title.text_frame
            if tf.paragraphs:
                para = tf.paragraphs[0]
                if para.runs:
                    para.runs[0].text = f"{trait_name}性状进展"
                    for r in para.runs[1:]:
                        r.text = ""
                else:
                    run = para.add_run()
                    run.text = f"{trait_name}性状进展"
            else:
                para = tf.add_paragraph()
                run = para.add_run()
                run.text = f"{trait_name}性状进展"
        except Exception as e:
            logger.warning("更新图表标题失败（%s）: %s", trait_name, e)

        logger.info(
            "✓ 更新年度性状折线图: %s，系列数=%s",
            trait_name,
            1 + len(comparison_series),
        )

    def _generate_trait_analysis(
        self,
        trait_name: str,
        present_df: pd.DataFrame,
        trait_col: str,
    ) -> str:
        """
        为单个性状生成分析文本

        Args:
            trait_name: 性状名称（如 NM$, TPI, MILK 等）
            present_df: 在群母牛年份汇总数据
            trait_col: 性状列名（如 平均NM$）

        Returns:
            分析文本字符串
        """
        parts = []

        # 1. 性状定义
        trait_def = self.TRAIT_DEFINITIONS.get(trait_name)
        if trait_def:
            cn_name, description = trait_def
            parts.append(f"{trait_name}（{cn_name}）：{description}。\n")  # 定义后换行
        else:
            parts.append(f"{trait_name}性状：\n")

        # 2. 数据趋势分析
        valid_data = present_df[
            ~present_df["出生年份"].astype(str).str.contains("总计|对比", na=False)
        ].copy()

        if valid_data.empty or trait_col not in valid_data.columns:
            parts.append(f"牧场{trait_name}性状数据暂无足够年份记录。")
            return "".join(parts)

        # 提取年份和数值
        years = valid_data["出生年份"].astype(str).tolist()
        values = pd.to_numeric(valid_data[trait_col], errors="coerce").tolist()

        # 过滤有效数据点
        valid_pairs = [(y, v) for y, v in zip(years, values) if pd.notna(v)]
        if len(valid_pairs) < 2:
            parts.append(f"牧场{trait_name}性状数据点不足，无法分析趋势。")
            return "".join(parts)

        years_clean = [p[0] for p in valid_pairs]
        values_clean = [p[1] for p in valid_pairs]

        # 计算趋势
        first_val = values_clean[0]
        last_val = values_clean[-1]
        change = last_val - first_val
        avg_val = sum(values_clean) / len(values_clean)

        # 计算年均变化（简单线性趋势）
        n = len(values_clean)
        if n >= 2:
            avg_change_per_year = change / (n - 1)
        else:
            avg_change_per_year = 0

        # 3. 趋势描述
        first_year = years_clean[0].replace("年", "").replace("及以前", "")
        last_year = years_clean[-1].replace("年", "").replace("及以前", "")

        # 根据性状类型判断趋势好坏
        # SCS、RFI 是负向指标（数值越低越好），其他多为正向指标
        is_negative_trait = trait_name in ("SCS", "RFI", "GL")

        if abs(change) < 0.01 * abs(avg_val) if avg_val != 0 else abs(change) < 0.1:
            trend_desc = "基本持平"
            trend_eval = "保持稳定"
        elif change > 0:
            if is_negative_trait:
                trend_desc = "呈上升趋势"
                trend_eval = "需关注改善"
            else:
                trend_desc = "呈上升趋势"
                trend_eval = "呈现良好的遗传进展"
        else:
            if is_negative_trait:
                trend_desc = "呈下降趋势"
                trend_eval = "呈现良好的改善趋势"
            else:
                trend_desc = "呈下降趋势"
                trend_eval = "需关注提升"

        parts.append(
            f"从{first_year}年到{last_year}年，牧场在群母牛{trait_name}性状{trend_desc}，"
            f"平均值从{first_val:.2f}变化到{last_val:.2f}，"
            f"年均变化约{avg_change_per_year:+.2f}，{trend_eval}。"
        )

        # 4. 当前水平评价
        parts.append(f"当前牛群{trait_name}平均值为{last_val:.2f}。")

        return "".join(parts)

    def _update_trait_analysis_textbox(self, slide, analysis_text: str):
        """
        更新性状进展页的分析文本框（文本框 13）

        Args:
            slide: 目标幻灯片
            analysis_text: 分析文本
        """
        text_box = None
        for shape in slide.shapes:
            if shape.name == "文本框 13":
                text_box = shape
                break

        if not text_box:
            # 尝试查找包含"分析"的文本框
            for shape in slide.shapes:
                if shape.has_text_frame and "分析" in shape.text and len(shape.text) < 100:
                    text_box = shape
                    break

        if not text_box:
            logger.warning("未找到分析文本框（文本框 13），跳过分析文本更新")
            return

        if not text_box.has_text_frame:
            logger.warning("文本框 13 没有文本框架")
            return

        # 设置文本，保持模板样式
        tf = text_box.text_frame
        if not tf.paragraphs:
            para = tf.add_paragraph()
        else:
            para = tf.paragraphs[0]

        # 清空段落并添加新run
        para.clear()
        run = para.add_run()

        # 避免重复"分析："前缀
        if analysis_text.startswith("分析：") or analysis_text.startswith("分析:"):
            run.text = analysis_text
        else:
            run.text = f"分析：{analysis_text}"

        # 设置字体：微软雅黑15号非加粗
        run.font.name = FONT_NAME_CN
        run.font.size = Pt(15)
        run.font.bold = False

        # 清理其他段落
        for extra_para in tf.paragraphs[1:]:
            for extra_run in extra_para.runs:
                extra_run.text = ""

        logger.debug("✓ 分析文本框更新完成")

    def _generate_distribution_analysis(
        self,
        index_name: str,
        distribution_data: list,
        total_count: int = 0,
    ) -> str:
        """
        为分布页面生成分析文本

        Args:
            index_name: 指数名称（NM$, TPI, 育种指数等）
            distribution_data: 分布数据列表 [{"interval": "0-100", "count": 50, "percent": 10.0}, ...]
            total_count: 总头数

        Returns:
            分析文本字符串
        """
        parts = []

        # 1. 指数定义（简短版）
        index_def = self.TRAIT_DEFINITIONS.get(index_name)
        if index_def:
            cn_name, _ = index_def
            parts.append(f"{index_name}（{cn_name}）")
        else:
            parts.append(f"{index_name}指数")

        if not distribution_data:
            parts.append("分布数据暂无。")
            return "".join(parts)

        # 2. 计算关键统计
        # 找出占比最高的区间
        max_interval = max(distribution_data, key=lambda x: x.get("percent", 0))
        max_interval_name = max_interval.get("interval", "")
        max_interval_pct = max_interval.get("percent", 0)

        # 计算低于0/负值的占比（如果有）
        negative_pct = 0
        low_pct = 0  # 低于某个阈值的占比
        high_pct = 0  # 高于某个阈值的占比

        for item in distribution_data:
            interval = str(item.get("interval", ""))
            pct = item.get("percent", 0) or 0

            # 检查是否为负值区间
            if interval.startswith("-") or interval.startswith("<0") or "以下" in interval:
                if "0" in interval or interval.startswith("-"):
                    negative_pct += pct

            # 根据指数类型判断低/高阈值
            if index_name == "NM$":
                if any(x in interval for x in ["-", "<0", "以下"]) and "0" in interval:
                    low_pct += pct
                elif any(x in interval for x in ["500", "600", "700", "800"]):
                    high_pct += pct
            elif index_name == "TPI":
                if any(x in interval for x in ["<2000", "2000以下", "1"]):
                    low_pct += pct
                elif any(x in interval for x in ["2800", "2900", "3000"]):
                    high_pct += pct

        # 3. 生成分析文本
        total_str = f"共{total_count}头" if total_count > 0 else ""
        parts.append(f"分布分析：牧场在群母牛{total_str}，")

        parts.append(f"其中{max_interval_name}区间占比最高（{max_interval_pct:.1f}%）。")

        # 评价
        if negative_pct > 20:
            parts.append(f"低于0的牛只占{negative_pct:.1f}%，建议关注低遗传水平牛只的淘汰或改良。")
        elif negative_pct > 0:
            parts.append(f"低于0的牛只占{negative_pct:.1f}%，整体遗传水平较好。")

        if high_pct > 30:
            parts.append(f"高遗传水平牛只占比{high_pct:.1f}%，牛群遗传基础优秀。")

        return "".join(parts)

    def _generate_normal_distribution_analysis(
        self,
        index_name: str,
        group_name: str,
        stats: Dict,
    ) -> str:
        """
        为正态分布页面生成分析文本

        Args:
            index_name: 指数名称（NM$, TPI等）
            group_name: 分组名称（整体、成母牛、后备牛、不同胎次等）
            stats: 统计数据 {"mean": 均值, "std": 标准差, "count": 头数, "median": 中位数}

        Returns:
            分析文本字符串
        """
        parts = []

        mean = stats.get("mean", 0)
        std = stats.get("std", 0)
        count = stats.get("count", 0)
        median = stats.get("median", 0)

        # 指数简介
        index_def = self.TRAIT_DEFINITIONS.get(index_name)
        if index_def:
            cn_name, _ = index_def
            parts.append(f"{index_name}（{cn_name}）")
        else:
            parts.append(f"{index_name}")

        parts.append(f"{group_name}分布分析：")

        if count > 0:
            parts.append(f"样本量{count}头，均值{mean:.1f}，标准差{std:.1f}，中位数{median:.1f}。")

            # 变异系数评价
            if mean != 0:
                cv = abs(std / mean) * 100
                if cv < 20:
                    parts.append("牛群遗传水平较为一致。")
                elif cv < 40:
                    parts.append("牛群遗传水平存在一定分化。")
                else:
                    parts.append("牛群遗传水平差异较大，建议针对性改良。")

            # 均值评价（基于指数类型）
            if index_name == "NM$":
                if mean > 400:
                    parts.append(f"平均NM$为{mean:.0f}，处于较高水平。")
                elif mean > 200:
                    parts.append(f"平均NM$为{mean:.0f}，处于中等水平。")
                else:
                    parts.append(f"平均NM$为{mean:.0f}，有较大提升空间。")
            elif index_name == "TPI":
                if mean > 2600:
                    parts.append(f"平均TPI为{mean:.0f}，处于较高水平。")
                elif mean > 2400:
                    parts.append(f"平均TPI为{mean:.0f}，处于中等水平。")
                else:
                    parts.append(f"平均TPI为{mean:.0f}，有较大提升空间。")

        else:
            parts.append("暂无足够数据进行分析。")

        return "".join(parts)

    def _update_distribution_analysis_textbox(self, slide, analysis_text: str):
        """
        更新分布页的分析文本框

        Args:
            slide: 目标幻灯片
            analysis_text: 分析文本
        """
        text_box = None

        # 方法1: 查找以"分析："开头的文本框（最可靠）
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip() if shape.text else ""
                if text.startswith("分析：") or text.startswith("分析:"):
                    text_box = shape
                    logger.debug(f"分布页找到分析文本框（通过'分析：'前缀）: {shape.name}")
                    break

        # 方法2: 按常见名称查找
        if not text_box:
            for name in ["文本框 13", "文本框 14", "文本框 15", "文本框 16"]:
                for shape in slide.shapes:
                    if shape.name == name and shape.has_text_frame:
                        text_box = shape
                        logger.debug(f"分布页找到分析文本框（通过名称）: {name}")
                        break
                if text_box:
                    break

        # 方法3: 查找页面下部的文本框（通常分析框在页面下方）
        if not text_box:
            candidates = []
            for shape in slide.shapes:
                if shape.has_text_frame and not getattr(shape, "has_table", False):
                    # 排除图表标题（通常在上方）
                    if hasattr(shape, "top") and shape.top > 4000000:  # 约4cm以下
                        text = shape.text.strip() if shape.text else ""
                        # 排除小标题和数据标签
                        if "分布" not in text and len(text) < 300:
                            candidates.append((shape.top, shape))

            if candidates:
                # 选择位置最靠下的
                candidates.sort(key=lambda x: x[0], reverse=True)
                text_box = candidates[0][1]
                logger.debug(f"分布页找到分析文本框（通过位置）: {text_box.name}")

        if not text_box:
            logger.warning("分布页未找到分析文本框，跳过")
            return

        # 设置文本（保留原格式）
        tf = text_box.text_frame
        if not tf.paragraphs:
            para = tf.add_paragraph()
        else:
            para = tf.paragraphs[0]

        # 保存原始格式
        font_size = None
        font_name = None
        if para.runs:
            run = para.runs[0]
            font_size = run.font.size
            font_name = run.font.name

        # 清空并设置新文本（避免重复"分析："前缀）
        para.clear()
        run = para.add_run()
        if analysis_text.startswith("分析：") or analysis_text.startswith("分析:"):
            run.text = analysis_text
        else:
            run.text = f"分析：{analysis_text}"

        # 设置字体：微软雅黑15号非加粗
        run.font.name = FONT_NAME_CN
        run.font.size = Pt(15)
        run.font.bold = False

        logger.info(f"✓ 分布页分析文本已更新: {text_box.name}")

    # ------------------------------------------------------------------ #
    def _delete_slide(self, index: int):
        """
        从演示文稿中删除指定索引的幻灯片

        python-pptx 未公开删除API，这里使用底层 _sldIdLst 操作，仅影响当前生成的PPT，不修改模板文件。
        """
        try:
            slides = self.prs.slides
            slide_id_list = slides._sldIdLst  # type: ignore[attr-defined]
            slide_id_list.remove(slide_id_list[index])
        except Exception as e:
            logger.error("删除Slide %s失败: %s", index + 1, e)

    def _duplicate_progress_slide(self, template_index: int, insert_index: int):
        """
        复制一页“在群牛关键性状年均进展”折线图页：
        - 完整复制该页所有内容（包括图表和文本）
        - 不修改任何数据和样式
        - 将新页插入到指定位置，以保持章节顺序
        """
        try:
            # 1) 复制模板页的所有shape（相当于PPT里的复制/粘贴）
            template_slide = self.prs.slides[template_index]
            new_slide = self.prs.slides.add_slide(template_slide.slide_layout)

            for shape in template_slide.shapes:
                new_el = deepcopy(shape.element)
                new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

            # 2) 将新幻灯片移动到指定位置（紧跟在Part4既有页面之后）
            try:
                slides = self.prs.slides
                slide_id_list = slides._sldIdLst  # type: ignore[attr-defined]
                new_id = slide_id_list[-1]
                slide_id_list.remove(new_id)
                slide_id_list.insert(insert_index, new_id)
            except Exception as e:
                logger.error("重排新增Slide位置失败: %s", e)

            return new_slide

        except Exception as e:
            logger.error("复制年度性状折线图页失败: %s", e, exc_info=True)
            return None

    # ------------------------------------------------------------------ #
    @staticmethod
    def _set_cell_text(cell, text: str, font_size: Optional[int] = None, bold: Optional[bool] = None):
        """更新单元格文字，并可选设置字体样式"""
        tf = cell.text_frame
        if not tf.paragraphs:
            para = tf.add_paragraph()
        else:
            para = tf.paragraphs[0]
        if para.runs:
            run = para.runs[0]
        else:
            run = para.add_run()
        run.text = text
        if font_size is not None or bold is not None:
            font = run.font
            font.name = "微软雅黑"
            if font_size is not None:
                font.size = Pt(font_size)
            if bold is not None:
                font.bold = bold
        # 清理其它 run/段落中的旧文本
        for extra_run in para.runs[1:]:
            extra_run.text = ""
        for extra_para in tf.paragraphs[1:]:
            for extra_run in extra_para.runs:
                extra_run.text = ""
