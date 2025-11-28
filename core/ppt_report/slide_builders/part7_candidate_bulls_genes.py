"""
Part 7: 备选公牛-隐性基因分析构建器
对应PPT模版的"备选公牛-隐性基因分析"页面（第126-148页，共23页模板）
"""

import logging
from typing import Dict, Any, List, Optional
from pathlib import Path

import pandas as pd
from openpyxl import load_workbook
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ..base_builder import BaseSlideBuilder

logger = logging.getLogger(__name__)


class Part7CandidateBullsGenesBuilder(BaseSlideBuilder):
    """Part 7: 备选公牛-隐性基因分析构建器"""

    # 16种隐性基因列表
    RECESSIVE_GENES = [
        "HH1", "HH2", "HH3", "HH4", "HH5", "HH6",
        "BLAD", "Brachyspina", "CVM", "Cholesterol deficiency",
        "Chondrodysplasia", "Citrullinemia", "DUMPS",
        "Factor XI", "MW", "Mulefoot"
    ]

    def __init__(self, prs, farm_name: str):
        super().__init__(prs, None)
        self.farm_name = farm_name

    def build(self, data: Dict[str, Any]):
        """
        构建 Part 7: 备选公牛-隐性基因分析

        Args:
            data: 包含 'excel_path' 的数据字典
        """
        logger.info("=" * 60)
        logger.info("开始构建Part 7: 备选公牛-隐性基因分析")
        logger.info("=" * 60)

        # 获取Excel路径
        excel_path = data.get("excel_path")
        if not excel_path:
            logger.error("❌ excel_path 未找到，跳过备选公牛-隐性基因分析")
            return

        try:
            # 1. 读取Excel数据
            bulls_data = self._read_bulls_genes_data(excel_path)
            if not bulls_data:
                logger.warning("⚠️  未找到备选公牛隐性基因数据，跳过")
                return

            bull_count = len(bulls_data)
            logger.info(f"✓ 读取到 {bull_count} 头备选公牛的隐性基因数据")

            # 2. 定位页面范围
            target_slides = self.find_slides_by_text("备选公牛-隐性基因分析", start_index=0)
            if not target_slides:
                logger.error("❌ 未找到备选公牛-隐性基因分析页面，跳过")
                return

            start_slide_index = target_slides[0]
            logger.info(f"✓ 定位到第{start_slide_index + 1}页（备选公牛-隐性基因分析起始页）")

            # 3. 动态页面管理
            self._manage_pages(start_slide_index, bull_count)

            # 4. 填充数据
            for i, bull_data in enumerate(bulls_data):
                slide_index = start_slide_index + i
                if slide_index < len(self.prs.slides):
                    slide = self.prs.slides[slide_index]
                    self._fill_bull_genes_page(slide, bull_data, i + 1)
                else:
                    logger.warning(f"⚠️  页面索引 {slide_index + 1} 超出范围，跳过公牛 {i + 1}")

            logger.info("✓ Part 7 备选公牛-隐性基因分析完成")

        except Exception as e:
            logger.error(f"❌ 构建备选公牛-隐性基因分析失败: {e}")
            logger.debug("错误详情:", exc_info=True)

    def _read_bulls_genes_data(self, excel_path: str) -> List[Dict]:
        """
        从Excel读取备选公牛隐性基因数据

        Args:
            excel_path: Excel文件路径

        Returns:
            List[Dict]: 公牛数据列表，每个字典包含公牛的隐性基因信息
        """
        try:
            wb = load_workbook(excel_path, data_only=True)

            # 检查sheet是否存在
            if "备选公牛-隐性基因分析" not in wb.sheetnames:
                logger.warning("⚠️  Excel中未找到'备选公牛-隐性基因分析' sheet")
                wb.close()
                return []

            ws = wb["备选公牛-隐性基因分析"]

            # 读取所有数据
            data = []
            for row in ws.iter_rows(values_only=False):
                row_data = [cell.value for cell in row]
                data.append(row_data)

            df = pd.DataFrame(data)
            logger.info(f"✓ 读取Excel: {df.shape[0]}行 × {df.shape[1]}列")

            # 解析每个公牛的数据块（每24行一个公牛）
            bulls_data = []
            row_idx = 0

            while row_idx < len(df):
                # 检查是否是公牛标识行（包含"【公牛"）
                cell_value = df.iloc[row_idx, 0]
                if pd.notna(cell_value) and "【公牛" in str(cell_value):
                    # 解析一个公牛数据块
                    bull_info = self._parse_bull_data_block(df, row_idx)
                    if bull_info:
                        bulls_data.append(bull_info)
                    row_idx += 24  # 每个公牛数据块占24行
                else:
                    row_idx += 1

            wb.close()
            logger.info(f"✓ 成功解析 {len(bulls_data)} 头公牛的隐性基因数据")
            return bulls_data

        except Exception as e:
            logger.error(f"❌ 读取备选公牛隐性基因数据失败: {e}")
            logger.debug("错误详情:", exc_info=True)
            return []

    def _parse_bull_data_block(self, df: pd.DataFrame, start_row: int) -> Optional[Dict]:
        """
        解析单个公牛的数据块（24行）

        Args:
            df: DataFrame
            start_row: 起始行索引

        Returns:
            Dict: 公牛隐性基因信息
        """
        try:
            # 行1: 公牛ID和原始号
            bull_id_line = str(df.iloc[start_row, 0])

            # 行2: 在群母牛统计
            stats_line = str(df.iloc[start_row + 1, 0]) if pd.notna(df.iloc[start_row + 1, 0]) else ""

            # 行4-19: 16种隐性基因数据（跳过行3表头）
            genes_data = []
            for i in range(16):
                gene_row_idx = start_row + 3 + i
                if gene_row_idx < len(df):
                    gene_info = {
                        'gene_name': str(df.iloc[gene_row_idx, 0]) if pd.notna(df.iloc[gene_row_idx, 0]) else "",
                        'adult_count': df.iloc[gene_row_idx, 1] if pd.notna(df.iloc[gene_row_idx, 1]) else 0,
                        'adult_ratio': str(df.iloc[gene_row_idx, 2]) if pd.notna(df.iloc[gene_row_idx, 2]) else "0.0%",
                        'heifer_count': df.iloc[gene_row_idx, 3] if pd.notna(df.iloc[gene_row_idx, 3]) else 0,
                        'heifer_ratio': str(df.iloc[gene_row_idx, 4]) if pd.notna(df.iloc[gene_row_idx, 4]) else "0.0%",
                        'total_count': df.iloc[gene_row_idx, 5] if pd.notna(df.iloc[gene_row_idx, 5]) else 0,
                        'total_ratio': str(df.iloc[gene_row_idx, 6]) if pd.notna(df.iloc[gene_row_idx, 6]) else "0.0%",
                    }
                    genes_data.append(gene_info)

            # 行21: 任意基因纯合小计
            summary_row_idx = start_row + 20
            summary_info = {}
            if summary_row_idx < len(df):
                summary_info = {
                    'adult_count': df.iloc[summary_row_idx, 1] if pd.notna(df.iloc[summary_row_idx, 1]) else 0,
                    'adult_ratio': str(df.iloc[summary_row_idx, 2]) if pd.notna(df.iloc[summary_row_idx, 2]) else "0.0%",
                    'heifer_count': df.iloc[summary_row_idx, 3] if pd.notna(df.iloc[summary_row_idx, 3]) else 0,
                    'heifer_ratio': str(df.iloc[summary_row_idx, 4]) if pd.notna(df.iloc[summary_row_idx, 4]) else "0.0%",
                    'total_count': df.iloc[summary_row_idx, 5] if pd.notna(df.iloc[summary_row_idx, 5]) else 0,
                    'total_ratio': str(df.iloc[summary_row_idx, 6]) if pd.notna(df.iloc[summary_row_idx, 6]) else "0.0%",
                }

            return {
                'bull_id': bull_id_line,
                'stats': stats_line,
                'genes': genes_data,
                'summary': summary_info
            }

        except Exception as e:
            logger.error(f"❌ 解析公牛数据块失败（起始行{start_row}）: {e}")
            return None

    def _manage_pages(self, start_index: int, bull_count: int):
        """
        动态管理页面数量

        Args:
            start_index: 起始页面索引
            bull_count: 实际公牛数量
        """
        template_pages = 23  # 模板提供的页面数

        if bull_count < template_pages:
            # 删除多余页面（从后往前删除）
            pages_to_delete = template_pages - bull_count
            logger.info(f"需要删除 {pages_to_delete} 个多余页面")

            for i in range(pages_to_delete):
                # 计算要删除的页面索引（从后往前）
                delete_index = start_index + template_pages - 1 - i
                if delete_index < len(self.prs.slides):
                    self._delete_slide(delete_index)
                    logger.debug(f"  ✓ 删除第{delete_index + 1}页")

            logger.info(f"✓ 已删除 {pages_to_delete} 个多余页面")

        elif bull_count > template_pages:
            # 需要复制模板页
            pages_to_copy = bull_count - template_pages
            logger.info(f"需要复制 {pages_to_copy} 个模板页")

            # 使用最后一个模板页作为复制源
            template_slide_index = start_index + template_pages - 1

            for i in range(pages_to_copy):
                insert_at = start_index + template_pages + i
                self._copy_slide(template_slide_index, insert_at)
                logger.debug(f"  ✓ 复制页面到第{insert_at + 1}页")

            logger.info(f"✓ 已复制 {pages_to_copy} 个模板页")
        else:
            logger.info(f"✓ 公牛数量（{bull_count}）正好匹配模板页数（{template_pages}），无需调整")

    def _delete_slide(self, slide_index: int):
        """
        删除指定索引的幻灯片

        Args:
            slide_index: 幻灯片索引
        """
        try:
            rId = self.prs.slides._sldIdLst[slide_index].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[slide_index]
        except Exception as e:
            logger.error(f"删除第{slide_index + 1}页失败: {e}")

    def _copy_slide(self, source_index: int, insert_at: int):
        """
        复制幻灯片到指定位置

        Args:
            source_index: 源幻灯片索引
            insert_at: 插入位置索引
        """
        try:
            from pptx.oxml.xmlchemy import OxmlElement
            import copy

            source = self.prs.slides[source_index]

            # 复制幻灯片XML
            blank_slide_layout = self.prs.slide_layouts[6]  # 空白布局
            dest = self.prs.slides.add_slide(blank_slide_layout)

            # 移动到指定位置
            xml_slides = self.prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[-1])
            xml_slides.insert(insert_at, slides[-1])

            # 复制所有形状
            for shape in source.shapes:
                el = shape.element
                newel = copy.deepcopy(el)
                dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

        except Exception as e:
            logger.error(f"复制第{source_index + 1}页到第{insert_at + 1}页失败: {e}")

    def _fill_bull_genes_page(self, slide, bull_data: Dict, bull_number: int):
        """
        填充单个公牛的隐性基因分析页面

        Args:
            slide: 幻灯片对象
            bull_data: 公牛数据
            bull_number: 公牛编号（用于日志）
        """
        try:
            logger.info(f"填充第{bull_number}头公牛的数据: {bull_data['bull_id'][:50]}...")

            # 1. 更新公牛标识（第一行）
            self._update_bull_id(slide, bull_data['bull_id'])

            # 2. 查找并更新主表格
            main_table = self._find_main_table(slide)
            if main_table:
                self._fill_genes_table(main_table, bull_data)
            else:
                logger.warning(f"  ⚠️  未找到第{bull_number}头公牛页面的主表格")

            # 3. 更新风险影响占比图表
            self._update_risk_chart(slide, bull_data)

            # 4. 更新分析文本（可选）
            self._update_analysis_text(slide, bull_data)

            logger.info(f"  ✓ 第{bull_number}头公牛数据填充完成")

        except Exception as e:
            logger.error(f"❌ 填充第{bull_number}头公牛数据失败: {e}")
            logger.debug("错误详情:", exc_info=True)

    def _update_bull_id(self, slide, bull_id: str):
        """更新公牛标识文本"""
        try:
            # 查找包含"【公牛"的文本框或表格单元格
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if "【公牛" in shape.text:
                        shape.text = bull_id
                        logger.debug(f"    ✓ 更新公牛标识: {bull_id[:50]}")
                        return
                elif shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if "【公牛" in cell.text:
                                cell.text = bull_id
                                # 设置字体格式：14号，白色
                                self._format_cell(cell, font_size=14, alignment=PP_ALIGN.LEFT,
                                                font_color=(255, 255, 255))
                                logger.debug(f"    ✓ 更新公牛标识（表格）: {bull_id[:50]}")
                                return
        except Exception as e:
            logger.debug(f"    更新公牛标识失败: {e}")

    def _find_main_table(self, slide):
        """查找主数据表格（20行×7列）"""
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # 查找20行×7列的表格
                if len(table.rows) == 20 and len(table.columns) == 7:
                    logger.debug(f"    ✓ 找到主表格: {len(table.rows)}行 × {len(table.columns)}列")
                    return table
        return None

    def _fill_genes_table(self, table, bull_data: Dict):
        """填充隐性基因数据表格"""
        try:
            # 行2: 在群母牛统计（如果有的话）
            if bull_data['stats']:
                table.cell(1, 0).text = bull_data['stats']
                self._format_cell(table.cell(1, 0), font_size=8, alignment=PP_ALIGN.LEFT)

            # 行4-19: 16种隐性基因数据
            for i, gene_info in enumerate(bull_data['genes'][:16]):
                row_idx = 3 + i  # 从第4行开始（索引3）

                # 填充7列数据
                table.cell(row_idx, 0).text = gene_info['gene_name']
                table.cell(row_idx, 1).text = str(gene_info['adult_count'])
                table.cell(row_idx, 2).text = gene_info['adult_ratio']
                table.cell(row_idx, 3).text = str(gene_info['heifer_count'])
                table.cell(row_idx, 4).text = gene_info['heifer_ratio']
                table.cell(row_idx, 5).text = str(gene_info['total_count'])
                table.cell(row_idx, 6).text = gene_info['total_ratio']

                # 格式化单元格
                for col_idx in range(7):
                    cell = table.cell(row_idx, col_idx)
                    if col_idx == 0:
                        self._format_cell(cell, font_size=9, alignment=PP_ALIGN.LEFT)
                    else:
                        self._format_cell(cell, font_size=9, alignment=PP_ALIGN.CENTER)

            # 行21: 任意基因纯合小计（索引19）
            if bull_data['summary']:
                summary = bull_data['summary']
                table.cell(19, 1).text = str(summary['adult_count'])
                table.cell(19, 2).text = summary['adult_ratio']
                table.cell(19, 3).text = str(summary['heifer_count'])
                table.cell(19, 4).text = summary['heifer_ratio']
                table.cell(19, 5).text = str(summary['total_count'])
                table.cell(19, 6).text = summary['total_ratio']

                for col_idx in range(1, 7):
                    self._format_cell(table.cell(19, col_idx), font_size=9, alignment=PP_ALIGN.CENTER, bold=True)

            logger.debug(f"    ✓ 表格数据填充完成")

        except Exception as e:
            logger.error(f"    ❌ 填充表格数据失败: {e}")
            logger.debug("错误详情:", exc_info=True)

    def _format_cell(self, cell, font_size: int = 8, alignment=PP_ALIGN.CENTER, bold: bool = False,
                     font_color: tuple = None, font_name: str = None):
        """
        格式化单元格

        Args:
            cell: 单元格对象
            font_size: 字号
            alignment: 对齐方式
            bold: 是否加粗
            font_color: 字体颜色RGB元组，如 (255, 255, 255) 代表白色
            font_name: 字体名称，如 "微软雅黑"
        """
        try:
            if not cell.text_frame or not cell.text_frame.paragraphs:
                return

            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = alignment
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    if bold:
                        run.font.bold = True
                    if font_color:
                        run.font.color.rgb = RGBColor(*font_color)
                    if font_name:
                        run.font.name = font_name
        except Exception as e:
            logger.debug(f"    格式化单元格失败: {e}")

    def _update_analysis_text(self, slide, bull_data: Dict):
        """更新分析文本框（可选）"""
        try:
            # 查找包含"分析："的文本框
            for shape in slide.shapes:
                if shape.has_text_frame and "分析：" in shape.text:
                    # 生成分析文本
                    analysis = self._generate_analysis_text(bull_data)
                    if analysis:
                        shape.text = f"分析：{analysis}"
                        logger.debug(f"    ✓ 更新分析文本")
                    return
        except Exception as e:
            logger.debug(f"    更新分析文本失败: {e}")

    def _generate_analysis_text(self, bull_data: Dict) -> str:
        """生成分析文本"""
        try:
            summary = bull_data.get('summary', {})
            total_count = summary.get('total_count', 0)
            total_ratio = summary.get('total_ratio', '0.0%')

            if total_count > 0:
                return f"该公牛可能造成全群{total_count}头母牛（{total_ratio}）产生隐性基因纯合后代，建议谨慎使用。"
            else:
                return "该公牛与牧场母牛群隐性基因匹配良好，无风险。"
        except Exception as e:
            logger.debug(f"    生成分析文本失败: {e}")
            return ""

    def _update_risk_chart(self, slide, bull_data: Dict):
        """
        更新"任意基因纯合风险影响占比"图表

        Args:
            slide: 幻灯片对象
            bull_data: 公牛数据
        """
        try:
            # 1. 查找图表
            chart_shape = None
            for shape in slide.shapes:
                if shape.has_chart:
                    chart = shape.chart
                    if chart.has_title and "任意基因纯合风险影响占比" in chart.chart_title.text_frame.text:
                        chart_shape = shape
                        break

            if not chart_shape:
                logger.debug("    未找到风险影响占比图表")
                return

            # 2. 提取数据
            summary = bull_data.get('summary', {})
            if not summary:
                logger.debug("    公牛数据中缺少summary信息")
                return

            # 3. 转换百分比为小数
            mature_ratio = self._parse_percentage(summary.get('adult_ratio', '0.0%'))
            heifer_ratio = self._parse_percentage(summary.get('heifer_ratio', '0.0%'))
            total_ratio = self._parse_percentage(summary.get('total_ratio', '0.0%'))

            logger.debug(f"    图表数据: 成母牛={mature_ratio:.4f}, 后备牛={heifer_ratio:.4f}, 全群={total_ratio:.4f}")

            # 4. 创建新图表数据
            from pptx.chart.data import CategoryChartData

            chart_data = CategoryChartData()
            chart_data.categories = ['成母牛', '后备牛', '全群']
            chart_data.add_series('风险占比', [mature_ratio, heifer_ratio, total_ratio])

            # 5. 替换图表数据
            chart = chart_shape.chart
            chart.replace_data(chart_data)

            logger.debug("    ✓ 图表更新完成")

        except Exception as e:
            logger.debug(f"    更新风险影响图表失败: {e}")

    def _parse_percentage(self, percent_value) -> float:
        """
        将百分比转换为小数

        Args:
            percent_value: 百分比值（可能是字符串"0.3%"或数值）

        Returns:
            float: 小数值（如0.003）
        """
        if isinstance(percent_value, (int, float)):
            return float(percent_value)

        try:
            percent_str = str(percent_value).strip().replace('%', '').strip()
            return float(percent_str) / 100.0
        except:
            return 0.0
