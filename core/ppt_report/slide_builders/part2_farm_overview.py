"""
Part 2: 牧场概况（直接基于完整PPT模板）
"""

import logging
from typing import Dict, List, Optional

from pptx.chart.data import CategoryChartData
from pptx.util import Pt

from ..base_builder import BaseSlideBuilder
from ..config import FONT_NAME_CN

logger = logging.getLogger(__name__)


class Part2FarmOverviewBuilder(BaseSlideBuilder):
    """基于模板的 Part2 填充器（Slide 3-6）"""

    SECTION_SLIDE_INDEX = 2
    BASIC_INFO_SLIDE_INDEX = 3
    STRUCTURE_SLIDE_INDEX = 4
    PARITY_SLIDE_INDEX = 5

    def __init__(self, prs, chart_creator, farm_name: str):
        super().__init__(prs, chart_creator)
        self.farm_name = farm_name

    def build(self, data: dict, versions: Optional[List[str]] = None):
        farm_info = data.get('farm_info_dict') or {}
        if not farm_info:
            logger.warning("farm_info_dict 为空，标记Part2页面待删除")
            # 标记固定索引页面待删除
            slides_to_delete = [
                self.SECTION_SLIDE_INDEX,
                self.BASIC_INFO_SLIDE_INDEX,
                self.STRUCTURE_SLIDE_INDEX,
                self.PARITY_SLIDE_INDEX
            ]
            self.mark_slides_for_deletion(slides_to_delete)
            return

        self._fill_section_slide()
        self._fill_basic_info_slide(farm_info)
        self._fill_structure_slide(farm_info)
        self._fill_parity_slide(farm_info)
        logger.info("✓ Part 2 模板页更新完成")

    # ------------------------------------------------------------------ #
    def _fill_section_slide(self):
        slide = self._get_slide(self.SECTION_SLIDE_INDEX)
        if not slide:
            logger.warning("模板缺少Part2章节页（Slide 3）")
            return
        logger.info("更新章节页：牧场概况")
        # 模板上的文字已固定，此处如需动态可扩展

    def _fill_basic_info_slide(self, farm_info: Dict):
        slide = self._get_slide(self.BASIC_INFO_SLIDE_INDEX)
        if not slide:
            logger.warning("模板缺少Part2基础信息页（Slide 4）")
            return

        logger.info("填充 Slide 4 - 牧场基本数据信息")
        self._fill_basic_info_table(slide, farm_info)
        self._fill_upload_table(slide, farm_info)
        self._update_basic_info_analysis(slide, farm_info)

    def _fill_structure_slide(self, farm_info: Dict):
        slide = self._get_slide(self.STRUCTURE_SLIDE_INDEX)
        if not slide:
            logger.warning("模板缺少Part2牛群结构页（Slide 5）")
            return

        logger.info("填充 Slide 5 - 牛群结构分析")
        cow_stats = farm_info.get('cow_type_distribution', [])
        self._fill_table(slide, "表格 2", cow_stats, columns=('type', 'count', 'percent'))
        avg_rows = [
            ('在群母牛平均胎次', self._format_float(farm_info.get('avg_lactation'))),
            ('在群母牛平均泌乳天数', self._format_days(farm_info.get('avg_dim'))),
        ]
        self._apply_table_rows(slide, "表格 6", avg_rows)
        self._update_charts(slide, cow_stats)
        self._update_structure_analysis(slide, farm_info, cow_stats)

    def _fill_parity_slide(self, farm_info: Dict):
        slide = self._get_slide(self.PARITY_SLIDE_INDEX)
        if not slide:
            logger.warning("模板缺少Part2胎次页（Slide 6）")
            return

        logger.info("填充 Slide 6 - 胎次分布分析")
        raw_stats = farm_info.get('parity_distribution', [])
        # 去掉表头“胎次”这一行，只保留实际数据与合计
        parity_stats = [item for item in raw_stats if item.get('type') != '胎次']
        self._fill_table(slide, "表格 5", parity_stats, columns=('type', 'count', 'percent'))
        # 图表中不需要“合计”这一类
        chart_stats = [item for item in parity_stats if item.get('type') != '合计']
        self._update_parity_charts(slide, chart_stats)
        self._update_parity_analysis(slide, farm_info, parity_stats)

    def _update_parity_charts(self, slide, stats: List[Dict]):
        """
        更新“按胎次分布数量/占比”图表：
        - 饼图：数量
        - 柱状图：占比
        图表中不包含“合计”这一类。
        """
        if not stats:
            return
        categories = [item["type"] for item in stats]
        counts = [item.get("count", 0) for item in stats]
        percents = [item.get("percent", 0.0) for item in stats]

        charts = [shape.chart for shape in slide.shapes if getattr(shape, "has_chart", False)]
        if not charts:
            return

        pie_chart = charts[0]
        bar_chart = charts[1] if len(charts) > 1 else charts[0]

        pie_data = CategoryChartData()
        pie_data.categories = categories
        pie_data.add_series("数量", counts)
        pie_chart.replace_data(pie_data)

        bar_data = CategoryChartData()
        bar_data.categories = categories
        bar_data.add_series("占比", percents)
        bar_chart.replace_data(bar_data)

    # ------------------------------------------------------------------ #
    def _fill_basic_info_table(self, slide, farm_info: Dict):
        rows = [
            ("牧场名称", farm_info.get('farm_name') or "-"),
            ("报告生成时间", farm_info.get('report_time') or "-"),
            ("牧场服务人员", farm_info.get('service_staff') or "-"),
        ]
        self._apply_table_rows(slide, "表格 5", rows)

    def _fill_upload_table(self, slide, farm_info: Dict):
        stats = farm_info.get('upload_stats', [])
        table = self._find_table(slide, "表格 6")
        if not table:
            logger.warning("未找到上传数据表（表格 6）")
            return

        header_offset = 1
        max_rows = len(table.rows) - header_offset
        for idx in range(max_rows):
            row_idx = idx + header_offset
            if idx < len(stats):
                item = stats[idx]
                self._set_cell_text(table.cell(row_idx, 0), item["type"])
                self._set_cell_text(table.cell(row_idx, 1), str(item["count"]))
                self._set_cell_text(table.cell(row_idx, 2), item["remark"])
            else:
                for col in range(len(table.columns)):
                    self._set_cell_text(table.cell(row_idx, col), "")
        if len(stats) > max_rows:
            logger.warning("上传数据条目超过模板容量，仅展示前 %s 条", max_rows)

    def _update_charts(self, slide, stats: List[Dict]):
        if not stats:
            return
        categories = [item['type'] for item in stats]
        values = [item['count'] for item in stats]
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series("数量", values)

        for shape in slide.shapes:
            if getattr(shape, "has_chart", False):
                shape.chart.replace_data(chart_data)

    def _fill_table(self, slide, table_name: str, stats: List[Dict], columns):
        table = self._find_table(slide, table_name)
        if not table:
            logger.warning("未找到表格 %s", table_name)
            return
        header_offset = 1
        max_rows = len(table.rows) - header_offset
        for idx in range(max_rows):
            row_idx = idx + header_offset
            if idx < len(stats):
                item = stats[idx]
                for col_idx, key in enumerate(columns):
                    value = item.get(key)
                    if key == 'percent' and value is not None:
                        text = f"{value:.1f}"
                    else:
                        text = "" if value is None else str(value)
                    self._set_cell_text(table.cell(row_idx, col_idx), text)
            else:
                for col in range(len(table.columns)):
                    self._set_cell_text(table.cell(row_idx, col), "")
        if len(stats) > max_rows:
            logger.warning("表格 %s 条目超过模板容量，仅展示前 %s 行", table_name, max_rows)

    def _apply_table_rows(self, slide, table_name: str, rows: List[tuple]):
        table = self._find_table(slide, table_name)
        if not table:
            logger.warning("未找到表格 %s", table_name)
            return
        for idx, row in enumerate(rows):
            if idx >= len(table.rows):
                break
            for col_idx, value in enumerate(row):
                if col_idx < len(table.columns):
                    self._set_cell_text(table.cell(idx, col_idx), str(value))

    def _update_basic_info_analysis(self, slide, farm_info: Dict):
        """第4页：牧场基本数据信息分析"""
        shape = self._find_shape(slide, "文本框 9")
        if not shape or not shape.has_text_frame:
            return

        total = farm_info.get('total_count') or 0
        lact = farm_info.get('lactating_count') or 0
        heifer = farm_info.get('heifer_count') or 0

        # 计算占比
        lact_pct = self._calc_percent(lact, total)
        heifer_pct = self._calc_percent(heifer, total)

        # 牛群规模评价
        if total >= 3000:
            scale = "大型牧场"
        elif total >= 1000:
            scale = "中型牧场"
        else:
            scale = "小型牧场"

        # 结构评价（正常比例：成母牛55%，后备牛45%）
        # 后备牛<35%过低，<40%偏低；成母牛<45%过低
        if heifer_pct < 35:
            structure_eval = "后备牛占比过低，建议增加性控冻精使用比例以补充后备牛"
        elif heifer_pct < 40:
            structure_eval = "后备牛占比偏低，建议适当增加性控冻精使用比例"
        elif lact_pct < 45:
            structure_eval = "后备牛占比较高，可考虑淘汰低遗传水平牛只"
        elif 50 <= lact_pct <= 60:
            structure_eval = "成母牛与后备牛比例合理"
        elif lact_pct > 60:
            structure_eval = "成母牛占比较高，可考虑淘汰低产牛"
        else:
            structure_eval = "牛群结构基本正常"

        # 数据完整性评价
        upload_stats = farm_info.get('upload_stats', [])
        has_genotype = any('基因组' in str(item.get('type', '')) for item in upload_stats)
        has_body_score = any('体型' in str(item.get('type', '')) for item in upload_stats)

        data_parts = []
        if has_genotype:
            data_parts.append("已有基因组检测数据")
        if has_body_score:
            data_parts.append("已有体型外貌测定数据")

        data_eval = "，".join(data_parts) if data_parts else "基础数据尚待完善"

        text = (
            f"分析：牧场为{scale}，在群牛共{total}头，其中成母牛{lact}头（{lact_pct:.1f}%），"
            f"后备牛{heifer}头（{heifer_pct:.1f}%）。{structure_eval}。{data_eval}，"
            f"为精准选配提供了良好的数据基础。"
        )
        self._set_text(shape, text)

    def _update_structure_analysis(self, slide, farm_info: Dict, stats: List[Dict]):
        """第5页：牛群结构分析"""
        shape = self._find_shape(slide, "文本框 5")
        if not shape:
            return

        avg_lact = farm_info.get('avg_lactation') or 0
        avg_dim = farm_info.get('avg_dim') or 0
        cow_dist = farm_info.get('cow_type_distribution', [])

        # 平均胎次评价
        if avg_lact < 2.0:
            lact_eval = "平均胎次偏低，牛群较年轻"
        elif avg_lact > 3.5:
            lact_eval = "平均胎次偏高，建议关注牛群更新"
        else:
            lact_eval = "平均胎次适中"

        # 泌乳天数评价
        if avg_dim < 150:
            dim_eval = "平均泌乳天数较短，产奶效率较高"
        elif avg_dim > 200:
            dim_eval = "平均泌乳天数偏长，建议关注繁殖管理"
        else:
            dim_eval = "平均泌乳天数正常"

        # 找出占比最大的牛群类型
        if cow_dist:
            main_type = max(cow_dist, key=lambda x: x.get('percent', 0))
            main_info = f"其中{main_type.get('type', '')}占比最高（{main_type.get('percent', 0):.1f}%）"
        else:
            main_info = ""

        text = (
            f"分析：牧场在群牛平均胎次{avg_lact:.2f}胎，{lact_eval}；"
            f"平均泌乳天数{int(avg_dim)}天，{dim_eval}。{main_info}。"
        )
        self._set_text(shape, text)

    def _update_parity_analysis(self, slide, farm_info: Dict, stats: List[Dict]):
        """
        第6页：胎次分布分析

        推荐比例（占成母牛）：
        - 1胎~2胎牛：40%
        - 3胎~5胎牛：40%
        - 6胎及以上：20%
        """
        shape = self._find_shape(slide, "文本框 2")
        if not shape:
            return

        # 过滤掉表头和合计
        valid_data = [x for x in stats if x.get("type") not in ("胎次", "合计")]

        if not valid_data:
            self._set_text(shape, "分析：暂无有效胎次分布数据。")
            return

        # 计算各组胎次占比（占成母牛比例）
        # 1-2胎
        low_parity = [x for x in valid_data
                      if str(x.get("type", "")).isdigit() and int(x.get("type", "0")) in (1, 2)]
        low_pct = sum(x.get("percent", 0) for x in low_parity)

        # 3-5胎
        mid_parity = [x for x in valid_data
                      if str(x.get("type", "")).isdigit() and int(x.get("type", "0")) in (3, 4, 5)]
        mid_pct = sum(x.get("percent", 0) for x in mid_parity)

        # 6胎及以上
        high_parity = [x for x in valid_data
                       if str(x.get("type", "")).isdigit() and int(x.get("type", "0")) >= 6]
        high_pct = sum(x.get("percent", 0) for x in high_parity)

        # 找占比最高的胎次
        highest = max(valid_data, key=lambda x: x.get("percent", 0))

        # 生成评价
        parts = []
        parts.append(
            f"牧场成母牛胎次分布：1-2胎占{low_pct:.1f}%，3-5胎占{mid_pct:.1f}%，6胎及以上占{high_pct:.1f}%。"
            f"其中{highest.get('type', '')}胎牛占比最高（{highest.get('percent', 0):.1f}%）。"
        )

        # 推荐比例：1-2胎40%、3-5胎40%、6胎及以上20%
        evals = []
        if low_pct < 30:
            evals.append("1-2胎牛占比偏低，建议加强后备牛培育和及时配种")
        elif low_pct > 50:
            evals.append("1-2胎牛占比偏高，牛群较年轻")

        if mid_pct < 30:
            evals.append("3-5胎牛占比偏低")
        elif mid_pct > 50:
            evals.append("3-5胎牛占比较高，牛群生产效率较好")

        if high_pct > 30:
            evals.append("6胎及以上牛占比偏高，建议加快高胎次低产牛淘汰")
        elif high_pct < 10:
            evals.append("高胎次牛占比较低")

        if not evals:
            evals.append("胎次结构较为均衡，接近推荐比例（1-2胎40%、3-5胎40%、6胎+20%）")

        parts.append("。".join(evals) + "。")

        text = "分析：" + "".join(parts)
        self._set_text(shape, text)

    # ------------------------------------------------------------------ #
    def _get_slide(self, index: int):
        if index < len(self.prs.slides):
            return self.prs.slides[index]
        return None

    def _find_shape(self, slide, name: str):
        for shape in slide.shapes:
            if shape.name == name:
                return shape
        return None

    def _find_table(self, slide, name: str):
        for shape in slide.shapes:
            if shape.name == name and getattr(shape, "has_table", False):
                return shape.table
        return None

    @staticmethod
    def _calc_percent(part: float, whole: float) -> float:
        return 0.0 if not whole else part * 100.0 / whole

    @staticmethod
    def _format_float(value) -> str:
        if value is None:
            return "-"
        return f"{float(value):.2f}"

    @staticmethod
    def _format_days(value) -> str:
        if value is None:
            return "-"
        return f"{int(value)}天"

    @staticmethod
    def _set_text(shape, text: str):
        """更新分析文本框内容，设置微软雅黑15号非加粗"""
        if not shape or not shape.has_text_frame:
            return
        tf = shape.text_frame
        # 清空并重新设置
        tf.clear()
        if not tf.paragraphs:
            para = tf.add_paragraph()
        else:
            para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text

        # 设置字体：微软雅黑15号非加粗
        run.font.name = FONT_NAME_CN
        run.font.size = Pt(15)
        run.font.bold = False

        # 清理多余段落
        for extra_para in tf.paragraphs[1:]:
            for extra_run in extra_para.runs:
                extra_run.text = ""

    @staticmethod
    def _set_cell_text(cell, text: str):
        """只改单元格文字，不动字体/颜色等样式"""
        tf = cell.text_frame
        if not tf.paragraphs:
            para = tf.add_paragraph()
        else:
            para = tf.paragraphs[0]
        if para.runs:
            run = para.runs[0]
        else:
            run = para.add_run()
        run.text = text
        # 清理多余 run/段落中的旧文本
        for extra_run in para.runs[1:]:
            extra_run.text = ""
        for extra_para in tf.paragraphs[1:]:
            for extra_run in extra_para.runs:
                extra_run.text = ""
