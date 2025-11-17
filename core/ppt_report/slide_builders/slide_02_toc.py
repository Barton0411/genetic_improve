"""
Slide 2: 目录页（间隔页）构建器
按照模板规则实现
"""

import logging
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


def build_slide_02_toc(prs):
    """
    构建Slide 2: 目录页

    使用Layout [1] 间隔页，手动添加38个元素：
    - 背景椭圆
    - "目录"标题和下划线
    - 7个章节条目（每个包含：垂直线、数字框、数字、标题、副标题）

    Args:
        prs: PowerPoint演示文稿对象
    """
    logger.info("开始构建Slide 2: 目录页")

    # 使用Layout [1] 间隔页
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # 章节数据
    sections = [
        {"num": "01", "title": "牧场概况", "subtitle": "基本信息 牛群结构", "top": 4.29},
        {"num": "02", "title": "系谱记录分析", "subtitle": "完整性 识别率", "top": 7.09},
        {"num": "03", "title": "牛群遗传评估", "subtitle": "性状 指数 进展", "top": 9.88},
        {"num": "04", "title": "配种记录分析", "subtitle": "概况 风险 近交", "top": 12.67},
        {"num": "05", "title": "公牛使用分析", "subtitle": "已用 备选 评估", "top": 5.66},
        {"num": "06", "title": "选配推荐方案", "subtitle": "推荐统计 质量", "top": 8.46},
        {"num": "07", "title": "项目总结建议", "subtitle": "发现 建议", "top": 11.25},
    ]

    # ============ 元素1: 背景椭圆 ============
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Cm(6.35), Cm(2.03),
        Cm(22.86), Cm(15.24)
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = RGBColor(240, 247, 252)  # 淡蓝色背景
    oval.line.fill.background()

    # ============ 元素2: "目录"标题 ============
    toc_title = slide.shapes.add_textbox(
        Cm(2.03), Cm(1.52),
        Cm(3.61), Cm(2.13)
    )
    toc_frame = toc_title.text_frame
    toc_frame.word_wrap = False
    toc_para = toc_frame.paragraphs[0]
    toc_para.text = "目录"
    toc_para.alignment = PP_ALIGN.LEFT

    for run in toc_para.runs:
        run.font.name = '微软雅黑'
        run.font.size = Pt(24)
        run.font.color.rgb = RGBColor(0, 112, 192)  # 深蓝色
        run.font.bold = True

    # ============ 元素3: "目录"下划线 ============
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(2.03), Cm(3.43),
        Cm(2.03), Cm(0.05)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = RGBColor(0, 112, 192)
    underline.line.fill.background()

    # ============ 左侧章节 (01-04) ============
    for sec in sections[:4]:
        _add_section_item(slide, sec, left_column=True)

    # ============ 右侧章节 (05-07) ============
    for sec in sections[4:]:
        _add_section_item(slide, sec, left_column=False)

    logger.info("✓ Slide 2（目录页）构建完成")
    return slide


def _add_section_item(slide, section_data, left_column=True):
    """
    添加一个章节条目

    Args:
        slide: 幻灯片对象
        section_data: 章节数据 {"num": "01", "title": "标题", "subtitle": "副标题", "top": 4.29}
        left_column: True表示左列(9.5cm起), False表示右列(18.14cm起)
    """
    num = section_data["num"]
    title = section_data["title"]
    subtitle = section_data["subtitle"]
    top = section_data["top"]

    # 左列和右列的起始位置不同
    if left_column:
        line_left = 9.5
        box_left = 10.01
        text_left = 12.17
    else:
        line_left = 18.14
        box_left = 18.62
        text_left = 20.8

    # 垂直线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(line_left), Cm(top),
        Cm(0.15), Cm(2.29)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 112, 192)  # 深蓝色
    line.line.fill.background()

    # 数字背景框
    num_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(box_left), Cm(top),
        Cm(1.93), Cm(1.8)
    )
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = RGBColor(217, 237, 247)  # 浅蓝色
    num_box.line.fill.background()

    # 数字文本
    num_text = slide.shapes.add_textbox(
        Cm(box_left), Cm(top),
        Cm(1.93), Cm(1.8)
    )
    num_frame = num_text.text_frame
    num_frame.word_wrap = False
    num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    num_para = num_frame.paragraphs[0]
    num_para.text = num
    num_para.alignment = PP_ALIGN.RIGHT

    for run in num_para.runs:
        run.font.name = '微软雅黑'
        run.font.size = Pt(48)
        run.font.color.rgb = RGBColor(0, 112, 192)
        run.font.bold = True

    # 标题文本
    title_text = slide.shapes.add_textbox(
        Cm(text_left), Cm(top),
        Cm(4.75), Cm(1.12)
    )
    title_frame = title_text.text_frame
    title_frame.word_wrap = False
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.alignment = PP_ALIGN.LEFT

    for run in title_para.runs:
        run.font.name = '微软雅黑'
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0, 112, 192)
        run.font.bold = True

    # 副标题文本
    subtitle_top = top + 1.27  # 副标题在标题下方
    subtitle_text = slide.shapes.add_textbox(
        Cm(text_left), Cm(subtitle_top),
        Cm(4.62), Cm(0.86)
    )
    subtitle_frame = subtitle_text.text_frame
    subtitle_frame.word_wrap = False
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = subtitle
    subtitle_para.alignment = PP_ALIGN.LEFT

    for run in subtitle_para.runs:
        run.font.name = '微软雅黑'
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(89, 89, 89)  # 深灰色
