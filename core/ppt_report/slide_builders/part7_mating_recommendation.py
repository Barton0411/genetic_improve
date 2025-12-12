"""
Part 7: 个体选配推荐结果构建器（页173）
"""

import logging
from typing import Dict, Optional

from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor

from ..base_builder import BaseSlideBuilder

logger = logging.getLogger(__name__)


class Part7MatingRecommendationBuilder(BaseSlideBuilder):
    """Part 7: 个体选配推荐结果构建器（页173 - 选配统计摘要）"""

    def __init__(self, prs, chart_creator, farm_name: str):
        super().__init__(prs, chart_creator)
        self.farm_name = farm_name

    def build(self, data: Dict):
        """
        构建 Part 7: 个体选配推荐结果（页173）

        Args:
            data: 包含 mating_summary 的数据字典
        """
        logger.info("=" * 60)
        logger.info("开始构建Part 7: 个体选配推荐结果（页173）")
        logger.info("=" * 60)

        # 使用动态查找定位页面（从索引50开始，避免找到目录页）
        target_slides = self.find_slides_by_text("选配统计摘要", start_index=50)

        if not target_slides:
            # 尝试备用关键字
            target_slides = self.find_slides_by_text("个体选配推荐结果", start_index=50)

        if not target_slides:
            logger.warning("未找到个体选配推荐结果页面（选配统计摘要），跳过")
            return

        mating_data = data.get("mating_summary")
        if not mating_data:
            logger.warning("mating_summary 数据为空，标记页面待删除")
            self.mark_slides_for_deletion(target_slides)
            return

        slide_index = target_slides[0]
        slide = self.prs.slides[slide_index]
        logger.info(f"✓ 定位到第{slide_index + 1}页（选配统计摘要）")

        # 更新基础统计表格（左侧 4×2 表格）
        self._update_basic_stats_table(slide, mating_data.get("basic_stats", {}))

        # 更新分组统计表格（右侧动态表格）
        self._update_group_stats_table(slide, mating_data.get("group_stats", []))

        # 更新Excel下载按钮
        self._update_excel_download_button(slide)

        # 更新分析文本
        self._update_analysis_text(slide, mating_data)

        logger.info("✓ Part 7 个体选配推荐结果页更新完成")

    # ------------------------------------------------------------------ #
    # 核心更新方法
    # ------------------------------------------------------------------ #

    def _update_basic_stats_table(self, slide, basic_stats: Dict):
        """更新基础统计表格（表格 2）"""
        table = self._find_table(slide, "表格 2")
        if not table:
            logger.warning("未找到基础统计表格（表格 2）")
            return

        # 表格结构：4行×2列
        # 行0: 总母牛数 | 值
        # 行1: 有性控推荐 | 值
        # 行2: 有常规推荐 | 值
        # 行3: 无推荐 | 值
        # 只更新数值列（列1），保留模板中的标签

        stats_values = [
            basic_stats.get("total_cows", 0),
            basic_stats.get("has_sexed", 0),
            basic_stats.get("has_regular", 0),
            basic_stats.get("no_recommendation", 0)
        ]

        for row_idx, value in enumerate(stats_values):
            if row_idx >= len(table.rows):
                break
            # 只更新列1的数值，保留列0的模板标签
            self._set_cell_text(table.cell(row_idx, 1), str(value))

        logger.info(f"✓ 基础统计表格更新完成：总数={basic_stats.get('total_cows', 0)}")

    def _update_group_stats_table(self, slide, group_stats: list):
        """更新分组统计表格（表格 6）"""
        table = self._find_table(slide, "表格 6")
        if not table:
            logger.warning("未找到分组统计表格（表格 6）")
            return

        # 表格结构：表头 + 数据行
        # 列: 分组 | 总数 | 性控 | 常规

        # 1) 确定表头行（假设第0行是表头）
        header_row_idx = 0
        data_start_row = 1

        # 2) 计算需要的数据行数
        num_groups = len(group_stats)
        template_data_rows = len(table.rows) - 1  # 减去表头行

        # 3) 如果模板行数不足，添加行
        if num_groups > template_data_rows:
            rows_to_add = num_groups - template_data_rows
            for _ in range(rows_to_add):
                self._add_table_row(table, copy_from_row=data_start_row)
            logger.info(f"添加了 {rows_to_add} 行到分组统计表格")

        # 4) 如果模板行数过多，删除多余行（从下往上删除）
        elif num_groups < template_data_rows:
            rows_to_delete = template_data_rows - num_groups
            for i in range(rows_to_delete):
                delete_row_idx = len(table.rows) - 1
                self._delete_row(table, delete_row_idx)
            logger.info(f"删除了 {rows_to_delete} 行从分组统计表格")

        # 5) 填充数据
        for i, group_data in enumerate(group_stats):
            row_idx = data_start_row + i
            if row_idx >= len(table.rows):
                break

            self._set_cell_text(table.cell(row_idx, 0), group_data.get("group", ""))
            self._set_cell_text(table.cell(row_idx, 1), str(group_data.get("total", 0)))
            self._set_cell_text(table.cell(row_idx, 2), str(group_data.get("sexed", 0)))
            self._set_cell_text(table.cell(row_idx, 3), str(group_data.get("regular", 0)))

        logger.info(f"✓ 分组统计表格更新完成：共 {num_groups} 个分组")

    def _update_excel_download_button(self, slide):
        """更新Excel下载按钮文本（圆角矩形 8）"""
        button_shape = None
        for shape in slide.shapes:
            if shape.name == "圆角矩形 8":
                button_shape = shape
                break

        if not button_shape:
            logger.warning("未找到Excel下载按钮（圆角矩形 8）")
            return

        if not button_shape.has_text_frame:
            logger.warning("Excel下载按钮没有文本框")
            return

        # 设置按钮文本
        button_text = f"点击下载{self.farm_name}选配推荐详情.xlsx"
        self._set_shape_text(button_shape, button_text)
        logger.info(f"✓ Excel下载按钮更新完成")

    def _update_analysis_text(self, slide, mating_data: Dict):
        """更新分析文本（文本框 14）"""
        text_box = None
        for shape in slide.shapes:
            if shape.name == "文本框 14":
                text_box = shape
                break

        if not text_box:
            logger.warning("未找到分析文本框（文本框 14）")
            return

        if not text_box.has_text_frame:
            logger.warning("文本框 14 没有文本框架")
            return

        # 生成智能分析文本
        basic_stats = mating_data.get("basic_stats", {})
        group_stats = mating_data.get("group_stats", [])

        total_cows = basic_stats.get("total_cows", 0)
        has_sexed = basic_stats.get("has_sexed", 0)
        has_regular = basic_stats.get("has_regular", 0)
        no_recommendation = basic_stats.get("no_recommendation", 0)

        # 计算覆盖率
        coverage_rate = 0
        if total_cows > 0:
            coverage_rate = ((has_sexed + has_regular) / total_cows) * 100

        # 生成分析文本
        analysis_parts = []

        analysis_parts.append(
            f"本期共为{total_cows}头应配母牛提供了选配推荐方案。"
        )

        if has_sexed > 0 and has_regular > 0:
            analysis_parts.append(
                f"其中{has_sexed}头获得性控公牛推荐，{has_regular}头获得常规公牛推荐，"
                f"推荐覆盖率达{coverage_rate:.1f}%。"
            )
        elif has_sexed > 0:
            analysis_parts.append(
                f"其中{has_sexed}头获得性控公牛推荐，推荐覆盖率达{coverage_rate:.1f}%。"
            )
        elif has_regular > 0:
            analysis_parts.append(
                f"其中{has_regular}头获得常规公牛推荐，推荐覆盖率达{coverage_rate:.1f}%。"
            )

        if no_recommendation > 0:
            no_rec_rate = (no_recommendation / total_cows) * 100 if total_cows > 0 else 0
            analysis_parts.append(
                f"另有{no_recommendation}头母牛（{no_rec_rate:.1f}%）暂无推荐方案，"
                f"建议人工审核后进行选配。"
            )

        # 分组统计分析（找出最大的分组）
        if group_stats:
            max_group = max(group_stats, key=lambda x: x.get("total", 0))
            max_group_name = max_group.get("group", "")
            max_group_total = max_group.get("total", 0)
            if max_group_total > 0:
                analysis_parts.append(
                    f"从分组来看，{max_group_name}的应配母牛数量最多，共{max_group_total}头。"
                )

        analysis_text = "".join(analysis_parts)
        self._set_shape_text(text_box, analysis_text)
        logger.info(f"✓ 分析文本更新完成")

    # ------------------------------------------------------------------ #
    # 辅助方法
    # ------------------------------------------------------------------ #

    def _find_table(self, slide, name: str):
        """查找指定名称的表格"""
        for shape in slide.shapes:
            if shape.name == name and getattr(shape, "has_table", False):
                return shape.table
        return None

    @staticmethod
    def _add_table_row(table, copy_from_row: int = 1):
        """
        添加新行到表格末尾，复制指定行的样式

        python-pptx 没有公开的添加行 API，这里通过底层 XML 复制行
        """
        try:
            from copy import deepcopy

            tbl = table._tbl
            source_row = table.rows[copy_from_row]._tr
            new_row = deepcopy(source_row)

            # 清空新行的文本
            for tc in new_row.tc_lst:
                for paragraph in tc.txBody.p_lst:
                    for run in paragraph.r_lst:
                        run.t.text = ""

            tbl.append(new_row)
        except Exception as e:
            logger.error(f"添加表格行失败: {e}", exc_info=True)

    @staticmethod
    def _delete_row(table, row_idx: int):
        """
        从表格中删除指定行

        python-pptx 没有公开的删除行 API，这里通过底层 XML 删除 <tr> 节点
        """
        try:
            row = table.rows[row_idx]
        except IndexError:
            return
        tbl = table._tbl
        tr = row._tr
        tbl.remove(tr)

    @staticmethod
    def _set_cell_text(cell, text: str):
        """设置单元格文本，保持模板样式"""
        tf = cell.text_frame
        if not tf.paragraphs:
            para = tf.add_paragraph()
        else:
            para = tf.paragraphs[0]
        if para.runs:
            run = para.runs[0]
        else:
            run = para.add_run()
        run.text = text
        # 清理其它 run/段落中的旧文本
        for extra_run in para.runs[1:]:
            extra_run.text = ""
        for extra_para in tf.paragraphs[1:]:
            for extra_run in extra_para.runs:
                extra_run.text = ""

    @staticmethod
    def _set_shape_text(shape, text: str):
        """设置形状文本，保持模板样式"""
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        if not tf.paragraphs:
            para = tf.add_paragraph()
        else:
            para = tf.paragraphs[0]
        if para.runs:
            run = para.runs[0]
        else:
            run = para.add_run()
        run.text = text
        # 清理其它内容
        for extra_run in para.runs[1:]:
            extra_run.text = ""
        for extra_para in tf.paragraphs[1:]:
            for extra_run in extra_para.runs:
                extra_run.text = ""
