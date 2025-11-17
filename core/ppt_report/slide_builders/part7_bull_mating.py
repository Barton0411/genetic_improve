"""
Part 7: 公牛性状与选配推荐
包含章节页 + 内容页(标准版+高级版)
"""

import logging
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..config import (
    FONT_NAME_CN, FONT_NAME_EN,
    FONT_SIZE_SECTION_TITLE, FONT_SIZE_SECTION_NUMBER,
    COLOR_BAR_GREEN, COLOR_NUMBER_BG
)

logger = logging.getLogger(__name__)


class Part7BullMatingBuilder:
    """Part 7: 公牛性状与选配推荐"""

    def __init__(self, prs, chart_creator=None, farm_name='示例牧场'):
        self.prs = prs
        self.chart_creator = chart_creator
        self.farm_name = farm_name

    def build(self, data: dict, versions: list = ['standard', 'enhanced']):
        """构建Part 7"""
        logger.info("=" * 60)
        logger.info("开始构建Part 7: 公牛性状与选配推荐")
        logger.info(f"生成版本: {versions}")
        logger.info("=" * 60)

        self._build_section_divider()

        if 'standard' in versions:
            logger.info("构建【标准版】内容页...")
            self._build_content_slide_standard(data)
            logger.info("✓ 【标准版】完成")

        if 'enhanced' in versions:
            logger.info("构建【高级版】内容页...")
            self._build_content_slide_enhanced(data)
            logger.info("✓ 【高级版】完成")

        logger.info("✓ Part 7构建完成")

    def _build_section_divider(self):
        """章节分隔页: 06 公牛性状与选配推荐"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(8.61), Cm(6.26), Cm(0.7), Cm(5.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_BAR_GREEN
        bar.line.fill.background()

        number_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(10.67), Cm(6.26), Cm(3.81), Cm(3.81))
        number_bg.fill.solid()
        number_bg.fill.fore_color.rgb = COLOR_NUMBER_BG
        number_bg.line.fill.background()

        number_box = slide.shapes.add_textbox(Cm(10.67), Cm(6.26), Cm(3.81), Cm(3.81))
        number_para = number_box.text_frame.paragraphs[0]
        number_para.text = "06"
        number_para.alignment = PP_ALIGN.CENTER
        number_para.font.name = FONT_NAME_EN
        number_para.font.size = FONT_SIZE_SECTION_NUMBER
        number_para.font.bold = True
        number_para.font.color.rgb = RGBColor(255, 255, 255)
        number_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        title_box = slide.shapes.add_textbox(Cm(15.8), Cm(6.26), Cm(17.78), Cm(2.03))
        title_para = title_box.text_frame.paragraphs[0]
        title_para.text = "公牛性状与选配推荐"
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = FONT_NAME_CN
        title_para.font.size = FONT_SIZE_SECTION_TITLE
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_box.text_frame.vertical_anchor = MSO_ANCHOR.TOP

        subtitle_box = slide.shapes.add_textbox(Cm(17.88), Cm(9.64), Cm(17.78), Cm(1.02))
        subtitle_para = subtitle_box.text_frame.paragraphs[0]
        subtitle_para.text = "已用公牛 备选公牛 选配推荐"
        subtitle_para.alignment = PP_ALIGN.LEFT
        subtitle_para.font.name = FONT_NAME_CN
        subtitle_para.font.size = Pt(20)
        subtitle_para.font.color.rgb = RGBColor(220, 235, 250)
        subtitle_box.text_frame.vertical_anchor = MSO_ANCHOR.TOP

        logger.info("✓ 章节分隔页创建完成")

    def _build_content_slide_standard(self, data: dict):
        """【标准版】内容页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        title_box = slide.shapes.add_textbox(Cm(2.74), Cm(1.37), Cm(27.94), Cm(1.52))
        title_para = title_box.text_frame.paragraphs[0]
        title_para.text = "公牛性状与选配推荐"
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = FONT_NAME_CN
        title_para.font.size = Pt(36)
        title_para.font.bold = False
        title_para.font.color.rgb = RGBColor(0, 139, 206)
        title_box.text_frame.vertical_anchor = MSO_ANCHOR.TOP

        logger.info("✓ 【标准版】公牛性状与选配推荐页创建完成（待添加数据内容）")

    def _build_content_slide_enhanced(self, data: dict):
        """【高级版】内容页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        title_box = slide.shapes.add_textbox(Cm(2.74), Cm(1.37), Cm(27.94), Cm(1.52))
        title_para = title_box.text_frame.paragraphs[0]
        title_para.text = "公牛性状与选配推荐"
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = FONT_NAME_CN
        title_para.font.size = Pt(36)
        title_para.font.bold = False
        title_para.font.color.rgb = RGBColor(0, 139, 206)
        title_box.text_frame.vertical_anchor = MSO_ANCHOR.TOP

        logger.info("✓ 【高级版】公牛性状与选配推荐页创建完成（待添加数据内容）")
