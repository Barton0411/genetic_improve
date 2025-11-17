"""
Part 1: 封面与目录构建器（新模板 ppt模版-2.pptx）
"""

import logging
from datetime import datetime
from typing import Optional

from pptx.enum.text import PP_ALIGN

from ..base_builder import BaseSlideBuilder

logger = logging.getLogger(__name__)


class Part1CoverBuilder(BaseSlideBuilder):
    """封面与目录构建器"""

    def __init__(self, prs, chart_creator, farm_info: Optional[dict] = None, fallback_farm_name: str = "牧场", fallback_reporter: str = "用户"):
        """
        初始化构建器

        Args:
            prs: PowerPoint演示文稿对象
            chart_creator: 图表创建器
            farm_info: 从Excel解析的牧场基础信息
            fallback_farm_name: 当Excel缺失时使用的默认牧场名
            fallback_reporter: 当Excel缺失时使用的默认汇报人
        """
        super().__init__(prs, chart_creator)
        self._farm_info = farm_info or {}
        self.farm_name = self._farm_info.get('farm_name') or fallback_farm_name
        self.reporter_name = self._farm_info.get('service_staff') or fallback_reporter
        self.report_date = self._farm_info.get('report_date')
        self.report_time_text = self._farm_info.get('report_time')

    def build(self):
        """填充模板中的封面与目录页"""
        logger.info("=" * 60)
        logger.info("开始构建Part 1: 封面与目录（套用模板）")
        logger.info("=" * 60)

        try:
            cover_slide = self.prs.slides[0]
            self._fill_cover_slide(cover_slide)
        except IndexError:
            logger.error("模板中缺少封面页（Slide 1）")

        try:
            toc_slide = self.prs.slides[1]
            self._fill_toc_slide(toc_slide)
        except IndexError:
            logger.error("模板中缺少目录页（Slide 2）")

        logger.info("✓ Part 1模板更新完成")

    def _fill_cover_slide(self, slide):
        """更新封面页文字"""
        logger.info("更新封面页")
        title_shape = self._find_shape_by_name(slide, "Title 6")
        info_shape = self._find_shape_by_name(slide, "Text Placeholder 1")

        title_text = f"{self.farm_name}牧场育种分析综合报告"
        if title_shape and title_shape.has_text_frame:
            tf = title_shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            para.text = title_text
            para.alignment = PP_ALIGN.LEFT if para.alignment is None else para.alignment
        else:
            logger.warning("未找到封面标题形状：Title 6")

        if info_shape and info_shape.has_text_frame:
            tf = info_shape.text_frame
            tf.clear()
            date_para = tf.paragraphs[0]
            date_para.text = self._format_report_date()
            date_para.alignment = PP_ALIGN.RIGHT

            reporter_para = tf.add_paragraph()
            reporter_para.text = f"服务人员：{self.reporter_name}"
            reporter_para.alignment = PP_ALIGN.RIGHT
        else:
            logger.warning("未找到封面信息形状：Text Placeholder 1")

    def _fill_toc_slide(self, slide):
        """目录页暂不做动态替换，保留模板"""
        logger.info("目录页暂时使用模板静态结构，后续可根据页码动态更新。")

    def _find_shape_by_name(self, slide, name: str):
        for shape in slide.shapes:
            if shape.name == name:
                return shape
        return None

    def _format_report_date(self) -> str:
        if self.report_date:
            return self.report_date.strftime("%Y年%m月")
        if self.report_time_text:
            return self.report_time_text
        return datetime.now().strftime("%Y年%m月")

        # 【装饰线】位置(0.80", 1.35")，尺寸0.80" × 0.02"
        title_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.80), Inches(1.35),
            Inches(0.80), Inches(0.02)
        )
        title_line.fill.solid()
        title_line.fill.fore_color.rgb = COLOR_DECORATION  # 蓝色RGB(0, 139, 206)
        title_line.line.fill.background()

        # 目录项配置
        toc_items = [
            ("01", "牧场概况", "基本信息 牛群结构"),
            ("02", "系谱记录分析", "完整性 识别率"),
            ("03", "牛群遗传评估", "性状 指数 进展"),
            ("04", "配种记录分析", "概况 风险 近交"),
            ("05", "公牛使用分析", "已用 备选 评估"),
            ("06", "选配推荐方案", "推荐统计 质量"),
            ("07", "项目总结建议", "发现 建议")
        ]

        # 左列和右列的精确位置（完全按照新模板）
        # 左列：色条3.74"，编号3.94"，文字4.79"，起始1.69"
        # 右列：色条7.14"，编号7.33"，文字8.19"，起始2.23"
        col1_bar_left = 3.74
        col1_num_left = 3.94
        col1_text_left = 4.79
        col1_start_top = 1.69

        col2_bar_left = 7.14
        col2_num_left = 7.33
        col2_text_left = 8.19
        col2_start_top = 2.23

        cell_height = 1.10  # 垂直间距

        for i, (number, title, desc) in enumerate(toc_items):
            # 确定列位置
            if i < 4:  # 左列（01-04）
                bar_left = col1_bar_left
                num_left = col1_num_left
                text_left = col1_text_left
                top = col1_start_top + i * cell_height
            else:  # 右列（05-07）
                bar_left = col2_bar_left
                num_left = col2_num_left
                text_left = col2_text_left
                top = col2_start_top + (i - 4) * cell_height

            # 【色条】0.06" × 0.90" - 黄绿色
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(bar_left), Inches(top),
                Inches(0.06), Inches(0.90)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = COLOR_BAR_GREEN  # 黄绿色RGB(138, 184, 74)
            bar.line.fill.background()

            # 【编号背景框】0.76"宽 × 0.71"高 - 浅蓝灰色背景
            num_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(num_left), Inches(top),
                Inches(0.76), Inches(0.71)
            )
            num_bg.fill.solid()
            num_bg.fill.fore_color.rgb = COLOR_NUMBER_BG  # 浅蓝灰色RGB(174, 200, 223)
            num_bg.line.fill.background()

            # 【编号文字】36pt Arial加粗白色，右对齐
            num_box = slide.shapes.add_textbox(
                Inches(num_left), Inches(top),
                Inches(0.76), Inches(0.71)
            )
            num_para = num_box.text_frame.paragraphs[0]
            num_para.text = number
            num_para.alignment = PP_ALIGN.RIGHT
            num_para.font.name = FONT_NAME_EN
            num_para.font.size = Pt(36)
            num_para.font.bold = True
            num_para.font.color.rgb = RGBColor(255, 255, 255)  # 白色（在浅色背景上）

            # 【标题】20pt微软雅黑加粗白色，左对齐（深蓝背景上）
            title_box = slide.shapes.add_textbox(
                Inches(text_left), Inches(top),
                Inches(1.87), Inches(0.44)
            )
            title_para = title_box.text_frame.paragraphs[0]
            title_para.text = title
            title_para.alignment = PP_ALIGN.LEFT
            title_para.font.name = FONT_NAME_CN
            title_para.font.size = Pt(20)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)  # 白色（深蓝背景）

            # 【描述】14pt微软雅黑浅色，左对齐（深蓝背景上）
            desc_box = slide.shapes.add_textbox(
                Inches(text_left), Inches(top + 0.50),
                Inches(1.82), Inches(0.34)
            )
            desc_para = desc_box.text_frame.paragraphs[0]
            desc_para.text = desc
            desc_para.alignment = PP_ALIGN.LEFT
            desc_para.font.name = FONT_NAME_CN
            desc_para.font.size = Pt(14)
            desc_para.font.bold = False
            desc_para.font.color.rgb = RGBColor(220, 235, 250)  # 浅蓝白色（深蓝背景）

        logger.info("✓ 目录页创建完成（用户修改版本：黄绿色条+浅色编号框+深蓝椭圆背景）")

    def _build_cover_slide_fallback(self, slide):
        """封面页降级方案（手动创建）"""
        # 中央标题区域
        title_box = slide.shapes.add_textbox(
            Inches(0.66), Inches(0.68),
            Inches(8.69), Inches(3.35)
        )
        title_frame = title_box.text_frame
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 牧场名称
        p1 = title_frame.paragraphs[0]
        p1.text = f"{self.farm_name}牧场"
        p1.alignment = PP_ALIGN.CENTER
        p1.font.name = FONT_NAME_CN
        p1.font.size = Pt(48)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_TITLE

        # 报告类型
        p2 = title_frame.add_paragraph()
        p2.text = "育种分析综合报告"
        p2.alignment = PP_ALIGN.CENTER
        p2.font.name = FONT_NAME_CN
        p2.font.size = Pt(28)
        p2.font.color.rgb = COLOR_TEXT_BODY
        p2.space_before = Pt(12)

        # 日期
        date_box = slide.shapes.add_textbox(
            Inches(0.78), Inches(3.42),
            Inches(8.44), Inches(0.52)
        )
        date_para = date_box.text_frame.paragraphs[0]
        date_para.text = datetime.now().strftime("%Y年%m月")
        date_para.alignment = PP_ALIGN.CENTER
        date_para.font.name = FONT_NAME_CN
        date_para.font.size = Pt(18)
        date_para.font.color.rgb = COLOR_TEXT_SECONDARY

        logger.info("✓ 封面页创建完成（降级方案）")
