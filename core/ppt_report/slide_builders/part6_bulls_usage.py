"""
Part 6: 配种记录分析-使用公牛分析构建器
对应PPT模版的"配种记录分析-使用公牛分析"页面（已用公牛性状汇总表）
"""

import logging
from typing import Dict, Optional

import pandas as pd
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from ..base_builder import BaseSlideBuilder
from ..config import FONT_NAME_CN

logger = logging.getLogger(__name__)


class Part6BullsUsageBuilder(BaseSlideBuilder):
    """Part 6: 配种记录分析-使用公牛分析构建器"""

    def __init__(self, prs, chart_creator, farm_name: str):
        super().__init__(prs, chart_creator)
        self.farm_name = farm_name

    def build(self, data: Dict):
        """
        构建 Part 6: 配种记录分析-使用公牛分析

        Args:
            data: 包含 'excel_path' 的数据字典
        """
        logger.info("=" * 60)
        logger.info("开始构建Part 6: 配种记录分析-使用公牛分析")
        logger.info("=" * 60)

        # 读取已用公牛性状汇总Sheet
        excel_path = data.get("excel_path")
        if excel_path is None:
            logger.error("excel_path 未找到，跳过使用公牛分析")
            return

        try:
            df_bulls = pd.read_excel(excel_path, sheet_name="已用公牛性状汇总", header=None)
            logger.info(f"✓ 读取bulls_usage数据（无header）: {len(df_bulls)}行 x {len(df_bulls.columns)}列")
        except Exception as e:
            logger.error(f"读取已用公牛性状汇总Sheet失败: {e}")
            return

        if df_bulls.empty:
            logger.warning("bulls_usage 数据为空DataFrame，跳过使用公牛分析")
            return

        # 查找"配种记录分析-使用公牛分析"页面（只处理第一页：汇总表）
        logger.info(f"开始查找包含'配种记录分析-使用公牛分析'的页面，模板共{len(self.prs.slides)}页")
        target_slides = self.find_slides_by_text("配种记录分析-使用公牛分析", start_index=0, max_count=1)

        if len(target_slides) == 0:
            logger.error("❌ 未找到配种记录分析-使用公牛分析页面，跳过")
            return

        logger.info(f"✓ 找到目标页面: 第{target_slides[0]+1}页")

        # 填充汇总表页面
        self._fill_summary_table(df_bulls, target_slides[0])

        logger.info("✓ Part 6 配种记录分析-使用公牛分析完成")

    def _fill_summary_table(self, df: pd.DataFrame, slide_index: int):
        """
        填充已用公牛性状汇总表（按年份）

        Excel结构:
        - 行0: 标题 "已用公牛性状汇总表（按年份）"
        - 行1: 表头（列名）
        - 行2-N: 年份数据（动态行数）
        - 最后一行: 总平均

        Args:
            df: 已用公牛性状汇总数据（header=None读取）
            slide_index: 幻灯片索引
        """
        try:
            slide = self.prs.slides[slide_index]
        except IndexError:
            logger.warning(f"无法访问第{slide_index + 1}页")
            return

        logger.info(f"填充第{slide_index + 1}页 - 已用公牛性状汇总表（按年份）")

        # 读取标题
        title = str(df.iloc[0, 0]).strip() if not pd.isna(df.iloc[0, 0]) else "已用公牛性状汇总表（按年份）"
        logger.info(f"标题: {title}")

        # 读取表头
        header_row = self._read_row_data(df, 1)
        logger.info(f"表头: {len(header_row)}列")

        # 读取数据行（从行2开始，直到遇到包含"总平均"或"总计"的行）
        data_rows = []
        summary_row = None
        for i in range(2, len(df)):
            cell_0 = str(df.iloc[i, 0]).strip() if not pd.isna(df.iloc[i, 0]) else ""

            if "总平均" in cell_0 or "总计" in cell_0:
                # 这是汇总行
                summary_row = self._read_row_data(df, i)
                logger.info(f"找到汇总行: 行{i}, 标识: {cell_0}")
                break

            if cell_0 == "" or pd.isna(df.iloc[i, 0]):
                # 遇到空行，停止读取
                break

            # 读取数据行
            row_data = self._read_row_data(df, i)
            data_rows.append(row_data)

        logger.info(f"读取数据: {len(data_rows)}行年份数据")
        if summary_row:
            logger.info(f"读取汇总行")

        # 填充PPT中的表格
        self._fill_ppt_table(slide, title, header_row, data_rows, summary_row)

        # 填充分析文本
        self._fill_analysis_text(slide, data_rows)

    def _read_row_data(self, df: pd.DataFrame, row_idx: int) -> list:
        """
        读取指定行的所有列数据

        Returns:
            list: [列0值, 列1值, ...]
        """
        row_data = []
        for j in range(df.shape[1]):
            cell_value = df.iloc[row_idx, j]
            if pd.isna(cell_value):
                row_data.append("")
            else:
                # 数值类型格式化
                if isinstance(cell_value, (int, float)):
                    # 前三列（年份、使用公牛数、配种头次）保留整数
                    if j in [0, 1, 2]:
                        row_data.append(str(int(cell_value)))
                    else:
                        # 其他列保留2位小数
                        row_data.append(f"{cell_value:.2f}")
                else:
                    row_data.append(str(cell_value).strip())
        return row_data

    def _fill_ppt_table(self, slide, title: str, header: list, data_rows: list, summary_row: list):
        """
        填充PPT中的表格

        PPT表格结构:
        - 表格12: 标题（1x1）
        - 表格2: 数据表（6行x17列，但行数和列数都是动态的）
        """
        # 查找表格
        table_title = None
        table_data = None

        for shape in slide.shapes:
            if shape.shape_type == 19:  # TABLE
                if shape.name == "表格 12":
                    table_title = shape.table
                elif shape.name == "表格 2":
                    table_data = shape.table

        # 填充标题表格
        if table_title:
            try:
                cell = table_title.cell(0, 0)
                cell.text = title
                logger.info(f"✓ 填充标题表格: {title}")
            except Exception as e:
                logger.error(f"填充标题表格失败: {e}")

        # 填充数据表格
        if table_data:
            self._fill_data_table(table_data, header, data_rows, summary_row)
        else:
            logger.warning("未找到数据表格（表格 2）")

    def _fill_data_table(self, table, header: list, data_rows: list, summary_row: list):
        """
        填充数据表格

        注意：
        1. 行数是动态的（需要根据年份数量调整）
        2. 列数也是动态的（需要根据Excel数据调整）
        """
        logger.info(f"数据表格原始大小: {len(table.rows)}行 x {len(table.columns)}列")
        logger.info(f"需要的列数: {len(header)}列")

        # 计算需要的列数
        required_columns = len(header)

        # 调整表格列数
        if len(table.columns) < required_columns:
            cols_to_add = required_columns - len(table.columns)
            logger.info(f"需要添加 {cols_to_add} 列")
            for _ in range(cols_to_add):
                self._add_table_column(table)
        elif len(table.columns) > required_columns:
            cols_to_delete = len(table.columns) - required_columns
            logger.info(f"需要删除 {cols_to_delete} 列")
            for _ in range(cols_to_delete):
                self._delete_table_column(table, len(table.columns) - 1)

        # 计算需要的行数
        header_rows = 1
        required_data_rows = len(data_rows)
        required_summary_rows = 1 if summary_row else 0
        required_total_rows = header_rows + required_data_rows + required_summary_rows

        # 调整表格行数
        if len(table.rows) < required_total_rows:
            rows_to_add = required_total_rows - len(table.rows)
            logger.info(f"需要添加 {rows_to_add} 行")
            for _ in range(rows_to_add):
                self._add_table_row(table)
        elif len(table.rows) > required_total_rows:
            rows_to_delete = len(table.rows) - required_total_rows
            logger.info(f"需要删除 {rows_to_delete} 行")
            for _ in range(rows_to_delete):
                self._delete_table_row(table, len(table.rows) - 1)

        # 填充数据
        try:
            # 填充表头（行0）
            self._fill_table_row(table, 0, header)

            # 填充数据行
            for i, row_data in enumerate(data_rows):
                self._fill_table_row(table, i + 1, row_data)

            # 填充汇总行
            if summary_row:
                self._fill_table_row(table, len(data_rows) + 1, summary_row)

            logger.info(f"✓ 填充数据表格: {len(data_rows)}行数据 + 1行汇总")

        except Exception as e:
            logger.error(f"填充数据表格失败: {e}", exc_info=True)

    def _fill_table_row(self, table, row_idx: int, row_data: list):
        """填充表格的一行"""
        for col_idx, cell_value in enumerate(row_data):
            if col_idx >= len(table.columns):
                # 超出表格列数，跳过
                break

            try:
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_value)

                # 设置格式
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.name = FONT_NAME_CN
                        run.font.size = Pt(10)
            except Exception as e:
                logger.debug(f"填充单元格 ({row_idx}, {col_idx}) 失败: {e}")

    def _fill_analysis_text(self, slide, data_rows: list):
        """
        填充分析文本

        根据年份数据生成详细分析
        """
        if not data_rows:
            return

        # 查找分析文本框（优先位置靠下的文本框）
        textbox = None
        candidate_textboxes = []

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                name = shape.name if shape.name else ""
                shape_text = shape.text if hasattr(shape, "text") else ""
                # 跳过标题
                if "配种记录分析" in shape_text or "使用公牛分析" in shape_text:
                    continue
                if "文本框" in name or "TextBox" in name:
                    candidate_textboxes.append(shape)

        # 按位置排序，选择最靠下的
        if candidate_textboxes:
            candidate_textboxes.sort(key=lambda s: getattr(s, 'top', 0), reverse=True)
            textbox = candidate_textboxes[0]

        if not textbox:
            logger.warning("未找到分析文本框")
            return

        try:
            analysis = "分析：" + self._generate_usage_analysis(data_rows)

            # 清空文本框并设置新内容
            tf = textbox.text_frame
            tf.clear()
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            run = p.add_run()
            run.text = analysis

            # 设置字体：微软雅黑15号非加粗
            run.font.name = FONT_NAME_CN
            run.font.size = Pt(15)
            run.font.bold = False

            logger.info(f"✓ 填充使用公牛分析文本")

        except Exception as e:
            logger.error(f"填充分析文本失败: {e}")

    def _generate_usage_analysis(self, data_rows: list) -> str:
        """
        生成已用公牛汇总分析文本

        Args:
            data_rows: 年份数据行列表，每行格式：[年份, 公牛数, 配种次数, NM$, TPI, ...]

        Returns:
            分析文本字符串
        """
        if not data_rows:
            return "暂无使用公牛数据。"

        try:
            years = [row[0] for row in data_rows if row]
            if not years:
                return "暂无使用公牛数据。"

            parts = []

            # 基础信息
            if len(years) >= 2:
                parts.append(f"牧场近{len(years)}年（{years[0]}-{years[-1]}）冻精使用情况：")
            else:
                parts.append(f"牧场{years[0]}年冻精使用情况：")

            # 提取公牛数和配种次数趋势
            bull_counts = []
            breeding_counts = []
            nm_values = []
            tpi_values = []

            for row in data_rows:
                if len(row) > 1:
                    try:
                        bull_counts.append(int(row[1]) if row[1] else 0)
                    except (ValueError, TypeError):
                        pass
                if len(row) > 2:
                    try:
                        breeding_counts.append(int(row[2]) if row[2] else 0)
                    except (ValueError, TypeError):
                        pass
                if len(row) > 3:
                    try:
                        nm_values.append(float(row[3]) if row[3] else 0)
                    except (ValueError, TypeError):
                        pass
                if len(row) > 4:
                    try:
                        tpi_values.append(float(row[4]) if row[4] else 0)
                    except (ValueError, TypeError):
                        pass

            # 公牛使用数量分析
            if bull_counts:
                avg_bulls = sum(bull_counts) / len(bull_counts)
                if len(bull_counts) >= 2:
                    trend = bull_counts[-1] - bull_counts[0]
                    if trend > 0:
                        parts.append(f"使用公牛数量呈上升趋势，年均使用{avg_bulls:.0f}头公牛。")
                    elif trend < 0:
                        parts.append(f"使用公牛数量有所减少，年均使用{avg_bulls:.0f}头公牛。")
                    else:
                        parts.append(f"使用公牛数量保持稳定，年均{avg_bulls:.0f}头。")
                else:
                    parts.append(f"年使用公牛{bull_counts[0]}头。")

            # NM$趋势分析
            if nm_values and len(nm_values) >= 2:
                nm_change = nm_values[-1] - nm_values[0]
                if nm_change > 50:
                    parts.append(f"公牛NM$水平持续提升，近年平均达{nm_values[-1]:.0f}。")
                elif nm_change > 0:
                    parts.append(f"公牛NM$水平稳步提高。")
                elif nm_change < -50:
                    parts.append("公牛NM$水平有所下降，建议优化选种策略。")

            # TPI趋势分析
            if tpi_values and len(tpi_values) >= 2:
                tpi_change = tpi_values[-1] - tpi_values[0]
                if tpi_change > 100:
                    parts.append(f"TPI指数明显提升，表明公牛综合品质改善。")

            return "".join(parts)

        except Exception as e:
            logger.warning(f"生成使用公牛分析失败: {e}")
            return "牧场冻精使用情况稳定。"

    def _add_table_column(self, table):
        """添加表格列"""
        try:
            from copy import deepcopy

            tbl = table._tbl

            # 获取所有行（直接遍历）
            trs = []
            for child in tbl:
                if 'tr' in child.tag.lower():
                    trs.append(child)

            # 在每一行中添加一个新单元格
            for tr in trs:
                # 找到该行的所有单元格
                tcs = []
                for child in tr:
                    if 'tc' in child.tag.lower():
                        tcs.append(child)

                if tcs:
                    # 复制最后一个单元格的结构
                    last_tc = tcs[-1]
                    new_tc = deepcopy(last_tc)
                    # 添加到行的末尾
                    tr.append(new_tc)

            # 更新表格的gridCol定义
            tblGrid = None
            for child in tbl:
                if 'tblgrid' in child.tag.lower():
                    tblGrid = child
                    break

            if tblGrid is not None:
                # 获取所有gridCol
                gridCols = []
                for child in tblGrid:
                    if 'gridcol' in child.tag.lower():
                        gridCols.append(child)

                if gridCols:
                    # 复制最后一个gridCol
                    last_gridCol = gridCols[-1]
                    new_gridCol = deepcopy(last_gridCol)
                    tblGrid.append(new_gridCol)

        except Exception as e:
            logger.warning(f"添加表格列失败: {e}")

    def _delete_table_column(self, table, col_idx: int):
        """删除表格中的一列"""
        try:
            tbl = table._tbl

            # 获取所有行
            trs = []
            for child in tbl:
                if 'tr' in child.tag.lower():
                    trs.append(child)

            # 在每一行中删除指定索引的单元格
            for tr in trs:
                tcs = []
                for child in tr:
                    if 'tc' in child.tag.lower():
                        tcs.append(child)

                if col_idx < len(tcs):
                    tr.remove(tcs[col_idx])

            # 更新表格的gridCol定义
            tblGrid = None
            for child in tbl:
                if 'tblgrid' in child.tag.lower():
                    tblGrid = child
                    break

            if tblGrid is not None:
                gridCols = []
                for child in tblGrid:
                    if 'gridcol' in child.tag.lower():
                        gridCols.append(child)

                if col_idx < len(gridCols):
                    tblGrid.remove(gridCols[col_idx])

        except Exception as e:
            logger.warning(f"删除表格列{col_idx}失败: {e}")

    def _add_table_row(self, table):
        """添加表格行（复制普通数据行格式，而非汇总行）"""
        try:
            from copy import deepcopy
            tbl = table._tbl

            # 复制第2行（索引1，通常是第一个数据行）而不是最后一行（汇总行）
            # 这样可以确保新增行使用普通数据行格式（白色背景），而不是汇总行格式（灰色背景）
            template_row_idx = 1 if len(table.rows) > 1 else 0
            template_tr = table.rows[template_row_idx]._tr
            new_tr = deepcopy(template_tr)
            # 添加到表格
            tbl.append(new_tr)

            # 通过 python-pptx API 清空新行的文本
            new_row_idx = len(table.rows) - 1
            for col_idx in range(len(table.columns)):
                try:
                    cell = table.cell(new_row_idx, col_idx)
                    cell.text = ""
                except:
                    pass
        except Exception as e:
            logger.warning(f"添加表格行失败: {e}")

    def _delete_table_row(self, table, row_idx: int):
        """删除表格中的一行"""
        try:
            tbl = table._tbl
            tr = table.rows[row_idx]._tr
            tbl.remove(tr)
        except Exception as e:
            logger.warning(f"删除表格行{row_idx}失败: {e}")
