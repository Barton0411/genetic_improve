"""
模板幻灯片复制工具
"""

import logging
from pathlib import Path
from copy import deepcopy
from typing import Optional

from pptx import Presentation

logger = logging.getLogger(__name__)


def copy_template_slide(prs, template_slide_index: int) -> Optional[int]:
    """
    从模板文件复制指定幻灯片到当前演示文稿末尾

    Args:
        prs: 当前的Presentation对象
        template_slide_index: 模板中的幻灯片索引（0-based）

    Returns:
        新幻灯片的索引，失败返回None
    """
    try:
        # 查找模板文件
        program_root = Path(__file__).parent.parent.parent.parent
        preferred_template = program_root / "牧场牧场育种分析报告-模版.pptx"
        fallback_template = program_root / "PPT模版.pptx"

        template_path = preferred_template if preferred_template.exists() else fallback_template

        if not template_path.exists():
            logger.error("❌ 未找到PPT模板文件")
            return None

        # 加载模板
        template_prs = Presentation(str(template_path))

        if template_slide_index >= len(template_prs.slides):
            logger.error(f"❌ 模板中不存在索引{template_slide_index}的幻灯片")
            return None

        # 获取模板幻灯片
        source_slide = template_prs.slides[template_slide_index]

        # 获取布局
        blank_layout = source_slide.slide_layout

        # 创建新幻灯片
        new_slide = prs.slides.add_slide(blank_layout)

        # 复制所有形状
        for shape in source_slide.shapes:
            el = shape.element
            newel = deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

        new_index = len(prs.slides) - 1
        logger.info(f"✓ 从模板复制第{template_slide_index + 1}页到当前PPT第{new_index + 1}页")
        return new_index

    except Exception as e:
        logger.error(f"❌ 复制模板幻灯片失败: {e}", exc_info=True)
        return None
