"""
Part 5: 配种记录分析-近交分析构建器
对应PPT模版的"配种记录分析-近交分析"两页
"""

import logging
from typing import Dict, Optional

import pandas as pd
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData

from ..base_builder import BaseSlideBuilder
from ..config import COLOR_TEXT_MAIN, FONT_NAME_CN

logger = logging.getLogger(__name__)


class Part5InbreedingBuilder(BaseSlideBuilder):
    """Part 5: 配种记录分析-近交分析构建器"""

    def __init__(self, prs, chart_creator, farm_name: str):
        super().__init__(prs, chart_creator)
        self.farm_name = farm_name

    def build(self, data: Dict):
        """
        构建 Part 5: 配种记录分析-近交分析

        Args:
            data: 包含 'breeding_inbreeding' 的数据字典
        """
        logger.info("=" * 60)
        logger.info("开始构建Part 5: 配种记录分析-近交分析")
        logger.info("=" * 60)

        # 重新读取配种记录-近交系数分析Sheet，使用header=None以获取完整数据
        excel_path = data.get("excel_path")
        if excel_path is None:
            logger.error("excel_path 未找到，跳过近交分析")
            return

        try:
            df_inbreeding = pd.read_excel(excel_path, sheet_name="配种记录-近交系数分析", header=None)
            logger.info(f"✓ 读取breeding_inbreeding数据（无header）: {len(df_inbreeding)}行 x {len(df_inbreeding.columns)}列")
        except Exception as e:
            logger.error(f"重新读取配种记录-近交系数分析Sheet失败: {e}")
            return

        if df_inbreeding.empty:
            logger.warning("breeding_inbreeding 数据为空DataFrame，跳过近交分析")
            return

        # 查找标题为"配种记录分析-近交分析"的页面，应该有2个
        logger.info(f"开始查找包含'配种记录分析-近交分析'的页面，模板共{len(self.prs.slides)}页")
        target_slides = self.find_slides_by_text("配种记录分析-近交分析", start_index=0, max_count=2)

        if len(target_slides) == 0:
            logger.error("❌ 未找到配种记录分析-近交分析页面，跳过近交分析")
            return

        if len(target_slides) < 2:
            logger.warning(f"⚠️ 只找到{len(target_slides)}个配种记录分析-近交分析页面，预期2个")
        else:
            logger.info(f"✓ 找到2个目标页面: 第{target_slides[0]+1}页, 第{target_slides[1]+1}页")

        # 填充第一页 - 近交系数分布汇总表
        if len(target_slides) > 0:
            self._fill_distribution_summary(df_inbreeding, target_slides[0])

        # 填充第二页 - 按年份近交系数趋势
        if len(target_slides) > 1:
            self._fill_yearly_trend(df_inbreeding, target_slides[1])

        logger.info("✓ Part 5 配种记录分析-近交分析完成")

    def _fill_distribution_summary(self, df: pd.DataFrame, slide_index: int):
        """
        填充第一页：近交系数分布汇总表

        Args:
            df: 近交系数数据（header=None读取）
            slide_index: 幻灯片索引
        """
        try:
            slide = self.prs.slides[slide_index]
        except IndexError:
            logger.warning(f"无法访问第{slide_index + 1}页")
            return

        logger.info(f"填充第{slide_index + 1}页 - 近交系数分布汇总表")

        # Excel结构：
        # 行0: 标题1（全部年份）
        # 行1: 表头1
        # 行2-6: 5行数据（4个区间+总计）

        # 读取第一个表的标题
        title1 = str(df.iloc[0, 0]).strip() if not pd.isna(df.iloc[0, 0]) else "近交系数分布汇总表（全部配次）"
        logger.info(f"第一个表标题: {title1}")

        # 读取第一个表的数据（5行：4个区间+总计）
        table1_data = self._read_table_data(df, 2, num_rows=5)
        logger.info(f"读取第一个表数据: {len(table1_data)}行")

        # 查找第二个标题行
        title2_row_idx = self._find_second_inbreeding_title(df)
        if title2_row_idx is None:
            logger.warning("未找到第二个近交分布表标题")
            title2 = "近交系数分布汇总表（近12个月）"
            table2_data = []
        else:
            title2 = str(df.iloc[title2_row_idx, 0]).strip()
            logger.info(f"第二个表标题: {title2} (行{title2_row_idx})")
            # 读取第二个表的数据
            table2_data = self._read_table_data(df, title2_row_idx + 2, num_rows=5)
            logger.info(f"读取第二个表数据: {len(table2_data)}行")

        # 填充PPT中的表格
        # 根据PPT截图，有两个表格和两个图表
        # 需要填充：表格1（全部年份）、表格2（近12个月）
        self._fill_distribution_tables(slide, title1, table1_data, title2, table2_data)

    def _fill_yearly_trend(self, df: pd.DataFrame, slide_index: int):
        """
        填充第二页：按年份近交系数趋势

        Args:
            df: 近交系数数据
            slide_index: 幻灯片索引
        """
        try:
            slide = self.prs.slides[slide_index]
        except IndexError:
            logger.warning(f"无法访问第{slide_index + 1}页")
            return

        logger.info(f"填充第{slide_index + 1}页 - 按年份近交系数趋势")

        # 查找"按年份近交系数趋势"标题行
        title_row_idx = self._find_yearly_trend_title(df)
        if title_row_idx is None:
            logger.warning("未找到按年份近交系数趋势标题")
            return

        logger.info(f"找到年份趋势标题: 行{title_row_idx}")

        # 读取年份数据（从标题+2行开始，动态读取直到遇到空行）
        yearly_data = self._read_yearly_data(df, title_row_idx + 2)
        logger.info(f"读取年份数据: {len(yearly_data)}行")

        if not yearly_data:
            logger.warning("年份数据为空")
            return

        # 填充PPT中的年份趋势表格
        self._fill_yearly_table(slide, yearly_data)

    # ------------------------------------------------------------------ #
    # 辅助方法
    # ------------------------------------------------------------------ #

    def _find_second_inbreeding_title(self, df: pd.DataFrame) -> Optional[int]:
        """查找第二个近交分布表标题行"""
        found_first = False
        for i in range(len(df)):
            cell_value = str(df.iloc[i, 0]).strip() if not pd.isna(df.iloc[i, 0]) else ""
            if "近交系数分布汇总表" in cell_value:
                if found_first:
                    return i
                else:
                    found_first = True
        return None

    def _find_yearly_trend_title(self, df: pd.DataFrame) -> Optional[int]:
        """查找按年份近交系数趋势标题行"""
        for i in range(len(df)):
            cell_value = str(df.iloc[i, 0]).strip() if not pd.isna(df.iloc[i, 0]) else ""
            if "按年份近交系数趋势" in cell_value:
                return i
        return None

    def _read_table_data(self, df: pd.DataFrame, start_row: int, num_rows: int) -> list:
        """
        读取表格数据

        Returns:
            list of lists: [[区间, 配种次数, 占比, 风险等级], ...]
        """
        data = []
        for i in range(start_row, min(start_row + num_rows, len(df))):
            row_data = []
            for j in range(4):  # 4列：区间、配种次数、占比、风险等级
                if j < df.shape[1]:
                    cell_value = df.iloc[i, j]
                    if pd.isna(cell_value):
                        row_data.append("")
                    else:
                        row_data.append(str(cell_value).strip())
                else:
                    row_data.append("")
            data.append(row_data)
        return data

    def _read_yearly_data(self, df: pd.DataFrame, start_row: int) -> list:
        """
        读取年份数据（动态读取直到遇到空行）

        Returns:
            list of lists: [[年份, 总配种次数, 高风险配种数, 高风险占比, 极高风险配种数, 极高风险占比], ...]
        """
        data = []
        for i in range(start_row, len(df)):
            # 检查第一列是否为空
            cell_0 = df.iloc[i, 0]
            if pd.isna(cell_0):
                break  # 遇到空行，停止读取

            row_data = []
            for j in range(6):  # 6列
                if j < df.shape[1]:
                    cell_value = df.iloc[i, j]
                    if pd.isna(cell_value):
                        row_data.append("")
                    else:
                        row_data.append(str(cell_value).strip())
                else:
                    row_data.append("")
            data.append(row_data)

        return data

    def _fill_distribution_tables(self, slide, title1: str, data1: list, title2: str, data2: list):
        """
        填充分布汇总表（两个表格和两个图表）

        根据PPT模板：
        - 表格 12：标题
        - 表格 2：左侧表格（全部年份）
        - 表格 5：右侧表格（近12个月）
        - Chart 1 (左)：全部年份柱状图
        - Chart 1 (右)：近12个月柱状图
        """
        # 查找PPT中的表格和图表
        tables = []
        charts = []

        for shape in slide.shapes:
            if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                tables.append(shape)
            elif shape.has_chart:
                charts.append(shape)

        logger.info(f"找到{len(tables)}个表格, {len(charts)}个图表")

        # 按名称查找特定表格
        table_title = None
        table_left = None
        table_right = None

        for shape in slide.shapes:
            if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                if shape.name == "表格 12":
                    table_title = shape.table
                elif shape.name == "表格 2":
                    table_left = shape.table
                elif shape.name == "表格 5":
                    table_right = shape.table

        # 填充标题表格
        if table_title:
            try:
                cell = table_title.cell(0, 0)
                cell.text = "近交系数分布汇总表"
                logger.info("✓ 填充标题表格")
            except Exception as e:
                logger.error(f"填充标题表格失败: {e}")

        # 填充左侧表格（全部年份）
        if table_left:
            self._fill_single_distribution_table(table_left, title1, data1, "左侧")
        else:
            logger.warning("未找到左侧表格（表格 2）")

        # 填充右侧表格（近12个月）
        if table_right:
            self._fill_single_distribution_table(table_right, title2, data2, "右侧")
        else:
            logger.warning("未找到右侧表格（表格 5）")

        # 填充图表
        if len(charts) >= 2:
            # 第一个图表：全部年份
            self._fill_distribution_chart(charts[0], data1, "全部年份")
            # 第二个图表：近12个月
            self._fill_distribution_chart(charts[1], data2, "近12个月")
        else:
            logger.warning(f"图表数量不足: 找到{len(charts)}个，需要2个")

        # 填充分析文本（使用全部年份数据）
        self._fill_inbreeding_analysis_text(slide, data1, is_all_years=True)

    def _fill_single_distribution_table(self, table, title: str, data: list, table_index: int):
        """填充单个分布表格"""
        try:
            # 表格结构假设：
            # 第一行可能是合并单元格的标题，或者没有标题
            # 数据从某一行开始

            # 简单方法：假设表头已存在，从第2行开始填充数据
            start_row = 1  # 跳过表头
            for row_idx, row_data in enumerate(data):
                actual_row_idx = start_row + row_idx
                if actual_row_idx >= len(table.rows):
                    logger.warning(f"表格{table_index}行数不足: {actual_row_idx} >= {len(table.rows)}")
                    break

                for col_idx, cell_value in enumerate(row_data):
                    if col_idx >= len(table.columns):
                        break

                    cell = table.cell(actual_row_idx, col_idx)
                    cell.text = cell_value

                    # 设置格式
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.name = FONT_NAME_CN
                            run.font.size = Pt(11)

            logger.info(f"✓ 填充表格{table_index}: {len(data)}行数据")

        except Exception as e:
            logger.error(f"填充表格{table_index}失败: {e}")

    def _fill_yearly_table(self, slide, yearly_data: list):
        """
        填充年份趋势表格

        注意：表格行数是动态的，需要根据实际数据调整
        """
        # 查找PPT中的表格
        table = None
        for shape in slide.shapes:
            if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                table = shape.table
                break

        if table is None:
            logger.warning("未找到年份趋势表格")
            return

        logger.info(f"找到年份趋势表格: {len(table.rows)}行 x {len(table.columns)}列")
        logger.info(f"需要填充: {len(yearly_data)}行年份数据")

        # 检查表格行数是否足够
        header_rows = 1  # 假设有1行表头
        required_rows = header_rows + len(yearly_data)

        if len(table.rows) < required_rows:
            # 需要添加行
            rows_to_add = required_rows - len(table.rows)
            logger.info(f"需要添加 {rows_to_add} 行")
            for _ in range(rows_to_add):
                self._add_table_row(table)
        elif len(table.rows) > required_rows:
            # 需要删除多余的行
            rows_to_delete = len(table.rows) - required_rows
            logger.info(f"需要删除 {rows_to_delete} 行")
            # 从后往前删除
            for _ in range(rows_to_delete):
                self._delete_table_row(table, len(table.rows) - 1)

        # 填充数据
        try:
            start_row = header_rows
            for row_idx, row_data in enumerate(yearly_data):
                actual_row_idx = start_row + row_idx
                if actual_row_idx >= len(table.rows):
                    break

                for col_idx, cell_value in enumerate(row_data):
                    if col_idx >= len(table.columns):
                        break

                    cell = table.cell(actual_row_idx, col_idx)
                    cell.text = cell_value

                    # 设置格式
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.name = FONT_NAME_CN
                            run.font.size = Pt(12)

            logger.info(f"✓ 填充年份趋势表格: {len(yearly_data)}行数据")

        except Exception as e:
            logger.error(f"填充年份趋势表格失败: {e}")

    def _add_table_row(self, table):
        """添加表格行"""
        try:
            from copy import deepcopy
            tbl = table._tbl
            # 复制最后一行的结构
            last_row_idx = len(table.rows) - 1
            last_tr = table.rows[last_row_idx]._tr
            new_tr = deepcopy(last_tr)
            # 添加到表格
            tbl.append(new_tr)

            # 通过 python-pptx API 清空新行的文本
            new_row_idx = len(table.rows) - 1
            for col_idx in range(len(table.columns)):
                try:
                    cell = table.cell(new_row_idx, col_idx)
                    cell.text = ""
                except:
                    pass
        except Exception as e:
            logger.warning(f"添加表格行失败: {e}")

    def _delete_table_row(self, table, row_idx: int):
        """删除表格中的一行"""
        try:
            tbl = table._tbl
            tr = table.rows[row_idx]._tr
            tbl.remove(tr)
        except Exception as e:
            logger.warning(f"删除表格行{row_idx}失败: {e}")

    def _fill_distribution_chart(self, chart_shape, data: list, chart_label: str):
        """
        填充近交系数分布柱状图

        Args:
            chart_shape: 图表形状对象
            data: 表格数据 [[区间, 配种次数, 占比, 风险等级], ...]
            chart_label: 图表标签（用于日志）
        """
        try:
            if not data or len(data) < 4:
                logger.warning(f"图表{chart_label}数据不足: {len(data)}行，需要至少4行")
                return

            chart = chart_shape.chart

            # 提取占比数据（列索引2）
            percentages = []
            for row in data[:4]:  # 只取前4行（4个区间）
                if len(row) > 2:
                    percentage_str = str(row[2]).strip()
                    # 移除百分号并转换为浮点数
                    percentage_str = percentage_str.replace('%', '').strip()
                    try:
                        percentage_value = float(percentage_str)
                        percentages.append(percentage_value)
                    except ValueError:
                        logger.warning(f"无法解析占比值: {row[2]}")
                        percentages.append(0.0)
                else:
                    percentages.append(0.0)

            # 确保有4个数据点
            while len(percentages) < 4:
                percentages.append(0.0)

            logger.info(f"图表{chart_label}数据: {percentages}")

            # 更新图表数据
            # 获取图表的第一个系列（应该只有一个系列）
            if len(chart.series) == 0:
                logger.warning(f"图表{chart_label}没有数据系列")
                return

            series = chart.series[0]

            # 更新系列数据
            # python-pptx的chart_data需要重新构建
            # 创建新的图表数据
            chart_data = CategoryChartData()
            chart_data.categories = ['< 3.125%', '3.125% - 6.25%', '6.25% - 12.5%', '> 12.5%']
            chart_data.add_series('占比(%)', percentages)

            # 替换图表数据
            chart.replace_data(chart_data)

            logger.info(f"✓ 填充图表{chart_label}: {percentages}")

        except Exception as e:
            logger.error(f"填充图表{chart_label}失败: {e}", exc_info=True)

    def _fill_inbreeding_analysis_text(self, slide, data: list, is_all_years: bool = True):
        """
        填充近交系数分析文本框

        Args:
            slide: 幻灯片对象
            data: 分布数据 [[区间, 配种次数, 占比, 风险等级], ...]
            is_all_years: 是否为全部年份数据
        """
        # 查找分析文本框（避免匹配到标题）
        textbox = None
        candidate_textboxes = []

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                name = shape.name if shape.name else ""
                shape_text = shape.text if hasattr(shape, "text") else ""

                # 跳过标题：页面标题通常包含"配种记录分析"且位置靠上
                if "配种记录分析" in shape_text:
                    continue
                # 跳过表格
                if "表格" in name:
                    continue
                # 跳过图表
                if shape.shape_type == 3:  # CHART
                    continue

                # 优先查找命名的文本框
                if name in ["文本框 13", "文本框 14", "文本框 15"]:
                    candidate_textboxes.append((shape, 100))  # 高优先级
                elif "文本框" in name or "TextBox" in name:
                    candidate_textboxes.append((shape, 50))  # 中优先级

        # 按优先级排序，优先级相同时按位置（top）排序（选择最靠下的）
        if candidate_textboxes:
            candidate_textboxes.sort(key=lambda x: (-x[1], -getattr(x[0], 'top', 0)))
            textbox = candidate_textboxes[0][0]

        if textbox is None:
            logger.debug("未找到近交分析文本框，跳过")
            return

        try:
            # 计算统计数据
            total_count = 0
            high_risk_count = 0  # 6.25%-12.5%
            very_high_risk_count = 0  # >12.5%

            for row in data:
                if len(row) >= 2:
                    try:
                        count = int(str(row[1]).replace(',', '')) if row[1] else 0
                    except:
                        count = 0

                    total_count += count

                    interval = str(row[0]) if row[0] else ""
                    if "6.25%" in interval and "12.5%" in interval:
                        high_risk_count = count
                    elif "> 12.5%" in interval or ">12.5%" in interval or "12.5%以上" in interval:
                        very_high_risk_count = count

            # 计算风险占比
            risk_pct = 0
            if total_count > 0:
                risk_pct = (high_risk_count + very_high_risk_count) / total_count * 100

            # 生成分析文本
            period = "全部年份" if is_all_years else "近12个月"
            parts = [f"近交系数分析（{period}）："]

            if total_count > 0:
                parts.append(f"牧场配种记录共{total_count}次，")

                if very_high_risk_count > 0:
                    parts.append(
                        f"其中极高风险（>12.5%）配次{very_high_risk_count}次，"
                        f"高风险（6.25%-12.5%）配次{high_risk_count}次，"
                        f"合计风险配次占比{risk_pct:.1f}%。"
                    )
                    parts.append("建议加强近交系数筛查，避免近亲配种。")
                elif high_risk_count > 0:
                    parts.append(
                        f"高风险（6.25%-12.5%）配次{high_risk_count}次，占比{risk_pct:.1f}%。"
                    )
                    parts.append("建议关注近交系数，选择血缘关系较远的公牛配种。")
                else:
                    parts.append("近交系数控制良好，风险配次占比较低，继续保持。")
            else:
                parts.append("暂无足够配种记录数据。")

            analysis_text = "分析：" + "".join(parts)

            # 填充文本
            textbox.text_frame.clear()
            p = textbox.text_frame.paragraphs[0] if textbox.text_frame.paragraphs else textbox.text_frame.add_paragraph()
            run = p.add_run()
            run.text = analysis_text
            p.alignment = PP_ALIGN.LEFT

            # 设置字体：微软雅黑15号非加粗
            run.font.name = FONT_NAME_CN
            run.font.size = Pt(15)
            run.font.bold = False

            logger.info(f"✓ 填充近交分析文本: 总配次{total_count}, 风险占比{risk_pct:.1f}%")

        except Exception as e:
            logger.error(f"填充近交分析文本失败: {e}")
