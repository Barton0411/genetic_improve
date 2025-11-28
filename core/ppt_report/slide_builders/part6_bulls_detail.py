"""
Part 6: 配种记录分析-使用公牛分析-已用公牛明细构建器
对应PPT模版的"配种记录分析-使用公牛分析"多页（每年一页）
"""

import logging
from typing import Dict, List, Tuple
import re

import pandas as pd
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from ..base_builder import BaseSlideBuilder
from ..config import FONT_NAME_CN

logger = logging.getLogger(__name__)


class Part6BullsDetailBuilder(BaseSlideBuilder):
    """Part 6: 配种记录分析-使用公牛分析-已用公牛明细构建器"""

    def __init__(self, prs, chart_creator, farm_name: str):
        super().__init__(prs, chart_creator)
        self.farm_name = farm_name

    def build(self, data: Dict):
        """
        构建 Part 6: 配种记录分析-使用公牛分析-已用公牛明细

        Args:
            data: 包含 'excel_path' 的数据字典
        """
        logger.info("=" * 60)
        logger.info("开始构建Part 6: 配种记录分析-使用公牛分析-已用公牛明细")
        logger.info("=" * 60)

        # 读取已用公牛性状明细Sheet
        excel_path = data.get("excel_path")
        if excel_path is None:
            logger.error("excel_path 未找到，跳过已用公牛明细")
            return

        try:
            df_bulls_detail = pd.read_excel(excel_path, sheet_name="已用公牛性状明细", header=None)
            logger.info(f"✓ 读取bulls_detail数据（无header）: {len(df_bulls_detail)}行 x {len(df_bulls_detail.columns)}列")
        except Exception as e:
            logger.error(f"读取已用公牛性状明细Sheet失败: {e}")
            return

        if df_bulls_detail.empty:
            logger.warning("bulls_detail 数据为空DataFrame，跳过已用公牛明细")
            return

        # 解析Excel中每年的数据
        yearly_data = self._parse_yearly_data(df_bulls_detail)
        logger.info(f"✓ 解析到 {len(yearly_data)} 年的数据")

        if not yearly_data:
            logger.warning("未解析到任何年份数据")
            return

        # 查找PPT中所有"配种记录分析-使用公牛分析"页面
        logger.info(f"开始查找包含'配种记录分析-使用公牛分析'的页面，模板共{len(self.prs.slides)}页")
        all_target_slides = self.find_slides_by_text("配种记录分析-使用公牛分析", start_index=0)

        # 跳过第一页（汇总表），从第二页开始是明细页
        detail_slides = all_target_slides[1:] if len(all_target_slides) > 1 else []

        if len(detail_slides) == 0:
            logger.error("❌ 未找到已用公牛明细页面，跳过")
            return

        logger.info(f"✓ 找到 {len(detail_slides)} 个明细页面")

        # 确保有足够的页面
        if len(detail_slides) < len(yearly_data):
            logger.warning(f"⚠️ PPT页面数({len(detail_slides)})少于年份数({len(yearly_data)})")
            # 只处理现有页面数量的年份
            yearly_data = yearly_data[:len(detail_slides)]
        elif len(detail_slides) > len(yearly_data):
            logger.info(f"PPT页面数({len(detail_slides)})多于年份数({len(yearly_data)})")
            # 删除多余的页面
            slides_to_delete = detail_slides[len(yearly_data):]
            for slide_idx in reversed(slides_to_delete):  # 从后往前删除，避免索引变化
                self._delete_slide(slide_idx)
                logger.info(f"  删除多余页面: 第{slide_idx + 1}页")
            # 只使用需要的页面
            detail_slides = detail_slides[:len(yearly_data)]

        # 填充每一年的明细页面
        for i, (year, bulls_data) in enumerate(yearly_data):
            slide_idx = detail_slides[i]
            self._fill_yearly_detail(df_bulls_detail, slide_idx, year, bulls_data)

        logger.info("✓ Part 6 配种记录分析-使用公牛分析-已用公牛明细完成")

    def _parse_yearly_data(self, df: pd.DataFrame) -> List[Tuple[str, Dict]]:
        """
        解析Excel中每年的数据

        Returns:
            [(年份, {title_row, header_row, data_start_row, data_end_row, bull_count}), ...]
        """
        yearly_data = []

        i = 0
        while i < len(df):
            # 查找年份标题行
            cell_0 = str(df.iloc[i, 0]).strip() if not pd.isna(df.iloc[i, 0]) else ""

            if "年已用公牛明细" in cell_0:
                # 提取年份和公牛数量
                year_match = re.search(r'(\d{4})年', cell_0)
                bull_count_match = re.search(r'（(\d+)头公牛', cell_0)

                year = year_match.group(1) if year_match else "未知"
                bull_count = int(bull_count_match.group(1)) if bull_count_match else 0

                title_row = i
                header_row = i + 1

                # 查找数据结束行（小计行）
                data_start_row = i + 2
                data_end_row = data_start_row

                for j in range(data_start_row, len(df)):
                    cell_j0 = str(df.iloc[j, 0]).strip() if not pd.isna(df.iloc[j, 0]) else ""
                    if "小计" in cell_j0:
                        data_end_row = j
                        break
                    if cell_j0 == "":
                        data_end_row = j
                        break

                yearly_data.append((year, {
                    'title_row': title_row,
                    'header_row': header_row,
                    'data_start_row': data_start_row,
                    'data_end_row': data_end_row,
                    'bull_count': bull_count,
                    'title_text': cell_0
                }))

                logger.info(f"  {year}年: {bull_count}头公牛, 数据行{data_start_row}-{data_end_row}")

                i = data_end_row + 1
            else:
                i += 1

        return yearly_data

    def _fill_yearly_detail(self, df: pd.DataFrame, slide_index: int, year: str, info: Dict):
        """
        填充某一年的已用公牛明细页面

        Args:
            df: 已用公牛性状明细数据
            slide_index: 幻灯片索引
            year: 年份
            info: 年份数据信息
        """
        try:
            slide = self.prs.slides[slide_index]
        except IndexError:
            logger.warning(f"无法访问第{slide_index + 1}页")
            return

        logger.info(f"填充第{slide_index + 1}页 - {year}年已用公牛明细")

        # 读取标题
        title_text = info['title_text']

        # 读取表头
        header_row = self._read_row_data(df, info['header_row'])

        # 读取公牛数据
        bulls_data = []
        for i in range(info['data_start_row'], info['data_end_row']):
            row_data = self._read_row_data(df, i)
            bulls_data.append(row_data)

        # 读取小计行（如果存在）
        summary_row = None
        if info['data_end_row'] < len(df):
            cell_0 = str(df.iloc[info['data_end_row'], 0]).strip() if not pd.isna(df.iloc[info['data_end_row'], 0]) else ""
            if "小计" in cell_0:
                summary_row = self._read_row_data(df, info['data_end_row'])

        logger.info(f"  标题: {title_text}")
        logger.info(f"  表头: {len(header_row)}列")
        logger.info(f"  公牛数据: {len(bulls_data)}行")
        if summary_row:
            logger.info(f"  有小计行")

        # 填充PPT中的表格
        self._fill_ppt_table(slide, title_text, header_row, bulls_data, summary_row)

        # 生成并填充分析文本
        analysis_text = self._generate_yearly_analysis(year, bulls_data, header_row)
        if analysis_text:
            self._fill_analysis_text(slide, analysis_text)

    def _generate_yearly_analysis(self, year: str, bulls_data: list, header_row: list) -> str:
        """
        生成某一年的公牛使用分析文本

        Args:
            year: 年份
            bulls_data: 公牛数据行列表
            header_row: 表头行

        Returns:
            分析文本字符串
        """
        if not bulls_data:
            return ""

        try:
            parts = []
            bull_count = len(bulls_data)

            # 基础信息
            parts.append(f"{year}年共使用{bull_count}头公牛，")

            # 查找关键指标列索引
            nm_idx = None
            tpi_idx = None
            usage_idx = None  # 使用次数

            for i, col in enumerate(header_row):
                col_str = str(col).upper().strip()
                if 'NM$' in col_str or 'NM' == col_str:
                    nm_idx = i
                elif 'TPI' in col_str:
                    tpi_idx = i
                elif '次数' in col_str or '使用' in col_str:
                    usage_idx = i

            # 统计NM$分布
            if nm_idx is not None:
                nm_values = []
                for row in bulls_data:
                    if len(row) > nm_idx and row[nm_idx]:
                        try:
                            nm_values.append(float(row[nm_idx]))
                        except (ValueError, TypeError):
                            pass

                if nm_values:
                    avg_nm = sum(nm_values) / len(nm_values)
                    max_nm = max(nm_values)
                    high_nm_count = sum(1 for v in nm_values if v >= 800)

                    if high_nm_count > 0:
                        parts.append(f"其中NM$≥800的优秀公牛{high_nm_count}头，")
                    parts.append(f"平均NM$为{avg_nm:.0f}。")

            # 统计使用次数
            if usage_idx is not None:
                usage_values = []
                for row in bulls_data:
                    if len(row) > usage_idx and row[usage_idx]:
                        try:
                            usage_values.append(int(row[usage_idx]))
                        except (ValueError, TypeError):
                            pass

                if usage_values:
                    total_usage = sum(usage_values)
                    max_usage = max(usage_values)
                    # 找出使用最多的公牛
                    max_idx = usage_values.index(max_usage)
                    if max_idx < len(bulls_data) and len(bulls_data[max_idx]) > 0:
                        top_bull = bulls_data[max_idx][0]  # 第一列通常是公牛名/号
                        if top_bull:
                            parts.append(f"使用最多的公牛为{top_bull}（{max_usage}次）。")

            return "".join(parts)

        except Exception as e:
            logger.warning(f"生成{year}年公牛分析失败: {e}")
            return f"{year}年公牛使用情况稳定。"

    def _fill_analysis_text(self, slide, analysis_text: str):
        """
        填充分析文本到页面

        Args:
            slide: 幻灯片对象
            analysis_text: 分析文本
        """
        if not analysis_text:
            return

        # 查找分析文本框（优先位置靠下的文本框，避免匹配标题）
        textbox = None
        candidate_textboxes = []

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                name = shape.name if shape.name else ""
                shape_text = shape.text if hasattr(shape, "text") else ""
                # 跳过标题
                if "配种记录分析" in shape_text or "使用公牛分析" in shape_text:
                    continue
                if "文本框" in name or "TextBox" in name:
                    candidate_textboxes.append(shape)

        # 按位置排序，选择最靠下的
        if candidate_textboxes:
            candidate_textboxes.sort(key=lambda s: getattr(s, 'top', 0), reverse=True)
            textbox = candidate_textboxes[0]

        if textbox:
            try:
                # 清空文本框并设置新内容
                tf = textbox.text_frame
                tf.clear()
                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                run = p.add_run()
                # 避免重复"分析："前缀
                if analysis_text.startswith("分析：") or analysis_text.startswith("分析:"):
                    run.text = analysis_text
                else:
                    run.text = "分析：" + analysis_text

                # 设置字体：微软雅黑15号非加粗
                run.font.name = FONT_NAME_CN
                run.font.size = Pt(15)
                run.font.bold = False

                logger.info(f"✓ 填充明细页分析文本")
            except Exception as e:
                logger.warning(f"填充分析文本失败: {e}")
        else:
            logger.debug("未找到明细页分析文本框")

    def _read_row_data(self, df: pd.DataFrame, row_idx: int) -> list:
        """读取指定行的所有列数据"""
        row_data = []
        for j in range(df.shape[1]):
            cell_value = df.iloc[row_idx, j]
            if pd.isna(cell_value):
                row_data.append("")
            else:
                # 数值类型格式化
                if isinstance(cell_value, (int, float)):
                    # 列2（使用次数）保留整数
                    if j == 2:
                        row_data.append(str(int(cell_value)))
                    else:
                        # 其他列保留2位小数
                        row_data.append(f"{cell_value:.2f}")
                else:
                    row_data.append(str(cell_value).strip())
        return row_data

    def _fill_ppt_table(self, slide, title: str, header: list, data_rows: list, summary_row: list):
        """
        填充PPT中的表格

        PPT表格结构:
        - 表格12: 标题（1x1）
        - 表格1: 数据表（动态行数 x 19列）
        """
        # 查找表格
        table_title = None
        table_data = None

        for shape in slide.shapes:
            if shape.shape_type == 19:  # TABLE
                if shape.name == "表格 12":
                    table_title = shape.table
                elif shape.name == "表格 1":
                    table_data = shape.table

        # 填充标题表格
        if table_title:
            try:
                cell = table_title.cell(0, 0)
                cell.text = title
                logger.info(f"✓ 填充标题表格: {title}")
            except Exception as e:
                logger.error(f"填充标题表格失败: {e}")

        # 填充数据表格
        if table_data:
            self._fill_data_table(table_data, header, data_rows, summary_row)
        else:
            logger.warning("未找到数据表格（表格 1）")

    def _fill_data_table(self, table, header: list, data_rows: list, summary_row: list):
        """
        填充数据表格

        注意：
        1. 行数是动态的（需要根据公牛数量调整）
        2. 列数也是动态的（需要根据Excel数据调整）
        """
        logger.info(f"数据表格原始大小: {len(table.rows)}行 x {len(table.columns)}列")
        logger.info(f"需要的列数: {len(header)}列")

        # 计算需要的列数
        required_columns = len(header)

        # 调整表格列数
        if len(table.columns) < required_columns:
            cols_to_add = required_columns - len(table.columns)
            logger.info(f"需要添加 {cols_to_add} 列")
            for _ in range(cols_to_add):
                self._add_table_column(table)
        elif len(table.columns) > required_columns:
            cols_to_delete = len(table.columns) - required_columns
            logger.info(f"需要删除 {cols_to_delete} 列")
            for _ in range(cols_to_delete):
                self._delete_table_column(table, len(table.columns) - 1)

        # 计算需要的行数
        header_rows = 1
        required_data_rows = len(data_rows)
        required_summary_rows = 1 if summary_row else 0
        required_total_rows = header_rows + required_data_rows + required_summary_rows

        # 调整表格行数
        # 注意：最后一行通常是灰色的小计行，需要保留其格式
        if len(table.rows) < required_total_rows:
            rows_to_add = required_total_rows - len(table.rows)
            logger.info(f"需要添加 {rows_to_add} 行")
            for _ in range(rows_to_add):
                # 在倒数第二行位置插入新行（保留最后一行的小计行格式）
                self._insert_table_row_before_last(table)
        elif len(table.rows) > required_total_rows:
            rows_to_delete = len(table.rows) - required_total_rows
            logger.info(f"需要删除 {rows_to_delete} 行")
            for _ in range(rows_to_delete):
                # 从倒数第二行开始删除（保留最后一行的小计行格式）
                self._delete_table_row(table, len(table.rows) - 2)

        # 填充数据
        try:
            # 填充表头（行0）
            self._fill_table_row(table, 0, header)

            # 填充数据行
            for i, row_data in enumerate(data_rows):
                self._fill_table_row(table, i + 1, row_data)

            # 填充小计行
            if summary_row:
                self._fill_table_row(table, len(data_rows) + 1, summary_row)

            logger.info(f"✓ 填充数据表格: {len(data_rows)}行数据 + {'1行小计' if summary_row else '无小计'}")

            # 移除所有数据行的边框（除了表头和小计行）
            self._remove_all_data_row_borders(table, len(data_rows), summary_row is not None)

        except Exception as e:
            logger.error(f"填充数据表格失败: {e}", exc_info=True)

    def _remove_all_data_row_borders(self, table, data_row_count: int, has_summary: bool):
        """
        移除所有数据行的边框（保留表头和小计行的边框）

        Args:
            table: 表格对象
            data_row_count: 数据行数量
            has_summary: 是否有小计行
        """
        try:
            # 数据行范围：从行1到行data_row_count（不包括行0的表头和最后的小计行）
            end_row = data_row_count if not has_summary else data_row_count

            for row_idx in range(1, end_row + 1):
                for col_idx in range(len(table.columns)):
                    try:
                        cell = table.cell(row_idx, col_idx)
                        tc = cell._tc
                        tcPr = tc.get_or_add_tcPr()

                        # 移除所有边框相关元素
                        for child in list(tcPr):
                            tag_name = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                            # 移除所有ln开头的元素（lnL, lnR, lnT, lnB等）和边框元素
                            if tag_name.startswith('ln') or 'Borders' in tag_name or 'borders' in tag_name:
                                tcPr.remove(child)
                    except:
                        pass

            logger.debug(f"  移除了行1-{end_row}的所有边框")

        except Exception as e:
            logger.warning(f"移除数据行边框失败: {e}")

    def _fill_table_row(self, table, row_idx: int, row_data: list):
        """填充表格的一行"""
        for col_idx, cell_value in enumerate(row_data):
            if col_idx >= len(table.columns):
                # 超出表格列数，跳过
                break

            try:
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_value)

                # 设置格式
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.name = FONT_NAME_CN
                        run.font.size = Pt(9)  # 明细表用更小的字体
            except Exception as e:
                logger.debug(f"填充单元格 ({row_idx}, {col_idx}) 失败: {e}")

    def _add_table_column(self, table):
        """添加表格列"""
        try:
            from copy import deepcopy

            tbl = table._tbl

            # 获取所有行
            trs = []
            for child in tbl:
                if 'tr' in child.tag.lower():
                    trs.append(child)

            # 在每一行中添加一个新单元格
            for tr in trs:
                tcs = []
                for child in tr:
                    if 'tc' in child.tag.lower():
                        tcs.append(child)

                if tcs:
                    last_tc = tcs[-1]
                    new_tc = deepcopy(last_tc)
                    tr.append(new_tc)

            # 更新表格的gridCol定义
            tblGrid = None
            for child in tbl:
                if 'tblgrid' in child.tag.lower():
                    tblGrid = child
                    break

            if tblGrid is not None:
                gridCols = []
                for child in tblGrid:
                    if 'gridcol' in child.tag.lower():
                        gridCols.append(child)

                if gridCols:
                    last_gridCol = gridCols[-1]
                    new_gridCol = deepcopy(last_gridCol)
                    tblGrid.append(new_gridCol)

        except Exception as e:
            logger.warning(f"添加表格列失败: {e}")

    def _delete_table_column(self, table, col_idx: int):
        """删除表格中的一列"""
        try:
            tbl = table._tbl

            trs = []
            for child in tbl:
                if 'tr' in child.tag.lower():
                    trs.append(child)

            for tr in trs:
                tcs = []
                for child in tr:
                    if 'tc' in child.tag.lower():
                        tcs.append(child)

                if col_idx < len(tcs):
                    tr.remove(tcs[col_idx])

            tblGrid = None
            for child in tbl:
                if 'tblgrid' in child.tag.lower():
                    tblGrid = child
                    break

            if tblGrid is not None:
                gridCols = []
                for child in tblGrid:
                    if 'gridcol' in child.tag.lower():
                        gridCols.append(child)

                if col_idx < len(gridCols):
                    tblGrid.remove(gridCols[col_idx])

        except Exception as e:
            logger.warning(f"删除表格列{col_idx}失败: {e}")

    def _add_table_row(self, table):
        """添加表格行（复制普通数据行格式，而非小计行）"""
        try:
            from copy import deepcopy
            tbl = table._tbl

            # 复制第2行（索引1，通常是第一个数据行）而不是最后一行（小计行）
            # 这样可以确保新增行使用普通数据行格式（白色背景），而不是小计行格式（灰色背景）
            template_row_idx = 1 if len(table.rows) > 1 else 0
            template_tr = table.rows[template_row_idx]._tr
            new_tr = deepcopy(template_tr)
            tbl.append(new_tr)

            # 清空新行的文本
            new_row_idx = len(table.rows) - 1
            for col_idx in range(len(table.columns)):
                try:
                    cell = table.cell(new_row_idx, col_idx)
                    cell.text = ""
                except:
                    pass
        except Exception as e:
            logger.warning(f"添加表格行失败: {e}")

    def _insert_table_row_before_last(self, table):
        """在最后一行之前插入新行（保留最后一行的小计行格式）"""
        try:
            from copy import deepcopy
            tbl = table._tbl

            # 复制第2行（索引1，通常是第一个数据行）的格式
            template_row_idx = 1 if len(table.rows) > 1 else 0
            template_tr = table.rows[template_row_idx]._tr
            new_tr = deepcopy(template_tr)

            # 获取最后一行的位置
            last_tr = table.rows[len(table.rows) - 1]._tr

            # 在最后一行之前插入新行
            tbl.insert(tbl.index(last_tr), new_tr)

            # 清空新行的文本并移除边框（新行现在在倒数第二的位置）
            new_row_idx = len(table.rows) - 2
            for col_idx in range(len(table.columns)):
                try:
                    cell = table.cell(new_row_idx, col_idx)
                    cell.text = ""

                    # 移除单元格的所有边框（避免出现多余的横线）
                    tc = cell._tc
                    tcPr = tc.get_or_add_tcPr()

                    # 移除所有边框相关元素
                    removed_count = 0
                    for child in list(tcPr):
                        tag_name = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                        # 移除所有ln开头的元素（lnL, lnR, lnT, lnB等）和边框元素
                        if tag_name.startswith('ln') or 'Borders' in tag_name or 'borders' in tag_name:
                            tcPr.remove(child)
                            removed_count += 1

                    if col_idx == 0 and removed_count > 0:
                        logger.debug(f"  行{new_row_idx}移除了{removed_count}个边框元素")
                except Exception as e:
                    logger.debug(f"  清除行{new_row_idx}列{col_idx}边框失败: {e}")
        except Exception as e:
            logger.warning(f"在最后一行之前插入表格行失败: {e}")

    def _delete_table_row(self, table, row_idx: int):
        """删除表格中的一行"""
        try:
            tbl = table._tbl
            tr = table.rows[row_idx]._tr
            tbl.remove(tr)
        except Exception as e:
            logger.warning(f"删除表格行{row_idx}失败: {e}")

    def _delete_slide(self, slide_index: int):
        """
        删除指定索引的幻灯片

        Args:
            slide_index: 幻灯片索引（从0开始）
        """
        try:
            rId = self.prs.slides._sldIdLst[slide_index].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[slide_index]
        except Exception as e:
            logger.warning(f"删除幻灯片{slide_index}失败: {e}")
