"""
Part 6: 配种记录时间线分析构建器
对应PPT模版的"配种记录时间线"页面（两页：全部记录、近一年）
"""

import logging
from typing import Dict, List, Optional
from pathlib import Path
from io import BytesIO

from openpyxl import load_workbook
from openpyxl.drawing.image import Image as OpenpyxlImage
from pptx.util import Cm, Pt

from ..base_builder import BaseSlideBuilder
from ..config import FONT_NAME_CN

logger = logging.getLogger(__name__)


class Part6TimelineBuilder(BaseSlideBuilder):
    """Part 6: 配种记录时间线分析构建器"""

    def __init__(self, prs, chart_creator, farm_name: str):
        super().__init__(prs, chart_creator)
        self.farm_name = farm_name

    def build(self, data: Dict) -> bool:
        """
        构建 Part 6: 配种记录时间线分析

        Args:
            data: 包含 'excel_path' 的数据字典

        Returns:
            bool: 是否需要从Excel提取图片（用于判断是否需要加载data_only=False workbook）
        """
        logger.info("=" * 60)
        logger.info("开始构建Part 6: 配种记录时间线分析")
        logger.info("=" * 60)

        excel_path = data.get("excel_path")
        if excel_path is None:
            logger.error("excel_path 未找到，跳过配种记录时间线")
            return False

        # 优先查找预导出的时间线图片（在Excel生成阶段已导出到analysis_results）
        pregenerated_images = self._find_pregenerated_images(excel_path)

        if pregenerated_images and len(pregenerated_images) >= 2:
            # 使用预导出的图片（无需加载data_only=False workbook，节省约20秒）
            logger.info(f"✓ 找到预导出的时间线图片，跳过从Excel提取")
            image_all_path = pregenerated_images[0]
            image_recent_path = pregenerated_images[1]
            used_pregenerated = True
        else:
            # 回退到从Excel提取图片
            logger.info("未找到预导出图片，从Excel中提取...")
            try:
                cached_wb = data.get("_cached_workbook")
                images = self._extract_images_from_excel(excel_path, cached_wb)
                logger.info(f"✓ 从Excel中提取到 {len(images)} 个图片")
            except Exception as e:
                logger.error(f"提取Excel图片失败: {e}")
                return False

            if len(images) < 2:
                logger.warning(f"Excel中图片数量不足（需要2个，实际{len(images)}个），跳过配种记录时间线")
                return False

            # 导出图片到临时文件
            temp_dir = Path(self.chart_creator.output_dir)
            image_all_path = self._export_image_to_file(images[0], temp_dir / "timeline_all.png")
            image_recent_path = self._export_image_to_file(images[1], temp_dir / "timeline_recent.png") if len(images) > 1 else None
            used_pregenerated = False

        # 查找"配种记录时间线"页面
        logger.info(f"开始查找包含'配种记录时间线'的页面，模板共{len(self.prs.slides)}页")
        target_slides = self.find_slides_by_text("配种记录时间线", start_index=0)

        if len(target_slides) < 2:
            logger.error(f"❌ 未找到足够的配种记录时间线页面（需要2页，实际{len(target_slides)}页），跳过")
            return not used_pregenerated

        logger.info(f"✓ 找到 {len(target_slides)} 个配种记录时间线页面")

        # Image 12: 配种记录时间线-全部记录
        if image_all_path:
            self._add_image_to_slide(target_slides[0], image_all_path, "配种记录时间线-全部记录")
            self._fill_timeline_analysis(target_slides[0], "all")
            logger.info(f"✓ 添加图片到第{target_slides[0] + 1}页: 配种记录时间线-全部记录")

        # Image 13: 配种记录时间线-近一年
        if image_recent_path:
            self._add_image_to_slide(target_slides[1], image_recent_path, "配种记录时间线-近一年")
            self._fill_timeline_analysis(target_slides[1], "recent")
            logger.info(f"✓ 添加图片到第{target_slides[1] + 1}页: 配种记录时间线-近一年")

        logger.info("✓ Part 6 配种记录时间线分析完成")
        return not used_pregenerated  # 返回是否需要从Excel提取（用于判断是否需要workbook）

    def _fill_timeline_analysis(self, slide_index: int, timeline_type: str):
        """
        填充时间线分析文本

        Args:
            slide_index: 幻灯片索引
            timeline_type: 时间线类型 ('all' 或 'recent')
        """
        try:
            slide = self.prs.slides[slide_index]

            # 生成分析文本
            if timeline_type == "all":
                analysis = (
                    "分析：配种记录时间线展示了牧场历史配种活动的整体分布情况。"
                    "时间线可直观反映配种的季节性规律和年度变化趋势，"
                    "有助于优化配种计划安排和冻精库存管理。"
                )
            else:  # recent
                analysis = (
                    "分析：近一年配种记录时间线展示了最近12个月的配种活动分布。"
                    "通过分析近期配种密度和时间分布，可以评估当前配种计划执行情况，"
                    "及时调整配种策略以提高繁殖效率。"
                )

            # 查找分析文本框（避免匹配到标题）
            textbox = None
            candidate_textboxes = []

            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    name = shape.name if shape.name else ""
                    shape_text = shape.text if hasattr(shape, "text") else ""

                    # 跳过标题：页面标题通常包含"配种记录时间线"
                    if "配种记录时间线" in shape_text:
                        continue
                    # 跳过图片
                    if "图片" in name or "Picture" in name:
                        continue

                    # 查找文本框
                    if "文本框" in name or "TextBox" in name:
                        candidate_textboxes.append(shape)

            # 按位置（top）排序，选择最靠下的文本框作为分析文本框
            if candidate_textboxes:
                candidate_textboxes.sort(key=lambda s: getattr(s, 'top', 0), reverse=True)
                textbox = candidate_textboxes[0]

            if textbox:
                # 清空文本框并设置新内容
                tf = textbox.text_frame
                tf.clear()
                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                run = p.add_run()
                run.text = analysis

                # 设置字体：微软雅黑15号非加粗
                run.font.name = FONT_NAME_CN
                run.font.size = Pt(15)
                run.font.bold = False

                logger.info(f"✓ 填充时间线分析文本")
            else:
                logger.debug("未找到时间线分析文本框")

        except Exception as e:
            logger.warning(f"填充时间线分析失败: {e}")

    def _extract_images_from_excel(self, excel_path: str, cached_wb=None) -> List[OpenpyxlImage]:
        """
        从Excel中提取图片（Image 12和Image 13）

        Args:
            excel_path: Excel文件路径
            cached_wb: 缓存的workbook对象（可选，避免重复加载）

        Returns:
            List[OpenpyxlImage]: 图片对象列表
        """
        # 优先使用缓存的workbook
        if cached_wb is not None:
            wb = cached_wb
            logger.debug("使用缓存的workbook")
        else:
            logger.debug("加载新的workbook（这可能需要约21秒）")
            wb = load_workbook(str(excel_path))
        ws = wb['已用公牛性状汇总']

        images = []

        # 获取工作表中的所有图片
        if hasattr(ws, '_images') and ws._images:
            # 按图片的位置或名称排序（假设Image 12和13是最后两个）
            all_images = list(ws._images)
            logger.info(f"工作表中共有 {len(all_images)} 个图片对象")

            # 通常Excel中的图表和图片是按顺序编号的
            # Image 12和Image 13应该是倒数第2个和倒数第1个
            if len(all_images) >= 2:
                images = all_images[-2:]  # 取最后两个
                logger.debug(f"提取最后两个图片作为时间线图")
            else:
                images = all_images
        else:
            logger.warning("工作表中没有找到图片对象")

        return images

    def _export_image_to_file(self, image: OpenpyxlImage, output_path: Path) -> Path:
        """
        将Excel图片导出为文件

        Args:
            image: openpyxl图片对象
            output_path: 输出文件路径

        Returns:
            Path: 导出的图片文件路径
        """
        try:
            # 获取图片数据
            if hasattr(image, '_data'):
                image_data = image._data()
            elif hasattr(image, 'ref'):
                # 从图片引用中获取数据
                image_data = image.ref
            else:
                logger.warning("无法获取图片数据")
                return None

            # 保存图片
            with open(output_path, 'wb') as f:
                if isinstance(image_data, bytes):
                    f.write(image_data)
                else:
                    f.write(image_data.read())

            logger.debug(f"图片已导出: {output_path.name}")
            return output_path

        except Exception as e:
            logger.warning(f"导出图片失败: {e}")
            return None

    def _add_image_to_slide(self, slide_index: int, image_path: Path, title: str):
        """
        将图片添加到幻灯片（大小：32cm宽 × 10cm高）

        Args:
            slide_index: 幻灯片索引
            image_path: 图片文件路径
            title: 图片标题（用于日志）
        """
        try:
            slide = self.prs.slides[slide_index]

            # 图片大小：32cm宽 × 10cm高（用户要求）
            width = Cm(32)
            height = Cm(10)

            # 居中放置
            # PPT默认尺寸：33.867cm × 19.05cm (13.33" × 7.5")
            prs_width = self.prs.slide_width
            prs_height = self.prs.slide_height

            left = (prs_width - width) // 2
            top = Cm(4)  # 距离顶部4cm

            # 查找并删除现有的图片占位符
            for shape in slide.shapes:
                if hasattr(shape, 'name') and ('图片' in shape.name or 'Picture' in shape.name):
                    sp = shape._element
                    sp.getparent().remove(sp)
                    logger.debug(f"  删除占位符: {shape.name}")
                    break

            # 添加新图片
            slide.shapes.add_picture(
                str(image_path),
                left=left,
                top=top,
                width=width,
                height=height
            )

            logger.debug(f"  图片已添加: {title}, 大小={width.cm:.1f}cm×{height.cm:.1f}cm")

        except Exception as e:
            logger.warning(f"添加图片到幻灯片失败: {e}")
            logger.debug(f"错误详情: {e}", exc_info=True)

    def _find_pregenerated_images(self, excel_path: str) -> List[Path]:
        """
        查找Excel生成阶段预导出的时间线图片

        Args:
            excel_path: Excel文件路径

        Returns:
            List[Path]: 图片路径列表，顺序为 [全部记录, 近一年]；未找到则返回空列表
        """
        try:
            excel_path = Path(excel_path)
            # Excel通常在 project/reports/ 下，analysis_results在 project/analysis_results/
            project_folder = excel_path.parent.parent
            analysis_folder = project_folder / "analysis_results"

            if not analysis_folder.exists():
                logger.debug(f"analysis_results目录不存在: {analysis_folder}")
                return []

            # 查找两个关键的时间线图片
            all_pattern = "配种记录时间线_全部_*.png"
            recent_pattern = "配种记录时间线_近一年_*.png"

            all_files = list(analysis_folder.glob(all_pattern))
            recent_files = list(analysis_folder.glob(recent_pattern))

            if not all_files or not recent_files:
                logger.debug(f"未找到完整的时间线图片集合 (全部: {len(all_files)}, 近一年: {len(recent_files)})")
                return []

            # 按修改时间排序，取最新的
            latest_all = max(all_files, key=lambda f: f.stat().st_mtime)
            latest_recent = max(recent_files, key=lambda f: f.stat().st_mtime)

            logger.info(f"  ✓ 找到预导出图片: {latest_all.name}, {latest_recent.name}")
            return [latest_all, latest_recent]

        except Exception as e:
            logger.debug(f"查找预导出图片失败: {e}")
            return []
