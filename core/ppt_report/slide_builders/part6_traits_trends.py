"""
Part 6: 配种记录分析-性状进展折线图构建器
对应PPT模版的"性状进展折线图"页面（每页2个图表）
"""

import logging
from typing import Dict, List
from pathlib import Path

from openpyxl import load_workbook
from openpyxl.chart import LineChart
from pptx.util import Inches, Pt

from ..base_builder import BaseSlideBuilder
from ..config import FONT_NAME_CN

logger = logging.getLogger(__name__)


class Part6TraitsTrendsBuilder(BaseSlideBuilder):
    """Part 6: 配种记录分析-性状进展折线图构建器"""

    def __init__(self, prs, chart_creator, farm_name: str):
        super().__init__(prs, chart_creator)
        self.farm_name = farm_name

    def build(self, data: Dict):
        """
        构建 Part 6: 配种记录分析-性状进展折线图

        Args:
            data: 包含 'excel_path' 的数据字典
        """
        logger.info("=" * 60)
        logger.info("开始构建Part 6: 配种记录分析-性状进展折线图")
        logger.info("=" * 60)

        # 读取Excel文件中的图表
        excel_path = data.get("excel_path")
        if excel_path is None:
            logger.error("excel_path 未找到，跳过性状进展折线图")
            return

        # 查找"性状进展折线图"页面
        logger.info(f"开始查找包含'性状进展折线图'的页面，模板共{len(self.prs.slides)}页")
        target_slides = self.find_slides_by_text("性状进展折线图", start_index=0)

        if len(target_slides) == 0:
            logger.error("❌ 未找到性状进展折线图页面，跳过")
            return

        try:
            # 优先使用缓存的workbook（避免重复加载，每次加载需要约21秒）
            cached_wb = data.get("_cached_workbook_data_only")
            charts_data = self._extract_charts_from_excel(excel_path, cached_wb)
            logger.info(f"✓ 从Excel中提取到 {len(charts_data)} 个图表")
        except Exception as e:
            logger.error(f"提取Excel图表失败: {e}")
            # 数据提取失败，标记页面待删除
            self.mark_slides_for_deletion(target_slides)
            return

        if not charts_data:
            logger.warning("没有找到任何图表，标记页面待删除")
            self.mark_slides_for_deletion(target_slides)
            return

        logger.info(f"✓ 找到 {len(target_slides)} 个性状进展折线图页面")

        # 计算需要的页面数（每页2个图表）
        required_pages = (len(charts_data) + 1) // 2
        logger.info(f"需要 {required_pages} 页来显示 {len(charts_data)} 个图表")

        # 更新PPT中的图表
        self._update_charts_in_slides(charts_data, target_slides, required_pages)

        logger.info("✓ Part 6 配种记录分析-性状进展折线图完成")

    def _extract_charts_from_excel(self, excel_path: str, cached_wb=None) -> List[Dict]:
        """
        从Excel中提取图表信息和数据

        Args:
            excel_path: Excel文件路径
            cached_wb: 缓存的workbook对象（可选，避免重复加载）

        Returns:
            List[Dict]: [{'index': 1, 'title': '经济指数', 'data': {...}}, ...]
        """
        # 优先使用缓存的workbook
        if cached_wb is not None:
            wb = cached_wb
            logger.debug("使用缓存的workbook (data_only)")
        else:
            logger.debug("加载新的workbook（这可能需要约21秒）")
            wb = load_workbook(str(excel_path), data_only=True)  # data_only=True 获取计算后的值
        ws = wb['已用公牛性状汇总']

        charts_data = []
        for i, chart in enumerate(ws._charts, 1):
            try:
                # 提取标题文本
                title_text = f'图表{i}'
                if hasattr(chart, 'title') and chart.title:
                    if hasattr(chart.title, 'tx') and chart.title.tx:
                        if hasattr(chart.title.tx, 'rich') and chart.title.tx.rich:
                            if hasattr(chart.title.tx.rich, 'p') and chart.title.tx.rich.p:
                                for p in chart.title.tx.rich.p:
                                    if hasattr(p, 'r') and p.r:
                                        for r in p.r:
                                            if hasattr(r, 't'):
                                                title_text = r.t
                                                break

                # 立即提取图表数据（传入worksheet以便直接读取源数据）
                chart_data = self._extract_chart_data_for_update(chart, wb)

                # 检查是否成功提取到数据
                has_data = False
                if chart_data and chart_data.get('series'):
                    total_values = sum(len(s.get('values', [])) for s in chart_data['series'])
                    has_data = total_values > 0

                logger.debug(f"  提取图表 {i}: {title_text} - 有数据={has_data}")

                charts_data.append({
                    'index': i,
                    'title': title_text,
                    'data': chart_data if has_data else None  # 如果没有数据就保存None
                })

            except Exception as e:
                logger.warning(f"提取图表 {i} 失败: {e}")
                continue

        return charts_data

    def _export_charts_to_images(self, excel_path: str, charts_data: List[Dict]) -> List[Dict]:
        """
        将Excel图表导出为图片文件

        Args:
            excel_path: Excel文件路径
            charts_data: 图表数据列表

        Returns:
            List[Dict]: [{'index': 1, 'title': '经济指数', 'image_path': '/path/to/image.png'}, ...]
        """
        chart_images = []

        # 使用chart_creator的临时目录
        temp_dir = Path(self.chart_creator.output_dir)

        for chart_data in charts_data:
            index = chart_data['index']
            title = chart_data['title']
            chart = chart_data['chart']

            try:
                # 生成图片文件名
                image_filename = f"trait_trend_{index}_{title}.png"
                image_path = temp_dir / image_filename

                # 使用matplotlib重新生成图表
                chart_image_path = self._recreate_chart_as_image(chart, title, image_path)

                chart_images.append({
                    'index': index,
                    'title': title,
                    'image_path': str(chart_image_path)
                })

                logger.debug(f"  导出图表 {index}: {title} -> {chart_image_path.name}")

            except Exception as e:
                logger.warning(f"导出图表 {index} ({title}) 失败: {e}")
                logger.debug(f"错误详情: {e}", exc_info=True)
                continue

        return chart_images

    def _recreate_chart_as_image(self, chart: LineChart, title: str, output_path: Path) -> Path:
        """
        使用matplotlib重新创建图表并保存为图片

        Args:
            chart: openpyxl图表对象
            title: 图表标题
            output_path: 输出文件路径

        Returns:
            Path: 生成的图片路径
        """
        import matplotlib.pyplot as plt
        import matplotlib
        from openpyxl.chart.series import SeriesLabel

        # 提取图表数据
        series_data = []
        x_labels = []

        for series in chart.series:
            # 提取系列名称
            series_name = ""
            if hasattr(series, 'tx') and series.tx:
                if hasattr(series.tx, 'strRef') and series.tx.strRef:
                    if hasattr(series.tx.strRef, 'strCache') and series.tx.strRef.strCache:
                        if hasattr(series.tx.strRef.strCache, 'pt') and series.tx.strRef.strCache.pt:
                            for pt in series.tx.strRef.strCache.pt:
                                if hasattr(pt, 'v'):
                                    series_name = pt.v
                                    break

            # 提取Y轴数据
            y_values = []
            if hasattr(series, 'val') and series.val:
                if hasattr(series.val, 'numRef') and series.val.numRef:
                    if hasattr(series.val.numRef, 'numCache') and series.val.numRef.numCache:
                        if hasattr(series.val.numRef.numCache, 'pt') and series.val.numRef.numCache.pt:
                            for pt in series.val.numRef.numCache.pt:
                                if hasattr(pt, 'v'):
                                    try:
                                        y_values.append(float(pt.v))
                                    except:
                                        y_values.append(0)

            # 提取X轴标签（只需要一次）
            if not x_labels and hasattr(series, 'cat') and series.cat:
                if hasattr(series.cat, 'strRef') and series.cat.strRef:
                    if hasattr(series.cat.strRef, 'strCache') and series.cat.strRef.strCache:
                        if hasattr(series.cat.strRef.strCache, 'pt') and series.cat.strRef.strCache.pt:
                            for pt in series.cat.strRef.strCache.pt:
                                if hasattr(pt, 'v'):
                                    x_labels.append(pt.v)

            series_data.append({
                'name': series_name or f'系列{len(series_data)+1}',
                'values': y_values
            })

        # 使用matplotlib绘制图表
        fig, ax = plt.subplots(figsize=(8, 5))

        # 绘制每个系列
        for series in series_data:
            if series['values']:
                x_range = range(len(series['values']))
                ax.plot(x_range, series['values'], marker='o', label=series['name'], linewidth=2)

        # 设置标题和标签
        if self.chart_creator.cn_font:
            ax.set_title(title, fontsize=14, fontweight='bold', fontproperties=self.chart_creator.cn_font)
            ax.set_xlabel('年份', fontproperties=self.chart_creator.cn_font)
        else:
            ax.set_title(title, fontsize=14, fontweight='bold')
            ax.set_xlabel('年份')

        # 设置X轴标签
        if x_labels:
            ax.set_xticks(range(len(x_labels)))
            if self.chart_creator.cn_font:
                ax.set_xticklabels(x_labels, fontproperties=self.chart_creator.cn_font)
            else:
                ax.set_xticklabels(x_labels)

        # 设置图例
        if len(series_data) > 1:
            if self.chart_creator.cn_font:
                ax.legend(prop=self.chart_creator.cn_font)
            else:
                ax.legend()

        # 添加网格
        ax.grid(True, alpha=0.3)

        # 保存图片
        plt.tight_layout()
        plt.savefig(str(output_path), dpi=150, bbox_inches='tight')
        plt.close()

        return output_path

    def _update_charts_in_slides(self, charts_data: List[Dict], target_slides: List[int], required_pages: int):
        """
        更新PPT中的图表数据和标题

        Args:
            charts_data: Excel图表数据列表
            target_slides: 目标幻灯片索引列表
            required_pages: 需要的页面数
        """
        # 检查页面数是否足够
        available_pages = len(target_slides)

        if available_pages < required_pages:
            logger.warning(f"PPT页面数({available_pages})少于需要的页数({required_pages})")
        elif available_pages > required_pages:
            logger.info(f"PPT页面数({available_pages})多于需要的页数({required_pages})")
            # 删除多余的页面（但保护索引>=100的页面，避免误删Part 7的选配推荐页面）
            slides_to_delete = target_slides[required_pages:]
            for slide_idx in reversed(slides_to_delete):
                if slide_idx < 100:  # 只删除索引<100的页面
                    self._delete_slide(slide_idx)
                    logger.info(f"  删除多余页面: 第{slide_idx + 1}页")
                else:
                    logger.warning(f"  跳过删除第{slide_idx + 1}页（保护Part 7页面，索引>{100})")
            # 只使用前required_pages个页面
            target_slides = target_slides[:required_pages]

        # 每页更新2个图表
        for page_idx in range(min(len(target_slides), required_pages)):
            slide_index = target_slides[page_idx]
            slide = self.prs.slides[slide_index]

            # 计算当前页应该显示的图表索引
            chart_idx_1 = page_idx * 2
            chart_idx_2 = page_idx * 2 + 1

            logger.info(f"更新第{slide_index + 1}页 - 性状进展折线图")

            # 查找页面中的图表
            ppt_charts = self._find_charts_in_slide(slide)

            # 更新第一个图表（左侧）
            if chart_idx_1 < len(charts_data) and len(ppt_charts) > 0:
                self._update_chart_data(ppt_charts[0], charts_data[chart_idx_1])
                logger.info(f"  ✓ 更新图表1: {charts_data[chart_idx_1]['title']}")

            # 更新第二个图表（右侧）
            if chart_idx_2 < len(charts_data) and len(ppt_charts) > 1:
                self._update_chart_data(ppt_charts[1], charts_data[chart_idx_2])
                logger.info(f"  ✓ 更新图表2: {charts_data[chart_idx_2]['title']}")

            # 生成并填充分析文本
            chart1_data = charts_data[chart_idx_1] if chart_idx_1 < len(charts_data) else None
            chart2_data = charts_data[chart_idx_2] if chart_idx_2 < len(charts_data) else None
            analysis_text = self._generate_traits_trends_analysis(chart1_data, chart2_data)
            if analysis_text:
                self._fill_analysis_text(slide, analysis_text)

    def _find_charts_in_slide(self, slide):
        """查找幻灯片中的所有图表"""
        charts = []
        for shape in slide.shapes:
            if shape.shape_type == 3:  # CHART = 3
                charts.append(shape)
        return charts

    def _update_chart_data(self, ppt_chart_shape, excel_chart_data: Dict):
        """
        更新PPT图表的数据和标题

        Args:
            ppt_chart_shape: PPT中的图表shape对象
            excel_chart_data: Excel中的图表数据 {'index': 1, 'title': '经济指数', 'data': {...} or None}
        """
        try:
            chart_data = excel_chart_data.get('data')
            title = excel_chart_data.get('title', '未命名图表')

            # 获取PPT图表对象
            ppt_chart = ppt_chart_shape.chart

            # 更新图表标题
            if ppt_chart.has_title:
                ppt_chart.chart_title.text_frame.text = title

            # 如果Excel图表有缓存数据，则更新PPT图表数据
            # 否则保持PPT图表的原有数据不变（只更新标题）
            if chart_data:
                self._update_ppt_chart_data(ppt_chart, chart_data)
                logger.debug(f"    图表数据已更新")
            else:
                logger.info(f"    保持原有图表数据（Excel图表无缓存数据）")

        except Exception as e:
            logger.warning(f"更新图表数据失败: {e}")
            logger.debug(f"错误详情: {e}", exc_info=True)

    def _parse_cell_reference(self, formula: str):
        """
        解析Excel公式引用，提取sheet名称和单元格范围

        Args:
            formula: 如 "'已用公牛性状汇总'!$B$2:$B$10" 或 "Sheet1!$A$1:$A$5"

        Returns:
            (sheet_name, start_cell, end_cell) 或 None
        """
        import re
        try:
            # 匹配格式: 'Sheet Name'!$A$1:$A$10 或 Sheet1!$A$1:$A$10
            pattern = r"(?:'([^']+)'|([^!]+))!\$?([A-Z]+)\$?(\d+):\$?([A-Z]+)\$?(\d+)"
            match = re.match(pattern, formula)
            if match:
                sheet_name = match.group(1) or match.group(2)
                start_col, start_row = match.group(3), int(match.group(4))
                end_col, end_row = match.group(5), int(match.group(6))
                return sheet_name, (start_col, start_row), (end_col, end_row)
        except Exception as e:
            logger.debug(f"解析公式引用失败: {formula}, 错误: {e}")
        return None

    def _read_range_from_worksheet(self, wb, sheet_name: str, start: tuple, end: tuple) -> list:
        """
        从worksheet读取指定范围的数据

        Args:
            wb: workbook对象
            sheet_name: sheet名称
            start: (列字母, 行号)
            end: (列字母, 行号)

        Returns:
            数据列表
        """
        try:
            ws = wb[sheet_name]
            start_col, start_row = start
            end_col, end_row = end

            values = []
            for row in range(start_row, end_row + 1):
                cell_value = ws[f"{start_col}{row}"].value
                if cell_value is not None:
                    values.append(cell_value)
                else:
                    values.append(0)
            return values
        except Exception as e:
            logger.debug(f"从worksheet读取数据失败: {e}")
            return []

    def _extract_chart_data_for_update(self, excel_chart, wb=None) -> Dict:
        """
        从Excel图表中提取数据用于更新PPT图表

        Args:
            excel_chart: openpyxl图表对象
            wb: workbook对象（用于直接读取源数据）

        Returns:
            {
                'categories': ['2023', '2024', '2025'],
                'series': [
                    {'name': '点位', 'values': [533.50, 633.06, 739.54]},
                ]
            }
        """
        try:
            categories = []
            series_list = []

            logger.debug(f"        开始提取图表数据，共{len(excel_chart.series) if hasattr(excel_chart, 'series') else 0}个系列")

            for idx, series in enumerate(excel_chart.series):
                logger.debug(f"          处理系列{idx + 1}: val={hasattr(series, 'val')}, cat={hasattr(series, 'cat')}")
                # 提取系列名称
                series_name = ""
                if hasattr(series, 'tx') and series.tx:
                    if hasattr(series.tx, 'strRef') and series.tx.strRef:
                        if hasattr(series.tx.strRef, 'strCache') and series.tx.strRef.strCache:
                            if hasattr(series.tx.strRef.strCache, 'pt') and series.tx.strRef.strCache.pt:
                                for pt in series.tx.strRef.strCache.pt:
                                    if hasattr(pt, 'v'):
                                        series_name = pt.v
                                        break

                # 提取Y轴数据 - 优先从源数据读取
                y_values = []
                if hasattr(series, 'val') and series.val:
                    if hasattr(series.val, 'numRef') and series.val.numRef:
                        # 优先从公式引用读取源数据（修复：不再依赖可能过期的numCache）
                        if wb and hasattr(series.val.numRef, 'f') and series.val.numRef.f:
                            formula = series.val.numRef.f
                            parsed = self._parse_cell_reference(formula)
                            if parsed:
                                sheet_name, start, end = parsed
                                y_values = self._read_range_from_worksheet(wb, sheet_name, start, end)
                                # 转换为浮点数
                                y_values = [float(v) if v is not None else 0 for v in y_values]
                                logger.debug(f"            从源数据读取: {formula} -> {len(y_values)}个值")

                        # 如果从源数据读取失败，回退到numCache
                        if not y_values:
                            if hasattr(series.val.numRef, 'numCache') and series.val.numRef.numCache:
                                if hasattr(series.val.numRef.numCache, 'pt') and series.val.numRef.numCache.pt:
                                    for pt in series.val.numRef.numCache.pt:
                                        if hasattr(pt, 'v'):
                                            try:
                                                y_values.append(float(pt.v))
                                            except:
                                                y_values.append(0)
                                    logger.debug(f"            从numCache读取: {len(y_values)}个值")

                # 提取X轴标签（只需要一次）- 同样优先从源数据读取
                if not categories and hasattr(series, 'cat') and series.cat:
                    # 尝试从源数据读取categories
                    if wb:
                        # 尝试字符串引用
                        if hasattr(series.cat, 'strRef') and series.cat.strRef:
                            if hasattr(series.cat.strRef, 'f') and series.cat.strRef.f:
                                parsed = self._parse_cell_reference(series.cat.strRef.f)
                                if parsed:
                                    sheet_name, start, end = parsed
                                    categories = self._read_range_from_worksheet(wb, sheet_name, start, end)
                                    categories = [str(v) for v in categories]
                                    logger.debug(f"            categories从源数据读取: {len(categories)}个")

                        # 尝试数值引用
                        if not categories and hasattr(series.cat, 'numRef') and series.cat.numRef:
                            if hasattr(series.cat.numRef, 'f') and series.cat.numRef.f:
                                parsed = self._parse_cell_reference(series.cat.numRef.f)
                                if parsed:
                                    sheet_name, start, end = parsed
                                    cat_values = self._read_range_from_worksheet(wb, sheet_name, start, end)
                                    categories = [str(int(float(v))) if v else '' for v in cat_values]
                                    logger.debug(f"            categories从源数据(num)读取: {len(categories)}个")

                    # 如果从源数据读取失败，回退到缓存
                    if not categories:
                        if hasattr(series.cat, 'strRef') and series.cat.strRef:
                            if hasattr(series.cat.strRef, 'strCache') and series.cat.strRef.strCache:
                                if hasattr(series.cat.strRef.strCache, 'pt') and series.cat.strRef.strCache.pt:
                                    for pt in series.cat.strRef.strCache.pt:
                                        if hasattr(pt, 'v'):
                                            categories.append(pt.v)
                        elif hasattr(series.cat, 'numRef') and series.cat.numRef:
                            if hasattr(series.cat.numRef, 'numCache') and series.cat.numRef.numCache:
                                if hasattr(series.cat.numRef.numCache, 'pt') and series.cat.numRef.numCache.pt:
                                    for pt in series.cat.numRef.numCache.pt:
                                        if hasattr(pt, 'v'):
                                            categories.append(str(int(float(pt.v))))

                series_list.append({
                    'name': series_name or f'系列{len(series_list)+1}',
                    'values': y_values
                })

            return {
                'categories': categories,
                'series': series_list
            }

        except Exception as e:
            logger.warning(f"提取图表数据失败: {e}")
            return None

    def _update_ppt_chart_data(self, ppt_chart, chart_data: Dict):
        """
        更新PPT图表的嵌入数据

        Args:
            ppt_chart: PPT图表对象
            chart_data: 图表数据 {'categories': [...], 'series': [...]}
        """
        try:
            from pptx.chart.data import CategoryChartData

            # 创建新的图表数据
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = chart_data['categories']

            # 添加所有系列
            for series in chart_data['series']:
                chart_data_obj.add_series(series['name'], series['values'])

            # 替换图表数据
            ppt_chart.replace_data(chart_data_obj)

            # 设置数据标签格式为2位小数
            try:
                for series in ppt_chart.series:
                    if hasattr(series, 'data_labels'):
                        # 确保数据标签可见
                        series.data_labels.show_value = True
                        # 设置数字格式为2位小数
                        series.data_labels.number_format = '0.00'
                        logger.debug(f"  设置系列'{series.name}'的数据标签格式为2位小数")
            except Exception as e:
                logger.debug(f"设置数据标签格式时出错: {e}")

        except Exception as e:
            logger.warning(f"更新PPT图表数据失败: {e}")
            logger.debug(f"错误详情: {e}", exc_info=True)

    def _fill_charts_to_slides_old(self, chart_images: List[Dict], target_slides: List[int], required_pages: int):
        """
        将图表图片填充到PPT页面

        Args:
            chart_images: 图表图片列表
            target_slides: 目标幻灯片索引列表
            required_pages: 需要的页面数
        """
        # 检查页面数是否足够
        available_pages = len(target_slides)

        if available_pages < required_pages:
            logger.warning(f"PPT页面数({available_pages})少于需要的页数({required_pages})")
            # TODO: 可以考虑复制页面
        elif available_pages > required_pages:
            logger.info(f"PPT页面数({available_pages})多于需要的页数({required_pages})")
            # 删除多余的页面
            slides_to_delete = target_slides[required_pages:]
            for slide_idx in reversed(slides_to_delete):  # 从后往前删除，避免索引变化
                self._delete_slide(slide_idx)
                logger.info(f"  删除多余页面: 第{slide_idx + 1}页")
            # 只使用需要的页面
            target_slides = target_slides[:required_pages]

        # 每页放置2个图表
        for page_idx in range(min(len(target_slides), required_pages)):
            slide_index = target_slides[page_idx]
            slide = self.prs.slides[slide_index]

            # 计算当前页应该显示的图表索引
            chart_idx_1 = page_idx * 2
            chart_idx_2 = page_idx * 2 + 1

            logger.info(f"填充第{slide_index + 1}页 - 性状进展折线图")

            # 填充第一个图表（左侧）
            if chart_idx_1 < len(chart_images):
                self._add_chart_to_slide(slide, chart_images[chart_idx_1], position='left')
                logger.info(f"  ✓ 添加图表1: {chart_images[chart_idx_1]['title']}")

            # 填充第二个图表（右侧）
            if chart_idx_2 < len(chart_images):
                self._add_chart_to_slide(slide, chart_images[chart_idx_2], position='right')
                logger.info(f"  ✓ 添加图表2: {chart_images[chart_idx_2]['title']}")

    def _add_chart_to_slide(self, slide, chart_data: Dict, position: str):
        """
        将图表图片添加到幻灯片

        Args:
            slide: 幻灯片对象
            chart_data: 图表数据 {'index': 1, 'title': '经济指数', 'image_path': '/path/to/image.png'}
            position: 位置 'left' 或 'right'
        """
        image_path = chart_data['image_path']

        # 定义图片位置和大小
        if position == 'left':
            left = Inches(0.5)
            top = Inches(2.0)
        else:  # right
            left = Inches(7.0)
            top = Inches(2.0)

        width = Inches(6.0)
        height = Inches(4.5)

        # 查找并替换现有的图片占位符，或直接添加图片
        # 先尝试查找Picture占位符
        picture_found = False
        target_name = "图片 1" if position == 'left' else "图片 2"

        for shape in slide.shapes:
            if hasattr(shape, 'name') and target_name in shape.name:
                # 找到对应的图片占位符，删除它
                sp = shape._element
                sp.getparent().remove(sp)
                picture_found = True
                break

        # 添加新图片
        try:
            slide.shapes.add_picture(
                image_path,
                left=left,
                top=top,
                width=width,
                height=height
            )
        except Exception as e:
            logger.warning(f"添加图片失败: {e}")

    def _delete_slide(self, slide_index: int):
        """
        删除指定索引的幻灯片

        Args:
            slide_index: 幻灯片索引（从0开始）
        """
        try:
            rId = self.prs.slides._sldIdLst[slide_index].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[slide_index]
        except Exception as e:
            logger.warning(f"删除幻灯片{slide_index}失败: {e}")

    def _has_valid_data(self, chart_data: Dict) -> bool:
        """
        检查图表是否有有效数据（不是全部相同或变化太小）

        Args:
            chart_data: 图表数据 {'index': 1, 'title': '经济指数', 'data': {...}}

        Returns:
            bool: 是否有有效数据
        """
        try:
            # 获取已提取的图表数据
            extracted_data = chart_data.get('data')
            if not extracted_data or not extracted_data.get('series'):
                return False

            # 检查所有系列的数据
            has_valid_series = False
            for series in extracted_data['series']:
                values = series.get('values', [])
                series_name = series.get('name', '未命名系列')

                logger.debug(f"    检查系列'{series_name}': {len(values)}个数据点")

                if not values:
                    logger.debug(f"      -> 跳过：无数据")
                    continue

                # 过滤掉非数字值
                numeric_values = [v for v in values if isinstance(v, (int, float)) and not (isinstance(v, float) and (v != v))]  # 排除NaN

                if len(numeric_values) < 2:
                    logger.debug(f"      -> 跳过：数据点少于2个")
                    continue

                # 计算最大值和最小值
                min_val = min(numeric_values)
                max_val = max(numeric_values)
                value_range = max_val - min_val

                logger.debug(f"      范围: {min_val:.2f} - {max_val:.2f} = {value_range:.2f}")

                # 方法1：检查绝对差异（对于接近0的值）
                if value_range < 0.001:
                    logger.debug(f"      -> 跳过：变化范围太小(<0.001)")
                    continue

                # 方法2：检查相对变化率（对于大数值）
                # 计算平均值
                avg_val = sum(numeric_values) / len(numeric_values)
                if abs(avg_val) > 1:  # 如果平均值大于1，使用相对变化率
                    relative_change = (value_range / abs(avg_val)) * 100  # 百分比
                    logger.debug(f"      平均值: {avg_val:.2f}, 变化率: {relative_change:.1f}%")
                    # 如果变化率小于10%，认为变化太小（看起来像水平线）
                    if relative_change < 10:
                        logger.debug(f"      -> 跳过：变化率太小(<10%)")
                        continue
                    else:
                        logger.debug(f"      -> ✓ 有效：变化率{relative_change:.1f}% >= 10%")
                        has_valid_series = True
                else:
                    # 对于小于1的值，只要range > 0.001就认为有效
                    logger.debug(f"      -> ✓ 有效：小数值但有变化")
                    has_valid_series = True

            # 返回是否有有效系列
            if has_valid_series:
                return True

            # 所有系列的数据都相同或变化太小
            logger.debug(f"  图表'{chart_data['title']}'无有效数据（数据全部相同或变化太小）")
            return False

        except Exception as e:
            logger.debug(f"检查图表数据有效性失败: {e}")
            # 出错时假设有数据，避免误删
            return True

    def _remove_chart_shape(self, chart_shape):
        """
        从幻灯片中删除图表shape

        Args:
            chart_shape: 图表shape对象
        """
        try:
            sp = chart_shape._element
            sp.getparent().remove(sp)
            logger.debug(f"  已删除图表shape")
        except Exception as e:
            logger.warning(f"删除图表shape失败: {e}")

    def _center_chart_in_slide(self, chart_shape):
        """
        将图表在幻灯片中居中显示

        Args:
            chart_shape: 图表shape对象
        """
        try:
            # 获取幻灯片尺寸
            slide_width = self.prs.slide_width
            slide_height = self.prs.slide_height

            # 保持图表原有尺寸
            chart_width = chart_shape.width
            chart_height = chart_shape.height

            # 计算居中位置
            # 水平居中
            chart_shape.left = (slide_width - chart_width) // 2
            # 垂直方向稍微偏上一点（距离顶部40%）
            chart_shape.top = int(slide_height * 0.3)

            logger.debug(f"  图表已居中: left={chart_shape.left}, top={chart_shape.top}")

        except Exception as e:
            logger.warning(f"居中图表失败: {e}")

    def _generate_traits_trends_analysis(self, chart1_data: dict, chart2_data: dict) -> str:
        """
        生成性状进展趋势分析文本

        Args:
            chart1_data: 第一个图表数据 {'title': '...', 'data': {...}}
            chart2_data: 第二个图表数据

        Returns:
            分析文本字符串
        """
        parts = []

        try:
            # 分析第一个图表
            if chart1_data:
                title1 = chart1_data.get('title', '')
                data1 = chart1_data.get('data')
                analysis1 = self._analyze_single_chart(title1, data1)
                if analysis1:
                    parts.append(analysis1)

            # 分析第二个图表
            if chart2_data:
                title2 = chart2_data.get('title', '')
                data2 = chart2_data.get('data')
                analysis2 = self._analyze_single_chart(title2, data2)
                if analysis2:
                    parts.append(analysis2)

            if parts:
                return " ".join(parts)
            else:
                return "冻精使用的各性状指标保持稳定发展态势。"

        except Exception as e:
            logger.warning(f"生成性状趋势分析失败: {e}")
            return ""

    def _analyze_single_chart(self, title: str, data: dict) -> str:
        """
        分析单个图表的趋势

        Args:
            title: 图表标题
            data: 图表数据 {'categories': [...], 'series': [...]}

        Returns:
            分析文本
        """
        if not data or not data.get('series'):
            return ""

        try:
            # 获取类别（年份）
            categories = data.get('categories', [])
            series_list = data.get('series', [])

            if not series_list:
                return ""

            # 分析主要系列的趋势
            trends = []
            for series in series_list[:2]:  # 只分析前两个系列
                name = series.get('name', '')
                values = series.get('values', [])

                if len(values) >= 2:
                    first_val = values[0]
                    last_val = values[-1]

                    if first_val is not None and last_val is not None:
                        change = last_val - first_val

                        if abs(change) > 0.01:  # 有变化
                            direction = "上升" if change > 0 else "下降"
                            # 直接写数据，不写百分比
                            first_year = categories[0] if categories else "起始"
                            last_year = categories[-1] if categories else "最新"
                            trends.append(f"{name}从{first_year}年的{first_val:.2f}{direction}至{last_year}年的{last_val:.2f}")

            if trends:
                return f"{title}方面，" + "；".join(trends) + "。"
            else:
                return f"{title}保持稳定。"

        except Exception as e:
            logger.debug(f"分析图表'{title}'失败: {e}")
            return ""

    def _fill_analysis_text(self, slide, analysis_text: str):
        """
        填充分析文本到页面

        Args:
            slide: 幻灯片对象
            analysis_text: 分析文本
        """
        if not analysis_text:
            return

        # 查找分析文本框
        textbox = None
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                name = shape.name if shape.name else ""
                # 查找文本框（通常在页面下方）
                if "文本框" in name or "TextBox" in name:
                    if hasattr(shape, "top") and shape.top > 4000000:  # 约4cm以下
                        textbox = shape
                        break

        if not textbox:
            # 查找任何可用的文本框
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and "文本框" in (shape.name or ""):
                    textbox = shape
                    break

        if textbox:
            try:
                # 清空文本框并设置新内容
                tf = textbox.text_frame
                tf.clear()
                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                run = p.add_run()
                # 避免重复"分析："前缀
                if analysis_text.startswith("分析：") or analysis_text.startswith("分析:"):
                    run.text = analysis_text
                else:
                    run.text = "分析：" + analysis_text

                # 设置字体：微软雅黑15号非加粗
                run.font.name = FONT_NAME_CN
                run.font.size = Pt(15)
                run.font.bold = False

                logger.info(f"✓ 填充性状趋势分析文本")
            except Exception as e:
                logger.warning(f"填充分析文本失败: {e}")
        else:
            logger.debug("未找到性状趋势分析文本框")
