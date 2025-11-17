"""
Slide 1: 封面页构建器
完全按照模板规则实现
"""

import logging
from datetime import datetime
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


def build_slide_01_cover(prs, farm_name: str, reporter_name: str = "XXX"):
    """
    构建Slide 1: 封面

    完全按照模板规则实现：
    - 使用Layout [0] 封面
    - 元素1: 标题 AUTO_SHAPE
    - 元素2: 日期和服务人员 PLACEHOLDER

    Args:
        prs: PowerPoint演示文稿对象
        farm_name: 牧场名称
        reporter_name: 服务人员姓名
    """
    logger.info("开始构建Slide 1: 封面")

    # 使用Layout [0] 封面
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # ============ 元素1: 标题 ============
    # 规则: AUTO_SHAPE, 位置(13.83cm, 9.26cm), 尺寸(18.31cm × 3.1cm)
    # 填充: None, 垂直对齐: MIDDLE, 段落对齐: JUSTIFY
    # 字体: 44pt, 主题颜色: TEXT_1

    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(13.83), Cm(9.26),    # 位置（按规则）
        Cm(18.31), Cm(3.1)      # 尺寸（按规则）
    )

    # 无填充，无边框
    title_shape.fill.background()
    title_shape.line.fill.background()

    # 文本框设置
    title_frame = title_shape.text_frame
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # 垂直居中（规则: MIDDLE (3)）
    title_frame.word_wrap = None  # 规则: None
    title_frame.margin_left = Cm(0.254)    # 规则: 左0.254cm
    title_frame.margin_right = Cm(0.254)   # 规则: 右0.254cm
    title_frame.margin_top = Cm(0.127)     # 规则: 上0.127cm
    title_frame.margin_bottom = Cm(0.127)  # 规则: 下0.127cm

    # 段落1: 标题文本
    # 规则: 对齐JUSTIFY (4), 44.0pt, 主题颜色TEXT_1
    title_para = title_frame.paragraphs[0]
    title_para.text = f"{farm_name}育种分析综合报告"
    title_para.alignment = PP_ALIGN.JUSTIFY  # 两端对齐（规则: JUSTIFY (4)）

    # 设置字体 - 必须在每个run上设置
    # 实际颜色：蓝色RGB(0, 132, 213)，加粗，有文字阴影
    for run in title_para.runs:
        run.font.name = '微软雅黑'  # 字体：微软雅黑
        run.font.size = Pt(44)  # 规则: 44.0pt
        run.font.color.rgb = RGBColor(0, 132, 213)  # 蓝色
        run.font.bold = True  # 加粗

        # 添加文字阴影
        try:
            from lxml import etree
            # 获取run的底层XML元素
            r_elem = run._r  # run元素
            # 获取或创建rPr（run properties）
            rPr = r_elem.rPr
            if rPr is None:
                rPr = etree.SubElement(r_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')

            # DrawingML命名空间
            a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

            # 添加阴影效果
            effectLst = etree.Element(f'{{{a_ns}}}effectLst')
            outerShdw = etree.SubElement(effectLst, f'{{{a_ns}}}outerShdw')
            outerShdw.set('blurRad', '38100')  # 阴影模糊半径
            outerShdw.set('dist', '38100')     # 阴影距离
            outerShdw.set('dir', '2700000')    # 阴影方向（270度，向下）
            outerShdw.set('algn', 'ctr')       # 对齐方式

            # 阴影颜色（黑色半透明）
            srgbClr = etree.SubElement(outerShdw, f'{{{a_ns}}}srgbClr')
            srgbClr.set('val', '000000')
            alpha = etree.SubElement(srgbClr, f'{{{a_ns}}}alpha')
            alpha.set('val', '50000')  # 50% 透明度

            rPr.append(effectLst)
        except Exception as e:
            logger.warning(f"无法添加文字阴影: {e}")

    # ============ 元素2: 日期和服务人员 ============
    # 规则: PLACEHOLDER, 位置(23.35cm, 13.81cm), 尺寸(8.79cm × 4.21cm)
    # 段落1: "年月", 段落2: "服务人员：XXX"
    # 字体: 24.0pt, 颜色类型: None

    info_box = slide.shapes.add_textbox(
        Cm(23.35), Cm(13.81),   # 位置（按规则）
        Cm(8.79), Cm(4.21)      # 尺寸（按规则）
    )

    # 文本框设置
    info_frame = info_box.text_frame
    info_frame.vertical_anchor = None  # 规则: None
    info_frame.word_wrap = None        # 规则: None
    info_frame.margin_left = Cm(0.254)
    info_frame.margin_right = Cm(0.254)
    info_frame.margin_top = Cm(0.127)
    info_frame.margin_bottom = Cm(0.127)

    # 段落1: 日期
    # 规则: 对齐None, 24.0pt, 颜色None
    para1 = info_frame.paragraphs[0]
    current_date = datetime.now()
    para1.text = f"{current_date.year}年{current_date.month}月"
    para1.alignment = None  # 规则: None (继承默认)

    # 设置字体
    # 实际颜色：蓝色RGB(0, 132, 213)，不加粗，无阴影
    for run in para1.runs:
        run.font.name = '微软雅黑'  # 字体：微软雅黑
        run.font.size = Pt(24)  # 规则: 24.0pt
        run.font.color.rgb = RGBColor(0, 132, 213)  # 蓝色
        run.font.bold = False  # 不加粗
        # 不添加阴影

    # 段落2: 服务人员
    # 规则: 对齐None, 24.0pt, 颜色None
    para2 = info_frame.add_paragraph()
    para2.text = f"服务人员：{reporter_name}"
    para2.alignment = None  # 规则: None (继承默认)

    # 设置字体
    # 实际颜色：蓝色RGB(0, 132, 213)，不加粗，无阴影
    for run in para2.runs:
        run.font.name = '微软雅黑'  # 字体：微软雅黑
        run.font.size = Pt(24)  # 规则: 24.0pt
        run.font.color.rgb = RGBColor(0, 132, 213)  # 蓝色
        run.font.bold = False  # 不加粗
        # 不添加阴影

    logger.info("✓ Slide 1（封面）构建完成")
    return slide
