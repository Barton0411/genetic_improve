"""
Part 7: 备选公牛排名分析构建器
对应PPT模版的"备选公牛排名分析"页面（第125页）
"""

import logging
from typing import Dict, Any, List, Optional, Tuple
from pathlib import Path
import tempfile
import subprocess

import pandas as pd
from openpyxl import load_workbook
from openpyxl.styles import PatternFill
from openpyxl.utils import get_column_letter
from PIL import Image
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor

from ..base_builder import BaseSlideBuilder

logger = logging.getLogger(__name__)


class Part7CandidateBullsRankingBuilder(BaseSlideBuilder):
    """Part 7: 备选公牛排名分析构建器"""

    def __init__(self, prs, farm_name: str):
        super().__init__(prs, None)
        self.farm_name = farm_name

    def build(self, data: Dict[str, Any]):
        """
        构建 Part 7: 备选公牛排名分析

        Args:
            data: 包含 'excel_path' 和 'bull_ranking' 的数据字典
        """
        logger.info("=" * 60)
        logger.info("开始构建Part 7: 备选公牛排名分析")
        logger.info("=" * 60)

        # 获取数据
        excel_path = data.get("excel_path")
        bull_ranking_data = data.get("bull_ranking", {})

        if not excel_path:
            logger.error(f"❌ excel_path 未找到，跳过备选公牛排名分析")
            logger.error(f"   data中的键: {list(data.keys())}")
            return

        # 确保excel_path是字符串
        excel_path = str(excel_path)

        try:
            # 动态查找包含"备选公牛排名分析"的页面
            target_slides = self.find_slides_by_text("备选公牛排名分析", start_index=0)

            if not target_slides:
                logger.error("❌ 未找到备选公牛排名分析页面，跳过")
                return

            if not bull_ranking_data:
                logger.warning("⚠️  bull_ranking 数据为空，标记页面待删除")
                self.mark_slides_for_deletion(target_slides)
                return

            slide_index = target_slides[0]
            slide = self.prs.slides[slide_index]
            logger.info(f"✓ 定位到第{slide_index + 1}页（备选公牛排名分析）")

            # 1. 优先使用预生成的图片（Excel生成阶段已导出）
            ranking_image = self._find_pregenerated_image(excel_path)

            # 2. 如果没有预生成图片，回退到AppleScript方式
            if ranking_image is None:
                logger.info("  未找到预生成图片，使用AppleScript方式导出...")
                ranking_image = self._convert_excel_table_to_image(excel_path)

            if ranking_image is None:
                logger.error("❌ 无法获取排名表图片")
                return

            # 3. 将图片插入到PPT
            self._insert_table_image_to_slide(slide, ranking_image)

            logger.info("✓ Part 7 备选公牛排名分析完成")

        except Exception as e:
            logger.error(f"❌ 构建备选公牛排名分析失败: {e}")
            logger.debug("错误详情:", exc_info=True)

    def _find_pregenerated_image(self, excel_path: str) -> Optional[Path]:
        """
        查找Excel生成阶段预导出的排名表图片

        Args:
            excel_path: Excel文件路径（用于定位analysis_results目录）

        Returns:
            Path: 图片路径，未找到返回None
        """
        try:
            # 从excel_path推断analysis_results目录
            # Excel通常在 project/reports/ 下，analysis_results在 project/analysis_results/
            excel_path = Path(excel_path)
            project_folder = excel_path.parent.parent  # reports的父目录
            analysis_folder = project_folder / "analysis_results"

            if not analysis_folder.exists():
                logger.debug(f"  analysis_results目录不存在: {analysis_folder}")
                return None

            # 查找最新的排名表图片（按文件名中的时间戳排序）
            pattern = "备选公牛排名表_*.png"
            matching_files = list(analysis_folder.glob(pattern))

            if not matching_files:
                logger.debug(f"  未找到匹配的图片文件: {pattern}")
                return None

            # 按修改时间排序，取最新的
            latest_image = max(matching_files, key=lambda f: f.stat().st_mtime)
            logger.info(f"  ✓ 找到预生成图片: {latest_image.name}")
            return latest_image

        except Exception as e:
            logger.debug(f"查找预生成图片失败: {e}")
            return None

    def _convert_excel_table_to_image(self, excel_path: str) -> Optional[Path]:
        """
        使用xlwings将Excel排名表导出为高清图片

        Args:
            excel_path: Excel文件路径

        Returns:
            Path: 图片文件路径，失败返回None
        """
        try:
            # 尝试使用xlwings（支持Mac和Windows）
            try:
                import xlwings as xw
                return self._export_with_xlwings(excel_path)
            except ImportError:
                logger.warning("xlwings未安装，尝试使用截图方法")
                return self._export_with_screenshot(excel_path)
            except Exception as e:
                logger.warning(f"xlwings导出失败: {e}，尝试使用截图方法")
                return self._export_with_screenshot(excel_path)

        except Exception as e:
            logger.error(f"❌ 转换Excel表格为图片失败: {e}")
            logger.debug("错误详情:", exc_info=True)
            return None

    def _export_with_xlwings(self, excel_path: str) -> Path:
        """
        使用xlwings和openpyxl创建高质量表格图片

        Args:
            excel_path: Excel文件路径

        Returns:
            Path: 图片路径
        """
        logger.info("  使用Excel数据生成高质量图片...")

        # 使用openpyxl读取表格数据和格式
        wb_openpyxl = load_workbook(excel_path, data_only=True)
        ws_openpyxl = wb_openpyxl['备选公牛排名']

        # 查找表格起始行
        ranking_start_row = None
        for row_idx in range(1, 21):
            cell_value = ws_openpyxl.cell(row=row_idx, column=1).value
            if cell_value and "备选公牛排名" in str(cell_value):
                ranking_start_row = row_idx
                break

        if not ranking_start_row:
            wb_openpyxl.close()
            raise ValueError("未找到备选公牛排名表")

        # 查找最后一行
        last_row = ranking_start_row
        for row_idx in range(ranking_start_row, ws_openpyxl.max_row + 1):
            cell_value = ws_openpyxl.cell(row=row_idx, column=1).value
            if pd.isna(cell_value) or cell_value is None or str(cell_value).strip() == '':
                last_row = row_idx - 1
                break
            last_row = row_idx

        # 查找最后一列（从表头行查找）
        header_row = ranking_start_row + 2
        last_col = 1
        for col_idx in range(1, ws_openpyxl.max_column + 1):
            cell_value = ws_openpyxl.cell(row=header_row, column=col_idx).value
            if cell_value is not None:
                last_col = col_idx

        logger.info(f"  表格范围: A{ranking_start_row}:{get_column_letter(last_col)}{last_row} ({last_row - ranking_start_row + 1}行 × {last_col}列)")

        # 创建临时Excel文件，只包含该表格
        temp_dir = Path(tempfile.gettempdir())
        temp_excel_path = temp_dir / f"excel_table_{id(self)}.xlsx"

        # 创建新工作簿并复制数据
        from openpyxl import Workbook
        temp_wb = Workbook()
        temp_ws = temp_wb.active
        temp_ws.title = '表格'

        # 复制数据和格式
        for row_idx in range(ranking_start_row, last_row + 1):
            for col_idx in range(1, last_col + 1):
                source_cell = ws_openpyxl.cell(row=row_idx, column=col_idx)
                target_cell = temp_ws.cell(row=row_idx - ranking_start_row + 1, column=col_idx)

                # 复制值
                target_cell.value = source_cell.value

                # 复制样式
                if source_cell.has_style:
                    target_cell.font = source_cell.font.copy()
                    target_cell.border = source_cell.border.copy()
                    target_cell.fill = source_cell.fill.copy()
                    target_cell.number_format = source_cell.number_format
                    target_cell.alignment = source_cell.alignment.copy()

        # 复制列宽
        for col_idx in range(1, last_col + 1):
            col_letter = get_column_letter(col_idx)
            if ws_openpyxl.column_dimensions[col_letter].width:
                temp_ws.column_dimensions[col_letter].width = ws_openpyxl.column_dimensions[col_letter].width

        # 复制行高
        for row_idx in range(ranking_start_row, last_row + 1):
            target_row_idx = row_idx - ranking_start_row + 1
            if ws_openpyxl.row_dimensions[row_idx].height:
                temp_ws.row_dimensions[target_row_idx].height = ws_openpyxl.row_dimensions[row_idx].height

        wb_openpyxl.close()

        # 保存临时Excel文件
        temp_wb.save(str(temp_excel_path))
        temp_wb.close()
        logger.info(f"  ✓ 创建临时Excel文件: {temp_excel_path.name}")

        # 使用AppleScript调用Excel导出为PDF
        image_path = temp_dir / f"excel_table_{id(self)}.png"
        pdf_path = temp_dir / f"excel_table_{id(self)}.pdf"

        # 使用AppleScript让Excel打开并导出PDF
        applescript = f'''
        set excelFile to POSIX file "{temp_excel_path}" as text
        set pdfFile to POSIX file "{pdf_path}" as text

        tell application "Microsoft Excel"
            set wb to open excelFile
            tell active sheet
                save as wb filename pdfFile file format PDF file format
            end tell
            close wb saving no
        end tell
        '''

        try:
            subprocess.run(['osascript', '-e', applescript], check=True, capture_output=True, text=True)
            logger.info(f"  ✓ 导出PDF: {pdf_path.name}")

            # 使用sips转换PDF为PNG
            subprocess.run(
                ['sips', '-s', 'format', 'png', str(pdf_path), '--out', str(image_path), '-Z', '2400'],
                check=True,
                capture_output=True,
                text=True
            )
            logger.info(f"  ✓ 转换为PNG: {image_path.name}")

            # 清理临时文件
            temp_excel_path.unlink()
            pdf_path.unlink()

            return image_path

        except Exception as e:
            logger.error(f"  AppleScript导出失败: {e}")
            # 清理临时文件
            if temp_excel_path.exists():
                temp_excel_path.unlink()
            if pdf_path.exists():
                pdf_path.unlink()
            raise

    def _export_with_screenshot(self, excel_path: str) -> Path:
        """
        使用openpyxl读取数据并用PIL渲染为高质量图片（带中文支持）

        Args:
            excel_path: Excel文件路径

        Returns:
            Path: 图片路径
        """
        logger.info("  使用openpyxl和PIL渲染表格...")

        # 使用openpyxl读取表格数据和格式
        wb = load_workbook(excel_path, data_only=True)
        ws = wb['备选公牛排名']

        # 查找表格起始行
        ranking_start_row = None
        for row_idx in range(1, 21):
            cell_value = ws.cell(row=row_idx, column=1).value
            if cell_value and "备选公牛排名" in str(cell_value):
                ranking_start_row = row_idx
                break

        if not ranking_start_row:
            wb.close()
            raise ValueError("未找到备选公牛排名表")

        # 查找最后一行
        last_row = ranking_start_row
        for row_idx in range(ranking_start_row, ws.max_row + 1):
            cell_value = ws.cell(row=row_idx, column=1).value
            if pd.isna(cell_value) or cell_value is None or str(cell_value).strip() == '':
                last_row = row_idx - 1
                break
            last_row = row_idx

        # 查找最后一列（从表头行查找）
        header_row = ranking_start_row + 2
        last_col = 1
        for col_idx in range(1, ws.max_column + 1):
            cell_value = ws.cell(row=header_row, column=col_idx).value
            if cell_value is not None:
                last_col = col_idx

        logger.info(f"  表格范围: A{ranking_start_row}:{get_column_letter(last_col)}{last_row} ({last_row - ranking_start_row + 1}行 × {last_col}列)")

        # 使用PIL绘制表格
        image_path = self._draw_table_image(ws, ranking_start_row, last_row, 1, last_col)

        wb.close()
        return image_path

    def _draw_table_image(self, ws, start_row: int, end_row: int, start_col: int, end_col: int) -> Path:
        """
        使用PIL绘制Excel表格为高质量图片（支持中文）

        Args:
            ws: openpyxl工作表对象
            start_row, end_row: 行范围（从1开始）
            start_col, end_col: 列范围（从1开始）

        Returns:
            Path: 生成的图片路径
        """
        from PIL import Image, ImageDraw, ImageFont

        # 设置参数（高分辨率）
        dpi_scale = 3  # 3倍分辨率
        cell_height = 30 * dpi_scale  # 单元格高度（像素）
        min_cell_width = 40 * dpi_scale  # 最小单元格宽度
        max_cell_width = 180 * dpi_scale  # 最大单元格宽度
        padding = 8 * dpi_scale  # 单元格内边距
        font_size = 10 * dpi_scale  # 字体大小
        border_width = 1 * dpi_scale  # 边框宽度

        # 计算每列的宽度（根据Excel列宽）
        col_widths = {}
        for col in range(start_col, end_col + 1):
            col_letter = get_column_letter(col)
            excel_width = ws.column_dimensions[col_letter].width
            if excel_width:
                # Excel列宽转像素（approximate）
                pixel_width = int(excel_width * 8 * dpi_scale)
            else:
                # 根据内容计算宽度
                max_width = min_cell_width
                for row in range(start_row, end_row + 1):
                    cell = ws.cell(row=row, column=col)
                    if cell.value:
                        text = str(cell.value)
                        # 估算宽度（中文字符约15像素，英文约8像素）
                        text_width = sum(15*dpi_scale if ord(c) > 127 else 8*dpi_scale for c in text) + padding * 2
                        max_width = max(max_width, min(text_width, max_cell_width))
                pixel_width = max_width
            col_widths[col] = pixel_width

        # 计算总尺寸
        total_width = sum(col_widths.values())
        total_height = (end_row - start_row + 1) * cell_height

        logger.info(f"  生成高清图片: {total_width}×{total_height}像素 (DPI×{dpi_scale})")

        # 创建图片
        img = Image.new('RGB', (total_width, total_height), color='white')
        draw = ImageDraw.Draw(img)

        # 加载中文字体
        font = None
        font_paths = [
            "/System/Library/Fonts/PingFang.ttc",  # Mac中文字体
            "/System/Library/Fonts/STHeiti Light.ttc",  # Mac中文字体
            "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",  # Unicode字体
            "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",  # Linux中文字体
        ]

        for font_path in font_paths:
            try:
                font = ImageFont.truetype(font_path, font_size)
                logger.debug(f"  使用字体: {font_path}")
                break
            except:
                continue

        if font is None:
            logger.warning("  未找到中文字体，使用默认字体")
            font = ImageFont.load_default()

        # 绘制单元格
        y = 0
        for row in range(start_row, end_row + 1):
            x = 0
            for col in range(start_col, end_col + 1):
                cell = ws.cell(row=row, column=col)
                cell_width = col_widths[col]

                # 获取背景色
                bg_color = (255, 255, 255)  # 默认白色
                if cell.fill and cell.fill.fill_type == 'solid':
                    color = cell.fill.start_color
                    if color and color.type == 'rgb':
                        rgb_str = color.rgb
                        if rgb_str and len(rgb_str) >= 6:
                            if len(rgb_str) == 8:
                                rgb_str = rgb_str[2:]
                            try:
                                r = int(rgb_str[0:2], 16)
                                g = int(rgb_str[2:4], 16)
                                b = int(rgb_str[4:6], 16)
                                bg_color = (r, g, b)
                            except:
                                pass

                # 绘制背景
                draw.rectangle([x, y, x + cell_width, y + cell_height], fill=bg_color, outline=(180, 180, 180), width=border_width)

                # 绘制文本
                if cell.value:
                    text = str(cell.value)
                    # 简化换行符
                    text = text.replace('\n', ' ').replace('\r', '')

                    # 文本颜色（默认黑色）
                    text_color = (0, 0, 0)

                    # 获取文本字体颜色（如果有）
                    if cell.font and cell.font.color:
                        try:
                            font_color = cell.font.color
                            if font_color.type == 'rgb' and font_color.rgb:
                                rgb_str = font_color.rgb
                                if len(rgb_str) == 8:
                                    rgb_str = rgb_str[2:]
                                r = int(rgb_str[0:2], 16)
                                g = int(rgb_str[2:4], 16)
                                b = int(rgb_str[4:6], 16)
                                text_color = (r, g, b)
                        except:
                            pass

                    # 绘制文本（居中）
                    try:
                        bbox = draw.textbbox((0, 0), text, font=font)
                        text_width = bbox[2] - bbox[0]
                        text_height = bbox[3] - bbox[1]
                    except:
                        # 估算文本尺寸
                        text_width = sum(15*dpi_scale if ord(c) > 127 else 8*dpi_scale for c in text)
                        text_height = font_size

                    text_x = x + (cell_width - text_width) // 2
                    text_y = y + (cell_height - text_height) // 2

                    draw.text((text_x, text_y), text, fill=text_color, font=font)

                x += cell_width
            y += cell_height

        # 保存图片（高质量）
        temp_dir = Path(tempfile.gettempdir())
        image_path = temp_dir / f"table_ranking_{id(self)}.png"
        img.save(str(image_path), 'PNG', quality=95, dpi=(300, 300))

        logger.info(f"  ✓ 高清图片已保存: {image_path.name}")
        return image_path

    def _insert_table_image_to_slide(self, slide, image_path: Path):
        """
        将表格图片插入到幻灯片

        Args:
            slide: PPT幻灯片对象
            image_path: 图片文件路径
        """
        try:
            # 查找并删除旧的排名表（Table 4）
            old_table_shape = None
            for shape in slide.shapes:
                if shape.has_table and hasattr(shape, 'name'):
                    if '表格 4' in shape.name or shape.name == '表格 4':
                        old_table_shape = shape
                        break

            if old_table_shape:
                # 保存位置
                old_left = old_table_shape.left
                old_top = old_table_shape.top

                # 删除旧表格
                sp = old_table_shape._element
                sp.getparent().remove(sp)
                logger.info("  ✓ 删除旧排名表")
            else:
                # 默认位置
                old_left = Cm(2)
                old_top = Cm(8)

            # 插入图片
            # PPT页面宽度约33.867cm，留出边距后最大32cm
            max_width = Cm(32)
            max_height = Cm(6)

            # 插入图片并获取对象
            pic = slide.shapes.add_picture(
                str(image_path),
                left=old_left,
                top=old_top
            )

            # 按比例缩放以适应
            # 先尝试按高度缩放
            pic.height = max_height
            if pic.width > max_width:
                # 如果宽度超限，改为按宽度缩放
                pic.width = max_width

            logger.info(f"  ✓ 图片已插入: 宽{pic.width.cm:.1f}cm × 高{pic.height.cm:.1f}cm")

        except Exception as e:
            logger.error(f"❌ 插入图片失败: {e}")
            logger.debug("错误详情:", exc_info=True)

    def _read_excel_data(self, excel_path: str):
        """
        从Excel读取标准表和排名表数据及格式

        Returns:
            Tuple[pd.DataFrame, pd.DataFrame, Dict]: (标准表数据, 排名表数据, 排名表格式)
        """
        try:
            # 使用openpyxl读取工作簿（保留格式）
            wb = load_workbook(excel_path, data_only=True)
            ws = wb['备选公牛排名']

            # 读取所有数据到列表
            data = []
            for row in ws.iter_rows(values_only=False):
                row_data = [cell.value for cell in row]
                data.append(row_data)

            # 转换为DataFrame
            df = pd.DataFrame(data)
            logger.info(f"✓ 读取Excel: {df.shape[0]}行 × {df.shape[1]}列")

            # 标准表：行0-3
            standards_df = df.iloc[0:4].copy()
            logger.info(f"  标准表: {standards_df.shape[0]}行 × {standards_df.shape[1]}列")

            # 排名表：从行6开始查找"备选公牛排名"
            ranking_start_row = None
            for idx in range(4, len(df)):
                cell_value = df.iloc[idx, 0]
                if pd.notna(cell_value) and "备选公牛排名" in str(cell_value):
                    ranking_start_row = idx
                    break

            if ranking_start_row is None:
                logger.error("❌ 未找到备选公牛排名表")
                wb.close()
                return None, None, None

            # 排名表数据：从标题行到最后
            ranking_df = df.iloc[ranking_start_row:].copy()
            # 去掉尾部的空行
            ranking_df = ranking_df.dropna(how='all')
            logger.info(f"  排名表: {ranking_df.shape[0]}行 × {ranking_df.shape[1]}列，起始行={ranking_start_row}")

            # 提取排名表的格式信息（单元格背景色）
            ranking_formats = self._extract_cell_formats(ws, ranking_start_row, len(ranking_df), df.shape[1])
            logger.info(f"  提取到 {len(ranking_formats)} 个单元格格式")

            wb.close()
            return standards_df, ranking_df, ranking_formats

        except Exception as e:
            logger.error(f"❌ 读取Excel数据失败: {e}")
            logger.debug("错误详情:", exc_info=True)
            return None, None, None

    def _extract_cell_formats(self, ws, start_row: int, num_rows: int, num_cols: int) -> Dict:
        """
        提取Excel单元格的格式信息（背景色）

        Args:
            ws: openpyxl工作表对象
            start_row: 起始行（从0开始）
            num_rows: 行数
            num_cols: 列数

        Returns:
            Dict: {(row, col): {'fill_color': (R, G, B)}}
        """
        formats = {}

        try:
            for row_offset in range(num_rows):
                excel_row = start_row + row_offset + 1  # openpyxl行从1开始
                for col_offset in range(num_cols):
                    excel_col = col_offset + 1  # openpyxl列从1开始

                    # 获取单元格
                    cell = ws.cell(row=excel_row, column=excel_col)

                    # 检查填充色
                    if cell.fill and cell.fill.fill_type == 'solid':
                        # 获取RGB颜色
                        color = cell.fill.start_color
                        if color and color.type == 'rgb':
                            rgb_str = color.rgb
                            if rgb_str and len(rgb_str) >= 6:
                                # 移除alpha通道（前2位）
                                if len(rgb_str) == 8:
                                    rgb_str = rgb_str[2:]

                                # 转换为RGB整数
                                r = int(rgb_str[0:2], 16)
                                g = int(rgb_str[2:4], 16)
                                b = int(rgb_str[4:6], 16)

                                # 只保存非白色背景（白色通常是默认）
                                if not (r > 250 and g > 250 and b > 250):
                                    formats[(row_offset, col_offset)] = {
                                        'fill_color': (r, g, b)
                                    }
                                    logger.debug(f"  单元格({row_offset}, {col_offset})颜色: RGB({r}, {g}, {b})")

        except Exception as e:
            logger.warning(f"提取单元格格式时出错: {e}")
            logger.debug("错误详情:", exc_info=True)

        return formats

    def _update_standards_table(self, slide, standards_df: pd.DataFrame):
        """
        更新优质冻精技术标准表

        Args:
            slide: PPT幻灯片对象
            standards_df: 标准表数据（4行 × N列），行0是标题，行1-3是实际数据
        """
        logger.info("更新标准表...")

        try:
            # 查找标准表（Table 2，名称为"表格 2"）
            table = None
            for shape in slide.shapes:
                if shape.has_table and hasattr(shape, 'name'):
                    if '表格 2' in shape.name or shape.name == '表格 2':
                        table = shape.table
                        logger.info(f"  找到标准表: {shape.name}")
                        break

            if table is None:
                logger.warning("⚠️  未找到标准表（表格 2），跳过")
                return

            # Excel行0是"优质冻精技术标准"大标题，跳过
            # 使用行1-3作为表格数据（行1=表头，行2=美国后裔，行3=基因组）
            data_start_row = 1
            data_end_row = min(4, standards_df.shape[0])

            # 获取实际数据列数（去除空列）
            data_cols = []
            for col_idx in range(standards_df.shape[1]):
                col_data = standards_df.iloc[data_start_row:data_end_row, col_idx]
                if col_data.notna().any():
                    data_cols.append(col_idx)

            num_data_cols = len(data_cols)
            logger.info(f"  标准表数据: {data_end_row - data_start_row}行 × {num_data_cols}列（有效列）")

            # 检查表格列数是否匹配
            if num_data_cols > len(table.columns):
                logger.warning(f"  数据列数({num_data_cols}) > 表格列数({len(table.columns)})，将截断数据")
                data_cols = data_cols[:len(table.columns)]
            elif num_data_cols < len(table.columns):
                logger.info(f"  数据列数({num_data_cols}) < 表格列数({len(table.columns)})，将填充空白")

            # 填充数据（从Excel行1开始，即表头行）
            for ppt_row_idx in range(min(data_end_row - data_start_row, len(table.rows))):
                excel_row_idx = data_start_row + ppt_row_idx
                for col_idx, data_col_idx in enumerate(data_cols):
                    if col_idx < len(table.columns):
                        cell = table.cell(ppt_row_idx, col_idx)
                        value = standards_df.iloc[excel_row_idx, data_col_idx]
                        cell.text = str(value) if pd.notna(value) else ""

                        # 设置字体大小和对齐
                        if cell.text_frame and cell.text_frame.paragraphs:
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.alignment = PP_ALIGN.CENTER
                                for run in paragraph.runs:
                                    run.font.size = Pt(8)

            logger.info("✓ 标准表更新完成")

        except Exception as e:
            logger.error(f"❌ 更新标准表失败: {e}")
            logger.debug("错误详情:", exc_info=True)

    def _update_ranking_table(self, slide, ranking_df: pd.DataFrame, bull_ranking_data: Dict, ranking_formats: Dict):
        """
        更新备选公牛排名表（动态行列，保留Excel格式）

        Args:
            slide: PPT幻灯片对象
            ranking_df: 排名表原始数据
            bull_ranking_data: 从collector收集的数据（用于获取格式化后的数据）
            ranking_formats: Excel单元格格式（背景色等）
        """
        logger.info("更新排名表...")

        try:
            # 查找并删除旧的排名表（Table 4）
            old_table_shape = None
            for shape in slide.shapes:
                if shape.has_table and hasattr(shape, 'name'):
                    if '表格 4' in shape.name or shape.name == '表格 4':
                        old_table_shape = shape
                        logger.info(f"  找到旧排名表: {shape.name}")
                        break

            if old_table_shape is None:
                logger.warning("⚠️  未找到旧排名表（表格 4），跳过")
                return

            # 保存旧表格的位置
            old_left = old_table_shape.left
            old_top = old_table_shape.top

            # 删除旧表格
            sp = old_table_shape._element
            sp.getparent().remove(sp)
            logger.info("  ✓ 删除旧排名表")

            # 获取实际数据的行列数
            # 去除空列
            data_cols = []
            for col_idx in range(ranking_df.shape[1]):
                col_data = ranking_df.iloc[:, col_idx]
                if col_data.notna().any():
                    data_cols.append(col_idx)

            num_cols = len(data_cols)
            num_rows = len(ranking_df)

            logger.info(f"  创建新表格: {num_rows}行 × {num_cols}列")

            # 创建新表格（固定高度6cm，宽度自动适应）
            table_height = Cm(6)
            # 初始宽度：根据列数估算（每列约1cm，最小20cm，最大32cm）
            estimated_width = Cm(min(32, max(20, num_cols * 1)))

            new_table_shape = slide.shapes.add_table(
                rows=num_rows,
                cols=num_cols,
                left=old_left,
                top=old_top,
                width=estimated_width,
                height=table_height
            )

            new_table = new_table_shape.table
            logger.info(f"  ✓ 创建新表格: 初始宽度{estimated_width.cm}cm × 高度{table_height.cm}cm")

            # 根据列数和行数自适应字体
            font_size = self._calculate_font_size(num_cols, num_rows)
            logger.info(f"  自适应字体大小: {font_size}pt")

            # 获取列名（用于判断列类型）
            # 表头在第3行（索引2）
            column_names = []
            if num_rows > 2:
                for col_idx, data_col_idx in enumerate(data_cols):
                    col_name = str(ranking_df.iloc[2, data_col_idx]) if pd.notna(ranking_df.iloc[2, data_col_idx]) else ""
                    column_names.append(col_name)
            else:
                column_names = [''] * num_cols

            # 填充数据
            for row_idx in range(num_rows):
                for col_idx, data_col_idx in enumerate(data_cols):
                    cell = new_table.cell(row_idx, col_idx)
                    value = ranking_df.iloc[row_idx, data_col_idx]

                    # 格式化数值
                    col_name = column_names[col_idx] if col_idx < len(column_names) else ""
                    formatted_value = self._format_value(value, row_idx, col_name)
                    cell.text = formatted_value

                    # 应用Excel背景色（在设置文本后，格式化前）
                    # 注意：ranking_formats的键是(row_offset, col_offset)，对应Excel的原始位置
                    if ranking_formats and (row_idx, data_col_idx) in ranking_formats:
                        cell_format = ranking_formats[(row_idx, data_col_idx)]
                        if 'fill_color' in cell_format:
                            self._apply_cell_fill(cell, cell_format['fill_color'])

                    # 设置单元格格式（字体、对齐）- 在应用背景色之后
                    self._format_cell(cell, row_idx, font_size)

            # 特殊处理：第一行标题合并单元格
            if num_rows > 0 and num_cols > 1:
                try:
                    first_cell = new_table.cell(0, 0)
                    last_cell = new_table.cell(0, num_cols - 1)
                    first_cell.merge(last_cell)
                    logger.debug("  ✓ 合并第一行标题单元格")
                except Exception as e:
                    logger.debug(f"  合并标题单元格失败（非关键）: {e}")

            # 特殊处理：第二行总计合并单元格
            if num_rows > 1 and num_cols > 1:
                try:
                    first_cell = new_table.cell(1, 0)
                    # 只合并部分单元格，避免覆盖所有列
                    merge_to_col = min(num_cols - 1, 5)  # 最多合并到第5列
                    last_cell = new_table.cell(1, merge_to_col)
                    first_cell.merge(last_cell)
                    logger.debug(f"  ✓ 合并第二行总计单元格（0-{merge_to_col}列）")
                except Exception as e:
                    logger.debug(f"  合并总计单元格失败（非关键）: {e}")

            logger.info("✓ 排名表更新完成")

        except Exception as e:
            logger.error(f"❌ 更新排名表失败: {e}")
            logger.debug("错误详情:", exc_info=True)

    def _format_value(self, value, row_idx: int, col_name: str) -> str:
        """
        格式化单元格值

        Args:
            value: 原始值
            row_idx: 行索引
            col_name: 列名（用于判断列类型）

        Returns:
            str: 格式化后的字符串
        """
        # 前3行（标题、总计、表头）不格式化
        if row_idx < 3:
            return str(value) if pd.notna(value) else ""

        # 空值
        if pd.notna(value) is False:
            return ""

        # 尝试转换为数字
        try:
            num_value = float(value)

            # 排名列：整数
            if 'ranking' in col_name.lower() or '排名' in col_name:
                return str(int(num_value))

            # 育种指数列：整数
            if 'index' in col_name.lower() or '指数' in col_name:
                return str(int(num_value))

            # 其他数值列：1位小数
            return f"{num_value:.1f}"

        except (ValueError, TypeError):
            # 非数值，直接返回字符串
            return str(value)

    def _calculate_font_size(self, num_cols: int, num_rows: int) -> int:
        """
        根据列数和行数计算合适的字体大小

        Args:
            num_cols: 列数
            num_rows: 行数

        Returns:
            int: 字体大小（pt）
        """
        # 根据列数和行数调整字体
        if num_cols <= 15:
            base_font = 7
        elif num_cols <= 25:
            base_font = 6
        elif num_cols <= 35:
            base_font = 5
        else:
            base_font = 5

        if num_rows > 15:
            base_font = max(5, base_font - 1)

        # 最小字体5pt，最大字体8pt
        return max(5, min(8, base_font))

    def _apply_cell_fill(self, cell, rgb_color: Tuple[int, int, int]):
        """
        应用单元格填充色

        Args:
            cell: 表格单元格对象
            rgb_color: RGB颜色元组 (R, G, B)
        """
        try:
            # 获取单元格填充
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(rgb_color[0], rgb_color[1], rgb_color[2])
            logger.debug(f"  应用填充色: RGB{rgb_color}")
        except Exception as e:
            logger.debug(f"  应用填充色失败: {e}")

    def _format_cell(self, cell, row_idx: int, font_size: int):
        """
        格式化单元格

        Args:
            cell: 表格单元格对象
            row_idx: 行索引
            font_size: 字体大小（pt）
        """
        if not cell.text_frame or not cell.text_frame.paragraphs:
            return

        for paragraph in cell.text_frame.paragraphs:
            # 设置对齐方式
            if row_idx == 0 or row_idx == 1:  # 标题行和总计行
                paragraph.alignment = PP_ALIGN.CENTER
            elif row_idx == 2:  # 表头行
                paragraph.alignment = PP_ALIGN.CENTER
            else:  # 数据行
                # 判断是否为数字
                try:
                    float(cell.text)
                    paragraph.alignment = PP_ALIGN.RIGHT
                except (ValueError, AttributeError):
                    paragraph.alignment = PP_ALIGN.CENTER

            # 设置字体
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                if row_idx <= 2:  # 标题、总计、表头加粗
                    run.font.bold = True
