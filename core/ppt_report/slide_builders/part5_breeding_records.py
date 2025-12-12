"""
Part 5: 配种记录分析构建器
对应PPT模版第96-98页
"""

import logging
from typing import Dict, Optional

import pandas as pd
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from ..base_builder import BaseSlideBuilder
from ..config import COLOR_TEXT_MAIN, FONT_NAME_CN

logger = logging.getLogger(__name__)


class Part5BreedingRecordsBuilder(BaseSlideBuilder):
    """Part 5: 配种记录分析构建器"""

    def __init__(self, prs, chart_creator, farm_name: str):
        super().__init__(prs, chart_creator)
        self.farm_name = farm_name

    def build(self, data: Dict):
        """
        构建 Part 5: 配种记录分析

        Args:
            data: 包含 'breeding_genes' 的数据字典
        """
        logger.info("=" * 60)
        logger.info("开始构建Part 5: 配种记录分析")
        logger.info("=" * 60)

        # 使用 DataCollector 缓存读取配种记录-隐性基因分析Sheet（header=None）
        data_collector = data.get("data_collector")
        if data_collector is None:
            logger.error("data_collector 未找到，跳过 Part 5")
            return

        df_genes = data_collector.get_raw_sheet("配种记录-隐性基因分析", header=None)

        # 查找标题为"配种记录分析-隐性基因"的页面，应该有2个
        logger.info(f"开始查找包含'配种记录分析-隐性基因'的页面，模板共{len(self.prs.slides)}页")
        target_slides = self.find_slides_by_text("配种记录分析-隐性基因", start_index=0, max_count=2)

        if len(target_slides) == 0:
            logger.error("❌ 未找到配种记录分析-隐性基因页面，跳过 Part 5")
            logger.error("请检查PPT模板中是否存在标题为'配种记录分析-隐性基因'的页面")
            return

        # 检查数据是否为空，如果为空则标记页面待删除
        if df_genes is None or df_genes.empty:
            logger.warning("配种记录-隐性基因分析数据为空，标记页面待删除")
            self.mark_slides_for_deletion(target_slides)
            return

        # 获取配种明细数据，用于统计总配次
        df_detail = data.get("breeding_detail")
        total_breeding_count = len(df_detail) if df_detail is not None else 0
        logger.info(f"✓ 配种记录总数: {total_breeding_count} 配次")

        if len(target_slides) < 2:
            logger.warning(f"⚠️ 只找到{len(target_slides)}个配种记录分析-隐性基因页面（索引: {target_slides}），预期2个")
        else:
            logger.info(f"✓ 找到2个目标页面: 第{target_slides[0]+1}页, 第{target_slides[1]+1}页")

        # 填充第一页 - 全部年份隐性基因纯合汇总
        if len(target_slides) > 0:
            self._fill_all_years_summary(df_genes, target_slides[0], total_breeding_count)

        # 填充第二页 - 近12个月隐性基因纯合汇总
        if len(target_slides) > 1:
            self._fill_recent_year_summary(df_genes, target_slides[1])

        logger.info("✓ Part 5 配种记录分析完成")

    def _fill_all_years_summary(self, df: pd.DataFrame, slide_index: int, total_count: int):
        """
        填充全部年份隐性基因纯合汇总表

        PPT中有2个表格：
        - 表格 12 (1行1列): 标题
        - 表格 10 (17行10列): 数据表格

        Args:
            df: 配种记录隐性基因数据（header=None读取）
            slide_index: 幻灯片索引
            total_count: 总配次数
        """
        try:
            slide = self.prs.slides[slide_index]
        except IndexError:
            logger.warning(f"无法访问第{slide_index + 1}页")
            return

        logger.info(f"填充第{slide_index + 1}页 - 全部年份隐性基因纯合汇总")

        # Excel结构（header=None读取）：
        # 行0：标题 "隐性基因纯合汇总表（全部配种记录年份，共XX配次）"
        # 行1：表头 "基因名称、翻译、..."
        # 行2-17：16行基因数据

        # 读取标题
        title_text = str(df.iloc[0, 0]).strip() if not pd.isna(df.iloc[0, 0]) else f"隐性基因纯合汇总表（全部配种记录年份，共{total_count}配次）"
        logger.info(f"第一个表的标题: {title_text}")

        # 读取16行基因数据（从行2开始）
        gene_data = self._read_gene_data(df, 2, num_rows=16)
        logger.info(f"读取了 {len(gene_data)} 行基因数据")

        # 填充PPT标题表格
        self._fill_title_table(slide, "表格 12", title_text)

        # 填充PPT数据表格
        self._fill_data_table(slide, "表格 10", gene_data)

        # 填充分析文本框（全部年份）
        self._fill_analysis_textbox(slide, gene_data, title_text, is_all_years=True)

    def _fill_recent_year_summary(self, df: pd.DataFrame, slide_index: int):
        """
        填充近12个月隐性基因纯合汇总表

        PPT中有2个表格 + 1个文本框：
        - 表格 12 (1行1列): 标题
        - 表格 10 (17行10列): 数据表格
        - 文本框 14: 分析说明

        Args:
            df: 配种记录隐性基因数据（header=None读取）
            slide_index: 幻灯片索引
        """
        try:
            slide = self.prs.slides[slide_index]
        except IndexError:
            logger.warning(f"无法访问第{slide_index + 1}页")
            return

        logger.info(f"填充第{slide_index + 1}页 - 近12个月隐性基因纯合汇总")

        # 查找包含"隐性基因纯合汇总表"且不在第0行的标题行（第二个表格）
        title_row_idx, title_text = self._find_second_title_row(df)
        if title_row_idx is None:
            logger.warning("未找到近12个月的标题行")
            return

        logger.info(f"找到第二个表格标题: 行{title_row_idx}, 标题: {title_text}")

        # 读取16行基因数据（标题行+2，跳过表头）
        data_start_row = title_row_idx + 2
        gene_data = self._read_gene_data(df, data_start_row, num_rows=16)
        logger.info(f"读取了 {len(gene_data)} 行基因数据（从行{data_start_row}开始）")

        # 填充PPT标题表格
        self._fill_title_table(slide, "表格 12", title_text)

        # 填充PPT数据表格
        self._fill_data_table(slide, "表格 10", gene_data)

        # 填充分析文本框（近一年）
        self._fill_analysis_textbox(slide, gene_data, title_text, is_all_years=False)

    # ------------------------------------------------------------------ #
    # 辅助方法
    # ------------------------------------------------------------------ #

    def _find_second_title_row(self, df: pd.DataFrame) -> tuple:
        """
        查找第二个包含"隐性基因纯合汇总表"的标题行
        返回 (行索引, 标题文本)
        """
        found_first = False
        for i in range(len(df)):
            cell_value = str(df.iloc[i, 0]).strip() if not pd.isna(df.iloc[i, 0]) else ""
            if "隐性基因纯合汇总表" in cell_value:
                if found_first:
                    # 这是第二个标题
                    return i, cell_value
                else:
                    # 跳过第一个，继续找第二个
                    found_first = True
        return None, None

    def _find_header_row(self, df: pd.DataFrame, start_from: int = 0) -> Optional[int]:
        """
        查找数据表格的表头行（包含"基因名称"）
        """
        for i in range(start_from, len(df)):
            cell_value = str(df.iloc[i, 0]).strip() if not pd.isna(df.iloc[i, 0]) else ""
            if cell_value == "基因名称":
                return i
        return None

    def _read_gene_data(self, df: pd.DataFrame, start_row: int, num_rows: int = 16) -> list:
        """
        从Excel读取基因数据（16行 × 10列）

        返回格式: [
            [基因名称, 翻译, 纯合配次, 占比, 仅公牛, 占比, 仅母牛父亲, 占比, 数据缺少, 占比],
            ...
        ]
        """
        end_row = min(start_row + num_rows, len(df))
        max_cols = min(10, df.shape[1])

        # P1优化：批量转换为numpy数组，减少iloc访问开销
        data_slice = df.iloc[start_row:end_row, :max_cols].fillna("").astype(str)
        data_array = data_slice.to_numpy()

        gene_data = []
        for i in range(data_array.shape[0]):
            row_data = [v.strip() for v in data_array[i]]
            # 补齐到10列
            while len(row_data) < 10:
                row_data.append("")
            gene_data.append(row_data)

        return gene_data

    def _fill_title_table(self, slide, table_name: str, title_text: str):
        """填充标题表格（1行1列）"""
        table = self._find_table(slide, table_name)
        if table is None:
            logger.warning(f"未找到表格: {table_name}")
            return

        try:
            # 填充标题
            cell = table.cell(0, 0)
            cell.text = title_text

            # 设置格式
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = FONT_NAME_CN
                    run.font.size = Pt(14)
                    run.font.bold = True
                    run.font.color.rgb = COLOR_TEXT_MAIN

            logger.info(f"✓ 填充标题表格: {title_text}")
        except Exception as e:
            logger.error(f"填充标题表格失败: {e}")

    def _fill_data_table(self, slide, table_name: str, gene_data: list):
        """
        填充数据表格（17行10列）
        第一行是表头，后面16行是数据
        """
        table = self._find_table(slide, table_name)
        if table is None:
            logger.warning(f"未找到表格: {table_name}")
            return

        try:
            # 表头已在模版中，不需要修改
            # 填充16行基因数据（从第2行开始，索引1）
            for row_idx, row_data in enumerate(gene_data, start=1):
                if row_idx >= len(table.rows):
                    logger.warning(f"表格行数不足: {row_idx} >= {len(table.rows)}")
                    break

                for col_idx, cell_value in enumerate(row_data):
                    if col_idx >= len(table.columns):
                        break

                    cell = table.cell(row_idx, col_idx)
                    cell.text = cell_value

                    # 设置格式
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.name = FONT_NAME_CN
                            run.font.size = Pt(11)

            logger.info(f"✓ 填充数据表格: {len(gene_data)}行数据")
        except Exception as e:
            logger.error(f"填充数据表格失败: {e}")

    def _fill_analysis_textbox(self, slide, gene_data: list, title_text: str, is_all_years: bool = False):
        """
        填充隐性基因分析文本框

        Args:
            slide: 幻灯片对象
            gene_data: 基因数据列表
            title_text: 表格标题
            is_all_years: 是否为全部年份数据
        """
        # 查找分析文本框（优先查找位置靠下的文本框，避免匹配到标题）
        textbox = None
        candidate_textboxes = []

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                name = shape.name if shape.name else ""
                # 跳过标题表格
                if "表格" in name:
                    continue
                # 查找文本框
                if "文本框" in name or "TextBox" in name:
                    candidate_textboxes.append(shape)

        # 按位置（top）排序，选择最靠下的文本框
        if candidate_textboxes:
            candidate_textboxes.sort(key=lambda s: getattr(s, 'top', 0), reverse=True)
            textbox = candidate_textboxes[0]

        if textbox is None:
            logger.warning("未找到分析文本框")
            return

        try:
            # 统计纯合配次数（第3列）
            total_homozygous = 0
            for row_data in gene_data:
                if len(row_data) > 2:
                    try:
                        count = int(row_data[2]) if row_data[2] and row_data[2].isdigit() else 0
                        total_homozygous += count
                    except:
                        pass

            # 根据数据范围生成分析文本
            period = "全部配种记录" if is_all_years else "过往1年配种"

            if total_homozygous > 0:
                analysis_text = f"分析：牧场{period}过程中共计出现{total_homozygous}次隐性纯合配次，可能造成胎儿死亡或出生后发病，建议加强配种前的基因筛查。"
            else:
                analysis_text = f"分析：牧场{period}过程中未出现隐性纯合配次，配种管理良好，继续保持。"

            # 填充文本
            textbox.text_frame.clear()
            p = textbox.text_frame.paragraphs[0] if textbox.text_frame.paragraphs else textbox.text_frame.add_paragraph()
            run = p.add_run()
            run.text = analysis_text
            p.alignment = PP_ALIGN.LEFT

            # 设置字体：微软雅黑15号非加粗
            run.font.name = FONT_NAME_CN
            run.font.size = Pt(15)
            run.font.bold = False

            logger.info(f"✓ 填充隐性基因分析文本({period}): {total_homozygous}次纯合")
        except Exception as e:
            logger.error(f"填充分析文本框失败: {e}")

    def _find_table(self, slide, name: str):
        """查找指定名称的表格"""
        for shape in slide.shapes:
            if shape.name == name and hasattr(shape, "table"):
                return shape.table
        return None
