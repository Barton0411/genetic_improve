#!/usr/bin/env python3
"""
逐页提取所有58页的完整细节，生成PPT报告生成规则
"""

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
from pathlib import Path


def extract_shape_complete_details(shape, shape_idx):
    """提取形状的所有细节"""
    details = {
        'index': shape_idx,
        'name': shape.name,
        'type': str(shape.shape_type),
        'position': {
            'left_cm': round(shape.left / 360000, 2),
            'top_cm': round(shape.top / 360000, 2),
        },
        'size': {
            'width_cm': round(shape.width / 360000, 2),
            'height_cm': round(shape.height / 360000, 2),
        }
    }

    # 填充
    if hasattr(shape, 'fill'):
        details['fill'] = {
            'type': str(shape.fill.type)
        }

    # 线条
    if hasattr(shape, 'line'):
        details['line'] = {
            'width_pt': shape.line.width / 914400 if shape.line.width else None,
        }

    # 文本框详细信息
    if shape.has_text_frame:
        tf = shape.text_frame
        details['text_frame'] = {
            'text': tf.text[:200] if len(tf.text) > 200 else tf.text,
            'vertical_anchor': str(tf.vertical_anchor),
            'word_wrap': tf.word_wrap,
            'margin_left_cm': round(tf.margin_left / 360000, 3),
            'margin_right_cm': round(tf.margin_right / 360000, 3),
            'margin_top_cm': round(tf.margin_top / 360000, 3),
            'margin_bottom_cm': round(tf.margin_bottom / 360000, 3),
            'paragraphs': []
        }

        # 提取段落
        for para_idx, para in enumerate(tf.paragraphs):
            para_details = {
                'index': para_idx + 1,
                'text': para.text[:100] if len(para.text) > 100 else para.text,
                'alignment': str(para.alignment),
                'level': para.level,
                'runs': []
            }

            # 提取run
            for run_idx, run in enumerate(para.runs):
                if run.text:
                    run_details = {
                        'text': run.text[:50] if len(run.text) > 50 else run.text,
                        'font': {}
                    }

                    font = run.font
                    run_details['font']['name'] = font.name
                    run_details['font']['size_pt'] = font.size.pt if font.size else None
                    run_details['font']['bold'] = font.bold
                    run_details['font']['italic'] = font.italic
                    run_details['font']['underline'] = font.underline

                    # 颜色详细信息
                    if font.color:
                        color_info = {
                            'type': str(font.color.type)
                        }

                        try:
                            if hasattr(font.color, 'rgb') and font.color.rgb:
                                rgb = font.color.rgb
                                color_info['rgb'] = f'RGB({rgb[0]}, {rgb[1]}, {rgb[2]})'
                                color_info['hex'] = f'#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'
                        except:
                            pass

                        try:
                            if hasattr(font.color, 'theme_color'):
                                color_info['theme_color'] = str(font.color.theme_color)
                        except:
                            pass

                        try:
                            if hasattr(font.color, 'brightness'):
                                color_info['brightness'] = font.color.brightness
                        except:
                            pass

                        run_details['font']['color'] = color_info
                    else:
                        run_details['font']['color'] = None

                    para_details['runs'].append(run_details)

            details['text_frame']['paragraphs'].append(para_details)

    # 表格
    if shape.has_table:
        table = shape.table
        details['table'] = {
            'rows': len(table.rows),
            'columns': len(table.columns),
            'sample_cells': []
        }

        # 提取前几个单元格作为示例
        for row_idx in range(min(2, len(table.rows))):
            for col_idx in range(min(3, len(table.columns))):
                cell = table.cell(row_idx, col_idx)
                cell_details = {
                    'position': f'[{row_idx}, {col_idx}]',
                    'text': cell.text[:50] if cell.text else ''
                }

                # 单元格字体
                if cell.text_frame and cell.text_frame.paragraphs:
                    para = cell.text_frame.paragraphs[0]
                    if para.runs:
                        run = para.runs[0]
                        font = run.font
                        cell_details['font'] = {
                            'name': font.name,
                            'size_pt': font.size.pt if font.size else None,
                            'bold': font.bold
                        }

                        # 颜色
                        if font.color:
                            try:
                                if hasattr(font.color, 'rgb') and font.color.rgb:
                                    rgb = font.color.rgb
                                    cell_details['font']['color'] = f'RGB({rgb[0]}, {rgb[1]}, {rgb[2]})'
                            except:
                                pass

                    cell_details['alignment'] = str(para.alignment)

                details['table']['sample_cells'].append(cell_details)

    # 图表
    if shape.has_chart:
        chart = shape.chart
        details['chart'] = {
            'type': str(chart.chart_type),
            'has_title': chart.has_title,
            'title': chart.chart_title.text_frame.text if chart.has_title and chart.chart_title else None,
            'has_legend': chart.has_legend,
        }

    return details


def extract_all_slides_rules(template_path):
    """提取所有幻灯片的完整规则"""
    prs = Presentation(template_path)

    rules = {
        'template_info': {
            'path': str(template_path),
            'total_slides': len(prs.slides),
            'slide_width_inches': prs.slide_width.inches,
            'slide_height_inches': prs.slide_height.inches,
        },
        'layouts': [],
        'slides': []
    }

    # 提取Layouts
    print("提取Layouts...")
    for idx, layout in enumerate(prs.slide_layouts):
        layout_info = {
            'index': idx,
            'name': layout.name,
            'placeholders': []
        }

        for ph in layout.placeholders:
            try:
                layout_info['placeholders'].append({
                    'idx': ph.placeholder_format.idx,
                    'name': ph.name,
                    'type': str(ph.placeholder_format.type)
                })
            except:
                pass

        rules['layouts'].append(layout_info)

    # 逐页提取
    print(f"\n逐页提取 {len(prs.slides)} 张幻灯片的完整细节...")
    for slide_idx, slide in enumerate(prs.slides, 1):
        print(f"  提取 Slide {slide_idx}...")

        # 找到layout index
        layout_index = None
        for idx, layout in enumerate(prs.slide_layouts):
            if layout == slide.slide_layout:
                layout_index = idx
                break

        slide_info = {
            'slide_number': slide_idx,
            'layout_name': slide.slide_layout.name,
            'layout_index': layout_index,
            'shapes': []
        }

        # 提取所有形状
        for shape_idx, shape in enumerate(slide.shapes, 1):
            shape_details = extract_shape_complete_details(shape, shape_idx)
            slide_info['shapes'].append(shape_details)

        rules['slides'].append(slide_info)

    return rules


def generate_rules_document(rules, output_dir):
    """生成规则文档（Markdown格式）"""
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    # 主规则文档
    rules_md = output_dir / "PPT生成规则.md"

    with open(rules_md, 'w', encoding='utf-8') as f:
        f.write("# PPT报告自动生成规则\n\n")
        f.write("本文档包含所有58页PPT的完整生成规则和格式规范。\n\n")

        f.write("## 1. 基本信息\n\n")
        f.write(f"- 幻灯片尺寸: {rules['template_info']['slide_width_inches']:.2f}\" × {rules['template_info']['slide_height_inches']:.2f}\"\n")
        f.write(f"- 总页数: {rules['template_info']['total_slides']}\n\n")

        f.write("## 2. Layouts定义\n\n")
        for layout in rules['layouts']:
            f.write(f"### Layout [{layout['index']}]: {layout['name']}\n\n")
            if layout['placeholders']:
                f.write("**占位符**:\n")
                for ph in layout['placeholders']:
                    f.write(f"- [{ph['idx']}] {ph['name']} - {ph['type']}\n")
            f.write("\n")

        f.write("## 3. 页面分类\n\n")

        # 按类型分类
        dividers = []
        content_pages = []

        for slide in rules['slides']:
            if slide['layout_name'] == '间隔页':
                dividers.append(slide['slide_number'])
            elif slide['layout_name'] == '1_图片页':
                content_pages.append(slide['slide_number'])

        f.write(f"- **封面**: Slide 1\n")
        f.write(f"- **间隔页** ({len(dividers)}页): {', '.join(map(str, dividers[:10]))}...\n")
        f.write(f"- **内容页** ({len(content_pages)}页): {', '.join(map(str, content_pages[:10]))}...\n\n")

        f.write("## 4. 逐页规则\n\n")
        f.write("详见各页规则文档：\n\n")

        for slide in rules['slides']:
            slide_num = slide['slide_number']
            f.write(f"- [Slide {slide_num:02d}: {slide['layout_name']}](./rules/Slide_{slide_num:02d}.md)\n")

    print(f"\n✅ 主规则文档: {rules_md}")

    # 生成每页详细规则
    rules_detail_dir = output_dir / "rules"
    rules_detail_dir.mkdir(exist_ok=True)

    for slide in rules['slides']:
        slide_num = slide['slide_number']
        slide_md = rules_detail_dir / f"Slide_{slide_num:02d}.md"

        with open(slide_md, 'w', encoding='utf-8') as f:
            f.write(f"# Slide {slide_num}: {slide['layout_name']}\n\n")

            f.write("## 实现规则\n\n")
            f.write(f"- **使用Layout**: [{slide['layout_index']}] {slide['layout_name']}\n")
            f.write(f"- **元素数量**: {len(slide['shapes'])}\n\n")

            f.write("## 元素清单\n\n")

            for shape in slide['shapes']:
                f.write(f"### 元素{shape['index']}: {shape['name']}\n\n")
                f.write(f"```python\n")
                f.write(f"# 类型: {shape['type']}\n")
                f.write(f"# 位置: Cm({shape['position']['left_cm']}), Cm({shape['position']['top_cm']})\n")
                f.write(f"# 尺寸: Cm({shape['size']['width_cm']}), Cm({shape['size']['height_cm']})\n")

                if shape.get('fill'):
                    f.write(f"# 填充: {shape['fill']['type']}\n")

                if shape.get('text_frame'):
                    tf = shape['text_frame']
                    f.write(f"\n# 文本框:\n")
                    f.write(f"# - 垂直对齐: {tf['vertical_anchor']}\n")
                    f.write(f"# - 自动换行: {tf['word_wrap']}\n")
                    f.write(f"# - 边距: 左{tf['margin_left_cm']}cm, 右{tf['margin_right_cm']}cm, 上{tf['margin_top_cm']}cm, 下{tf['margin_bottom_cm']}cm\n")

                    for para in tf['paragraphs']:
                        f.write(f"\n# 段落{para['index']}: \"{para['text']}\"\n")
                        f.write(f"# - 对齐: {para['alignment']}\n")

                        for run in para['runs']:
                            font = run['font']
                            f.write(f"#   Run: \"{run['text']}\"\n")
                            f.write(f"#     字体: {font['name']}, {font['size_pt']}pt\n")
                            if font['bold']:
                                f.write(f"#     粗体: {font['bold']}\n")
                            if font.get('color'):
                                color = font['color']
                                f.write(f"#     颜色类型: {color['type']}\n")
                                if color.get('rgb'):
                                    f.write(f"#     RGB: {color['rgb']}\n")
                                if color.get('theme_color'):
                                    f.write(f"#     主题颜色: {color['theme_color']}\n")

                if shape.get('table'):
                    table = shape['table']
                    f.write(f"\n# 表格: {table['rows']}行 × {table['columns']}列\n")
                    for cell in table['sample_cells']:
                        f.write(f"# 单元格{cell['position']}: \"{cell['text']}\"\n")
                        if cell.get('font'):
                            f.write(f"#   字体: {cell['font']['name']}, {cell['font']['size_pt']}pt, 粗体={cell['font']['bold']}\n")

                if shape.get('chart'):
                    chart = shape['chart']
                    f.write(f"\n# 图表:\n")
                    f.write(f"# - 类型: {chart['type']}\n")
                    f.write(f"# - 标题: {chart['title']}\n")

                f.write(f"```\n\n")

                f.write("---\n\n")

    print(f"✅ 生成 {len(rules['slides'])} 个详细规则文档")


if __name__ == "__main__":
    template_path = "/Users/bozhenwang/Desktop/自动生成报告PPT模板/牧场牧场育种分析报告_20251106_112926.pptx"

    print("="*80)
    print("开始提取所有58页的完整规则...")
    print("="*80)

    # 提取规则
    rules = extract_all_slides_rules(template_path)

    # 保存JSON
    json_output = "/Users/bozhenwang/projects/mating/genetic_improve/core/ppt_report/ppt_generation_rules.json"
    with open(json_output, 'w', encoding='utf-8') as f:
        json.dump(rules, f, ensure_ascii=False, indent=2)
    print(f"\n✅ 完整规则JSON: {json_output}")

    # 生成规则文档
    docs_dir = "/Users/bozhenwang/projects/mating/genetic_improve/core/ppt_report/docs/generation_rules"
    generate_rules_document(rules, docs_dir)

    print("\n" + "="*80)
    print("✅ 所有规则提取完成！")
    print("="*80)
    print(f"\n规则文档位置: {docs_dir}")
    print("下一步：基于这些规则实现PPT生成代码")
